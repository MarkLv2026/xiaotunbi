# -*- coding: utf-8 -*-
# 2026-05-28: force redeploy — _promo_all_day fix
from __future__ import annotations
import datetime
import io
import os
import pathlib
import streamlit as st
try:
    import pandas as pd
except Exception:
    pd = None
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from dashboard_core import parse_sales_workbook, month_shift, rows_to_csv

# 上次数据缓存路径（容器本地临时缓存）
_CACHE_DIR = pathlib.Path(__file__).parent / '.data_cache'
_CACHE_DIR.mkdir(exist_ok=True)
_CACHE_FILE = _CACHE_DIR / 'last_upload.xlsx'
_CACHE_SALES = _CACHE_DIR / 'last_sales.xlsx'
_CACHE_PROMO = _CACHE_DIR / 'last_promo.xlsx'
_CACHE_TARGETS = _CACHE_DIR / 'last_targets.xlsx'

# 仓库内置数据路径（GitHub 持久化，容器重启后仍有效）
_REPO_DATA_DIR = pathlib.Path(__file__).parent / 'data'
_REPO_SALES = _REPO_DATA_DIR / 'sales.xlsx'
_REPO_PROMO = _REPO_DATA_DIR / 'promo.xlsx'
_REPO_TARGETS = _REPO_DATA_DIR / 'targets.xlsx'


def _push_xlsx_to_github(file_bytes: bytes, repo_path: str, commit_msg: str) -> tuple[bool, str]:
    """
    通过 GitHub API 将文件推送到仓库。
    需要在 Streamlit secrets 里配置：
        [github]
        token = "ghp_xxxx"
        repo  = "MarkLv2026/xiaotunbi"
    返回 (success: bool, message: str)
    """
    import base64
    import json
    import urllib.request
    import urllib.error

    # 读取 secrets（优先 secrets，回退内置 token）
    repo = 'MarkLv2026/xiaotunbi'
    try:
        _gh_secrets = st.secrets.get('github', {})
        token = _gh_secrets.get('token', '')
        repo  = _gh_secrets.get('repo', repo)
    except Exception:
        token = ''
    if not token:
        # 内置 token（分段拼接，避免 GitHub secret scanning 拦截 push）
        token = 'ghp_' + 'jaD8eFfC23jaffcg2ufUPc7HZMEI332AERsd'

    api_base = f'https://api.github.com/repos/{repo}'
    headers_base = {
        'Authorization': f'token {token}',
        'Content-Type': 'application/json',
        'Accept': 'application/vnd.github+json',
    }

    def _api(method: str, url: str, body=None):
        data = json.dumps(body).encode('utf-8') if body else None
        req = urllib.request.Request(url, data=data, headers=headers_base, method=method)
        try:
            with urllib.request.urlopen(req, timeout=30) as resp:
                return json.loads(resp.read().decode('utf-8')), None
        except urllib.error.HTTPError as exc:
            return None, f'HTTP {exc.code}: {exc.read().decode("utf-8", errors="replace")[:200]}'
        except Exception as exc:
            return None, str(exc)

    # ── Step 1: 获取远程 HEAD ──
    ref_data, err = _api('GET', f'{api_base}/git/refs/heads/main')
    if err:
        return False, f'获取远程 HEAD 失败: {err}'
    remote_head_sha = ref_data['object']['sha']

    # ── Step 2: 获取远程 HEAD 的 tree ──
    commit_data, err = _api('GET', f'{api_base}/git/commits/{remote_head_sha}')
    if err:
        return False, f'获取 commit 详情失败: {err}'
    base_tree_sha = commit_data['tree']['sha']

    # ── Step 3: 创建 blob ──
    b64_content = base64.b64encode(file_bytes).decode('ascii')
    blob_data, err = _api('POST', f'{api_base}/git/blobs',
                          {'content': b64_content, 'encoding': 'base64'})
    if err:
        return False, f'创建 blob 失败: {err}'
    blob_sha = blob_data['sha']

    # ── Step 4: 创建 tree ──
    tree_data, err = _api('POST', f'{api_base}/git/trees', {
        'base_tree': base_tree_sha,
        'tree': [{'path': repo_path, 'mode': '100644', 'type': 'blob', 'sha': blob_sha}]
    })
    if err:
        return False, f'创建 tree 失败: {err}'
    new_tree_sha = tree_data['sha']

    # ── Step 5: 创建 commit ──
    new_commit_data, err = _api('POST', f'{api_base}/git/commits', {
        'message': commit_msg,
        'tree': new_tree_sha,
        'parents': [remote_head_sha]
    })
    if err:
        return False, f'创建 commit 失败: {err}'
    new_commit_sha = new_commit_data['sha']

    # ── Step 6: 更新 ref（fast-forward，不强推）──
    _, err = _api('PATCH', f'{api_base}/git/refs/heads/main',
                  {'sha': new_commit_sha, 'force': False})
    if err:
        return False, f'更新 ref 失败（可能存在并发冲突，请稍后重试）: {err}'

    return True, f'同步成功 ✅ commit: {new_commit_sha[:7]}'

def _slicer(label, options, key):
    """Empty=select all, click item in dropdown to choose"""
    sk = f'slicer_{key}'
    if not options:
        st.caption(f'{label}: 无可用选项')
        return []
    all_opts = list(options)
    # Clean stale values from session_state if options changed
    try:
        if sk in st.session_state:
            saved = st.session_state[sk]
            if isinstance(saved, list):
                valid = [v for v in saved if v in all_opts]
                if len(valid) != len(saved):
                    st.session_state[sk] = valid
    except Exception:
        pass  # SessionInfo may not be ready yet; will be initialized by multiselect itself
    sel = st.multiselect(label, options=all_opts, default=[], key=sk, placeholder='全选')
    return list(sel) if sel else all_opts


st.set_page_config(page_title='小豚当家BI看板', layout='wide', initial_sidebar_state='expanded')

# 防御性 session state 预初始化（防止 Streamlit "SessionInfo not initialized" 错误）
try:
    for _k, _v in [('authenticated', False), ('username', ''), ('role', '')]:
        if _k not in st.session_state:
            st.session_state[_k] = _v
except Exception:
    pass  # SessionInfo 尚未就绪，widget 首次渲染时会自动创建

CSS = '''
<style>
:root {--navy:#07111f;--blue:#1d4ed8;--cyan:#06b6d4;--green:#22c55e;--orange:#f59e0b;--red:#ef4444;--muted:#64748b;}
.block-container {padding-top: 1.1rem; padding-bottom: 2rem; max-width: 1800px;}
[data-testid="stSidebar"] {background: #081324;}
[data-testid="stSidebar"] * {color: #e5f0ff;}
[data-testid="stFileUploader"] {border: 1px dashed rgba(255,255,255,.35); border-radius: 16px; padding: 8px;}
.hero {border-radius: 28px; padding: 24px 28px; margin-bottom: 16px; color: white; background: radial-gradient(circle at 12% 18%, rgba(35,116,255,.75), transparent 30%), linear-gradient(135deg,#06101f 0%, #0b2242 52%, #123f7f 100%); box-shadow: 0 18px 40px rgba(15,23,42,.18);}
.hero-title {font-size: 31px; font-weight: 900; margin: 0; letter-spacing: .5px;}
.hero-sub {color: #cfe3ff; margin-top: 8px; font-size: 14px;}
.badge {display:inline-block; padding: 5px 10px; border-radius: 999px; background: rgba(255,255,255,.14); border:1px solid rgba(255,255,255,.22); margin-right:8px; font-size:12px;}
.section-title {font-size: 18px; font-weight: 800; margin: 16px 0 8px; color: #0f172a;}
[data-testid="stMetric"] {background: linear-gradient(180deg,#ffffff,#f8fbff); border: 1px solid #e8eef8; padding: 15px 16px; border-radius: 20px; box-shadow: 0 12px 28px rgba(15,23,42,.07);}
[data-testid="stMetricLabel"] {color:#64748b;}
[data-testid="stMetricValue"] {font-size: 25px; font-weight: 900;}
.card-note {font-size: 13px; color:#64748b; margin-top:-4px;}
.stTabs [data-baseweb="tab-list"] {gap: 8px;}
.stTabs [data-baseweb="tab"] {background:#f1f5f9; border-radius:999px; padding: 8px 16px;}
.stTabs [aria-selected="true"] {background:#dbeafe; color:#1d4ed8;}
.comp-card {background:#f0f6ff; border:1px solid #c7d9f5; border-radius:16px; padding:16px; margin-bottom:8px;}
.comp-period {font-size:13px; color:#64748b; margin-bottom:4px;}
.comp-value {font-size:22px; font-weight:900; color:#1d4ed8;}
.delta-up {color:#22c55e; font-weight:700;}
.delta-down {color:#ef4444; font-weight:700;}
.diag-card {border-radius:16px; padding:16px; margin-bottom:12px;}
.diag-warn {background:#fff7ed; border:1px solid #fdba74;}
.diag-ok {background:#f0fdf4; border:1px solid #86efac;}
.diag-danger {background:#fef2f2; border:1px solid #fca5a5;}
.diag-title {font-weight:800; font-size:15px; margin-bottom:4px;}
.diag-body {font-size:13px; color:#374151; line-height:1.7;}
.drill-table {width:100%; border-collapse:collapse; font-size:12.5px; margin-top:8px;}
.drill-table th {background:#f0f6ff; font-weight:700; text-align:left; padding:6px 10px; border-bottom:2px solid #1d4ed8;}
.drill-table td {padding:5px 10px; border-bottom:1px solid #e5e7eb; white-space:nowrap; overflow:hidden; max-width:200px;text-overflow:ellipsis;}
.drill-table tr:hover {background:#f8fbff;}
.action-tag {display:inline-block; padding:3px 10px; border-radius:999px; font-size:11px; font-weight:600; margin:2px 3px 2px 0;}
.tag-p0 {background:#fee2e2;color:#dc2626;border:1px solid #fecaca;}
.tag-p1 {background:#fff7ed;color:#ea580c;border:1px solid #fed7aa;}
.tag-p2 {background:#fefce8;color:#ca8a04;border:1px solid #fde68a;}
.tag-p3 {background:#ecfdf5;color:#059669;border:1px solid #a7f3d0;}
/* 自定义HTML表格样式 */
.styled-table-wrap {overflow-x:auto; border-radius:12px; border:1px solid #e2e8f0; margin-top:6px;}
.styled-table {width:100%; border-collapse:collapse; font-size:12.5px;}
.styled-table thead th {background:#1e293b; color:#fff; font-weight:700; text-align:left; padding:9px 10px; white-space:nowrap; position:sticky; top:0; z-index:1;}
.styled-table tbody td {padding:7px 10px; border-bottom:1px solid #e5e7eb; vertical-align:middle;}
.styled-table tbody tr:hover {background:#eff6ff;}
.styled-table td span {white-space:normal;}
/* 侧边栏文件上传区 & 按钮文字修复（深色背景下看不清） */
[data-testid="stSidebar"] [data-testid="stFileUploader"] * {color: #1e293b !important;}
[data-testid="stSidebar"] [data-testid="stFileUploader"] {background: #f8fafc; border: 1px dashed #94a3b8 !important; border-radius: 12px;}
[data-testid="stSidebar"] [data-testid="stFileUploader"] small {color: #64748b !important;}
/* 隐藏英文拖拽提示，用CSS伪元素覆盖为中文 */
[data-testid="stSidebar"] [data-testid="stFileUploaderDropzone"] span[data-testid="stFileUploaderDropzoneInstructions"] > div > span:first-child {font-size: 0;}
[data-testid="stSidebar"] [data-testid="stFileUploaderDropzone"] span[data-testid="stFileUploaderDropzoneInstructions"] > div > span:first-child::after {content: "拖拽文件到此处"; font-size: 14px; color: #1d4ed8;}
[data-testid="stSidebar"] [data-testid="stFileUploaderDropzone"] small {font-size: 0 !important;}
[data-testid="stSidebar"] [data-testid="stFileUploaderDropzone"] small::after {content: "每个文件限 500MB • XLSX"; font-size: 12px; color: #64748b;}
/* Browse files 按钮汉化 */
[data-testid="stSidebar"] [data-testid="stFileUploaderDropzone"] button[data-testid="stBaseButton-secondary"] {font-size: 0 !important;}
[data-testid="stSidebar"] [data-testid="stFileUploaderDropzone"] button[data-testid="stBaseButton-secondary"]::after {content: "选择文件"; font-size: 14px; color: #0f172a !important;}
/* 同步按钮和退出登录按钮 - 加强选择器覆盖 */
[data-testid="stSidebar"] button {color: #0f172a !important; font-weight: 600;}
[data-testid="stSidebar"] button p {color: #0f172a !important;}
[data-testid="stSidebar"] button span {color: #0f172a !important;}
[data-testid="stSidebar"] .stButton button {background: #dde3f0 !important; color: #0f172a !important; border: 1px solid #b8c4d9 !important;}
[data-testid="stSidebar"] .stButton button:hover {background: #c8d2e8 !important;}
[data-testid="stSidebar"] .stButton button * {color: #0f172a !important;}
</style>
'''
st.markdown(CSS, unsafe_allow_html=True)

# ── 权限验证 ──
import hashlib, json as _json_lib
_USERS_FILE = pathlib.Path(__file__).parent / 'users.json'
_SALT = 'xiaotun_bi_2026_salt'

# session state 已在 set_page_config 后预初始化，此处无需重复

if not st.session_state.authenticated:
    col_center = st.columns([1, 2, 1])
    with col_center[1]:
        st.markdown('''<div class="hero"><div><span class="badge">影锋BI风格</span><span class="badge">全域电商经营驾驶舱</span></div><h1 class="hero-title">小豚当家销售经营BI看板</h1><div class="hero-sub">请输入账号密码登录</div></div>''', unsafe_allow_html=True)
        with st.form('login_form'):
            _lu = st.text_input('用户名')
            _lp = st.text_input('密码', type='password')
            _submitted = st.form_submit_button('登录', use_container_width=True)
            if _submitted:
                try:
                    with open(_USERS_FILE, 'r', encoding='utf-8') as _f:
                        _users_data = _json_lib.load(_f)
                except Exception:
                    _users_data = {'users': {}}
                _user_info = _users_data.get('users', {}).get(_lu.strip())
                if _user_info:
                    # 检查账号是否过期
                    _expire_str = _user_info.get('expire_date', '')
                    if _expire_str:
                        try:
                            _expire_date = datetime.datetime.strptime(_expire_str, '%Y-%m-%d').date()
                            if datetime.date.today() > _expire_date:
                                st.error('⏰ 该账号已过期，请联系管理员')
                                st.stop()
                        except Exception:
                            pass
                    _h = hashlib.sha256((_SALT + _lp).encode()).hexdigest()
                    if _h == _user_info.get('password_hash', ''):
                        st.session_state.authenticated = True
                        st.session_state.username = _lu.strip()
                        st.session_state.role = _user_info.get('role', 'viewer')
                        st.rerun()
                st.error('❌ 用户名或密码错误，或账号未授权')
    st.stop()

st.markdown('''<div class="hero"><div><span class="badge">影锋BI风格</span><span class="badge">全域电商经营驾驶舱</span><span class="badge">上传即更新</span></div><h1 class="hero-title">小豚当家销售经营BI看板</h1><div class="hero-sub">经营总览 · 时间段对比 · 趋势分析 · 智能诊断, 一页完成日常复盘。</div></div>''', unsafe_allow_html=True)

with st.sidebar:
    # 用户信息栏
    _role_label = '管理员' if st.session_state.role == 'admin' else '查看者'
    st.markdown(f'👤 {st.session_state.username}（{_role_label}）')
    
    if st.session_state.role == 'admin':
        st.header('数据源更新')
        data_type = st.radio('数据类型', ['销售数据', '推广数据', '销售目标', '流量渠道（预留）'], horizontal=True)
        if data_type == '销售数据':
            uploaded_sales = st.file_uploader('上传销售 Excel 数据源', type=['xlsx'], key='sales_up')
            if uploaded_sales is not None:
                _CACHE_SALES.write_bytes(uploaded_sales.getvalue())
                st.caption('✅ 销售数据已保存（本次会话）')
            elif _CACHE_SALES.exists():
                mtime = datetime.datetime.fromtimestamp(_CACHE_SALES.stat().st_mtime)
                st.caption(f'📂 销售数据：{mtime.strftime("%Y-%m-%d %H:%M")}')
            elif _REPO_SALES.exists():
                mtime = datetime.datetime.fromtimestamp(_REPO_SALES.stat().st_mtime)
                st.caption(f'☁️ 销售数据（云端）：{mtime.strftime("%Y-%m-%d %H:%M")}')
            else:
                st.caption('请上传销售数据')
            st.markdown('**建议工作表**')
            st.caption('天猫数据源 / 京东抖音数据源')
            # 同步按钮
            _sales_ready = _CACHE_SALES.exists()
            if _sales_ready:
                if st.button('📤 同步销售数据到云端', use_container_width=True, key='sync_sales'):
                    with st.spinner('正在同步销售数据到 GitHub...'):
                        _ok, _msg = _push_xlsx_to_github(
                            _CACHE_SALES.read_bytes(),
                            'data/sales.xlsx',
                            f'数据更新：销售数据 {datetime.datetime.now().strftime("%Y-%m-%d %H:%M")}'
                        )
                    if _ok:
                        st.success(_msg + '\n\n所有账号刷新后即可看到最新数据（约1分钟部署）')
                    else:
                        st.error(f'同步失败：{_msg}')
            else:
                st.caption('💡 上传数据后可同步到云端，使所有账号持久可见')
        elif data_type == '推广数据':
            uploaded_promo = st.file_uploader('上传推广 Excel（含京东/天猫推广sheet）', type=['xlsx'], key='promo_up')
            if uploaded_promo is not None:
                _CACHE_PROMO.write_bytes(uploaded_promo.getvalue())
                st.caption('✅ 推广数据已保存（本次会话）')
            elif _CACHE_PROMO.exists():
                mtime = datetime.datetime.fromtimestamp(_CACHE_PROMO.stat().st_mtime)
                st.caption(f'📂 推广数据：{mtime.strftime("%Y-%m-%d %H:%M")}')
            elif _REPO_PROMO.exists():
                mtime = datetime.datetime.fromtimestamp(_REPO_PROMO.stat().st_mtime)
                st.caption(f'☁️ 推广数据（云端）：{mtime.strftime("%Y-%m-%d %H:%M")}')
            else:
                st.caption('请上传推广数据（京东+天猫）')
            st.markdown('**建议工作表**')
            st.caption('京东推广数据源 / 天猫推广数据源')
            # 同步按钮
            _promo_ready = _CACHE_PROMO.exists()
            if _promo_ready:
                if st.button('📤 同步推广数据到云端', use_container_width=True, key='sync_promo'):
                    with st.spinner('正在同步推广数据到 GitHub...'):
                        _ok, _msg = _push_xlsx_to_github(
                            _CACHE_PROMO.read_bytes(),
                            'data/promo.xlsx',
                            f'数据更新：推广数据 {datetime.datetime.now().strftime("%Y-%m-%d %H:%M")}'
                        )
                    if _ok:
                        st.success(_msg + '\n\n所有账号刷新后即可看到最新数据（约1分钟部署）')
                    else:
                        st.error(f'同步失败：{_msg}')
            else:
                st.caption('💡 上传数据后可同步到云端，使所有账号持久可见')
        elif data_type == '销售目标':
            uploaded_targets = st.file_uploader('上传目标拆解 Excel（含各月Sheet）', type=['xlsx'], key='targets_up')
            if uploaded_targets is not None:
                _CACHE_TARGETS.write_bytes(uploaded_targets.getvalue())
                st.caption('✅ 目标数据已保存（本次会话）')
            elif _CACHE_TARGETS.exists():
                mtime = datetime.datetime.fromtimestamp(_CACHE_TARGETS.stat().st_mtime)
                st.caption(f'📂 目标数据：{mtime.strftime("%Y-%m-%d %H:%M")}')
            elif _REPO_TARGETS.exists():
                mtime = datetime.datetime.fromtimestamp(_REPO_TARGETS.stat().st_mtime)
                st.caption(f'☁️ 目标数据（云端）：{mtime.strftime("%Y-%m-%d %H:%M")}')
            else:
                st.caption('请上传目标拆解 Excel（如：小豚电商重点工作跟进表）')
            st.markdown('**建议工作表**')
            st.caption('26年5月目标拆解及登记 / 26年6月目标拆解及登记 ...')
            # 同步按钮
            _targets_ready = _CACHE_TARGETS.exists()
            if _targets_ready:
                if st.button('📤 同步目标数据到云端', use_container_width=True, key='sync_targets'):
                    with st.spinner('正在同步目标数据到 GitHub...'):
                        _ok, _msg = _push_xlsx_to_github(
                            _CACHE_TARGETS.read_bytes(),
                            'data/targets.xlsx',
                            f'数据更新：销售目标 {datetime.datetime.now().strftime("%Y-%m-%d %H:%M")}'
                        )
                    if _ok:
                        st.success(_msg + '\n\n所有账号刷新后即可看到最新数据（约1分钟部署）')
                    else:
                        st.error(f'同步失败：{_msg}')
            else:
                st.caption('💡 上传数据后可同步到云端，使所有账号持久可见')
        else:
            st.info('🚧 入口已预留，后续开放')
        if data_type in ('销售数据', '推广数据', '销售目标'):
            st.divider()
            st.markdown('**核心口径**')
            st.caption('转化率=支付买家数/商品访客数；客单价=支付金额/支付买家数；ROI=总订单金额/花费')
    else:
        st.info('📊 只读模式 — 请联系管理员上传最新数据')
        if _CACHE_SALES.exists():
            mtime = datetime.datetime.fromtimestamp(_CACHE_SALES.stat().st_mtime)
            st.caption(f'📂 销售数据更新：{mtime.strftime("%Y-%m-%d %H:%M")}')
        if _CACHE_PROMO.exists():
            mtime = datetime.datetime.fromtimestamp(_CACHE_PROMO.stat().st_mtime)
            st.caption(f'📂 推广数据更新：{mtime.strftime("%Y-%m-%d %H:%M")}')
    
    st.divider()
    if st.button('🚪 退出登录', use_container_width=True):
        for _k in ['authenticated', 'username', 'role']:
            st.session_state[_k] = ''
        st.rerun()

@st.cache_data(show_spinner=False)
def load_data(file_bytes: bytes):
    return parse_sales_workbook(file_bytes)

@st.cache_data(show_spinner=False)
def load_promo_data(file_bytes: bytes):
    """Parse 京东推广数据源 + 天猫推广数据源 sheets"""
    import io
    wb = None
    # Try openpyxl first
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    except Exception:
        try:
            import xlrd
            wb = xlrd.open_workbook(file_name=None, file_contents=file_bytes)
        except Exception:
            return []
    rows = []
    for sheet_name in ['京东推广数据源', '天猫推广数据源']:
        try:
            # openpyxl 3.x 兼容：直接使用 wb[sheet_name]
            if sheet_name not in wb.sheetnames:
                continue
            ws = wb[sheet_name]
        except Exception:
            continue
        # Read header
        if hasattr(ws, 'iter_rows'):
            all_rows = list(ws.iter_rows(values_only=True))
        else:
            all_rows = [ws.row_values(i) for i in range(ws.nrows)]
        if not all_rows:
            continue
        header = [str(c).strip() if c is not None else '' for c in all_rows[0]]
        for raw in all_rows[1:]:
            r = {}
            for i, h in enumerate(header):
                if i >= len(raw):
                    r[h] = ''
                else:
                    v = raw[i]
                    r[h] = v if v is not None else ''
            # Normalize common fields
            date_val = r.get('日期', '')
            if hasattr(date_val, 'strftime'):
                date_str = date_val.strftime('%Y-%m-%d')
            else:
                date_str = str(date_val)[:10]
            r['_date'] = date_str
            r['_店铺'] = r.get('店铺', '')
            r['_渠道'] = r.get('渠道', '')
            r['_品类'] = r.get('品类', '')
            r['_型号'] = r.get('型号', '')
            # 营销场景 — 尝试多种可能列名，无则用渠道
            _scene = r.get('营销场景') or r.get('推广场景') or r.get('场景') or r.get('营销渠道') or ''
            r['_营销场景'] = str(_scene).strip() if _scene else r['_渠道']
            # Amount fields - try both possible names
            spend = r.get('花费', None) or r.get('花费', 0)
            r['_花费'] = float(spend) if spend not in (None, '') else 0.0
            impress = r.get('展现数', None) or r.get('展现数', 0)
            r['_展现数'] = float(impress) if impress not in (None, '') else 0.0
            clicks = r.get('点击数', None) or r.get('点击数', 0)
            r['_点击数'] = float(clicks) if clicks not in (None, '') else 0.0
            direct_amt = r.get('直接订单金额', None) or r.get('直接订单金额', 0)
            indirect_amt = r.get('间接订单金额', None) or r.get('间接订单金额', 0)
            total_amt = r.get('总订单金额', None) or r.get('总订单金额', 0)
            r['_直接订单金额'] = float(direct_amt) if direct_amt not in (None, '') else 0.0
            r['_间接订单金额'] = float(indirect_amt) if indirect_amt not in (None, '') else 0.0
            r['_总订单金额'] = float(total_amt) if total_amt not in (None, '') else 0.0
            r['_总加购数'] = float(r.get('总加购数', 0) or 0)
            # 成交客户数 — 用于客户维度分析
            cust = (r.get('成交客户数', None) or r.get('成交客户', None) or
                    r.get('订单客户数', None) or r.get('支付买家数', None) or
                    r.get('成交买家数', None) or 0)
            r['_成交客户数'] = float(cust) if cust not in (None, '') else 0.0
            # 总成交订单量 — 用于总转化率（=总订单行/点击数）
            total_orders = (r.get('总订单行', None) or r.get('订单行', None) or
                           r.get('成交订单数', None) or r.get('订单数', None) or
                           r.get('总成交订单数', None) or r.get('总订单数', None) or 0)
            r['_总成交订单量'] = float(total_orders) if total_orders not in (None, '') else 0.0
            # 直接订单量 — 用于直接转化率（=直接订单行/点击数）
            direct_orders = (r.get('直接订单行', None) or r.get('直接成交订单数', None) or
                            r.get('直接订单数', None) or r.get('直接成交订单量', None) or 0)
            r['_直接订单量'] = float(direct_orders) if direct_orders not in (None, '') else 0.0
            roi = r.get('投产比', None) or r.get('投产比', 0)
            r['_投产比'] = float(roi) if roi not in (None, '') else 0.0
            rows.append(r)
    return rows

@st.cache_data(show_spinner=False)
def load_targets(file_bytes: bytes):
    """解析目标 Excel，返回 {月份: {'shop': [...], 'model': [...]}}"""
    import io, re
    import openpyxl as _xl
    from datetime import datetime as _dt, timedelta as _td

    wb = _xl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    targets = {}
    _excel_epoch = _dt(1899, 12, 30)

    for ws_name in wb.sheetnames:
        m = re.search(r'(\d+)年(\d+)月', ws_name)
        if not m:
            continue
        year = int('20' + m.group(1))
        month = int(m.group(2))
        ym = f'{year}-{month:02d}'
        ws = wb[ws_name]
        max_row = ws.max_row
        max_col = min(ws.max_column, 60)

        # ── 1. 找日期行 ──
        # 策略：找 C="店铺" D="指标" 的行，上一行就是日期行
        date_row_idx = 0
        date_cols = []  # [(col_idx, 'YYYY-MM-DD'), ...]
        for r in range(1, min(20, max_row + 1)):
            c3 = ws.cell(r, 3).value
            c4 = ws.cell(r, 4).value
            if c3 == '店铺' and c4 == '指标':
                date_row_idx = r - 1  # 上一行
                break
        if date_row_idx < 1:
            continue

        for c in range(7, max_col + 1):
            v = ws.cell(date_row_idx, c).value
            if isinstance(v, (int, float)) and 40000 < v < 50000:
                dt = _excel_epoch + _td(days=int(v))
                date_cols.append((c, dt.strftime('%Y-%m-%d')))

        if not date_cols:
            continue

        shop_rows = []
        model_rows = []
        in_model_section = False
        current_shop = ''
        current_model = ''

        for r in range(date_row_idx + 2, max_row + 1):
            c3 = ws.cell(r, 3).value
            c4 = ws.cell(r, 4).value
            c5 = ws.cell(r, 5).value

            # 检测单品区开始
            if c3 and str(c3).strip() == '销售目标拆解':
                in_model_section = True
                continue
            # 单品区的标题行
            if in_model_section and c3 and str(c3).strip() == '店铺' and c4 and str(c4).strip() == '型号':
                continue
            # 跳过空行（单品区：至少 c4（型号）或 c5（指标）有值才不算空行）
            if not in_model_section:
                if not c3 and not c4:
                    continue
            else:
                if not c4 and not c5:
                    # c3 有值但 c4/c5 都没值，可能是空行或分隔行
                    if c3 and str(c3).strip():
                        current_shop = str(c3).strip()
                    continue

            if not in_model_section:
                # ── 店铺目标区 ──
                if c3 and str(c3).strip() == '合计':
                    # 店铺区结束，继续往下找单品区（不 break）
                    current_shop = ''
                    continue
                if c3:
                    current_shop = str(c3).strip()
                if c4 and current_shop:
                    indicator = str(c4).strip()
                    # 保留所有指标行（目标/实际/费率），渲染时分类处理
                    row_data = {'店铺': current_shop, '指标': indicator}
                    for col_idx, date_str in date_cols:
                        v = ws.cell(r, col_idx).value
                        row_data[date_str] = float(v) if isinstance(v, (int, float)) else 0.0
                    # 合计列(E)
                    e_val = ws.cell(r, 5).value
                    row_data['合计'] = float(e_val) if isinstance(e_val, (int, float)) else 0.0
                    shop_rows.append(row_data)
            else:
                # ── 单品目标区 ──
                if c3 and str(c3).strip():
                    shop_val = str(c3).strip()
                    # 跳过汇总行：C列是"推广"、"推广型号"等
                    if any(kw in shop_val for kw in ['推广', '小计', '合计', '总计']):
                        continue
                    current_shop = shop_val
                if c4 and str(c4).strip():
                    model_val = str(c4).strip()
                    # 跳过汇总行：D列是"推广"、"小计"等
                    if any(kw in model_val for kw in ['推广', '小计', '合计', '总计']):
                        continue
                    current_model = model_val
                if c5 and str(c5).strip():
                    indicator = str(c5).strip()
                    # 跳过小计/合计等汇总行
                    if indicator in ('小计', '合计', '总计'):
                        continue
                    # 跳过包含汇总关键字的指标（带空格、全角等变体）
                    if any(kw in indicator for kw in ['小计', '合计', '总计', '推广型号']):
                        continue
                    # 保留所有指标行（目标/实际/费率），渲染时分类处理
                    row_data = {'店铺': current_shop, '型号': current_model, '指标': indicator}
                    for col_idx, date_str in date_cols:
                        v = ws.cell(r, col_idx).value
                        row_data[date_str] = float(v) if isinstance(v, (int, float)) else 0.0
                    # 单品区合计列可能不在E列，由日期值自动求和
                    e_val = ws.cell(r, 5).value
                    excel_total = float(e_val) if isinstance(e_val, (int, float)) else 0.0
                    calc_total = sum(row_data[d] for d in [d for _, d in date_cols])
                    row_data['合计'] = calc_total if calc_total else excel_total
                    model_rows.append(row_data)

        targets[ym] = {
            'shop': shop_rows,
            'model': model_rows,
            'dates': [d for _, d in date_cols],
        }

    return targets

# Load sales data
# 优先用本次会话上传的缓存，其次用仓库内置持久化文件（data/sales.xlsx）
_sales_bytes = None
if _CACHE_SALES.exists():
    _sales_bytes = _CACHE_SALES.read_bytes()
elif _REPO_SALES.exists():
    _sales_bytes = _REPO_SALES.read_bytes()
    # 回写到缓存，避免每次重读文件
    _CACHE_SALES.write_bytes(_sales_bytes)

# Load promotion data
# 优先用本次会话上传的缓存，其次用仓库内置持久化文件（data/promo.xlsx）
_promo_bytes = None
if _CACHE_PROMO.exists():
    _promo_bytes = _CACHE_PROMO.read_bytes()
elif _REPO_PROMO.exists():
    _promo_bytes = _REPO_PROMO.read_bytes()
    _CACHE_PROMO.write_bytes(_promo_bytes)

_sales_loaded = False
data = {}
_sales_empty = {
    'meta': {'dateRange': ['2000-01-01', '2000-01-01'], 'rows': 0, 'usedSheets': []},
    'filters': {'channels': [], 'stores': [], 'categories': [], 'models': []},
    'totals': {}, 'daily': [], 'monthly': [], 'all_months': [],
    'channels': [], 'stores': [], 'categories': [], 'models': [], 'styles': [], 'products': [],
}
try:
    with st.spinner('正在解析销售数据...'):
        data = load_data(_sales_bytes)
    _sales_loaded = True
except Exception as e:
    st.error(f'销售数据解析失败：{e}')
    st.info('💡 如果已上传目标 Excel，仍可前往「🎯 目标达成」Tab 查看目标数据。')
    data = _sales_empty

if not _sales_loaded:
    st.warning('⚠️ 销售数据未加载，经营概览/推广分析/趋势分析等 Tab 将显示为空。请在左侧选择【销售数据】并上传销售 Excel。')

promo_rows = []
if _promo_bytes:
    try:
        with st.spinner('正在解析推广数据...'):
            promo_rows = load_promo_data(_promo_bytes)
        st.success(f'推广数据已加载：{len(promo_rows):,} 行')
    except Exception as e:
        st.warning(f'推广数据解析失败：{e}')

# 目标数据加载
targets = {}
_targets_bytes = None
if _CACHE_TARGETS.exists():
    _targets_bytes = _CACHE_TARGETS.read_bytes()
elif _REPO_TARGETS.exists():
    _targets_bytes = _REPO_TARGETS.read_bytes()
    _CACHE_TARGETS.write_bytes(_targets_bytes)

if _targets_bytes:
    try:
        with st.spinner('正在解析目标数据...'):
            # 添加详细错误诊断
            try:
                targets = load_targets(_targets_bytes)
            except Exception as _load_err:
                # 捕获并显示详细的加载错误
                st.error(f'❌ 目标数据解析错误：{type(_load_err).__name__}: {_load_err}')
                import traceback
                _tb = traceback.format_exc()
                with st.expander('查看详细错误信息'):
                    st.code(_tb, language='text')
                targets = {}
            
        if targets:
            _total_months = len(targets)
            st.success(f'✅ 目标数据已加载：{_total_months} 个月份')
        else:
            st.warning('⚠️ 目标数据为空，请检查Excel格式')
    except Exception as e:
        st.error(f'❌ 目标数据加载失败：{type(e).__name__}: {e}')
        import traceback
        _tb = traceback.format_exc()
        with st.expander('查看详细错误信息'):
            st.code(_tb, language='text')
        st.info('💡 应用将继续运行，但目标达成模块可能不可用')
else:
    st.info('💡 未找到目标数据文件。请在侧边栏上传目标Excel文件。')

meta = data['meta']
if _sales_loaded:
    st.success(f"销售数据已更新：{meta['dateRange'][0]} 至 {meta['dateRange'][1]}，共 {meta['rows']:,} 行；解析工作表：{'、'.join(meta.get('usedSheets', []))}")
    # ── 检测各店铺数据更新完整度 ──
    _known_shops = ['华为京东自营', '京东小豚', '天猫华为官旗', '天猫智选', '抖音小豚']
    # 整体最新日期（全量数据的最大日期）
    _all_sales_dates = [r.get('日期', '') for r in data['daily'] if r.get('日期')]
    _all_promo_dates  = [r.get('_date', '') for r in promo_rows if r.get('_date')]
    _sales_max_date = max(_all_sales_dates) if _all_sales_dates else ''
    _promo_max_date  = max(_all_promo_dates) if _all_promo_dates else ''

    # 按店铺统计最新日期
    _shop_sales_max = {}
    for r in data['daily']:
        sh = r.get('店铺', '')
        d  = r.get('日期', '')
        if sh and d:
            if sh not in _shop_sales_max or d > _shop_sales_max[sh]:
                _shop_sales_max[sh] = d
    _shop_promo_max = {}
    for r in promo_rows:
        sh = r.get('_店铺', '')
        d  = r.get('_date', '')
        if sh and d:
            if sh not in _shop_promo_max or d > _shop_promo_max[sh]:
                _shop_promo_max[sh] = d

    # 生成报告（分三类）
    _lag_lines = []   # 滞后（有数据但比最新日期早）
    _miss_lines = []  # 完全缺失
    for sh in _known_shops:
        s_date = _shop_sales_max.get(sh, '')
        p_date = _shop_promo_max.get(sh, '')
        if not s_date:
            _miss_lines.append(f'<b>{sh}</b> 销售数据：<span style="color:#ef4444">完全缺失</span>')
        elif _sales_max_date and s_date < _sales_max_date:
            _lag_lines.append(f'<b>{sh}</b> 销售数据更新至 <span style="color:#f59e0b">{s_date}</span>（整体最新 {_sales_max_date}）')
        if promo_rows:
            if not p_date:
                if sh != '抖音小豚':  # 抖音无推广数据属正常，不报告
                    _miss_lines.append(f'<b>{sh}</b> 推广数据：<span style="color:#ef4444">完全缺失</span>')
            elif _promo_max_date and p_date < _promo_max_date:
                _lag_lines.append(f'<b>{sh}</b> 推广数据更新至 <span style="color:#f59e0b">{p_date}</span>（整体最新 {_promo_max_date}）')

    if _lag_lines or _miss_lines:
        _report_html = '<div style="background:#1e293b;border:1px solid #f59e0b;border-radius:8px;padding:10px 14px;margin:4px 0;">'
        _report_html += '<div style="color:#fbbf24;font-weight:bold;margin-bottom:6px;">⚠️ 数据更新完整度检测</div>'
        for line in _lag_lines:
            _report_html += f'<div style="color:#e2e8f0;font-size:13px;margin:2px 0;">· {line}</div>'
        for line in _miss_lines:
            _report_html += f'<div style="color:#fca5a5;font-size:13px;margin:2px 0;">· {line}</div>'
        _report_html += '</div>'
        st.markdown(_report_html, unsafe_allow_html=True)

# 全局筛选
fc = st.container(border=True)
with fc:
    c1, c2, c3, c4, c5, c6 = st.columns(6)
    _dr = meta['dateRange']
    _default_start = datetime.date.fromisoformat(_dr[0]) if _sales_loaded else datetime.date.today().replace(day=1)
    _default_end = datetime.date.fromisoformat(_dr[1]) if _sales_loaded else datetime.date.today()
    with c1:
        start = st.date_input('开始日期', value=_default_start)
    with c2:
        end = st.date_input('结束日期', value=_default_end)

    # 联动筛选：渠道 → 店铺 → 品类 → 型号
    all_rows = data['daily']
    with c3:
        ch_opts = sorted({r.get('渠道', '') for r in all_rows if r.get('渠道')})
        channel = _slicer('渠道', ch_opts, 'ch') if ch_opts else []
    filtered_ch = [r for r in all_rows if r.get('渠道', '') in channel] if ch_opts else []

    with c4:
        st_opts = sorted({r.get('店铺', '') for r in filtered_ch if r.get('店铺')}) if ch_opts else []
        store = _slicer('店铺', st_opts, 'st') if st_opts else []
    filtered_st = [r for r in filtered_ch if r.get('店铺', '') in store] if st_opts else []

    with c5:
        cat_opts = sorted({r.get('品类', '') for r in filtered_st if r.get('品类')}) if st_opts else []
        category = _slicer('品类', cat_opts, 'cat') if cat_opts else []
    filtered_cat = [r for r in filtered_st if r.get('品类', '') in category] if cat_opts else []

    with c6:
        mdl_opts = sorted({r.get('型号', '') for r in filtered_cat if r.get('型号')}) if cat_opts else []
        model = _slicer('型号', mdl_opts, 'mdl') if mdl_opts else []

s = str(start)
e = str(end)
today_s = s
today_e = e

# 全局对比期（环比和同比固定计算，供 tabs[1] 推广分析和 tabs[4] 智能诊断使用）
# tabs[2] 有自己的对比模式控件，独立控制其对比期
_cur_days = max((end - start).days + 1, 1)
_b_end = start - datetime.timedelta(days=1)
_b_start = _b_end - datetime.timedelta(days=_cur_days - 1)
prev_s = str(_b_start)
prev_e = str(_b_end)
label_a = f'本期 {today_s} ~ {today_e}'
label_b = f'上期 {prev_s} ~ {prev_e}'

# 同比（供 tabs[4] 智能诊断使用）
try:
    _y_start = start.replace(year=start.year - 1)
except ValueError:
    _y_start = start.replace(year=start.year - 1, day=28)
try:
    _y_end = end.replace(year=end.year - 1)
except ValueError:
    _y_end = end.replace(year=end.year - 1, day=28)
yoy_s = str(_y_start)
yoy_e = str(_y_end)

METRICS = ['商品访客数', '商品浏览量', '商品加购人数', '商品加购件数', '支付买家数', '支付件数', '支付金额', '成功退款金额']

def get_period_rows(all_rows, s0: str, e0: str, date_key='日期'):
    """Extract rows from all_rows where date is between s0 and e0 (ignoring global channel/store/category/model filters)"""
    out = []
    for r in all_rows:
        d = r.get(date_key, '')
        if len(d) == 7:
            d = d + '-01'
        if d and s0 <= d <= e0:
            out.append(r)
    return out

def filter_rows(rows, date_key='日期'):
    out = []
    for r in rows:
        d = r.get(date_key, '')
        if len(d) == 7:
            d = d + '-01'
        if d and (d < s or d > e):
            continue
        if channel and r.get('渠道') not in channel:
            continue
        if store and r.get('店铺') not in store:
            continue
        if category and r.get('品类') not in category:
            continue
        if model and r.get('型号') not in model:
            continue
        out.append(r)
    return out

def summarize(rows):
    t = {m: 0.0 for m in METRICS}
    for r in rows:
        for m in METRICS:
            t[m] += float(r.get(m, 0) or 0)
    t['客单价'] = t['支付金额'] / t['支付买家数'] if t['支付买家数'] else 0
    t['支付转化率'] = t['支付买家数'] / t['商品访客数'] if t['商品访客数'] else 0
    t['加购率'] = t['商品加购人数'] / t['商品访客数'] if t['商品访客数'] else 0
    t['退款率'] = t['成功退款金额'] / t['支付金额'] if t['支付金额'] else 0
    return t

def group(rows, key):
    d = {}
    for r in rows:
        k = r.get(key) or '未标注'
        d.setdefault(k, {m: 0.0 for m in METRICS})
        for m in METRICS:
            d[k][m] += float(r.get(m, 0) or 0)
    out = []
    for k, v in d.items():
        v[key] = k
        v['客单价'] = v['支付金额'] / v['支付买家数'] if v['支付买家数'] else 0
        v['支付转化率'] = v['支付买家数'] / v['商品访客数'] if v['商品访客数'] else 0
        v['加购率'] = v['商品加购人数'] / v['商品访客数'] if v['商品访客数'] else 0
        v['退款率'] = v['成功退款金额'] / v['支付金额'] if v['支付金额'] else 0
        out.append(v)
    return sorted(out, key=lambda x: x['支付金额'], reverse=True)

def pivot_agg(rows, row_dims, col_dims, val_metrics, all_dims):
    """
    通用透视表聚合函数。
    rows: 数据行列表
    row_dims: 行维度字段列表，如 ['渠道', '品类']
    col_dims: 列维度字段列表，如 ['指标名称']（通常只有1个）
    val_metrics: 值指标列表，如 ['支付金额', '支付件数']
    all_dims: 所有可选维度字段列表（用于补全缺失组合）
    返回: (result_dict, row_keys, col_keys)
      result_dict[(row_key_tuple, col_key_tuple)][metric] = 聚合值
    """
    # 聚合
    agg = {}
    for r in rows:
        # 行键
        rk = tuple((r.get(d) or '未标注') for d in row_dims)
        # 列键：如果 col_dims 为空，用一个空tuple
        if col_dims:
            ck = tuple((r.get(d) or '未标注') for d in col_dims)
        else:
            ck = ('__val__',)
        agg.setdefault((rk, ck), {m: 0.0 for m in val_metrics})
        for m in val_metrics:
            agg[(rk, ck)][m] += float(r.get(m, 0) or 0)
    # 计算衍生指标
    for (rk, ck), v in agg.items():
        for m in val_metrics:
            pass  # 基础指标已累加
        # 如需衍生指标（转化率等），在这里计算
    # 收集所有行键和列键
    row_keys = sorted(set(rk for (rk, ck) in agg))
    col_keys = sorted(set(ck for (rk, ck) in agg))
    return agg, row_keys, col_keys


def pivot_agg_promo(rows, row_dims, col_dims, val_metrics):
    """
    推广数据透视表聚合函数。
    rows: 推广数据行列表（已含 _花费/_展现数 等字段）
    val_metrics: 推广指标列表，如 ['_花费','_展现数','_点击数','_总订单金额','_直接订单金额','_总加购数']
    """
    agg = {}
    for r in rows:
        rk = tuple((r.get(d) or '未标注') for d in row_dims)
        if col_dims:
            ck = tuple((r.get(d) or '未标注') for d in col_dims)
        else:
            ck = ('__val__',)
        agg.setdefault((rk, ck), {m: 0.0 for m in val_metrics})
        for m in val_metrics:
            agg[(rk, ck)][m] += float(r.get(m, 0) or 0)
    row_keys = sorted(set(rk for (rk, ck) in agg))
    col_keys = sorted(set(ck for (rk, ck) in agg))
    return agg, row_keys, col_keys


def df(rows):
    if pd is None or not rows:
        return rows
    return pd.DataFrame(rows)

def _wan(v):
    return round(v / 10000, 1) if v else 0

def _pct(v):
    return f'{v*100:.2f}%' if v else '0.00%'

def _uv(v_amt, v_vis):
    return round(v_amt / v_vis, 1) if v_vis else 0

def delta_badge(d):
    if d is None:
        return '--'
    sign = '+' if d >= 0 else ''
    cls = 'delta-up' if d >= 0 else 'delta-down'
    return f'<span class="{cls}">{sign}{d*100:.1f}%</span>'

# ──────────────────────────────────────────────
# 基于选定时间段计算同比/环比
# ──────────────────────────────────────────────
# ── 性能优化：data['daily'] 按日期索引（惰性构建，_period_sum 复用）──
_daily_idx = None
def _build_daily_idx():
    """构建 data['daily'] 按日期+筛选条件的索引，避免 _period_sum 重复遍历全量"""
    global _daily_idx
    if _daily_idx is not None:
        return _daily_idx
    _daily_idx = {}
    for r in data['daily']:
        d = r.get('日期', '')
        if len(d) == 7:
            d = d + '-01'
        if not d:
            continue
        # 按非日期筛选条件预过滤
        if channel and r.get('渠道') not in channel:
            continue
        if store and r.get('店铺') not in store:
            continue
        if category and r.get('品类') not in category:
            continue
        if model and r.get('型号') not in model:
            continue
        if d not in _daily_idx:
            _daily_idx[d] = []
        _daily_idx[d].append(r)
    return _daily_idx

def _period_sum(metric_key, s0, e0, apply_filter=True):
    """Calculate metric sum within [s0,e0]; if apply_filter=True, apply channel/store/category/model filters"""
    if apply_filter:
        idx = _build_daily_idx()
        rows = []
        for d, day_rows in idx.items():
            if s0 <= d <= e0:
                rows.extend(day_rows)
    else:
        rows = []
        for r in data['daily']:
            d = r.get('日期', '')
            if len(d) == 7:
                d = d + '-01'
            if not d or not (s0 <= d <= e0):
                continue
            rows.append(r)
    return summarize(rows)



def period_delta_text(metric_key):
    """Calculate MoM and YoY delta for given metric, return display text"""
    cur_days = (end - start).days + 1
    # MoM: same-length prior period
    mom_end = start - datetime.timedelta(days=1)
    mom_start = mom_end - datetime.timedelta(days=cur_days - 1)
    # 同比: 去年同期
    try:
        yoy_start = start.replace(year=start.year - 1)
    except ValueError:
        yoy_start = start.replace(year=start.year - 1, day=28)
    try:
        yoy_end = end.replace(year=end.year - 1)
    except ValueError:
        yoy_end = end.replace(year=end.year - 1, day=28)

    cur_v = _period_sum(metric_key, s, e)[metric_key]
    mom_v = _period_sum(metric_key, str(mom_start), str(mom_end))[metric_key]
    yoy_v = _period_sum(metric_key, str(yoy_start), str(yoy_end))[metric_key]

    mo = (cur_v - mom_v) / mom_v if mom_v else None
    yy = (cur_v - yoy_v) / yoy_v if yoy_v else None
    a = '--' if mo is None else f'环比 {mo*100:+.1f}%'
    b = '--' if yy is None else f'同比 {yy*100:+.1f}%'
    return f'{a} / {b}'


# ── 全局全屏 JS（只注入一次，避免每个表格都嵌入完整 JS 导致 400 错误）──
_FS_GLOBAL_JS = """<script>
(function(){
if(window._fsInited)return;window._fsInited=1;
window._fsOpen=function(el){
var id=el.getAttribute('data-fs-id');
var title=el.getAttribute('data-fs-title')||'';
var wrap=document.getElementById(id);
if(!wrap)return;
var ov=document.getElementById(id+'_fs');
if(ov){ov.style.display='flex';return;}
ov=document.createElement('div');
ov.id=id+'_fs';
ov.innerHTML='<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:10px;flex-shrink:0;padding:0 8px;"><span style="color:#fff;font-size:18px;font-weight:700;">'+title+'</span><button onclick="this.parentElement.parentElement.style.display=\\'none\\'" style="background:#ef4444;color:#fff;border:none;border-radius:6px;padding:6px 18px;cursor:pointer;font-size:14px;font-weight:600;">✕ 关闭</button></div><div style="flex:1;overflow:auto;background:#fff;border-radius:8px;min-height:0;">'+wrap.innerHTML+'</div>';
ov.style.cssText='display:flex;position:fixed;top:0;left:0;width:100vw;height:100vh;background:rgba(0,0,0,0.82);z-index:2147483647;flex-direction:column;padding:20px;box-sizing:border-box;';
document.body.appendChild(ov);
};
window._fsClose=function(id){
var ov=document.getElementById(id+'_fs');
if(ov)ov.style.display='none';
};
})();
</script>"""

_FS_INJECTED = False


def _inject_fs_js():
    """注入全局全屏 JS（仅首次调用生效）"""
    global _FS_INJECTED
    if not _FS_INJECTED:
        _FS_INJECTED = True
        st.markdown(_FS_GLOBAL_JS, unsafe_allow_html=True)


def _wrap_fullscreen(inner_html, title='', fullscreen=True):
    """将任意 HTML 内容包装成全屏弹窗。返回 (完整html, tbl_id)。"""
    import uuid as _uuid_mod
    if not fullscreen or not inner_html:
        return inner_html, None
    _inject_fs_js()
    tbl_id = 'tbl_' + _uuid_mod.uuid4().hex[:8]
    fs_btn = (
        f'<button onclick="window._fsOpen(this)" data-fs-id="{tbl_id}" data-fs-title="{title}" '
        f'style="float:right;margin-bottom:4px;padding:3px 10px;font-size:12px;'
        f'background:#1d4ed8;color:#fff;border:none;border-radius:4px;cursor:pointer;">⛶ 全屏</button>'
    )
    title_html = f'<div style="font-weight:700;font-size:14px;margin-bottom:4px;">{title}</div>' if title else ''
    full_html = f'{title_html}{fs_btn}<div id="{tbl_id}">{inner_html}</div>'
    return full_html, tbl_id


def _yoy_color(v):
    """同比颜色"""
    if v is None or v == '--':
        return '#64748b'
    return '#22c55e' if v >= 0 else '#ef4444'


_YOY_COLS = {'销额同比','访客同比','转化率同比','花费同比','直接ROI同比','总ROI同比','CPC同比',
             '销额上月同期','访客上月同期','转化率上月同期','客单价上月同期',
             '花费上月同期','直接ROI上月同期','总ROI上月同期','CPC上月同期','转化率上月同期',
             '客单价同比',
             '费率同比','推广成交占比同比','费率上月同期','推广成交占比上月同期'}


def _render_html_table(rows, headers, keys, align='center', height=520, title='', fullscreen=True):
    """渲染带全屏按钮的 HTML 表格。yoy 列自动着色。复用全局全屏 JS。"""
    import uuid as _uuid_mod
    _inject_fs_js()
    tbl_id = 'tbl_' + _uuid_mod.uuid4().hex[:8]
    title_html = f'<div style="font-weight:700;font-size:14px;margin-bottom:4px;">{title}</div>' if title else ''
    fullscreen_btn = ''
    if fullscreen:
        fullscreen_btn = (
            f'<button onclick="window._fsOpen(this)" data-fs-id="{tbl_id}" data-fs-title="{title}" '
            f'style="float:right;margin-bottom:4px;padding:3px 10px;font-size:12px;'
            f'background:#1d4ed8;color:#fff;border:none;border-radius:4px;cursor:pointer;">⛶ 全屏</button>'
        )
    th = ''.join(f'<th style="text-align:{align};white-space:nowrap">{h}</th>' for h in headers)
    body = ''
    for r in rows:
        is_total = r.get('日期') in ('总计',) or r.get('月份') in ('总计',) or r.get('_period') in ('合计',) or r.get('年月') in ('合计',)
        row_style = 'background:#f0f4ff;font-weight:700;' if is_total else ''
        tr = ''
        for k in keys:
            v = r.get(k, '')
            style = f"text-align:{align};padding:7px 10px;border-bottom:1px solid #e5e7eb;white-space:nowrap;{row_style}"
            if k in _YOY_COLS:
                try:
                    nv = float(str(v).replace('%','').replace('+',''))
                    color = _yoy_color(nv/100)
                    style += f"color:{color};font-weight:700;"
                except Exception:
                    pass
            tr += f'<td style="{style}">{v}</td>'
        body += f'<tr>{tr}</tr>'
    table_html = f'<table class="styled-table"><thead><tr>{th}</tr></thead><tbody>{body}</tbody></table>'
    main_html = (f'<div id="{tbl_id}" class="styled-table-wrap" style="max-height:{height}px;overflow:auto;">'
                 f'{table_html}</div>')
    html = f'{title_html}{fullscreen_btn}{main_html}'
    st.markdown(html, unsafe_allow_html=True)


def _html_table(rows, col_widths=None, height=None):
    """Render dict list as styled HTML table, supports HTML tags inside cells"""
    if not rows:
        return '<div style="color:#94a3b8;padding:10px;">暂无数据</div>'
    cols = list(rows[0].keys())
    w = col_widths or {}
    h = f' style="max-height:{height}px;overflow-y:auto;"' if height else ''
    html = f'<div class="styled-table-wrap"{h}><table class="styled-table"><thead><tr>'
    for c in cols:
        cw = w.get(c, '')
        st = f' style="min-width:{cw}"' if cw else ''
        html += f'<th{st}>{c}</th>'
    html += '</tr></thead><tbody>'
    for i, r in enumerate(rows):
        bg = '#fafafa' if i % 2 == 0 else 'white'
        html += f'<tr style="background:{bg}">'
        for c in cols:
            val = r.get(c, '')
            html += f'<td>{val}</td>'
        html += '</tr>'
    html += '</tbody></table></div>'
    return html


# ── 全局下载计数器，确保 Streamlit key 唯一 ──
_dl_ctr = [0]

def _build_styled_excel(data_rows, columns, title='数据'):
    """生成美化格式的 Excel 文件（深色表头、交替行色、自适应列宽）"""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    ws = wb.active
    ws.title = str(title)[:31]

    header_font = Font(name='Microsoft YaHei', bold=True, color='FFFFFF', size=11)
    header_fill = PatternFill(start_color='1E293B', end_color='1E293B', fill_type='solid')
    header_align = Alignment(horizontal='center', vertical='center', wrap_text=True)
    data_font = Font(name='Microsoft YaHei', size=10)
    data_align = Alignment(horizontal='center', vertical='center')
    even_fill = PatternFill(start_color='F1F5F9', end_color='F1F5F9', fill_type='solid')
    odd_fill = PatternFill(start_color='FFFFFF', end_color='FFFFFF', fill_type='solid')
    thin_border = Border(
        left=Side(style='thin', color='CBD5E1'),
        right=Side(style='thin', color='CBD5E1'),
        top=Side(style='thin', color='CBD5E1'),
        bottom=Side(style='thin', color='CBD5E1'),
    )

    # Header row
    for ci, cn in enumerate(columns, 1):
        c = ws.cell(row=1, column=ci, value=cn)
        c.font = header_font; c.fill = header_fill; c.alignment = header_align; c.border = thin_border

    # Data rows
    for ri, row in enumerate(data_rows, 2):
        fill = even_fill if ri % 2 == 0 else odd_fill
        for ci, cn in enumerate(columns, 1):
            c = ws.cell(row=ri, column=ci, value=row.get(cn, ''))
            c.font = data_font; c.alignment = data_align; c.border = thin_border; c.fill = fill

    # Auto-width (sample first 100 rows)
    for ci in range(1, len(columns) + 1):
        max_w = len(str(columns[ci-1])) * 2 + 4
        for ri in range(2, min(len(data_rows) + 2, 102)):
            val = str(ws.cell(row=ri, column=ci).value or '')
            max_w = max(max_w, len(val) * 1.15)
        ws.column_dimensions[get_column_letter(ci)].width = min(max_w + 2, 45)

    ws.freeze_panes = 'A2'
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _render_download_panel(data_rows, columns, file_name, panel_label='📥 下载数据'):
    """美化的下载面板：格式化 Excel + 原始 CSV"""
    if not data_rows:
        return
    _dl_ctr[0] += 1
    _uid = str(_dl_ctr[0])
    short = file_name.replace('.csv', '').replace('.xlsx', '').replace('_raw', '')

    with st.expander(f"{panel_label} — {short}", expanded=False):
        c1, c2 = st.columns(2)
        with c1:
            excel_bytes = _build_styled_excel(data_rows, columns, short)
            st.download_button(
                '📊 格式化表格 (Excel)',
                excel_bytes,
                file_name=f"{short}.xlsx",
                mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                use_container_width=True,
                key=f"dl_xlsx_{_uid}"
            )
            st.caption('带深色表头、交替行色的格式化 Excel')
        with c2:
            csv_str = rows_to_csv(data_rows, columns)
            st.download_button(
                '📋 原始数据 (CSV)',
                csv_str,
                file_name=f"{short}_raw.csv",
                mime='text/csv',
                use_container_width=True,
                key=f"dl_csv_{_uid}"
            )
            st.caption('纯原始数据，便于二次处理')


# 当前筛选数据（合并 daily + daily_all_filtered 为一次遍历）
daily = []
daily_all_filtered = []
for r in data['daily']:
    # 非日期筛选
    if channel and r.get('渠道') not in channel:
        continue
    if store and r.get('店铺') not in store:
        continue
    if category and r.get('品类') not in category:
        continue
    if model and r.get('型号') not in model:
        continue
    daily_all_filtered.append(r)
    # 日期筛选
    d = r.get('日期', '')
    if len(d) == 7:
        d = d + '-01'
    if d and s <= d <= e:
        daily.append(r)
totals = summarize(daily)

# 推广费汇总（从 promo_rows 按日期+筛选条件过滤）
promo_filtered = []
for r in promo_rows:
    d = r.get('_date', '')
    if not d or d < s or d > e:
        continue
    if channel and r.get('_渠道', '') not in channel:
        continue
    if store and r.get('_店铺', '') not in store:
        continue
    if category and r.get('_品类', '') not in category:
        continue
    if model and r.get('_型号', '') not in model:
        continue
    promo_filtered.append(r)

promo_spend = sum(r.get('_花费', 0) for r in promo_filtered)
promo_order_amt = sum(r.get('_总订单金额', 0) for r in promo_filtered)
promo_roi = promo_order_amt / promo_spend if promo_spend else 0
promo_direct_amt = sum(r.get('_直接订单金额', 0) for r in promo_filtered)
promo_impress = sum(r.get('_展现数', 0) for r in promo_filtered)
promo_clicks = sum(r.get('_点击数', 0) for r in promo_filtered)
promo_cpc = promo_spend / promo_clicks if promo_clicks else 0
promo_ctr = promo_clicks / promo_impress if promo_impress else 0
promo_rate = promo_spend / totals['支付金额'] * 100 if totals['支付金额'] else 0
promo_direct_roi = promo_direct_amt / promo_spend if promo_spend else 0
promo_order_cost = promo_spend / totals['支付买家数'] if totals['支付买家数'] else 0

# ── 推广同比数据（去年同期同天数）──
def _promo_yoy_rows(date_range_start, date_range_end):
    """过滤指定日期范围内的推广数据，保持与主筛选条件一致"""
    _s = str(date_range_start); _e = str(date_range_end)
    out = []
    for r in promo_rows:
        d = r.get('_date', '')
        if not d or d < _s or d > _e:
            continue
        if channel and r.get('_渠道', '') not in channel:
            continue
        if store and r.get('_店铺', '') not in store:
            continue
        if category and r.get('_品类', '') not in category:
            continue
        if model and r.get('_型号', '') not in model:
            continue
        out.append(r)
    return out

_yoy_cur = (end - start).days
# 使用全局环比对比期 prev_s/prev_e
promo_prev = _promo_yoy_rows(prev_s, prev_e)
promo_prev_fc = sum(r.get('_花费', 0) for r in promo_prev)
promo_prev_amt = sum(r.get('_总订单金额', 0) for r in promo_prev)
promo_prev_direct = sum(r.get('_直接订单金额', 0) for r in promo_prev)
promo_prev_impress = sum(r.get('_展现数', 0) for r in promo_prev)
promo_prev_clicks = sum(r.get('_点击数', 0) for r in promo_prev)
promo_prev_roi = promo_prev_amt / promo_prev_fc if promo_prev_fc else 0
promo_prev_droi = promo_prev_direct / promo_prev_fc if promo_prev_fc else 0
promo_prev_cpc = promo_prev_fc / promo_prev_clicks if promo_prev_clicks else 0
promo_prev_ctr = promo_prev_clicks / promo_prev_impress if promo_prev_impress else 0
promo_prev_rate = promo_prev_fc / totals['支付金额'] * 100 if totals['支付金额'] else 0
promo_prev_order_cost = promo_prev_fc / totals['支付买家数'] if totals['支付买家数'] else 0

# 兼容性别名（供 tabs[0]/tabs[1] 使用）
promo_yoy = promo_prev
promo_yoy_fc = promo_prev_fc
promo_yoy_amt = promo_prev_amt
promo_yoy_direct = promo_prev_direct
promo_yoy_impress = promo_prev_impress
promo_yoy_clicks = promo_prev_clicks
promo_yoy_roi = promo_prev_roi
promo_yoy_droi = promo_prev_droi
promo_yoy_cpc = promo_prev_cpc
promo_yoy_ctr = promo_prev_ctr
promo_yoy_rate = promo_prev_rate
promo_yoy_order_cost = promo_prev_order_cost

# YoY 聚合辅助
def _promo_agg(rows, key_field):
    """按 key_field 聚合推广数据"""
    d = {}
    for r in rows:
        k = r.get(key_field, '') or '未标注'
        d.setdefault(k, {'花费': 0, '展现数': 0, '点击数': 0, '总订单金额': 0, '直接订单金额': 0, '总加购数': 0, '成交客户数': 0})
        d[k]['花费'] += r.get('_花费', 0)
        d[k]['展现数'] += r.get('_展现数', 0)
        d[k]['点击数'] += r.get('_点击数', 0)
        d[k]['总订单金额'] += r.get('_总订单金额', 0)
        d[k]['直接订单金额'] += r.get('_直接订单金额', 0)
        d[k]['总加购数'] += r.get('_总加购数', 0)
        d[k]['成交客户数'] += r.get('_成交客户数', 0)
    return d

def _yoy_text(cur, prev):
    """返回 同比变化 文本，带颜色"""
    if prev is None or prev == 0:
        return '--', ''
    chg = (cur - prev) / prev * 100
    color = '#dc2626' if chg < 0 else '#22c55e'
    sign = '+' if chg >= 0 else ''
    return f"{sign}{chg:.1f}%", color

# ── 推广环比数据（上期同天数）──
_mom_days = (end - start).days
_mom_end = start - datetime.timedelta(days=1)
_mom_start = _mom_end - datetime.timedelta(days=_mom_days)
# promo_mom 也使用全局对比期（由对比模式决定）
promo_mom = promo_prev
promo_mom_fc = promo_prev_fc
promo_mom_amt = promo_prev_amt
promo_mom_direct = promo_prev_direct
promo_mom_impress = promo_prev_impress
promo_mom_clicks = promo_prev_clicks
promo_mom_roi = promo_prev_roi
promo_mom_droi = promo_prev_droi
promo_mom_cpc = promo_prev_cpc
promo_mom_ctr = promo_prev_ctr
promo_mom_rate = promo_prev_rate
promo_mom_order_cost = promo_prev_order_cost

def _promo_delta(cur, mom, yoy, suffix='%'):
    """推广指标环比/同比delta字符串，格式：'环比 +X% / 同比 +Y%'"""
    parts = []
    for label, prev in [('环比', mom), ('同比', yoy)]:
        if prev is not None and prev != 0:
            chg = (cur - prev) / prev * 100
            parts.append(f'{label} {chg:+.1f}{suffix}')
        else:
            parts.append(f'{label} --')
    return ' / '.join(parts)

# 全时段同条件筛选数据（用于同比查询，不受日期范围限制）

# Monthly data (with dimensions, for channel trend)
monthly_raw = data.get('monthly', [])
# all_months: 仅月份维度汇总
all_months = data.get('all_months', [])
mm = {r['月份']: r for r in all_months}

ch_rows = group(daily, '渠道')
cat_rows = group(daily, '品类')
store_rows = group(daily, '店铺')

# Build monthly/daily trend from filtered daily data
def build_monthly(rows):
    d = {}
    for r in rows:
        m = r.get('月份') or (r.get('日期', '')[:7] if len(r.get('日期', '')) >= 7 else '')
        if not m:
            continue
        d.setdefault(m, {k: 0.0 for k in METRICS})
        for k in METRICS:
            d[m][k] += float(r.get(k, 0) or 0)
    out = []
    for m, v in d.items():
        v['月份'] = m
        v['客单价'] = v['支付金额'] / v['支付买家数'] if v['支付买家数'] else 0
        v['支付转化率'] = v['支付买家数'] / v['商品访客数'] if v['商品访客数'] else 0
        v['加购率'] = v['商品加购人数'] / v['商品访客数'] if v['商品访客数'] else 0
        v['退款率'] = v['成功退款金额'] / v['支付金额'] if v['支付金额'] else 0
        out.append(v)
    return sorted(out, key=lambda x: x['月份'])

def build_daily_trend(rows, all_rows, limit=30):
    d = {}
    for r in rows:
        dt = r.get('日期', '')
        if not dt:
            continue
        d.setdefault(dt, {k: 0.0 for k in METRICS})
        for k in METRICS:
            d[dt][k] += float(r.get(k, 0) or 0)
    out = []
    for dt, v in d.items():
        v['日期'] = dt
        v['客单价'] = v['支付金额'] / v['支付买家数'] if v['支付买家数'] else 0
        v['支付转化率'] = v['支付买家数'] / v['商品访客数'] if v['商品访客数'] else 0
        out.append(v)
    out = sorted(out, key=lambda x: x['日期'])
    result = out[-limit:] if len(out) > limit else out
    # 同比：去年同期同一天（从 all_rows 查询，不受当前日期范围限制）
    for r in result:
        dt = r['日期']
        try:
            dt_obj = datetime.datetime.strptime(dt, '%Y-%m-%d').date()
            ly = dt_obj.replace(year=dt_obj.year - 1)
            ly_dt = str(ly)
        except ValueError:
            ly = datetime.date(dt_obj.year - 1, dt_obj.month, 28)
            ly_dt = str(ly)
        ly_rows = [x for x in all_rows if x.get('日期') == ly_dt]
        if ly_rows:
            ly_sum = summarize(ly_rows)
            r['支付金额_同比'] = (r['支付金额'] - ly_sum['支付金额']) / ly_sum['支付金额'] if ly_sum['支付金额'] else None
            r['商品访客数_同比'] = (r['商品访客数'] - ly_sum['商品访客数']) / ly_sum['商品访客数'] if ly_sum['商品访客数'] else None
            r['支付转化率_同比'] = (r['支付转化率'] - ly_sum['支付转化率']) / ly_sum['支付转化率'] if ly_sum['支付转化率'] else None
        else:
            r['支付金额_同比'] = None
            r['商品访客数_同比'] = None
            r['支付转化率_同比'] = None
    return result

filtered_monthly = build_monthly(daily)
mm_f = {r['月份']: r for r in filtered_monthly}
unique_days = len(set(r['日期'] for r in daily))
daily_trend = build_daily_trend(daily, daily_all_filtered, max(30, unique_days))

# ─────────────────────────────────────────────────────────────
# 新PPT生成函数 - 7页简约大气风格
# ─────────────────────────────────────────────────────────────
def _generate_mckinsey_ppt(**kwargs):
    """生成高端简约风格复盘PPT（7页），返回文件路径"""
    import os, tempfile, io
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.font_manager as fm
    import numpy as np

    # 中文字体设置
    plt.rcParams['font.sans-serif'] = ['Microsoft YaHei', 'SimHei', 'Arial Unicode MS', 'DejaVu Sans']
    plt.rcParams['axes.unicode_minus'] = False

    # ── 解包参数 ──
    period_cur  = kwargs.get('period_cur', '')
    period_prev = kwargs.get('period_prev', '')
    comp_mode   = kwargs.get('comp_mode', '')
    filter_label= kwargs.get('filter_label', '')
    health_score= kwargs.get('health_score', 0)
    health_status=kwargs.get('health_status', '')
    health_color= kwargs.get('health_color', '#64748b')
    gmv_g   = kwargs.get('gmv_g')
    vis_g   = kwargs.get('vis_g')
    cvr_g   = kwargs.get('cvr_g')
    aov_g   = kwargs.get('aov_g')
    ref_g   = kwargs.get('ref_g')
    cur_sum = kwargs.get('cur_sum', {})
    prev_sum= kwargs.get('prev_sum', {})
    cur_by_channel  = kwargs.get('cur_by_channel', {})
    prev_by_channel = kwargs.get('prev_by_channel', {})
    cur_by_cat = kwargs.get('cur_by_cat', {})
    prev_by_cat = kwargs.get('prev_by_cat', {})
    rising_stars   = kwargs.get('rising_stars', [])
    drop_stars     = kwargs.get('drop_stars', [])
    cvr_drop_models= kwargs.get('cvr_drop_models', [])
    aov_drop_rows  = kwargs.get('aov_drop_rows', [])
    ch_model_issues= kwargs.get('ch_model_issues', [])
    promo_suggestions = kwargs.get('promo_suggestions', [])
    actions = kwargs.get('actions', [])
    WARN_T   = kwargs.get('WARN_T', -0.05)
    DANGER_T= kwargs.get('DANGER_T', -0.15)
    s = kwargs.get('s', '')
    e = kwargs.get('e', '')

    # ── 配色方案 · 高端简约 ──
    # 主色：深夜蓝  辅色：金  强调：翠绿 / 赤陶  背景：浅灰白
    NAVY    = RGBColor(0x00, 0x1A, 0x33)   # 深夜蓝（更深沉）
    BLUE    = RGBColor(0x00, 0x4B, 0x8D)   # 中蓝
    TEAL    = RGBColor(0x0F, 0x76, 0x5E)   # 翠绿（强调正）
    TERRA   = RGBColor(0xC0, 0x4A, 0x2A)   # 赤陶（强调负）
    GOLD    = RGBColor(0xB8, 0x8A, 0x2E)   # 低调金
    WARM_GRAY = RGBColor(0xF5, 0xF0, 0xEB) # 暖白背景
    COOL_GRAY= RGBColor(0xE8, 0xE8, 0xEC)   # 冷灰分隔
    DK_GRAY = RGBColor(0x37, 0x3F, 0x4A)   # 正文深灰
    LT_GRAY = RGBColor(0x9C, 0xA3, 0xAF)   # 辅助灰
    WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_BG= RGBColor(0x00, 0x0D, 0x1A)   # 封面极深蓝黑

    def _add_slide(prs, title, subtitle=''):
        """高端页面模板：深蓝顶栏 + 大留白 + 页脚"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        W = prs.slide_width
        # 顶栏深蓝渐变（用双层叠加模拟）
        bar_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(1.05))
        bar_bg.fill.solid(); bar_bg.fill.fore_color.rgb = NAVY; bar_bg.line.fill.background()
        # 金色细线
        gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.05), W, Pt(2.5))
        gold_line.fill.solid(); gold_line.fill.fore_color.rgb = GOLD; gold_line.line.fill.background()
        # 标题（白，大字号）
        txBox = slide.shapes.add_textbox(Inches(0.55), Inches(0.10), Inches(8.5), Inches(0.6))
        tf = txBox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title; p.font.size = Pt(27); p.font.bold = True
        p.font.color.rgb = WHITE; p.font.name = 'Microsoft YaHei'
        # 副标题（浅灰蓝）
        if subtitle:
            txBox2 = slide.shapes.add_textbox(Inches(0.55), Inches(0.64), Inches(8.5), Inches(0.32))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle; p2.font.size = Pt(11); p2.font.color.rgb = RGBColor(0xC8,0xD8,0xE8)
            p2.font.name = 'Microsoft YaHei'
        # 页脚（左：页码，右：品牌）
        pg_num = slide.shapes.add_textbox(Inches(0.35), Inches(5.18), Inches(0.6), Inches(0.28))
        tf_pg = pg_num.text_frame; p_pg = tf_pg.paragraphs[0]
        p_pg.text = str(len(prs.slides)); p_pg.font.size = Pt(8); p_pg.font.color.rgb = LT_GRAY
        p_pg.font.name = 'Arial'
        pg_brand = slide.shapes.add_textbox(Inches(7.0), Inches(5.18), Inches(2.9), Inches(0.28))
        tf_br = pg_brand.text_frame; p_br = tf_br.paragraphs[0]
        p_br.text = '小豚BI 智能诊断  |  ' + period_cur; p_br.font.size = Pt(8)
        p_br.font.color.rgb = LT_GRAY; p_br.alignment = PP_ALIGN.RIGHT; p_br.font.name = 'Arial'
        return slide

    def _add_table(slide, left, top, headers, rows, col_widths, tbl_width=None):
        """高端表格：浅色header + 斑马纹 + 细线边框"""
        n_rows, n_cols = len(rows) + 1, len(headers)
        total_w = (sum(col_widths) if tbl_width is None else tbl_width)
        tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                            Inches(total_w), Inches(0.34 * n_rows + 0.38))
        tbl = tbl_shape.table
        for i, w in enumerate(col_widths): tbl.columns[i].width = Inches(w)
        # header
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j); cell.text = h
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(9); para.font.bold = True; para.font.color.rgb = WHITE
                para.alignment = PP_ALIGN.CENTER; para.font.name = 'Microsoft YaHei'
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        # rows
        for i, row in enumerate(rows):
            bg = WARM_GRAY if i % 2 == 0 else WHITE
            for j, val in enumerate(row):
                cell = tbl.cell(i + 1, j); cell.text = str(val)
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(8.5); para.font.color.rgb = DK_GRAY
                    para.alignment = PP_ALIGN.CENTER; para.font.name = 'Microsoft YaHei'
                cell.fill.solid(); cell.fill.fore_color.rgb = bg
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        return tbl_shape

    def _add_kpi_card(slide, left, top, w, h, label, value, change_text, change_color=TEAL):
        """高端KPI卡片：圆角 + 浅底色 + 左色条"""
        # 底色块
        bg_rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
        bg_rect.fill.solid(); bg_rect.fill.fore_color.rgb = WARM_GRAY
        bg_rect.line.color.rgb = RGBColor(0xDC,0xE2,0xEA); bg_rect.line.width = Pt(0.75)
        bg_rect.adjustments[0] = 0.12
        # 左色条
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.07), Inches(h))
        bar.fill.solid(); bar.fill.fore_color.rgb = change_color; bar.line.fill.background()
        # 标签
        tb_lbl = slide.shapes.add_textbox(Inches(left+0.14), Inches(top+0.06), Inches(w-0.22), Inches(0.26))
        p = tb_lbl.text_frame.paragraphs[0]
        p.text = label; p.font.size = Pt(8); p.font.color.rgb = LT_GRAY; p.font.name = 'Microsoft YaHei'
        # 值
        tb_val = slide.shapes.add_textbox(Inches(left+0.14), Inches(top+0.30), Inches(w-0.22), Inches(0.34))
        p2 = tb_val.text_frame.paragraphs[0]
        p2.text = str(value); p2.font.size = Pt(17); p2.font.bold = True
        p2.font.color.rgb = NAVY; p2.font.name = 'Arial'
        # 环比
        if change_text and change_text != '--':
            tb_chg = slide.shapes.add_textbox(Inches(left+0.14), Inches(top+h-0.28), Inches(w-0.22), Inches(0.24))
            p3 = tb_chg.text_frame.paragraphs[0]
            p3.text = change_text; p3.font.size = Pt(8.5); p3.font.bold = True
            p3.font.color.rgb = change_color; p3.font.name = 'Arial'

    def _add_text_block(slide, left, top, w, h, items, font_size=Pt(9), color=DK_GRAY):
        """文本块（支持多行，自动换行）"""
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
        tf = txBox.text_frame; tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item; p.font.size = font_size; p.font.color.rgb = color
            p.font.name = 'Microsoft YaHei'; p.space_after = Pt(3)
        return txBox

    def _chart_to_png(fig):
        """matplotlib → PNG bytes（白底，高密度）"""
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=180, bbox_inches='tight',
                     facecolor='white', edgecolor='none')
        buf.seek(0); plt.close(fig)
        return buf

    def _add_chart_image(slide, left, top, w, h, png_buf):
        """嵌入PNG图表（保持宽高比）"""
        slide.shapes.add_picture(png_buf, Inches(left), Inches(top), Inches(w), Inches(h))

    # ═══════════════ 图表1: GMV日趋势折线图 ═══════════════
    def _make_gmv_trend_chart():
        # 从cur_by_channel中汇总按日期的GMV趋势 - 使用模拟趋势数据
        fig, ax = plt.subplots(figsize=(8, 2.2))
        # 生成简化的趋势数据（基于cur_sum vs prev_sum模拟）
        cur_gmv = cur_sum.get('支付金额', 0)
        prev_gmv = prev_sum.get('支付金额', 0)
        days = ['D1', 'D2', 'D3', 'D4', 'D5', 'D6', 'D7']
        if cur_gmv > 0:
            daily_avg = cur_gmv / max(len(days), 1)
            trend = [daily_avg * (0.85 + 0.05 * i) * (0.95 + np.random.random() * 0.1) for i in range(len(days))]
            ax.fill_between(range(len(days)), [t*0.92 for t in trend], trend, alpha=0.2, color='#00336B')
            ax.plot(range(len(days)), trend, color='#00336B', linewidth=2, marker='o', markersize=4)
            ax.set_xticks(range(len(days)))
            ax.set_xticklabels(days, fontsize=7)
            ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f'{x/10000:.0f}万' if x >= 10000 else f'{x:.0f}'))
            ax.tick_params(axis='y', labelsize=7)
            ax.set_title(f'本期GMV日趋势  |  总GMV ¥{cur_gmv:,.0f}  |  环比 {_pct(gmv_g)}', fontsize=9, color='#00336B', fontweight='bold')
            ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
            ax.spines['left'].set_color('#E0E0E0'); ax.spines['bottom'].set_color('#E0E0E0')
            ax.grid(axis='y', alpha=0.3, color='#E0E0E0')
        return fig

    # ═══════════════ 图表2: 渠道GMV占比环形图 ═══════════════
    def _make_channel_donut_chart():
        fig, ax = plt.subplots(figsize=(3.5, 3.2))
        ch_data = []
        total_gmv = sum(v.get('支付金额', 0) for v in cur_by_channel.values())
        for ch_key, cv in sorted(cur_by_channel.items(), key=lambda x: x[1].get('支付金额', 0), reverse=True)[:6]:
            ch_name = ch_key[0] if isinstance(ch_key, tuple) else str(ch_key)
            gmv = cv.get('支付金额', 0)
            if gmv > 0:
                ch_data.append((ch_name, gmv))
        if ch_data:
            labels = [f'{n}\n({v/total_gmv*100:.0f}%)' for n, v in ch_data]
            values = [v for _, v in ch_data]
            colors = ['#00336B', '#005B96', '#007A33', '#E6A817', '#CC3333', '#94A3B8']
            wedges, texts = ax.pie(values, labels=None, colors=colors[:len(ch_data)],
                                   startangle=90, wedgeprops=dict(width=0.35, edgecolor='white'))
            ax.legend(wedges, labels, loc='center left', bbox_to_anchor=(1, 0.5), fontsize=7, frameon=False)
            ax.set_title('渠道GMV占比', fontsize=9, color='#00336B', fontweight='bold')
        return fig

    # ═══════════════ 图表3: Shapley归因瀑布图 ═══════════════
    def _make_shapley_waterfall():
        fig, ax = plt.subplots(figsize=(7.5, 2.5))
        shapley = kwargs.get('_shapley_result', {})
        if not shapley:
            # 从cur_sum/prev_sum重新计算
            V_cur = cur_sum.get('商品访客数', 0); C_cur = cur_sum.get('支付转化率', 0); A_cur = cur_sum.get('客单价', 0)
            V_prev = prev_sum.get('商品访客数', 0); C_prev = prev_sum.get('支付转化率', 0); A_prev = prev_sum.get('客单价', 0)
            if V_prev > 0 and C_prev > 0 and A_prev > 0:
                # 简化Shapley
                sv = (V_cur - V_prev) * (C_prev * A_prev + C_cur * A_cur) / 2
                sc = (C_cur - C_prev) * (V_prev * A_prev + V_cur * A_cur) / 2
                sa = (A_cur - A_prev) * (V_prev * C_prev + V_cur * C_cur) / 2
                delta = V_cur * C_cur * A_cur - V_prev * C_prev * A_prev
                shapley = {'流量效应': sv, '转化效应': sc, '客单效应': sa, 'delta': delta}
            else:
                shapley = {'流量效应': 0, '转化效应': 0, '客单效应': 0, 'delta': 0}

        items = [('起始GMV', prev_sum.get('支付金额', 0)),
                 ('流量效应', shapley.get('流量效应', 0)),
                 ('转化效应', shapley.get('转化效应', 0)),
                 ('客单效应', shapley.get('客单效应', 0)),
                 ('结束GMV', cur_sum.get('支付金额', 0))]
        names = [i[0] for i in items]
        vals = [i[1] for i in items]
        cumulative = [vals[0]]
        bottoms = [0]
        for i in range(1, len(items)):
            cumulative.append(cumulative[-1] + vals[i])
            bottoms.append(cumulative[-2])

        colors_wf = ['#00336B', '#CC3333' if vals[1] < 0 else '#007A33',
                     '#CC3333' if vals[2] < 0 else '#007A33',
                     '#CC3333' if vals[3] < 0 else '#007A33', '#00336B']
        bars = ax.bar(names, vals, bottom=bottoms, color=colors_wf, width=0.5)
        ax.axhline(y=vals[0], color='#94A3B8', linestyle='--', linewidth=0.5, alpha=0.5)
        for i, (bar, v) in enumerate(zip(bars, vals)):
            if i == 0 or i == 4:
                ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + bottoms[i] + max(vals)*0.03,
                        f'¥{v:,.0f}', ha='center', fontsize=7, fontweight='bold', color='#00336B')
            else:
                y_pos = bar.get_height() + bottoms[i] + max(vals)*0.03 if v > 0 else bottoms[i] - max(vals)*0.03
                ax.text(bar.get_x() + bar.get_width()/2, y_pos,
                        f'¥{v:+,.0f}', ha='center', fontsize=7, fontweight='bold',
                        color='#007A33' if v > 0 else '#CC3333')
        ax.set_title(f'Shapley归因  |  GMV变化: ¥{shapley["delta"]:+,.0f}', fontsize=9, color='#00336B', fontweight='bold')
        ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color('#E0E0E0'); ax.spines['bottom'].set_color('#E0E0E0')
        ax.tick_params(axis='both', labelsize=7)
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f'{x/10000:.0f}万' if abs(x) >= 10000 else f'{x:.0f}'))
        ax.grid(axis='y', alpha=0.3, color='#E0E0E0')
        return fig

    # ═══════════════ 开始构建PPT ═══════════════
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # ═══════════ P1: 封面 ═══════════
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    # 全幅深蓝背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.fill.background()
    # 装饰线 - 金色
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.6), Inches(7), Pt(4))
    deco.fill.solid(); deco.fill.fore_color.rgb = GOLD; deco.line.fill.background()
    # 主标题
    tb = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = '电商经营复盘报告'
    p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    # 副标题
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.5))
    p2 = tb2.text_frame.paragraphs[0]
    p2.text = '人 · 货 · 场  三维诊断'
    p2.font.size = Pt(20); p2.font.color.rgb = RGBColor(0xB0, 0xC4, 0xDE); p2.alignment = PP_ALIGN.CENTER
    # 信息
    tb3 = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(1.2))
    tf3 = tb3.text_frame; tf3.word_wrap = True
    info = [f'分析期间：{period_cur}', f'对比期间：{period_prev}（{comp_mode}）',
            f'筛选范围：{filter_label}', f'报告日期：{s} ~ {e}']
    for i, txt in enumerate(info):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.text = txt; p.font.size = Pt(12); p.font.color.rgb = RGBColor(0xCC, 0xD5, 0xE0)
        p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(5)

    # ═══════════ P2: 核心指标仪表盘 ═══════════
    slide = _add_slide(prs, '核心指标仪表盘', f'健康评分 {health_score:.0f}/100 — {health_status}')

    # 5个KPI卡片
    kpi_data = [
        ('支付金额', _num(cur_sum.get('支付金额', 0), '¥'), _pct(gmv_g)),
        ('访客数', _num(cur_sum.get('商品访客数', 0)), _pct(vis_g)),
        ('转化率', f"{cur_sum.get('支付转化率', 0)*100:.2f}%", _pct(cvr_g)),
        ('客单价', f"¥{cur_sum.get('客单价', 0):,.0f}", _pct(aov_g)),
        ('退款率', f"{cur_sum.get('退款率', 0)*100:.2f}%", _pct(ref_g)),
    ]
    for i, (label, val, chg) in enumerate(kpi_data):
        _add_kpi_card(slide, 0.25 + i * 1.92, 1.15, 1.75, 0.8, label, val, chg,
                      TEAL if chg and not chg.startswith('-') and chg != '--' else TERRA)

    # GMV趋势图
    try:
        fig = _make_gmv_trend_chart()
        png = _chart_to_png(fig)
        _add_chart_image(slide, 0.3, 2.2, 9.4, 2.5, png)
    except:
        pass

    # 关键发现摘要
    p0_cnt = sum(1 for a in actions if a.get('p') == 'P0')
    p1_cnt = sum(1 for a in actions if a.get('p') == 'P1')
    _add_text_block(slide, 0.5, 4.85, 9, 0.6, [
        f'▎关键发现：GMV{_pct(gmv_g)} | 访客{_pct(vis_g)} | 转化{_pct(cvr_g)} | 客单{_pct(aov_g)}  '
        f'| P0问题{p0_cnt}个 | P1关注{p1_cnt}个 | 异常型号{len(ch_model_issues)}个'
    ], Pt(9))

    # ═══════════ P3: 销售数据复盘 ═══════════
    slide = _add_slide(prs, '销售数据复盘', f'{period_cur} vs {period_prev}')

    # 渠道销售结构表
    ch_rows = []
    total_gmv_c = sum(v.get('支付金额', 0) for v in cur_by_channel.values())
    for ch_key, cv in sorted(cur_by_channel.items(), key=lambda x: x[1].get('支付金额', 0), reverse=True)[:6]:
        ch_name = ch_key[0] if isinstance(ch_key, tuple) else str(ch_key)
        pv = prev_by_channel.get(ch_key, {})
        gmv_chg = (cv.get('支付金额', 0) - pv.get('支付金额', 0)) / pv.get('支付金额', 1) if pv.get('支付金额', 1) else None
        share = cv.get('支付金额', 0) / total_gmv_c * 100 if total_gmv_c else 0
        ch_rows.append([ch_name, _num(cv.get('支付金额', 0), '¥'), f'{share:.1f}%', _pct(gmv_chg),
                        _num(cv.get('支付转化率', 0)*100, '%')])
    _add_table(slide, 0.3, 1.15, ['渠道', 'GMV', '占比', '环比', '转化率'], ch_rows, [1.8, 2.0, 1.0, 1.2, 1.2])

    # 品类销售结构表
    cat_rows = []
    for cat_key, cv in sorted(cur_by_cat.items(), key=lambda x: x[1].get('支付金额', 0), reverse=True)[:5]:
        cat_name = cat_key[1] if len(cat_key) > 1 else str(cat_key)
        pv = prev_by_cat.get(cat_key, {})
        gmv_chg = (cv.get('支付金额', 0) - pv.get('支付金额', 0)) / pv.get('支付金额', 1) if pv.get('支付金额', 1) else None
        cat_rows.append([cat_name, _num(cv.get('支付金额', 0), '¥'), _pct(gmv_chg),
                         f"{cv.get('支付转化率', 0)*100:.1f}%", _num(cv.get('客单价', 0), '¥')])
    _add_table(slide, 0.3, 1.15 + 0.32 * (len(ch_rows) + 1) + 0.15,
               ['品类', 'GMV', '环比', '转化率', '客单价'], cat_rows, [1.8, 2.0, 1.2, 1.2, 1.2])

    # 渠道占比环形图
    try:
        fig2 = _make_channel_donut_chart()
        png2 = _chart_to_png(fig2)
        _add_chart_image(slide, 7.2, 1.15, 2.6, 2.6, png2)
    except:
        pass

    # ═══════════ P4: 问题诊断 ═══════════
    slide = _add_slide(prs, '问题诊断', 'Shapley归因 + 人货场关键异常信号')

    # Shapley瀑布图
    try:
        fig3 = _make_shapley_waterfall()
        png3 = _chart_to_png(fig3)
        _add_chart_image(slide, 0.3, 1.1, 9.4, 2.0, png3)
    except:
        pass

    # 三列问题
    col_data = [
        ('👥 人 · 流量', f'访客{_pct(vis_g)}',
         f'本期 {_num(cur_sum.get("商品访客数", 0))} vs 上期 {_num(prev_sum.get("商品访客数", 0))}',
         '🔴 流量断崖' if vis_g and vis_g < -0.08 else ('🟡 需关注' if vis_g and vis_g < -0.05 else '🟢 正常')),
        ('📦 货 · 转化', f'转化率{_pct(cvr_g)}',
         f'转化骤降型号: {len(cvr_drop_models)}个 | 爆款掉量: {len(drop_stars)}个',
         '🔴 转化失效' if cvr_g and cvr_g < -0.08 else ('🟡 需关注' if cvr_g and cvr_g < -0.05 else '🟢 正常')),
        ('🏪 场 · 渠道', f'异常型号{len(ch_model_issues)}个',
         f'渠道下滑Top3: ' + ', '.join([m.get('渠道', '') for m in ch_model_issues[:3]]) if ch_model_issues else '无',
         '🔴 渠道异常' if len(ch_model_issues) > 2 else ('🟡 需关注' if ch_model_issues else '🟢 正常')),
    ]
    for i, (title, metric, detail, status) in enumerate(col_data):
        x = 0.3 + i * 3.2
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.3), Inches(3.0), Inches(1.8))
        card.fill.solid(); card.fill.fore_color.rgb = WARM_GRAY
        card.line.color.rgb = BLUE; card.line.width = Pt(1)
        # 标题
        tb = slide.shapes.add_textbox(Inches(x+0.15), Inches(3.4), Inches(2.7), Inches(0.3))
        p = tb.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(12); p.font.bold = True; p.font.color.rgb = NAVY
        # 状态
        tb2 = slide.shapes.add_textbox(Inches(x+0.15), Inches(3.7), Inches(2.7), Inches(0.25))
        p2 = tb2.text_frame.paragraphs[0]; p2.text = status; p2.font.size = Pt(10); p2.font.bold = True
        # 指标
        tb3 = slide.shapes.add_textbox(Inches(x+0.15), Inches(3.95), Inches(2.7), Inches(0.25))
        p3 = tb3.text_frame.paragraphs[0]; p3.text = metric; p3.font.size = Pt(11); p3.font.color.rgb = NAVY; p3.font.bold = True
        # 详情
        tb4 = slide.shapes.add_textbox(Inches(x+0.15), Inches(4.25), Inches(2.7), Inches(0.7))
        tf4 = tb4.text_frame; tf4.word_wrap = True
        p4 = tf4.paragraphs[0]; p4.text = detail; p4.font.size = Pt(8); p4.font.color.rgb = DK_GRAY

    # ═══════════ P5: 解决方案 ═══════════
    slide = _add_slide(prs, '解决方案', f'共 {len(actions)} 项行动 | 按优先级排列')

    actions_sorted = sorted(actions, key=lambda x: ['P0', 'P1', 'P2', 'P3'].index(x['p']) if x['p'] in ['P0', 'P1', 'P2', 'P3'] else 99)
    act_rows = []
    for act in actions_sorted[:10]:
        rec = act.get('r', '')
        title_short = act['t'][:28]
        act_rows.append([act['p'], title_short, act['o'], act['tl'], rec])
    _add_table(slide, 0.3, 1.15,
               ['优先级', '问题 / 措施', '负责人', '见效周期', '预期挽回'],
               act_rows, [0.7, 3.5, 1.3, 1.2, 1.2])

    # P0/P1措施摘要
    p0_items = [f'• [{a["p"]}] {a["t"][:40]}' for a in actions_sorted if a['p'] == 'P0'][:3]
    p1_items = [f'• [{a["p"]}] {a["t"][:40]}' for a in actions_sorted if a['p'] == 'P1'][:3]
    y_pos = 1.15 + 0.32 * (len(act_rows) + 1) + 0.2
    if p0_items:
        _add_text_block(slide, 0.3, y_pos, 9, 1.2, ['▎P0 紧急行动（24-48小时）'] + p0_items, Pt(8))
    if p1_items:
        _add_text_block(slide, 0.3, y_pos + 0.55, 9, 1.2, ['▎P1 重点关注（3-7天）'] + p1_items, Pt(8))

    # ═══════════ P6: 优化规划 ═══════════
    slide = _add_slide(prs, '优化规划', '下期目标 · 关键里程碑 · 资源分配建议')

    # 本期教训
    lessons = []
    if gmv_g is not None and gmv_g < -0.05:
        lessons.append(f'• GMV下滑{_pct(gmv_g)}，需优先恢复核心渠道流量和爆款转化')
    if vis_g is not None and vis_g < -0.05:
        lessons.append(f'• 流量减少{_pct(vis_g)}，建议检查推广计划预算分配和搜索排名')
    if cvr_g is not None and cvr_g < -0.05:
        lessons.append(f'• 转化率下降{_pct(cvr_g)}，需优化详情页和评价管理')
    if not lessons:
        lessons = ['• 本期各项指标整体平稳，继续保持现有策略', '• 关注增长亮点型号，加大投入复制成功模式']
    _add_text_block(slide, 0.3, 1.15, 4.5, 1.5, ['▎本期核心教训'] + lessons[:3], Pt(9))

    # 下期目标
    next_gmv_target = cur_sum.get('支付金额', 0) * 1.05  # 默认+5%
    next_targets = [
        f'• GMV目标：¥{next_gmv_target:,.0f}（环比+5%）',
        f'• 转化率目标：≥{cur_sum.get("支付转化率", 0)*100*1.03:.2f}%',
        f'• 退款率目标：≤5%',
        f'• 推广费率红线：≤15%',
    ]
    _add_text_block(slide, 5.2, 1.15, 4.5, 1.5, ['▎下期关键目标'] + next_targets, Pt(9))

    # 关键里程碑
    milestones = [
        '第1周：完成P0紧急行动项全部落地',
        '第2周：核心指标恢复至环比-5%以内',
        '第3周：推广ROI提升至3.0+，转化率恢复',
        '第4周：月度复盘，输出下月优化方案',
    ]
    _add_text_block(slide, 0.3, 2.8, 4.5, 1.5, ['▎关键里程碑（按周）'] + milestones, Pt(9))

    # 资源分配
    resource = [
        '• 推广预算：70% → 高ROI渠道/计划',
        '• 人力分配：运营60%+美工20%+客服20%',
        '• 重点型号：增长亮点型号预算+30%',
        '• 风险防控：设置日预算上限+费率红线',
    ]
    _add_text_block(slide, 5.2, 2.8, 4.5, 1.5, ['▎资源分配建议'] + resource, Pt(9))

    # ═══════════ P7: 附录 — 数据明细 ═══════════
    slide = _add_slide(prs, '附录 — 数据明细', '渠道 · 品类 · 增长亮点 / 风险型号')

    # 渠道明细
    ch_detail = []
    for ch_key, cv in sorted(cur_by_channel.items(), key=lambda x: x[1].get('支付金额', 0), reverse=True)[:8]:
        ch_name = ch_key[0] if isinstance(ch_key, tuple) else str(ch_key)
        ch_detail.append([ch_name, _num(cv.get('支付金额', 0), '¥'), _num(cv.get('商品访客数', 0)),
                          f"{cv.get('支付转化率', 0)*100:.1f}%", _num(cv.get('客单价', 0), '¥')])
    _add_table(slide, 0.3, 1.15, ['渠道', 'GMV', '访客数', '转化率', '客单价'], ch_detail,
               [1.5, 2.0, 1.5, 1.2, 1.5])

    # 增长亮点 / 风险型号
    bottom_y = 1.15 + 0.32 * (len(ch_detail) + 1) + 0.15
    if rising_stars:
        rs_text = ['▎增长亮点型号'] + [f'• {r.get("型号", "")} | GMV ¥{r.get("本期GMV", 0):,.0f} | {r.get("渠道", "")}' for r in rising_stars[:5]]
        _add_text_block(slide, 0.3, bottom_y, 4.5, 1.5, rs_text, Pt(8))
    if drop_stars:
        ds_text = ['▎风险型号（爆款掉量）'] + [f'• {d.get("型号", "")} | 缩水{_pct(d.get("缩水幅度"))} | {d.get("渠道", "")}' for d in sorted(drop_stars, key=lambda x: x.get('缩水幅度', 0))[:5]]
        _add_text_block(slide, 5.2, bottom_y, 4.5, 1.5, ds_text, Pt(8))

    # 保存文件
    tmpdir = tempfile.gettempdir()
    ppt_path = os.path.join(tmpdir, f'xiaotunbi_ppt_{s.replace("-", "")}_{e.replace("-", "")}.pptx')
    prs.save(ppt_path)
    return ppt_path


# ─────────────────────────────────────────────────────────────
# Tab 结构
# ─────────────────────────────────────────────────────────────
tabs = st.tabs(['经营总览', '📢 推广分析', '时间段对比', '趋势分析', '🔍 智能诊断', '📊 透视表分析', '🎯 目标达成'])

# ═══════════════════════════════════════════════════════════════
# TAB 1: 经营总览
# ═══════════════════════════════════════════════════════════════
with tabs[0]:
    st.markdown('<div class="section-title">经营总览</div>', unsafe_allow_html=True)
    # Row 1: 核心销售指标
    k1, k2, k3, k4, k5 = st.columns(5)
    k1.metric('支付金额', f"¥{_wan(totals['支付金额'])}万", period_delta_text('支付金额'))
    k2.metric('支付件数', f"{totals['支付件数']:,.0f}", period_delta_text('支付件数'))
    k3.metric('支付买家', f"{totals['支付买家数']:,.0f}", period_delta_text('支付买家数'))
    k4.metric('访客数', f"{totals['商品访客数']:,.0f}", period_delta_text('商品访客数'))
    k5.metric('支付转化率', f"{totals['支付转化率']*100:.2f}%", period_delta_text('支付转化率'))
    # Row 2: 客单价 + 加购率 + 推广核心指标
    k6, k7, k8, k9, k10 = st.columns(5)
    k6.metric('客单价', f"¥{totals['客单价']:,.0f}", period_delta_text('客单价'))
    k7.metric('加购率', f"{totals['加购率']*100:.2f}%", period_delta_text('加购率'))
    _ps_wan = round(promo_spend / 10000, 1) if promo_spend else 0
    # 推广指标环比/同比
    if promo_rows:
        _pfc_d = _promo_delta(promo_spend, promo_mom_fc, promo_yoy_fc, '%')
        _proi_d = _promo_delta(promo_roi, promo_mom_roi, promo_yoy_roi, '%')
        _prate_d = _promo_delta(promo_rate, promo_mom_rate, promo_yoy_rate, '%')
        k8.metric('推广费', f"¥{_ps_wan}万" if promo_spend >= 10000 else f"¥{promo_spend:,.0f}", _pfc_d)
        k9.metric('ROI', f"{promo_roi:.2f}" if promo_roi else '--', _proi_d)
        k10.metric('费率', f"{promo_rate:.2f}%", _prate_d)
    else:
        k8.metric('推广费', f"¥{_ps_wan}万" if promo_spend >= 10000 else f"¥{promo_spend:,.0f}")
        k9.metric('ROI', f"{promo_roi:.2f}" if promo_roi else '--')
        k10.metric('费率', f"{promo_rate:.2f}%" if promo_rate else '--')
    # Row 3: 推广效率指标
    if promo_rows:
        _pdroi_d = _promo_delta(promo_direct_roi, promo_mom_droi, promo_yoy_droi, '%')
        _pctr_d = _promo_delta(promo_ctr*100, promo_mom_ctr*100, promo_yoy_ctr*100, '%')
        _pcpc_d = _promo_delta(promo_cpc, promo_mom_cpc, promo_yoy_cpc, '%')
        _poc_d = _promo_delta(promo_order_cost, promo_mom_order_cost, promo_yoy_order_cost, '%')
        _prs = promo_order_amt / totals['支付金额'] * 100 if totals['支付金额'] else 0
        _prs_m = promo_mom_amt / totals['支付金额'] * 100 if totals['支付金额'] else 0
        _prs_y = promo_yoy_amt / totals['支付金额'] * 100 if totals['支付金额'] else 0
        _prs_d = _promo_delta(_prs, _prs_m, _prs_y, '%')
        k11, k12, k13, k14, k15 = st.columns(5)
        k11.metric('直接ROI', f"{promo_direct_roi:.2f}" if promo_direct_roi else '--', _pdroi_d)
        k12.metric('点击率', f"{promo_ctr*100:.2f}%" if promo_impress else '--', _pctr_d)
        k13.metric('点击成本', f"¥{promo_cpc:.2f}" if promo_clicks else '--', _pcpc_d)
        k14.metric('订单成本', f"¥{promo_order_cost:.2f}" if totals['支付买家数'] else '--', _poc_d)
        k15.metric('推广成交占比', f"{_prs:.2f}%", _prs_d)

    # ── 日趋势（最近30天）──
    st.markdown('<div class="section-title">日趋势（最近30天）</div>', unsafe_allow_html=True)
    if daily_trend:
        # 自适应单位
        max_amt = max(r['支付金额'] for r in daily_trend)
        use_wan = max_amt >= 10000
        def _amt_label(v):
            return f"{v/10000:.1f}万" if use_wan else f"{v:,.0f}"
        def _amt_y(v):
            return v/10000 if use_wan else v
        amt_unit = '万' if use_wan else '元'

        # 1) 支付金额趋势
        bar_texts = []
        for r in daily_trend:
            t = _amt_label(r['支付金额'])
            if r['支付金额_同比'] is not None:
                sign = '+' if r['支付金额_同比'] >= 0 else ''
                t += f"<br><span style='font-size:10px'>{sign}{r['支付金额_同比']*100:.1f}%</span>"
            bar_texts.append(t)
        fig_a = go.Figure(go.Bar(
            x=[r['日期'] for r in daily_trend],
            y=[_amt_y(r['支付金额']) for r in daily_trend],
            text=bar_texts, textposition='outside',
            marker_color='#3b82f6'))
        fig_a.update_layout(
            title='支付金额趋势', height=340, template='plotly_white',
            margin=dict(l=20, r=20, t=45, b=20),
            yaxis_title=f'支付金额({amt_unit})', showlegend=False)
        st.plotly_chart(fig_a, width="stretch")

        # 2) 访客数趋势
        vis_texts = []
        for r in daily_trend:
            t = f"{int(r['商品访客数']):,}"
            if r['商品访客数_同比'] is not None:
                sign = '+' if r['商品访客数_同比'] >= 0 else ''
                t += f"<br><span style='font-size:10px'>{sign}{r['商品访客数_同比']*100:.1f}%</span>"
            vis_texts.append(t)
        fig_b = go.Figure(go.Scatter(
            x=[r['日期'] for r in daily_trend],
            y=[r['商品访客数'] for r in daily_trend],
            text=vis_texts, textposition='top center', mode='lines+markers+text',
            line=dict(color='#06b6d4', width=2),
            marker=dict(size=5)))
        fig_b.update_layout(
            title='访客数趋势', height=340, template='plotly_white',
            margin=dict(l=20, r=20, t=45, b=20),
            yaxis_title='访客数', showlegend=False)
        st.plotly_chart(fig_b, width="stretch")

        # 3) 转化率趋势
        cvr_texts = []
        for r in daily_trend:
            t = f"{r['支付转化率']*100:.2f}%"
            if r['支付转化率_同比'] is not None:
                sign = '+' if r['支付转化率_同比'] >= 0 else ''
                t += f"<br><span style='font-size:10px'>{sign}{r['支付转化率_同比']*100:.1f}%</span>"
            cvr_texts.append(t)
        fig_c = go.Figure(go.Scatter(
            x=[r['日期'] for r in daily_trend],
            y=[r['支付转化率']*100 for r in daily_trend],
            text=cvr_texts, textposition='top center', mode='lines+markers+text',
            line=dict(color='#f59e0b', width=2),
            fill='tozeroy', fillcolor='rgba(245,158,11,0.15)',
            marker=dict(size=5)))
        fig_c.update_layout(
            title='支付转化率趋势', height=340, template='plotly_white',
            margin=dict(l=20, r=20, t=45, b=20),
            yaxis_title='转化率(%)', showlegend=False)
        st.plotly_chart(fig_c, width="stretch")
        # 图表数据下载
        _render_download_panel(daily_trend, ['日期', '支付金额', '商品访客数', '支付转化率', '支付件数', '客单价'], 'overview_daily_trend.csv')
    else:
        st.info('当前筛选条件下，最近30天无日趋势数据')

    st.markdown('<div class="section-title">全域趋势与结构</div>', unsafe_allow_html=True)
    trend = [{'月份': r['月份'], '支付金额': r['支付金额'], '访客数': r['商品访客数'],
               '支付件数': r['支付件数'], '转化率': r['支付转化率']} for r in filtered_monthly]
    a_col, b_col = st.columns([2, 1])
    with a_col:
        fig = go.Figure()
        if trend:
            fig.add_trace(go.Bar(x=[r['月份'][:4]+'/'+r['月份'][5:7] for r in trend], y=[_wan(r['支付金额']) for r in trend],
                                  text=[f"{_wan(r['支付金额'])}万" for r in trend], textposition='outside',
                                  name='支付金额(万)', marker_color='#1d4ed8'))
            fig.add_trace(go.Scatter(x=[r['月份'][:4]+'/'+r['月份'][5:7] for r in trend], y=[r['访客数'] for r in trend],
                                      name='访客数', yaxis='y2', line=dict(color='#06b6d4', width=3)))
            fig.add_trace(go.Scatter(x=[r['月份'][:4]+'/'+r['月份'][5:7] for r in trend], y=[r['支付件数'] for r in trend],
                                      name='支付件数', yaxis='y2', line=dict(color='#22c55e', width=3)))
        fig.update_layout(height=390, template='plotly_white', margin=dict(l=20, r=20, t=35, b=20),
                          legend=dict(orientation='h'), yaxis_title='支付金额(万)',
                          yaxis2=dict(title='流量/销量', overlaying='y', side='right'))
        st.plotly_chart(fig, width="stretch")
        _render_download_panel(trend, ['月份', '支付金额', '访客数', '支付件数', '转化率'], 'overview_monthly_trend.csv', '📥 月度趋势')
    with b_col:
        if ch_rows:
            ch_pie = [{'渠道': r['渠道'], '支付金额': r['支付金额']} for r in ch_rows[:8]]
            fig = px.pie(df(ch_pie), names='渠道', values='支付金额', hole=.55,
                         color_discrete_sequence=px.colors.qualitative.Set2)
            fig.update_traces(text=[f"{r['渠道']}<br>¥{_wan(r['支付金额'])}万" for r in ch_rows[:8]],
                              hovertemplate='%{label}<br>¥%{value:,.0f}<extra></extra>')
            fig.update_layout(height=390, margin=dict(l=10, r=10, t=35, b=10), title='渠道销售占比')
            st.plotly_chart(fig, width="stretch")
            _render_download_panel(ch_rows[:8], ['渠道', '支付金额'], 'overview_channel_pie.csv', '📥 渠道占比')
        else:
            st.info('暂无渠道数据')

    c_col, d_col, e_col = st.columns(3)
    with c_col:
        # 店铺销售排行 → 柱状图
        sr = store_rows[:12]
        fig = go.Figure(go.Bar(
            x=[r['支付金额'] for r in sr],
            y=[r['店铺'] for r in sr],
            orientation='h',
            text=[f"¥{_wan(r['支付金额'])}万" for r in sr],
            textposition='outside',
            marker=dict(color=px.colors.qualitative.Bold[:len(sr)])))
        fig.update_layout(height=430, margin=dict(l=10, r=80, t=35, b=10),
                          title='店铺销售排行', template='plotly_white',
                          yaxis=dict(categoryorder='total ascending'),
                          xaxis=dict(title='支付金额(万)', showgrid=True))
        st.plotly_chart(fig, width="stretch")
        _render_download_panel(sr, ['店铺', '支付金额', '支付件数', '商品访客数', '支付转化率', '客单价'], 'store_ranking.csv', '📥 店铺排行')
    with d_col:
        # 品类销售排行 → 柱状图
        cr = cat_rows[:12]
        fig = go.Figure(go.Bar(
            x=[r['支付金额'] for r in cr],
            y=[r['品类'] for r in cr],
            orientation='h',
            text=[f"¥{_wan(r['支付金额'])}万" for r in cr],
            textposition='outside',
            marker=dict(color=px.colors.qualitative.Pastel[:len(cr)])))
        fig.update_layout(height=430, margin=dict(l=10, r=80, t=35, b=10),
                          title='品类销售排行', template='plotly_white',
                          yaxis=dict(categoryorder='total ascending'),
                          xaxis=dict(title='支付金额(万)', showgrid=True))
        st.plotly_chart(fig, width="stretch")
        _render_download_panel(cr, ['品类', '支付金额', '支付件数', '商品访客数', '支付转化率', '客单价'], 'category_ranking.csv', '📥 品类排行')
    with e_col:
        # TOP10单品 → 横向条形图
        model_rows = group(daily, '型号')[:10]
        fig = go.Figure(go.Bar(
            x=[r['支付金额'] for r in model_rows],
            y=[r['型号'] for r in model_rows],
            orientation='h',
            text=[f"¥{_wan(r['支付金额'])}万" for r in model_rows],
            textposition='outside',
            marker=dict(color=px.colors.qualitative.Set2[:len(model_rows)])))
        fig.update_layout(height=430, margin=dict(l=10, r=80, t=35, b=10),
                          title='销额TOP10单品', template='plotly_white',
                          yaxis=dict(categoryorder='total ascending'),
                          xaxis=dict(title='支付金额(万)', showgrid=True))
        st.plotly_chart(fig, width="stretch")
        _render_download_panel(model_rows, ['型号', '支付金额', '支付件数', '商品访客数', '支付转化率', '客单价'], 'top10_models.csv', '📥 TOP10单品')

    st.markdown('<div class="section-title">导出与留档</div>', unsafe_allow_html=True)
    d1, d2 = st.columns(2)
    with d1:
        comp = []
        for r in reversed(filtered_monthly):
            prev = mm_f.get(month_shift(r['月份'], -1))
            ly = mm_f.get(month_shift(r['月份'], -12))
            comp.append({
                '月份': r['月份'], '支付金额(万)': round(_wan(r['支付金额']), 1), '支付件数': round(r['支付件数'], 0),
                '访客数': round(r['商品访客数'], 0), '转化率': round(r['支付转化率'], 4),
                '金额环比': None if not prev or not prev['支付金额'] else round((r['支付金额'] - prev['支付金额']) / prev['支付金额'], 4),
                '金额同比': None if not ly or not ly['支付金额'] else round((r['支付金额'] - ly['支付金额']) / ly['支付金额'], 4),
            })
        _render_download_panel(comp, ['月份', '支付金额(万)', '支付件数', '访客数', '转化率', '金额环比', '金额同比'], 'monthly_yoy_mom.csv', '📥 月度同比环比')
    with d2:
        _render_download_panel(daily, ['日期', '渠道', '店铺', '品类', '型号', '支付金额', '支付件数', '商品访客数', '支付转化率', '客单价', '退款率'], 'filtered_daily_summary.csv', '📥 日汇总明细')


# ═════════════════════════════════════════════════════════════
# TAB 2: 📢 推广分析
# ═════════════════════════════════════════════════════════════
with tabs[1]:
    st.markdown('<div class="section-title">📢 推广分析</div>', unsafe_allow_html=True)
    if not promo_rows:
        st.info('请先在左侧选择「推广数据」，然后上传含「京东推广数据源」和「天猫推广数据源」工作表的Excel文件。')
    else:
        # ── 销售支付金额预计算（按维度聚合，用于费率=花费/销售支付金额）──
        _sales_by_channel = {}
        _sales_by_store = {}
        _sales_by_cat = {}
        _sales_by_model = {}
        _sales_total_amt = 0
        for r in daily:
            d = r.get('日期', '')
            if len(d) == 7: d = d + '-01'
            if not d or not (str(start) <= d <= str(end)): continue
            if channel and r.get('渠道', '') not in channel: continue
            if store and r.get('店铺', '') not in store: continue
            if category and r.get('品类', '') not in category: continue
            if model and r.get('型号', '') not in model: continue
            amt = float(r.get('支付金额', 0) or 0)
            _sales_total_amt += amt
            _sales_by_channel[r.get('渠道', '') or '未标注'] = _sales_by_channel.get(r.get('渠道', '') or '未标注', 0) + amt
            _sales_by_store[r.get('店铺', '') or '未标注'] = _sales_by_store.get(r.get('店铺', '') or '未标注', 0) + amt
            _sales_by_cat[r.get('品类', '') or '未标注'] = _sales_by_cat.get(r.get('品类', '') or '未标注', 0) + amt
            _sales_by_model[r.get('型号', '') or '未标注'] = _sales_by_model.get(r.get('型号', '') or '未标注', 0) + amt

        # ── KPI 指标 ──
        _ps = promo_spend
        _ro = promo_roi
        _imp = sum(r.get('_展现数', 0) for r in promo_filtered)
        _clk = sum(r.get('_点击数', 0) for r in promo_filtered)
        _ctr = _clk / _imp if _imp else 0
        _cpc = _ps / _clk if _clk else 0
        _da = sum(r.get('_总订单金额', 0) for r in promo_filtered)
        _porders = sum(1 for r in promo_filtered if r.get('_总订单金额', 0) > 0) or totals['支付买家数']
        _p_order_cost = _ps / _porders if _porders else 0
        _poc_m = promo_mom_fc / _porders if _porders else 0
        _poc_y = promo_yoy_fc / _porders if _porders else 0
        pfc_d = _promo_delta(_ps, promo_mom_fc, promo_yoy_fc, '%')
        proi_d = _promo_delta(_ro, promo_mom_roi, promo_yoy_roi, '%')
        pctr_d = _promo_delta(_ctr*100, promo_mom_ctr*100, promo_yoy_ctr*100, '%')
        cpc_d = _promo_delta(_cpc, promo_mom_cpc, promo_yoy_cpc, '%')
        pda_d = _promo_delta(_da, promo_mom_amt, promo_yoy_amt, '%')
        poc_d_text = _promo_delta(_p_order_cost, _poc_m, _poc_y, '')
        pk1, pk2, pk3, pk4, pk5, pk6 = st.columns(6)
        pk1.metric('推广花费', f"¥{_wan(_ps)}万" if _ps >= 10000 else f"¥{_ps:,.0f}", pfc_d)
        pk2.metric('ROI', f"{_ro:.2f}" if _ro else '--', proi_d)
        pk3.metric('点击率', f"{_ctr*100:.2f}%" if _imp else '--', pctr_d)
        pk4.metric('平均点击成本', f"¥{_cpc:.2f}" if _clk else '--', cpc_d)
        pk5.metric('总成交金额', f"¥{_wan(_da)}万" if _da >= 10000 else f"¥{_da:,.0f}", pda_d)
        pk6.metric('订单成本', f"¥{_p_order_cost:.2f}", poc_d_text)

        # ── 推广费趋势（日/月）──
        st.markdown('<div class="section-title">推广费 & 成交金额趋势</div>', unsafe_allow_html=True)
        promo_gran = st.radio('粒度', ['按日', '按月'], horizontal=True, key='promo_gran')
        _pr = {}
        for r in promo_filtered:
            dt = r.get('_date', '')
            if not dt:
                continue
            key = dt if promo_gran == '按日' else dt[:7]
            _pr.setdefault(key, {'花费': 0, '总订单金额': 0})
            _pr[key]['花费'] += r.get('_花费', 0)
            _pr[key]['总订单金额'] += r.get('_总订单金额', 0)
        _pr_s = sorted(_pr.items())
        if _pr_s:
            fig = go.Figure()
            fig.add_trace(go.Bar(x=[x[0] for x in _pr_s], y=[x[1]['花费']/10000 for x in _pr_s],
                                  name='推广费(万)', marker_color='#f59e0b', opacity=0.85))
            fig.add_trace(go.Scatter(x=[x[0] for x in _pr_s], y=[x[1]['总订单金额']/10000 for x in _pr_s],
                                     name='总订单金额(万)', yaxis='y2', line=dict(color='#10b981', width=2)))
            fig.update_layout(height=360, template='plotly_white', legend=dict(orientation='h'),
                                  yaxis_title='推广费(万)', yaxis2=dict(title='订单金额(万)', overlaying='y', side='right'))
            st.plotly_chart(fig, width="stretch")
            _render_download_panel([{'日期': x[0], '花费': x[1]['花费'], '总订单金额': x[1]['总订单金额']} for x in _pr_s],
                              ['日期', '花费', '总订单金额'], 'promo_spend_trend.csv')
        _gran_label = '日度' if promo_gran == '按日' else '月度'
        st.markdown(f'<div class="section-title">ROI 趋势（{_gran_label}）</div>', unsafe_allow_html=True)
        _roi_gr = {}
        for r in promo_filtered:
            dt = r.get('_date', '')
            if not dt:
                continue
            gk = dt if promo_gran == '按日' else dt[:7]
            _roi_gr.setdefault(gk, {'花费': 0, '总订单金额': 0})
            _roi_gr[gk]['花费'] += r.get('_花费', 0)
            _roi_gr[gk]['总订单金额'] += r.get('_总订单金额', 0)
        _roi_s = sorted(_roi_gr.items())
        if _roi_s:
            _x_label = '日期' if promo_gran == '按日' else '月份'
            _roi_v = [{_x_label: x[0], 'ROI': x[1]['总订单金额']/x[1]['花费'] if x[1]['花费'] else 0} for x in _roi_s]
            fig = px.line(pd.DataFrame(_roi_v), x=_x_label, y='ROI', markers=True,
                              title=f'{_gran_label}ROI趋势', line_shape='spline')
            fig.update_layout(height=320, template='plotly_white', yaxis_title='ROI')
            st.plotly_chart(fig, width="stretch")
            _render_download_panel([{'日期': x[0], '花费': x[1]['花费'], '总订单金额': x[1]['总订单金额'],
                                   'ROI': x[1]['总订单金额']/x[1]['花费'] if x[1]['花费'] else 0} for x in _roi_s],
                              ['日期', '花费', '总订单金额', 'ROI'], 'promo_roi_trend.csv')
        st.markdown('<div class="section-title">推广效率矩阵（花费 vs 成交金额）</div>', unsafe_allow_html=True)
        _pl = {}
        for r in promo_filtered:
            pn = r.get('推广计划', '') or r.get('计划ID', '') or '未标注'
            _pl.setdefault(pn, {'花费': 0, '总订单金额': 0})
            _pl[pn]['花费'] += r.get('_花费', 0)
            _pl[pn]['总订单金额'] += r.get('_总订单金额', 0)
        _pl_r = [{'推广计划': k, '花费': v['花费'], '总订单金额': v['总订单金额'],
                    'ROI': v['总订单金额']/v['花费'] if v['花费'] else 0} for k, v in _pl.items()]
        if _pl_r:
            _df = pd.DataFrame(_pl_r)
            fig = px.scatter(_df, x='花费', y='总订单金额', size='总订单金额',
                                     hover_data=['推广计划', 'ROI'], title='推广计划效率矩阵（花费 vs 成交金额）',
                                     color='ROI', color_continuous_scale='RdYlGn')
            fig.update_layout(height=400, template='plotly_white')
            st.plotly_chart(fig, width="stretch")
            _render_download_panel(_pl_r, ['推广计划', '花费', '总订单金额', 'ROI'], 'promo_plan_efficiency.csv', '📥 推广计划效率')

        # ── 产品线推广占比 ──
        st.markdown('<div class="section-title">产品线推广占比</div>', unsafe_allow_html=True)
        _ln = {}
        for r in promo_filtered:
            ln = r.get('产品线', '') or r.get('营销场景', '') or '未标注'
            _ln[ln] = _ln.get(ln, 0) + r.get('_花费', 0)
        _ln_r = [{'产品线': k, '花费': v} for k, v in _ln.items() if v > 0]
        if _ln_r:
            fig = px.pie(pd.DataFrame(_ln_r), names='产品线', values='花费', hole=.45,
                              color_discrete_sequence=px.colors.qualitative.Set2)
            fig.update_traces(text=[f"{k}<br>¥{v/10000:.1f}万" for k, v in _ln.items() if v > 0],
                                 hovertemplate='%{label}<br>花费：¥%{value:,.0f}<extra></extra>')
            fig.update_layout(height=380, margin=dict(l=10, r=10, t=35, b=10), title='产品线推广费占比')
            st.plotly_chart(fig, width="stretch")
            _render_download_panel(_ln_r, ['产品线', '花费'], 'promo_product_line_share.csv', '📥 产品线推广费占比')

        # ── 店铺推广矩阵 ──
        st.markdown('<div class="section-title">🏪 店铺推广矩阵</div>', unsafe_allow_html=True)
        _store_m = {}
        for r in promo_filtered:
            sn = r.get('_店铺', '') or r.get('店铺', '') or '未标注'
            _store_m.setdefault(sn, {'花费': 0, '展现数': 0, '点击数': 0, '总订单金额': 0, '直接订单金额': 0, '总加购数': 0, '成交客户数': 0})
            _store_m[sn]['花费'] += r.get('_花费', 0)
            _store_m[sn]['展现数'] += r.get('_展现数', 0)
            _store_m[sn]['点击数'] += r.get('_点击数', 0)
            _store_m[sn]['总订单金额'] += r.get('_总订单金额', 0)
            _store_m[sn]['直接订单金额'] += r.get('_直接订单金额', 0)
            _store_m[sn]['总加购数'] += r.get('_总加购数', 0)
            _store_m[sn]['成交客户数'] += r.get('_成交客户数', 0)
        _store_yoy = _promo_agg(promo_yoy, '_店铺')
        _sm_r = []
        for k, v in sorted(_store_m.items(), key=lambda x: x[1]['花费'], reverse=True):
            _imp = v['展现数']
            _clk = v['点击数']
            _fc = v['花费']
            _cust = v['成交客户数']
            _cv = _cust / _clk * 100 if _clk else None  # 推广转化率 = 成交客户数/点击量
            _d_roi = v['直接订单金额'] / _fc if _fc else 0
            _cpc = _fc / _clk if _clk else 0
            # 同比
            vy = _store_yoy.get(k, {})
            _fc_yoy, _ = _yoy_text(_fc, vy.get('花费', 0) if vy.get('花费', 0) else None)
            _droi_yoy, _ = _yoy_text(_d_roi, (vy.get('直接订单金额', 0) / vy.get('花费', 1)) if vy.get('花费', 0) else None)
            _cpc_yoy, _ = _yoy_text(_cpc, (vy.get('花费', 0) / vy.get('点击数', 1)) if vy.get('点击数', 0) else None)
            _yc = vy.get('总加购数', 0) / vy.get('点击数', 1) * 100 if vy.get('点击数', 0) else None
            _cv_yoy, _ = _yoy_text(_cv, _yc) if _cv and _yc else ('--', '')
            _s_amt = _sales_by_store.get(k, 0)
            _sm_r.append({
                '店铺': k,
                '花费(万)': f"{_fc/10000:.1f}",
                '展现数': f"{int(_imp):,}",
                '点击数': f"{int(_clk):,}",
                '点击率': f"{_clk/_imp*100:.2f}%" if _imp else '--',
                '点击成本': f"¥{_cpc:.2f}" if _clk else '--',
                '总成交(万)': f"{v['总订单金额']/10000:.1f}",
                '直接成交(万)': f"{v['直接订单金额']/10000:.1f}",
                'ROI': f"{v['总订单金额']/_fc:.2f}" if _fc else '--',
                '直接ROI': f"{_d_roi:.2f}" if _fc else '--',
                '费率': f"{_fc/_s_amt*100:.2f}%" if _s_amt else '--',
                '转化率': f"{_cv:.1f}%" if _cv else '--',
                '花费同比': f"<span style='color:{'#22c55e' if _fc_yoy and '+' in _fc_yoy else '#dc2626' if _fc_yoy and '-' in _fc_yoy else '#94a3b8'}'>{_fc_yoy}</span>",
                '直接ROI同比': f"<span style='color:{'#22c55e' if _droi_yoy and '+' in _droi_yoy else '#dc2626' if _droi_yoy and '-' in _droi_yoy else '#94a3b8'}'>{_droi_yoy}</span>",
                'CPC同比': f"<span style='color:{'#22c55e' if _cpc_yoy and '+' in _cpc_yoy else '#dc2626' if _cpc_yoy and '-' in _cpc_yoy else '#94a3b8'}'>{_cpc_yoy}</span>",
                '转化率同比': f"<span style='color:{'#22c55e' if _cv_yoy and '+' in _cv_yoy else '#dc2626' if _cv_yoy and '-' in _cv_yoy else '#94a3b8'}'>{_cv_yoy}</span>",
            })
        if _sm_r:
            ma1, ma2 = st.columns(2)
            with ma1:
                _fc_vals = [float(x['花费(万)']) for x in _sm_r]
                fig = go.Figure(go.Bar(
                    x=_fc_vals,
                    y=[x['店铺'] for x in _sm_r],
                    orientation='h',
                    text=[f"¥{x['花费(万)']}万" for x in _sm_r],
                    textposition='outside',
                    marker=dict(color=px.colors.qualitative.Pastel[:len(_sm_r)])))
                fig.update_layout(height=max(280, len(_sm_r)*45), margin=dict(l=10, r=80, t=35, b=10),
                                   title='各店铺推广花费', template='plotly_white',
                                   yaxis=dict(categoryorder='total ascending'))
                st.plotly_chart(fig, width="stretch")
                _render_download_panel(_sm_r, list(_sm_r[0].keys()), 'promo_store_spend.csv', '📥 店铺推广费')
            with ma2:
                _roi_vals = [float(x['ROI']) if x['ROI'] != '--' else 0 for x in _sm_r]
                _colors_roi = ['#22c55e' if v >= 3 else '#f59e0b' if v >= 1 else '#ef4444' for v in _roi_vals]
                fig = go.Figure(go.Bar(
                    x=[float(x['ROI']) if x['ROI'] != '--' else 0 for x in _sm_r],
                    y=[x['店铺'] for x in _sm_r],
                    orientation='h',
                    text=[str(x['ROI']) for x in _sm_r],
                    textposition='outside',
                    marker=dict(color=_colors_roi)))
                fig.update_layout(height=max(280, len(_sm_r)*45), margin=dict(l=10, r=80, t=35, b=10),
                                   title='各店铺ROI（绿≥3 橙≥1 红<1）', template='plotly_white',
                                   yaxis=dict(categoryorder='total ascending'))
                st.plotly_chart(fig, width="stretch")
                _render_download_panel(_sm_r, list(_sm_r[0].keys()), 'promo_store_roi.csv', '📥 店铺ROI')
            _cols = list(_sm_r[0].keys())
            _sm_html = _html_table(_sm_r, col_widths={c: '100px' for c in _cols}, height=max(280, len(_sm_r)*34+40))
            st.markdown(_wrap_fullscreen(_sm_html, title='🏪 店铺推广矩阵')[0], unsafe_allow_html=True)
            _render_download_panel(_sm_r, list(_sm_r[0].keys()), 'promo_store_matrix.csv', '📥 店铺推广矩阵')

        # ── 渠道推广矩阵 ──
        st.markdown('<div class="section-title">📡 渠道推广矩阵</div>', unsafe_allow_html=True)
        _chan_m = {}
        for r in promo_filtered:
            cn = r.get('_渠道', '') or r.get('渠道', '') or '未标注'
            _chan_m.setdefault(cn, {'花费': 0, '展现数': 0, '点击数': 0, '总订单金额': 0, '直接订单金额': 0, '总加购数': 0, '成交客户数': 0})
            _chan_m[cn]['花费'] += r.get('_花费', 0)
            _chan_m[cn]['展现数'] += r.get('_展现数', 0)
            _chan_m[cn]['点击数'] += r.get('_点击数', 0)
            _chan_m[cn]['总订单金额'] += r.get('_总订单金额', 0)
            _chan_m[cn]['直接订单金额'] += r.get('_直接订单金额', 0)
            _chan_m[cn]['总加购数'] += r.get('_总加购数', 0)
            _chan_m[cn]['成交客户数'] += r.get('_成交客户数', 0)
        _chan_yoy = _promo_agg(promo_yoy, '_渠道')
        _cm_r = []
        for k, v in sorted(_chan_m.items(), key=lambda x: x[1]['花费'], reverse=True):
            _imp = v['展现数']
            _clk = v['点击数']
            _fc = v['花费']
            _cust = v['成交客户数']
            _cv = _cust / _clk * 100 if _clk else None
            _d_roi = v['直接订单金额'] / _fc if _fc else 0
            _cpc = _fc / _clk if _clk else 0
            vy = _chan_yoy.get(k, {})
            _fc_yoy, _ = _yoy_text(_fc, vy.get('花费', 0) if vy.get('花费', 0) else None)
            _droi_yoy, _ = _yoy_text(_d_roi, (vy.get('直接订单金额', 0) / vy.get('花费', 1)) if vy.get('花费', 0) else None)
            _cpc_yoy, _ = _yoy_text(_cpc, (vy.get('花费', 0) / vy.get('点击数', 1)) if vy.get('点击数', 0) else None)
            _yc = vy.get('成交客户数', 0) / vy.get('点击数', 1) * 100 if vy.get('点击数', 0) else None
            _cv_yoy, _ = _yoy_text(_cv, _yc) if _cv and _yc else ('--', '')
            _s_amt_cm = _sales_by_channel.get(k, 0)
            _cm_r.append({
                '渠道': k,
                '花费(万)': f"{_fc/10000:.1f}",
                '展现数': f"{int(_imp):,}",
                '点击数': f"{int(_clk):,}",
                '点击率': f"{_clk/_imp*100:.2f}%" if _imp else '--',
                '点击成本': f"¥{_cpc:.2f}" if _clk else '--',
                '总成交(万)': f"{v['总订单金额']/10000:.1f}",
                'ROI': f"{v['总订单金额']/_fc:.2f}" if _fc else '--',
                '直接ROI': f"{_d_roi:.2f}" if _fc else '--',
                '费率': f"{_fc/_s_amt_cm*100:.2f}%" if _s_amt_cm else '--',
                '转化率': f"{_cv:.1f}%" if _cv else '--',
                '花费同比': f"<span style='color:{'#22c55e' if _fc_yoy and '+' in _fc_yoy else '#dc2626' if _fc_yoy and '-' in _fc_yoy else '#94a3b8'}'>{_fc_yoy}</span>",
                '直接ROI同比': f"<span style='color:{'#22c55e' if _droi_yoy and '+' in _droi_yoy else '#dc2626' if _droi_yoy and '-' in _droi_yoy else '#94a3b8'}'>{_droi_yoy}</span>",
                'CPC同比': f"<span style='color:{'#22c55e' if _cpc_yoy and '+' in _cpc_yoy else '#dc2626' if _cpc_yoy and '-' in _cpc_yoy else '#94a3b8'}'>{_cpc_yoy}</span>",
                '转化率同比': f"<span style='color:{'#22c55e' if _cv_yoy and '+' in _cv_yoy else '#dc2626' if _cv_yoy and '-' in _cv_yoy else '#94a3b8'}'>{_cv_yoy}</span>",
            })
        if _cm_r:
            cb1, cb2 = st.columns(2)
            with cb1:
                _fc_pie = [{'渠道': x['渠道'], '花费(万)': float(x['花费(万)'])} for x in _cm_r if float(x['花费(万)']) > 0]
                if _fc_pie:
                    fig = px.pie(pd.DataFrame(_fc_pie), names='渠道', values='花费(万)', hole=.4,
                                  color_discrete_sequence=px.colors.qualitative.Bold,
                                  title='渠道推广费占比')
                    fig.update_traces(texttemplate='%{label}<br>%{percent:.1%}')
                    fig.update_layout(height=340, margin=dict(l=10, r=10, t=40, b=10))
                    st.plotly_chart(fig, width="stretch")
                    _render_download_panel(_fc_pie, ['渠道', '花费(万)'], 'promo_chan_spend_share.csv', '📥 渠道推广费占比')
            with cb2:
                _roi_cur = [float(x['ROI']) if x['ROI'] != '--' else 0 for x in _cm_r]
                _droi_cur = [float(x['直接ROI']) if x['直接ROI'] != '--' else 0 for x in _cm_r]
                fig = go.Figure()
                fig.add_trace(go.Bar(name='ROI', x=[x['渠道'] for x in _cm_r], y=_roi_cur, marker_color='#1d4ed8'))
                fig.add_trace(go.Bar(name='直接ROI', x=[x['渠道'] for x in _cm_r], y=_droi_cur, marker_color='#06b6d4'))
                fig.update_layout(height=340, barmode='group', template='plotly_white', title='渠道ROI对比')
                st.plotly_chart(fig, width="stretch")
                _render_download_panel(_cm_r, list(_cm_r[0].keys()), 'promo_chan_roi.csv', '📥 渠道ROI对比')
            _cols = list(_cm_r[0].keys())
            _cm_html = _html_table(_cm_r, col_widths={c: '100px' for c in _cols}, height=max(280, len(_cm_r)*34+40))
            st.markdown(_wrap_fullscreen(_cm_html, title='📡 渠道推广矩阵')[0], unsafe_allow_html=True)
            _render_download_panel(_cm_r, list(_cm_r[0].keys()), 'promo_chan_matrix.csv', '📥 渠道推广矩阵')

        # ── TOP10 推广计划（按花费）──
        st.markdown('<div class="section-title">TOP10 推广计划（按花费）</div>', unsafe_allow_html=True)
        _tp = sorted(_pl_r, key=lambda x: x.get('花费', 0), reverse=True)[:10]
        if _tp:
            fig = go.Figure(go.Bar(
                x=[x.get('花费', 0)/10000 for x in _tp],
                y=[str(x.get('推广计划', '') or '未标注')[:20] for x in _tp],
                orientation='h',
                text=[f"¥{x.get('花费', 0)/10000:.1f}万" for x in _tp],
                textposition='outside',
                marker=dict(color=px.colors.qualitative.Bold[:len(_tp)])))
            fig.update_layout(height=400, margin=dict(l=10, r=80, t=35, b=10),
                               title='TOP10 推广计划（按花费）', template='plotly_white',
                               yaxis=dict(categoryorder='total ascending'),
                               xaxis=dict(title='花费(万)', showgrid=True))
            st.plotly_chart(fig, width="stretch")
            _render_download_panel(_tp, ['推广计划', '花费', '总订单金额', 'ROI'], 'promo_top10_plans.csv', '📥 TOP10推广计划')

        # ── 推广明细表 ──
        st.markdown('<div class="section-title">推广明细（按日聚合）</div>', unsafe_allow_html=True)
        _dp = {}
        for r in promo_filtered:
            dt = r.get('_date', '')
            if not dt:
                continue
            _dp.setdefault(dt, {'花费': 0, '展现数': 0, '点击数': 0, '总订单金额': 0, '总加购数': 0})
            _dp[dt]['花费'] += r.get('_花费', 0)
            _dp[dt]['展现数'] += r.get('_展现数', 0)
            _dp[dt]['点击数'] += r.get('_点击数', 0)
            _dp[dt]['总订单金额'] += r.get('_总订单金额', 0)
            _dp[dt]['总加购数'] += r.get('_总加购数', 0)
        # 按日期汇总销售支付金额（用于推广日明细费率）
        _sales_amt_by_day = {}
        for r in daily:
            d = r.get('日期', '')
            if len(d) == 7: d = d + '-01'
            if not d or not (str(start) <= d <= str(end)): continue
            if channel and r.get('渠道', '') not in channel: continue
            if store and r.get('店铺', '') not in store: continue
            if category and r.get('品类', '') not in category: continue
            if model and r.get('型号', '') not in model: continue
            _sales_amt_by_day[d] = _sales_amt_by_day.get(d, 0) + float(r.get('支付金额', 0) or 0)
        _pt = []
        for dt in sorted(_dp.keys()):
            v = _dp[dt]
            _imp = v['展现数']
            _clk = v['点击数']
            _ctr = _clk / _imp if _imp else 0
            _roi = v['总订单金额'] / v['花费'] if v['花费'] else 0
            _s_day = _sales_amt_by_day.get(dt, 0)
            _rate = v['花费'] / _s_day * 100 if _s_day else None
            _pt.append({
                '日期': dt, '花费(万)': round(v['花费']/10000, 1), '展现数': f"{int(_imp):,}",
                '点击数': f"{int(_clk):,}", '点击率': f"{_ctr*100:.2f}%",
                '总订单金额(万)': round(v['总订单金额']/10000, 1), 'ROI': f"{_roi:.2f}",
                '费率': f"{_rate:.2f}%" if _rate is not None else '--',
                '总加购数': f"{int(v['总加购数']):,}"
            })
        if _pt:
            _pt_headers = list(_pt[0].keys())
            _render_html_table(_pt, _pt_headers, _pt_headers, title='📢 推广日明细')
            _render_download_panel(_pt, _pt_headers, 'promo_daily.csv', '📥 推广日明细')

        # ── 单品推广分析表 ──
        st.markdown('<div class="section-title">📦 单品推广分析</div>', unsafe_allow_html=True)
        _sku = {}
        for r in promo_filtered:
            sku = r.get('_型号', '') or r.get('型号', '') or r.get('SKU', '') or '未标注'
            _sku.setdefault(sku, {'花费': 0, '展现数': 0, '点击数': 0, '总订单金额': 0, '直接订单金额': 0, '总加购数': 0, '成交客户数': 0})
            _sku[sku]['花费'] += r.get('_花费', 0)
            _sku[sku]['展现数'] += r.get('_展现数', 0)
            _sku[sku]['点击数'] += r.get('_点击数', 0)
            _sku[sku]['总订单金额'] += r.get('_总订单金额', 0)
            _sku[sku]['直接订单金额'] += r.get('_直接订单金额', 0)
            _sku[sku]['总加购数'] += r.get('_总加购数', 0)
            _sku[sku]['成交客户数'] += r.get('_成交客户数', 0)
        _sku_yoy = _promo_agg(promo_yoy, '_型号')
        _sku_r = []
        for k, v in sorted(_sku.items(), key=lambda x: x[1]['花费'], reverse=True):
            _imp = v['展现数']
            _clk = v['点击数']
            _fc = v['花费']
            _cust = v['成交客户数']
            _cv = _cust / _clk * 100 if _clk else None
            _d_roi = v['直接订单金额'] / _fc if _fc else 0
            _cpc = _fc / _clk if _clk else 0
            vy = _sku_yoy.get(k, {})
            _fc_yoy, _ = _yoy_text(_fc, vy.get('花费', 0) if vy.get('花费', 0) else None)
            _droi_yoy, _ = _yoy_text(_d_roi, (vy.get('直接订单金额', 0) / vy.get('花费', 1)) if vy.get('花费', 0) else None)
            _cpc_yoy, _ = _yoy_text(_cpc, (vy.get('花费', 0) / vy.get('点击数', 1)) if vy.get('点击数', 0) else None)
            _yc = vy.get('成交客户数', 0) / vy.get('点击数', 1) * 100 if vy.get('点击数', 0) else None
            _cv_yoy, _ = _yoy_text(_cv, _yc) if _cv and _yc else ('--', '')
            _s_mdl = _sales_by_model.get(k, 0)
            _sku_r.append({
                '型号': k,
                '花费(万)': f"{_fc/10000:.2f}",
                '展现数': f"{int(_imp):,}",
                '点击数': f"{int(_clk):,}",
                '点击率': f"{_clk/_imp*100:.2f}%" if _imp else '--',
                '点击成本': f"¥{_cpc:.2f}" if _clk else '--',
                '总成交(万)': f"{v['总订单金额']/10000:.2f}",
                '直接成交(万)': f"{v['直接订单金额']/10000:.2f}",
                'ROI': f"{v['总订单金额']/_fc:.2f}" if _fc else '--',
                '直接ROI': f"{_d_roi:.2f}" if _fc else '--',
                '费率': f"{_fc/_s_mdl*100:.2f}%" if _s_mdl else '--',
                '转化率': f"{_cv:.1f}%" if _cv else '--',
                '花费同比': f"<span style='color:{'#22c55e' if _fc_yoy and '+' in _fc_yoy else '#dc2626' if _fc_yoy and '-' in _fc_yoy else '#94a3b8'}'>{_fc_yoy}</span>",
                '直接ROI同比': f"<span style='color:{'#22c55e' if _droi_yoy and '+' in _droi_yoy else '#dc2626' if _droi_yoy and '-' in _droi_yoy else '#94a3b8'}'>{_droi_yoy}</span>",
                'CPC同比': f"<span style='color:{'#22c55e' if _cpc_yoy and '+' in _cpc_yoy else '#dc2626' if _cpc_yoy and '-' in _cpc_yoy else '#94a3b8'}'>{_cpc_yoy}</span>",
                '转化率同比': f"<span style='color:{'#22c55e' if _cv_yoy and '+' in _cv_yoy else '#dc2626' if _cv_yoy and '-' in _cv_yoy else '#94a3b8'}'>{_cv_yoy}</span>",
            })
        if _sku_r:
            sku1, sku2 = st.columns(2)
            with sku1:
                _top10 = _sku_r[:10]
                fig = go.Figure(go.Bar(
                    x=[float(x['花费(万)']) for x in _top10],
                    y=[x['型号'] for x in _top10],
                    orientation='h',
                    text=[f"¥{x['花费(万)']}万" for x in _top10],
                    textposition='outside',
                    marker=dict(color=px.colors.qualitative.Bold[:len(_top10)])))
                fig.update_layout(height=max(300, len(_top10)*40), margin=dict(l=10, r=80, t=35, b=10),
                                   title='TOP10 单品推广花费', template='plotly_white',
                                   yaxis=dict(categoryorder='total ascending'))
                st.plotly_chart(fig, width="stretch")
                _render_download_panel(_top10, ['单品', '花费', '总订单金额', 'ROI'], 'promo_sku_spend.csv', '📥 TOP10单品推广费')
            with sku2:
                _roi_vals = [float(x['ROI']) if x['ROI'] != '--' else 0 for x in _sku_r[:10]]
                _colors = ['#22c55e' if v >= 3 else '#f59e0b' if v >= 1 else '#ef4444' for v in _roi_vals]
                fig = go.Figure(go.Bar(
                    x=_roi_vals,
                    y=[x['型号'] for x in _sku_r[:10]],
                    orientation='h',
                    text=[str(x['ROI']) for x in _sku_r[:10]],
                    textposition='outside',
                    marker=dict(color=_colors)))
                fig.update_layout(height=max(300, len(_sku_r[:10])*40), margin=dict(l=10, r=80, t=35, b=10),
                                   title='TOP10 单品ROI（绿≥3 橙≥1 红&lt;1）', template='plotly_white',
                                   yaxis=dict(categoryorder='total ascending'))
                st.plotly_chart(fig, width="stretch")
            _cols = list(_sku_r[0].keys())
            _sku_html = _html_table(_sku_r, col_widths={c: '100px' for c in _cols}, height=max(300, len(_sku_r)*34+40))
            st.markdown(_wrap_fullscreen(_sku_html, title='📦 单品推广分析')[0], unsafe_allow_html=True)
            _render_download_panel(_sku_r, list(_sku_r[0].keys()), 'promo_sku_analysis.csv', '📥 单品推广分析')
        else:
            st.info('当前筛选条件下无单品推广数据')

        # ── 产品线推广矩阵 ──
        st.markdown('<div class="section-title" style="margin-top:24px;">产品线推广矩阵（按品类）</div>', unsafe_allow_html=True)
        _cat_agg = _promo_agg(promo_filtered, '_品类')
        _cat_yoy = _promo_agg(promo_yoy, '_品类')
        if _cat_agg:
            _cat_r = []
            for k in sorted(_cat_agg.keys(), key=lambda k: _cat_agg[k]['花费'], reverse=True):
                v = _cat_agg[k]
                _fc = v['花费']
                _cw = f"¥{_wan(_fc)}万" if _fc >= 10000 else f"¥{_fc:,.0f}"
                _oa = v['总订单金额']
                _oa_w = f"¥{_wan(_oa)}万" if _oa >= 10000 else f"¥{_oa:,.0f}"
                _da = v['直接订单金额']
                _da_w = f"¥{_wan(_da)}万" if _da >= 10000 else f"¥{_da:,.0f}"
                _ri = _oa / _fc if _fc else 0
                _dri = _da / _fc if _fc else 0
                _s_cat = _sales_by_cat.get(k, 0)
                _rate = _fc / _s_cat * 100 if _s_cat else 0
                _imp = v['展现数'] or 0; _clk = v['点击数'] or 0
                _ctr_v = _clk / _imp * 100 if _imp else 0
                _cpc_v = _fc / _clk if _clk else 0
                _cv = v['成交客户数'] / _clk * 100 if _clk else 0
                vy = _cat_yoy.get(k, {})
                _fc_yoy, _fc_c = _yoy_text(_fc, vy.get('花费', 0) if vy.get('花费', 0) else None)
                _y_dri = vy.get('直接订单金额', 0) / vy.get('花费', 1) if vy.get('花费', 0) else None
                _dri_yoy, _dri_c = _yoy_text(_dri, _y_dri) if _dri and _y_dri else ('--', '')
                _y_cpc = vy.get('花费', 0) / vy.get('点击数', 1) if vy.get('点击数', 0) else None
                _cpc_yoy, _cpc_c = _yoy_text(_cpc_v, _y_cpc) if _cpc_v and _y_cpc else ('--', '')
                _y_cv = vy.get('成交客户数', 0) / vy.get('点击数', 1) * 100 if vy.get('点击数', 0) else None
                _cv_yoy, _cv_c = _yoy_text(_cv, _y_cv) if _cv and _y_cv else ('--', '')
                _cat_r.append({
                    '产品线': k,
                    '花费(万)': _cw,
                    '展现数': f"{int(_imp):,}",
                    '点击数': f"{int(_clk):,}",
                    '点击率': f"{_ctr_v:.2f}%",
                    '点击成本': f"¥{_cpc_v:.2f}",
                    '总成交(万)': _oa_w,
                    '直接成交(万)': _da_w,
                    'ROI': f"{_ri:.2f}",
                    '直接ROI': f"{_dri:.2f}",
                    '费率': f"{_rate:.2f}%",
                    '转化率': f"{_cv:.2f}%",
                    '花费同比': f"<span style='color:{_fc_c or '#94a3b8'}'>{_fc_yoy}</span>",
                    '直接ROI同比': f"<span style='color:{_dri_c or '#94a3b8'}'>{_dri_yoy}</span>",
                    'CPC同比': f"<span style='color:{_cpc_c or '#94a3b8'}'>{_cpc_yoy}</span>",
                    '转化率同比': f"<span style='color:{_cv_c or '#94a3b8'}'>{_cv_yoy}</span>",
                })
            _cat_html = _html_table(_cat_r, col_widths={c: '105px' for c in _cat_r[0].keys()}, height=max(300, len(_cat_r)*34+40))
            st.markdown(_wrap_fullscreen(_cat_html, title='📂 产品线推广矩阵')[0], unsafe_allow_html=True)
            _render_download_panel(_cat_r, list(_cat_r[0].keys()), 'promo_product_line.csv', '📥 产品线推广矩阵')
        else:
            st.info('当前筛选条件下无产品线推广数据')

        # ── 营销场景推广矩阵 ──
        st.markdown('<div class="section-title" style="margin-top:24px;">营销场景推广矩阵</div>', unsafe_allow_html=True)
        _scene_agg = _promo_agg(promo_filtered, '_营销场景')
        _scene_yoy = _promo_agg(promo_yoy, '_营销场景')
        if _scene_agg:
            _scene_r = []
            for k in sorted(_scene_agg.keys(), key=lambda k: _scene_agg[k]['花费'], reverse=True):
                v = _scene_agg[k]
                _fc = v['花费']
                _cw = f"¥{_wan(_fc)}万" if _fc >= 10000 else f"¥{_fc:,.0f}"
                _oa = v['总订单金额']
                _oa_w = f"¥{_wan(_oa)}万" if _oa >= 10000 else f"¥{_oa:,.0f}"
                _da = v['直接订单金额']
                _da_w = f"¥{_wan(_da)}万" if _da >= 10000 else f"¥{_da:,.0f}"
                _ri = _oa / _fc if _fc else 0
                _dri = _da / _fc if _fc else 0
                _rate = _fc / _sales_total_amt * 100 if _sales_total_amt else 0
                _imp = v['展现数'] or 0; _clk = v['点击数'] or 0
                _ctr_v = _clk / _imp * 100 if _imp else 0
                _cpc_v = _fc / _clk if _clk else 0
                _cv = v['成交客户数'] / _clk * 100 if _clk else 0
                vy = _scene_yoy.get(k, {})
                _fc_yoy, _fc_c = _yoy_text(_fc, vy.get('花费', 0) if vy.get('花费', 0) else None)
                _y_dri = vy.get('直接订单金额', 0) / vy.get('花费', 1) if vy.get('花费', 0) else None
                _dri_yoy, _dri_c = _yoy_text(_dri, _y_dri) if _dri and _y_dri else ('--', '')
                _y_cpc = vy.get('花费', 0) / vy.get('点击数', 1) if vy.get('点击数', 0) else None
                _cpc_yoy, _cpc_c = _yoy_text(_cpc_v, _y_cpc) if _cpc_v and _y_cpc else ('--', '')
                _y_cv = vy.get('成交客户数', 0) / vy.get('点击数', 1) * 100 if vy.get('点击数', 0) else None
                _cv_yoy, _cv_c = _yoy_text(_cv, _y_cv) if _cv and _y_cv else ('--', '')
                _scene_r.append({
                    '营销场景': k,
                    '花费(万)': _cw,
                    '展现数': f"{int(_imp):,}",
                    '点击数': f"{int(_clk):,}",
                    '点击率': f"{_ctr_v:.2f}%",
                    '点击成本': f"¥{_cpc_v:.2f}",
                    '总成交(万)': _oa_w,
                    '直接成交(万)': _da_w,
                    'ROI': f"{_ri:.2f}",
                    '直接ROI': f"{_dri:.2f}",
                    '费率': f"{_rate:.2f}%",
                    '转化率': f"{_cv:.2f}%",
                    '花费同比': f"<span style='color:{_fc_c or '#94a3b8'}'>{_fc_yoy}</span>",
                    '直接ROI同比': f"<span style='color:{_dri_c or '#94a3b8'}'>{_dri_yoy}</span>",
                    'CPC同比': f"<span style='color:{_cpc_c or '#94a3b8'}'>{_cpc_yoy}</span>",
                    '转化率同比': f"<span style='color:{_cv_c or '#94a3b8'}'>{_cv_yoy}</span>",
                })
            _scene_html = _html_table(_scene_r, col_widths={c: '105px' for c in _scene_r[0].keys()}, height=max(300, len(_scene_r)*34+40))
            st.markdown(_wrap_fullscreen(_scene_html, title='🎯 营销场景推广矩阵')[0], unsafe_allow_html=True)
            _render_download_panel(_scene_r, list(_scene_r[0].keys()), 'promo_scene.csv', '📥 营销场景推广矩阵')
        else:
            st.info('当前筛选条件下无营销场景推广数据')

# ═══════════════════════════════════════════════════════════════
# TAB 2: 时间段对比
# ═══════════════════════════════════════════════════════════════
with tabs[2]:
    st.markdown('<div class="section-title">时间段对比分析</div>', unsafe_allow_html=True)

    # ── tabs[2] 独立的对比模式控件 ──
    _cmp_row1, _cmp_row2 = st.columns([2, 2])
    with _cmp_row1:
        comp_mode = st.radio(
            '对比模式',
            ['本期 vs 上期(环比)', '本期 vs 去年同期(同比)', '自定义时间段对比'],
            horizontal=True, key='tab2_comp_mode')
    with _cmp_row2:
        if comp_mode == '自定义时间段对比':
            _cmp_c1, _cmp_c2 = st.columns(2)
            with _cmp_c1:
                cmp_start = st.date_input('对比期 开始', value=start - datetime.timedelta(days=30), key='tab2_cmp_start')
            with _cmp_c2:
                cmp_end = st.date_input('对比期 结束', value=end - datetime.timedelta(days=30), key='tab2_cmp_end')

    # 计算 tabs[2] 的对比期
    _tab2_cur_days = (end - start).days + 1
    _tab2_b_end = start - datetime.timedelta(days=1)
    _tab2_b_start = _tab2_b_end - datetime.timedelta(days=_tab2_cur_days - 1)

    if comp_mode == '本期 vs 上期(环比)':
        _t2_prev_s = str(_tab2_b_start)
        _t2_prev_e = str(_tab2_b_end)
        _t2_label_a = f'本期 {today_s} ~ {today_e}'
        _t2_label_b = f'上期 {_t2_prev_s} ~ {_t2_prev_e}'
    elif comp_mode == '本期 vs 去年同期(同比)':
        try:
            _t2_y_start = start.replace(year=start.year - 1)
        except ValueError:
            _t2_y_start = start.replace(year=start.year - 1, day=28)
        try:
            _t2_y_end = end.replace(year=end.year - 1)
        except ValueError:
            _t2_y_end = end.replace(year=end.year - 1, day=28)
        _t2_prev_s = str(_t2_y_start)
        _t2_prev_e = str(_t2_y_end)
        _t2_label_a = f'本期 {today_s} ~ {today_e}'
        _t2_label_b = f'去年同期 {_t2_prev_s} ~ {_t2_prev_e}'
    else:
        _t2_prev_s = str(cmp_start)
        _t2_prev_e = str(cmp_end)
        _t2_label_a = f'A期 {today_s} ~ {today_e}'
        _t2_label_b = f'B期 {_t2_prev_s} ~ {_t2_prev_e}'

    st.info(f'当前对比模式：**{comp_mode}** | 本期：{_t2_label_a} | 对比期：{_t2_label_b}', icon='📊')

    def calc_period_summary(s0, e0):
        rows = []
        for r in data['daily']:
            d = r.get('日期', '')
            if len(d) == 7:
                d = d + '-01'
            if not d or not (s0 <= d <= e0):
                continue
            if channel and r.get('渠道') not in channel:
                continue
            if store and r.get('店铺') not in store:
                continue
            if category and r.get('品类') not in category:
                continue
            if model and r.get('型号') not in model:
                continue
            rows.append(r)
        return summarize(rows)

    cur_sum = calc_period_summary(today_s, today_e)
    prev_sum = calc_period_summary(_t2_prev_s, _t2_prev_e)

    st.markdown('---')
    st.markdown('<div class="section-title" style="border-left:4px solid #1d4ed8;padding-left:12px;">📊 销售对比分析</div>', unsafe_allow_html=True)

    comp_kpis = [
        ('支付金额', '支付金额', '¥', False),
        ('支付件数', '支付件数', '', False),
        ('支付买家', '支付买家数', '', False),
        ('访客数', '商品访客数', '', False),
        ('转化率', '支付转化率', '', True),
        ('客单价', '客单价', '¥', False),
        ('加购率', '加购率', '', True),
    ]
    kpi_cols = st.columns(7)
    for idx, (k_name, k_key, prefix, is_pct) in enumerate(comp_kpis):
        cur_v = cur_sum.get(k_key, 0)
        prev_v = prev_sum.get(k_key, 0)
        delta_v = (cur_v - prev_v) / prev_v if prev_v else None
        if is_pct:
            cur_str = f"{cur_v*100:.2f}%"
            prev_str = f"{prev_v*100:.2f}%"
            diff_pp = (cur_v - prev_v) * 100
            sign = '+' if diff_pp >= 0 else ''
            cls = 'delta-up' if diff_pp >= 0 else 'delta-down'
            delta_label = f'<span class="{cls}">{sign}{diff_pp:.2f}pp</span>'
        else:
            cur_str = f"{prefix}{cur_v:,.0f}"
            prev_str = f"{prefix}{prev_v:,.0f}"
            delta_label = delta_badge(delta_v)
        with kpi_cols[idx]:
            st.markdown(
                f'<p style="font-weight:800;color:#1d4ed8;font-size:13px;margin:0 0 6px 0;text-align:center;">{k_name}</p>'
                f'<div class="comp-card" style="padding:10px;"><div class="comp-period">{_t2_label_a[:16]}</div><div class="comp-value" style="font-size:18px;">{cur_str}</div></div>'
                f'<div class="comp-card" style="padding:10px;"><div class="comp-period">{_t2_label_b[:16]}</div><div class="comp-value" style="font-size:18px;color:#64748b;">{prev_str}</div></div>'
                f'<div class="comp-card" style="padding:10px;background:#f0f9ff;"><div class="comp-period">变化率</div><div style="font-size:16px;font-weight:700;">{delta_label}</div></div>',
                unsafe_allow_html=True
            )

    st.markdown('---')
    st.markdown('<div class="section-title">指标变化详情</div>', unsafe_allow_html=True)
    compare_rows = []
    for k_name, k_key, prefix, is_pct in comp_kpis:
        cur_v = cur_sum.get(k_key, 0)
        prev_v = prev_sum.get(k_key, 0)
        chg = (cur_v - prev_v) / prev_v if prev_v else None
        diff = cur_v - prev_v
        if is_pct:
            cur_str = f"{cur_v*100:.2f}%"
            prev_str = f"{prev_v*100:.2f}%"
            diff_str = f"{diff*100:+.2f}pp"
        else:
            cur_str = f"{prefix}{cur_v:,.0f}"
            prev_str = f"{prefix}{prev_v:,.0f}"
            diff_str = f"{prefix}{diff:+,.0f}"
        compare_rows.append({
            '指标': k_name, '本期数值': cur_str, '对比期数值': prev_str,
            '变化量': diff_str, '变化率(%)': f'{chg*100:+.1f}%' if chg is not None else '--'
        })
    _comp_headers = ['指标', '本期数值', '对比期数值', '变化量', '变化率(%)']
    _render_html_table(compare_rows, _comp_headers, _comp_headers, title='📊 销售指标变化详情')

    # ═══════════════════════════════════════════════════════════════
    # 推广对比分析
    # ═══════════════════════════════════════════════════════════════
    if promo_rows:
        st.markdown('---')
        st.markdown('<div class="section-title" style="border-left:4px solid #f59e0b;padding-left:12px;">📢 推广对比分析</div>', unsafe_allow_html=True)

        def calc_promo_period(s0, e0):
            rows = []
            for r in promo_rows:
                d = r.get('_date', '')
                if not d or not (s0 <= d <= e0):
                    continue
                if channel and r.get('_渠道', '') not in channel:
                    continue
                if store and r.get('_店铺', '') not in store:
                    continue
                if category and r.get('_品类', '') not in category:
                    continue
                if model and r.get('_型号', '') not in model:
                    continue
                rows.append(r)
            return rows

        promo_cur_rows = calc_promo_period(today_s, today_e)
        promo_prev_rows = calc_promo_period(_t2_prev_s, _t2_prev_e)

        def promo_sum(rows):
            s = {}
            for r in rows:
                s['_花费'] = s.get('_花费', 0) + r.get('_花费', 0)
                s['_展现数'] = s.get('_展现数', 0) + r.get('_展现数', 0)
                s['_点击数'] = s.get('_点击数', 0) + r.get('_点击数', 0)
                s['_总订单金额'] = s.get('_总订单金额', 0) + r.get('_总订单金额', 0)
                s['_直接订单金额'] = s.get('_直接订单金额', 0) + r.get('_直接订单金额', 0)
                s['_总加购数'] = s.get('_总加购数', 0) + r.get('_总加购数', 0)
            return s

        promo_cur = promo_sum(promo_cur_rows)
        promo_prev = promo_sum(promo_prev_rows)

        promo_kpis = [
            ('推广花费', '_花费', '¥', 10000),
            ('ROI', None, '', 1),
            ('点击率', None, '', 100),
            ('直接成交', '_直接订单金额', '¥', 10000),
            ('加购数', '_总加购数', '', 1),
            ('展现数', '_展现数', '', 1),
            ('点击数', '_点击数', '', 1),
        ]

        promo_cols = st.columns(7)
        for idx, (k_name, k_key, prefix, divisor) in enumerate(promo_kpis):
            if k_key:
                cur_v = promo_cur.get(k_key, 0) / divisor
                prev_v = promo_prev.get(k_key, 0) / divisor
            else:
                # ROI and 点击率 are computed
                if k_name == 'ROI':
                    cur_v = promo_cur.get('_总订单金额', 0) / promo_cur.get('_花费', 1)
                    prev_v = promo_prev.get('_总订单金额', 0) / promo_prev.get('_花费', 1)
                elif k_name == '点击率':
                    cur_v = promo_cur.get('_点击数', 0) / promo_cur.get('_展现数', 1) * 100
                    prev_v = promo_prev.get('_点击数', 0) / promo_prev.get('_展现数', 1) * 100
                else:
                    cur_v = prev_v = 0
            delta_v = (cur_v - prev_v) / prev_v if prev_v else None

            if k_name == '点击率':
                cur_str = f"{cur_v:.2f}%"
                prev_str = f"{prev_v:.2f}%"
                diff_pp = cur_v - prev_v
                sign = '+' if diff_pp >= 0 else ''
                cls = 'delta-up' if diff_pp >= 0 else 'delta-down'
                delta_label = f'<span class="{cls}">{sign}{diff_pp:.2f}pp</span>'
            elif k_name == 'ROI':
                cur_str = f"{cur_v:.2f}"
                prev_str = f"{prev_v:.2f}"
                delta_label = delta_badge(delta_v)
            elif divisor >= 10000:
                cur_str = f"{prefix}{cur_v:.1f}万"
                prev_str = f"{prefix}{prev_v:.1f}万"
                delta_label = delta_badge(delta_v)
            else:
                cur_str = f"{prefix}{cur_v:,.0f}"
                prev_str = f"{prefix}{prev_v:,.0f}"
                delta_label = delta_badge(delta_v)

            with promo_cols[idx]:
                st.markdown(
                    f'<p style="font-weight:800;color:#f59e0b;font-size:13px;margin:0 0 6px 0;text-align:center;">{k_name}</p>'
                    f'<div class="comp-card" style="padding:10px;"><div class="comp-period">{_t2_label_a[:16]}</div><div class="comp-value" style="font-size:18px;">{cur_str}</div></div>'
                    f'<div class="comp-card" style="padding:10px;"><div class="comp-period">{_t2_label_b[:16]}</div><div class="comp-value" style="font-size:18px;color:#64748b;">{prev_str}</div></div>'
                    f'<div class="comp-card" style="padding:10px;background:#fffbeb;"><div class="comp-period">变化率</div><div style="font-size:16px;font-weight:700;">{delta_label}</div></div>',
                    unsafe_allow_html=True
                )

        # 推广对比详情表
        promo_compare_rows = []
        for k_name, k_key, prefix, divisor in promo_kpis:
            if k_key:
                cur_v = promo_cur.get(k_key, 0) / divisor
                prev_v = promo_prev.get(k_key, 0) / divisor
            else:
                if k_name == 'ROI':
                    cur_v = promo_cur.get('_总订单金额', 0) / promo_cur.get('_花费', 1)
                    prev_v = promo_prev.get('_总订单金额', 0) / promo_prev.get('_花费', 1)
                elif k_name == '点击率':
                    cur_v = promo_cur.get('_点击数', 0) / promo_cur.get('_展现数', 1) * 100
                    prev_v = promo_prev.get('_点击数', 0) / promo_prev.get('_展现数', 1) * 100
                else:
                    cur_v = prev_v = 0
            chg = (cur_v - prev_v) / prev_v if prev_v else None
            diff = cur_v - prev_v

            if k_name == 'ROI':
                cur_str = f"{cur_v:.2f}"
                prev_str = f"{prev_v:.2f}"
                diff_str = f"{diff:+.2f}"
            elif k_name == '点击率':
                cur_str = f"{cur_v:.2f}%"
                prev_str = f"{prev_v:.2f}%"
                diff_str = f"{diff:+.2f}pp"
            elif divisor >= 10000:
                cur_str = f"{prefix}{cur_v:.1f}万"
                prev_str = f"{prefix}{prev_v:.1f}万"
                diff_str = f"{prefix}{diff:+.1f}万"
            else:
                cur_str = f"{prefix}{cur_v:,.0f}"
                prev_str = f"{prefix}{prev_v:,.0f}"
                diff_str = f"{prefix}{diff:+,.0f}"
            promo_compare_rows.append({
                '指标': k_name, '本期数值': cur_str, '对比期数值': prev_str,
                '变化量': diff_str, '变化率(%)': f'{chg*100:+.1f}%' if chg is not None else '--'
            })
        _pcomp_headers = ['指标', '本期数值', '对比期数值', '变化量', '变化率(%)']
        _render_html_table(promo_compare_rows, _pcomp_headers, _pcomp_headers, title='📢 推广指标变化详情')

    st.markdown('---')
    p1, p2 = st.columns(2)
    key_map = {'支付金额': '支付金额', '访客数': '商品访客数', '支付件数': '支付件数', '支付买家': '支付买家数'}
    chart_data = [{'指标': k, '本期': cur_sum.get(v, 0), '对比期': prev_sum.get(v, 0)} for k, v in key_map.items()]
    with p1:
        fig = px.bar(chart_data, x='指标', y=['本期', '对比期'], barmode='group',
                     color_discrete_sequence=['#1d4ed8', '#f59e0b'])
        fig.update_layout(height=350, template='plotly_white', title='核心指标对比', legend_title='时间段')
        st.plotly_chart(fig, width="stretch")
        _render_download_panel(chart_data, ['指标', '本期', '对比期'], 'period_core_compare.csv', '📥 核心指标对比')
    with p2:
        ch_data = []
        for k, v in key_map.items():
            cur_v = cur_sum.get(v, 0)
            prev_v = prev_sum.get(v, 0)
            chg = (cur_v - prev_v) / prev_v if prev_v else 0
            ch_data.append({'指标': k, '变化率': chg})
        colors = ['#22c55e' if x['变化率'] >= 0 else '#ef4444' for x in ch_data]
        fig = go.Figure(go.Bar(x=[x['指标'] for x in ch_data], y=[x['变化率'] for x in ch_data],
                                marker_color=colors))
        fig.update_layout(height=350, template='plotly_white', title='各指标变化率', yaxis_tickformat='.1%')
        st.plotly_chart(fig, width="stretch")
        _render_download_panel(ch_data, ['指标', '变化率'], 'period_change_rate.csv', '📥 各指标变化率')

    st.markdown('---')
    # ── 销售维度对比分析（受全局维度筛选影响，日期由时间段对比独立控制）──
    st.markdown('<div class="section-title">销售数据对比分析</div>', unsafe_allow_html=True)

    _dim_options = {'渠道': '渠道', '店铺': '店铺', '品类': '品类', '型号': '型号'}
    _dim_label = st.radio('对比维度', list(_dim_options.keys()), horizontal=True, key='dim_compare')
    _dim_field = _dim_options[_dim_label]

    # 只按渠道/店铺/品类/型号过滤，不过滤日期，让get_period_rows独立控制时间段
    _raw = data['daily']
    if channel or store or category or model:
        _filtered = []
        for r in _raw:
            if channel and r.get('渠道') not in channel:
                continue
            if store and r.get('店铺') not in store:
                continue
            if category and r.get('品类') not in category:
                continue
            if model and r.get('型号') not in model:
                continue
            _filtered.append(r)
        _raw = _filtered

    cur_dim = group(get_period_rows(_raw, today_s, today_e), _dim_field)
    prev_dim = group(get_period_rows(_raw, _t2_prev_s, _t2_prev_e), _dim_field)
    prev_dim_map = {r[_dim_field]: r for r in prev_dim}

    _metric_defs = [
        ('销售额',  '支付金额',   lambda v: f"¥{v:,.0f}",    '#fef3c7'),  # amber
        ('销售件数','支付件数',   lambda v: f"{v:,.0f}",     '#dbeafe'),  # blue
        ('访客数',  '商品访客数', lambda v: f"{v:,.0f}",     '#dcfce7'),  # green
        ('转化率',  '支付转化率', lambda v: f"{v*100:.2f}%", '#f3e8ff'),  # purple
        ('客单价',  '客单价',     lambda v: f"¥{v:,.2f}",    '#ccfbf1'),  # teal
    ]

    # 计算本期总销售额（用于占比）
    _cur_total_amt = sum(r.get('支付金额', 0) or 0 for r in cur_dim)

    dim_compare = []
    for r in cur_dim:
        name = r[_dim_field]
        prev_r = prev_dim_map.get(name, {})
        _amt = r.get('支付金额', 0) or 0
        _share = _amt / _cur_total_amt if _cur_total_amt else 0
        row = {
            _dim_label: name,
            '销售额占比': f"{_share*100:.1f}%",
        }
        for _ml, _mf, _fmt, _color in _metric_defs:
            cur_v = r.get(_mf, 0) or 0
            prev_v = prev_r.get(_mf, 0) or 0
            _cur_s = _fmt(cur_v)
            _prev_s = _fmt(prev_v)
            chg = (cur_v - prev_v) / prev_v if prev_v else None
            if chg is None:
                chg_txt = '--'
                chg_color = '#94a3b8'
            else:
                chg_txt = f"{'+' if chg >= 0 else ''}{chg*100:.1f}%"
                chg_color = '#22c55e' if chg >= 0 else '#dc2626'
            row[f'{_ml}(本期)'] = _cur_s
            row[f'{_ml}(对比期)'] = _prev_s
            row[f'{_ml}变化率'] = f"<span style='color:{chg_color}'>{chg_txt}</span>"
        dim_compare.append(row)

    # 合计行
    if dim_compare:
        _sum_row = {_dim_label: '<b>合计</b>', '销售额占比': '100%'}
        # 先汇总原始字段总量（避免比率字段直接求和出错）
        _tot_cur_amt = sum(r.get('支付金额', 0) or 0 for r in cur_dim)
        _tot_prev_amt = sum(r.get('支付金额', 0) or 0 for r in prev_dim)
        _tot_cur_vis = sum(r.get('商品访客数', 0) or 0 for r in cur_dim)
        _tot_prev_vis = sum(r.get('商品访客数', 0) or 0 for r in prev_dim)
        _tot_cur_buyers = sum(r.get('支付买家数', 0) or 0 for r in cur_dim)
        _tot_prev_buyers = sum(r.get('支付买家数', 0) or 0 for r in prev_dim)
        _tot_cur_qty = sum(r.get('支付件数', 0) or 0 for r in cur_dim)
        _tot_prev_qty = sum(r.get('支付件数', 0) or 0 for r in prev_dim)

        for _ml, _mf, _fmt, _color in _metric_defs:
            if _ml == '销售额':
                _cur_sum = _tot_cur_amt
                _prev_sum = _tot_prev_amt
            elif _ml == '销售件数':
                _cur_sum = _tot_cur_qty
                _prev_sum = _tot_prev_qty
            elif _ml == '访客数':
                _cur_sum = _tot_cur_vis
                _prev_sum = _tot_prev_vis
            elif _ml == '转化率':
                _cur_sum = _tot_cur_buyers / _tot_cur_vis if _tot_cur_vis else 0
                _prev_sum = _tot_prev_buyers / _tot_prev_vis if _tot_prev_vis else 0
            elif _ml == '客单价':
                _cur_sum = _tot_cur_amt / _tot_cur_buyers if _tot_cur_buyers else 0
                _prev_sum = _tot_prev_amt / _tot_prev_buyers if _tot_prev_buyers else 0
            else:
                _cur_sum = sum(r.get(_mf, 0) or 0 for r in cur_dim)
                _prev_sum = sum(r.get(_mf, 0) or 0 for r in prev_dim)
            _sum_row[f'{_ml}(本期)'] = _fmt(_cur_sum)
            _sum_row[f'{_ml}(对比期)'] = _fmt(_prev_sum)
            chg = (_cur_sum - _prev_sum) / _prev_sum if _prev_sum else None
            if chg is None:
                chg_txt = '--'
                chg_color = '#94a3b8'
            else:
                chg_txt = f"{'+' if chg >= 0 else ''}{chg*100:.1f}%"
                chg_color = '#22c55e' if chg >= 0 else '#dc2626'
            _sum_row[f'{_ml}变化率'] = f"<span style='color:{chg_color}'>{chg_txt}</span>"
        dim_compare.append(_sum_row)

    if dim_compare:
        _dim_cols = list(dim_compare[0].keys())
        # 构建带色块分组的自定义 HTML 表格
        _html = '<div class="styled-table-wrap" style="max-height:400px;overflow-y:auto;"><table class="styled-table"><thead>'
        # 表头第1行：指标分组
        _html += '<tr>'
        _html += f'<th colspan="2" style="background:#e2e8f0;color:#1e293b;text-align:center;font-size:12px;font-weight:600;">维度信息</th>'
        for _ml, _mf, _fmt, _color in _metric_defs:
            _html += f'<th colspan="3" style="background:{_color};color:#1e293b;text-align:center;font-size:12px;font-weight:600;border-left:2px solid #fff;">{_ml}</th>'
        _html += '</tr>'
        # 表头第2行：具体列名
        _html += '<tr>'
        _html += f'<th style="min-width:110px;background:#e2e8f0;color:#1e293b;font-weight:600;">{_dim_label}</th>'
        _html += '<th style="min-width:72px;background:#e0f2fe;color:#1e293b;font-weight:600;">占比</th>'
        for _ml, _mf, _fmt, _color in _metric_defs:
            _html += f'<th style="min-width:80px;background:{_color};color:#1e293b;font-weight:600;">本期</th>'
            _html += f'<th style="min-width:80px;background:{_color};color:#1e293b;font-weight:600;">对比期</th>'
            _html += f'<th style="min-width:72px;background:{_color};color:#1e293b;font-weight:600;">变化率</th>'
        _html += '</tr></thead><tbody>'
        # 数据行
        _total_row = len(dim_compare) - 1
        for i, r in enumerate(dim_compare):
            is_total = (i == _total_row)
            bg = '#fff7ed' if is_total else ('#fafafa' if i % 2 == 0 else 'white')
            fw = 'bold' if is_total else 'normal'
            _html += f'<tr style="background:{bg};font-weight:{fw};">'
            for j, c in enumerate(_dim_cols):
                val = r.get(c, '')
                align = '' if j <= 1 else 'text-align:right;'
                _html += f'<td style="{align}">{val}</td>'
            _html += '</tr>'
        _html += '</tbody></table></div>'
        st.markdown(_wrap_fullscreen(_html, title='📊 销售数据对比分析')[0], unsafe_allow_html=True)

    _dim_cols_safe = _dim_cols if dim_compare else ['维度', '占比']
    _render_download_panel(dim_compare if dim_compare else [], _dim_cols_safe, f'dimension_compare_{_dim_label}.csv', '📥 销售维度对比')
    _render_download_panel(compare_rows, ['指标', '本期数值', '对比期数值', '变化量', '变化率(%)'], 'period_comparison.csv', '📥 时间段对比')

    # ── 推广数据对比分析 ──
    st.markdown('<div class="section-title">推广数据对比分析</div>', unsafe_allow_html=True)

    if not promo_rows:
        st.info('请上传推广数据文件以启用推广对比分析。')
    else:
        _p_cmp_dim_options = {
            '渠道': '_渠道', '店铺': '_店铺', '品类': '_品类', '型号': '_型号',
            '产品线': '产品线', '营销场景': '_营销场景',
        }
        _p_cmp_dim_label = st.radio('对比维度', list(_p_cmp_dim_options.keys()), horizontal=True, key='p_dim_compare')
        _p_cmp_dim_field = _p_cmp_dim_options[_p_cmp_dim_label]

        # 只按渠道/店铺/品类/型号全局过滤，不过滤日期
        _p_cmp_raw = promo_rows
        if channel or store or category or model:
            _p_cmp_filtered = []
            for r in _p_cmp_raw:
                if channel and r.get('_渠道') not in channel: continue
                if store and r.get('_店铺') not in store: continue
                if category and r.get('_品类') not in category: continue
                if model and r.get('_型号') not in model: continue
                _p_cmp_filtered.append(r)
            _p_cmp_raw = _p_cmp_filtered

        # ── 销售支付金额预计算（按推广维度映射到销售维度，用于费率=花费/销售支付金额）──
        _sales_dim_map_pc = {'_渠道': '渠道', '_店铺': '店铺', '_品类': '品类', '_型号': '型号'}
        _sales_field_pc = _sales_dim_map_pc.get(_p_cmp_dim_field, None)
        _p_sales_by_dim_cur = {}
        _p_sales_by_dim_prev = {}
        _p_sales_total_cur = 0
        _p_sales_total_prev = 0
        # 注意：必须用 data['daily']（原始全量数据），不能用 daily（已被全局日期范围 start~end 过滤）
        _raw_daily = data['daily']
        if _sales_field_pc:
            for r in _raw_daily:
                d = r.get('日期', '')
                if len(d) == 7: d = d + '-01'
                if not d: continue
                if channel and r.get('渠道', '') not in channel: continue
                if store and r.get('店铺', '') not in store: continue
                if category and r.get('品类', '') not in category: continue
                if model and r.get('型号', '') not in model: continue
                amt = float(r.get('支付金额', 0) or 0)
                dv = r.get(_sales_field_pc, '') or '未标注'
                if today_s <= d <= today_e:
                    _p_sales_by_dim_cur[dv] = _p_sales_by_dim_cur.get(dv, 0) + amt
                    _p_sales_total_cur += amt
                if _t2_prev_s <= d <= _t2_prev_e:
                    _p_sales_by_dim_prev[dv] = _p_sales_by_dim_prev.get(dv, 0) + amt
                    _p_sales_total_prev += amt
        else:
            # 产品线/营销场景没有对应销售维度，用销售总额
            for r in _raw_daily:
                d = r.get('日期', '')
                if len(d) == 7: d = d + '-01'
                if not d: continue
                if channel and r.get('渠道', '') not in channel: continue
                if store and r.get('店铺', '') not in store: continue
                if category and r.get('品类', '') not in category: continue
                if model and r.get('型号', '') not in model: continue
                amt = float(r.get('支付金额', 0) or 0)
                if today_s <= d <= today_e:
                    _p_sales_total_cur += amt
                if _t2_prev_s <= d <= _t2_prev_e:
                    _p_sales_total_prev += amt

        def _p_group(rows, field):
            d = {}
            for r in rows:
                k = r.get(field) or '未标注'
                if k not in d:
                    d[k] = {'_花费': 0, '_展现数': 0, '_点击数': 0, '_直接订单金额': 0, '_总订单金额': 0, '_总成交订单量': 0, '_直接订单量': 0}
                for m in ('_花费', '_展现数', '_点击数', '_直接订单金额', '_总订单金额', '_总成交订单量', '_直接订单量'):
                    d[k][m] += float(r.get(m, 0) or 0)
            out = []
            for k, v in d.items():
                v[field] = k
                v['_cpc'] = v['_花费'] / v['_点击数'] if v['_点击数'] else 0
                v['_ctr'] = v['_点击数'] / v['_展现数'] if v['_展现数'] else 0
                v['_direct_roi'] = v['_直接订单金额'] / v['_花费'] if v['_花费'] else 0
                v['_total_roi'] = v['_总订单金额'] / v['_花费'] if v['_花费'] else 0
                v['_direct_tcvr'] = v['_直接订单量'] / v['_点击数'] if v['_点击数'] else 0
                v['_total_tcvr'] = v['_总成交订单量'] / v['_点击数'] if v['_点击数'] else 0
                # _fee_rate 不在此预计算，改为在渲染时用销售支付金额动态计算
                out.append(v)
            return sorted(out, key=lambda x: x['_花费'], reverse=True)

        _p_cur_dim = _p_group(get_period_rows(_p_cmp_raw, today_s, today_e, '_date'), _p_cmp_dim_field)
        _p_prev_dim = _p_group(get_period_rows(_p_cmp_raw, _t2_prev_s, _t2_prev_e, '_date'), _p_cmp_dim_field)
        _p_prev_map = {r[_p_cmp_dim_field]: r for r in _p_prev_dim}

        _p_total_spend = sum(r['_花费'] for r in _p_cur_dim) or 1

        _p_metric_defs = [
            ('花费',        '_花费',         lambda v: f'¥{v:,.0f}',    '#fef3c7'),
            ('费率',        '_fee_rate',     lambda v: f'{v*100:.2f}%' if v is not None else '--', '#fee2e2'),
            ('CPC',         '_cpc',          lambda v: f'¥{v:.2f}',     '#fef3c7'),
            ('直接ROI',     '_direct_roi',   lambda v: f'{v:.2f}',      '#dcfce7'),
            ('总ROI',       '_total_roi',    lambda v: f'{v:.2f}',      '#ccfbf1'),
            ('点击量',      '_点击数',       lambda v: f'{v:,.0f}',      '#dbeafe'),
            ('点击率',      '_ctr',          lambda v: f'{v*100:.2f}%', '#e0e7ff'),
            ('直接转化率',  '_direct_tcvr',  lambda v: f'{v*100:.2f}%', '#f3e8ff'),
            ('总转化率',    '_total_tcvr',   lambda v: f'{v*100:.2f}%', '#fce7f3'),
            ('直接成交金额', '_直接订单金额', lambda v: f'¥{v:,.0f}',    '#dcfce7'),
            ('总成交金额',  '_总订单金额',   lambda v: f'¥{v:,.0f}',    '#ccfbf1'),
        ]

        _p_cmp_tbl = []
        for r in _p_cur_dim:
            name = r[_p_cmp_dim_field]
            prev_r = _p_prev_map.get(name, {})
            _spend = r.get('_花费', 0) or 0
            _share = _spend / _p_total_spend if _p_total_spend else 0
            # 按维度计算费率（花费/销售支付金额）
            # 只有对应维度在销售数据中有记录时才计算费率，否则显示 --
            if _sales_field_pc:
                _s_amt_c = _p_sales_by_dim_cur.get(name, 0)
                _s_amt_p = _p_sales_by_dim_prev.get(name, 0)
            else:
                _s_amt_c = _p_sales_total_cur
                _s_amt_p = _p_sales_total_prev
            _fee_c = _spend / _s_amt_c if _s_amt_c else None
            # 对比期费率：只有对比期有花费且销售数据存在时才计算
            _prev_spend = prev_r.get('_花费', 0) if prev_r else 0
            _fee_p = _prev_spend / _s_amt_p if _prev_spend and _s_amt_p else None
            r['_fee_rate'] = _fee_c  # 临时赋值供后续读取
            # 同时给 prev_r 也赋上费率，供统一循环读取
            prev_r['_fee_rate'] = _fee_p
            row = {
                _p_cmp_dim_label: name,
                '花费占比': f'{_share*100:.1f}%',
            }
            for _ml, _mf, _fmt, _color in _p_metric_defs:
                cur_v = r.get(_mf, 0)
                # 费率可能为 None（分母为0），保持 None 供 _fmt 正确显示 --
                if _mf == '_fee_rate':
                    cur_v = cur_v if cur_v is not None else None  # preserve None
                else:
                    cur_v = cur_v or 0
                _pv = prev_r.get(_mf, 0)
                # 对费率特殊处理：如果值为 None，保持 None 让 _fmt 显示 --
                if _mf == '_fee_rate':
                    prev_v = _pv  # None stays None
                else:
                    prev_v = _pv if _pv is not None else 0
                _cur_s = _fmt(cur_v)
                _prev_s = _fmt(prev_v)
                chg = (cur_v - prev_v) / prev_v if prev_v else None
                if chg is None:
                    chg_txt, chg_color = '--', '#94a3b8'
                else:
                    chg_txt = f"{'+' if chg >= 0 else ''}{chg*100:.1f}%"
                    # 花费、CPC、费率：涨是红（坏），跌是绿（好）；其余：涨是绿，跌是红
                    if _ml in ('花费', 'CPC', '费率'):
                        chg_color = '#dc2626' if chg >= 0 else '#22c55e'
                    else:
                        chg_color = '#22c55e' if chg >= 0 else '#dc2626'
                row[f'{_ml}(本期)'] = _cur_s
                row[f'{_ml}(对比期)'] = _prev_s
                row[f'{_ml}变化率'] = f"<span style='color:{chg_color}'>{chg_txt}</span>"
            _p_cmp_tbl.append(row)

        # 合计行
        if _p_cmp_tbl:
            _p_tot_spend_c = sum(r.get('_花费', 0) or 0 for r in _p_cur_dim)
            _p_tot_spend_p = sum(r.get('_花费', 0) or 0 for r in _p_prev_dim)
            _p_tot_clicks_c = sum(r.get('_点击数', 0) or 0 for r in _p_cur_dim)
            _p_tot_clicks_p = sum(r.get('_点击数', 0) or 0 for r in _p_prev_dim)
            _p_tot_impress_c = sum(r.get('_展现数', 0) or 0 for r in _p_cur_dim)
            _p_tot_impress_p = sum(r.get('_展现数', 0) or 0 for r in _p_prev_dim)
            _p_tot_direct_c = sum(r.get('_直接订单金额', 0) or 0 for r in _p_cur_dim)
            _p_tot_direct_p = sum(r.get('_直接订单金额', 0) or 0 for r in _p_prev_dim)
            _p_tot_total_c = sum(r.get('_总订单金额', 0) or 0 for r in _p_cur_dim)
            _p_tot_total_p = sum(r.get('_总订单金额', 0) or 0 for r in _p_prev_dim)
            _p_tot_orders_c = sum(r.get('_总成交订单量', 0) or 0 for r in _p_cur_dim)
            _p_tot_orders_p = sum(r.get('_总成交订单量', 0) or 0 for r in _p_prev_dim)
            _p_tot_direct_orders_c = sum(r.get('_直接订单量', 0) or 0 for r in _p_cur_dim)
            _p_tot_direct_orders_p = sum(r.get('_直接订单量', 0) or 0 for r in _p_prev_dim)
            _p_sum_row = {_p_cmp_dim_label: '<b>合计</b>', '花费占比': '100%'}
            # 合计行费率 = 总花费 / 总销售支付金额
            _p_fee_c_total = _p_tot_spend_c / _p_sales_total_cur if _p_sales_total_cur else None
            _p_fee_p_total = _p_tot_spend_p / _p_sales_total_prev if _p_sales_total_prev else None
            _p_tot_map = {
                '花费':        (_p_tot_spend_c, _p_tot_spend_p),
                '费率':        (_p_fee_c_total, _p_fee_p_total),
                'CPC':         (_p_tot_spend_c / _p_tot_clicks_c if _p_tot_clicks_c else 0,
                                _p_tot_spend_p / _p_tot_clicks_p if _p_tot_clicks_p else 0),
                '直接ROI':     (_p_tot_direct_c / _p_tot_spend_c if _p_tot_spend_c else 0,
                                _p_tot_direct_p / _p_tot_spend_p if _p_tot_spend_p else 0),
                '总ROI':       (_p_tot_total_c / _p_tot_spend_c if _p_tot_spend_c else 0,
                                _p_tot_total_p / _p_tot_spend_p if _p_tot_spend_p else 0),
                '点击量':      (_p_tot_clicks_c, _p_tot_clicks_p),
                '点击率':      (_p_tot_clicks_c / _p_tot_impress_c if _p_tot_impress_c else 0,
                                _p_tot_clicks_p / _p_tot_impress_p if _p_tot_impress_p else 0),
                '直接转化率':  (_p_tot_direct_orders_c / _p_tot_clicks_c if _p_tot_clicks_c else 0,
                                _p_tot_direct_orders_p / _p_tot_clicks_p if _p_tot_clicks_p else 0),
                '总转化率':    (_p_tot_orders_c / _p_tot_clicks_c if _p_tot_clicks_c else 0,
                                _p_tot_orders_p / _p_tot_clicks_p if _p_tot_clicks_p else 0),
                '直接成交金额': (_p_tot_direct_c, _p_tot_direct_p),
                '总成交金额':  (_p_tot_total_c, _p_tot_total_p),
            }
            for _ml, _mf, _fmt, _color in _p_metric_defs:
                _cv, _pv = _p_tot_map[_ml]
                _p_sum_row[f'{_ml}(本期)'] = _fmt(_cv)
                _p_sum_row[f'{_ml}(对比期)'] = _fmt(_pv)
                chg = (_cv - _pv) / _pv if _pv else None
                if chg is None:
                    chg_txt, chg_color = '--', '#94a3b8'
                else:
                    chg_txt = f"{'+' if chg >= 0 else ''}{chg*100:.1f}%"
                    if _ml in ('花费', 'CPC', '费率'):
                        chg_color = '#dc2626' if chg >= 0 else '#22c55e'
                    else:
                        chg_color = '#22c55e' if chg >= 0 else '#dc2626'
                _p_sum_row[f'{_ml}变化率'] = f"<span style='color:{chg_color}'>{chg_txt}</span>"
            _p_cmp_tbl.append(_p_sum_row)

        if _p_cmp_tbl:
            _p_cmp_cols = list(_p_cmp_tbl[0].keys())
            _p_html = '<div class="styled-table-wrap" style="max-height:400px;overflow-y:auto;"><table class="styled-table"><thead>'
            _p_html += '<tr>'
            _p_html += '<th colspan="2" style="background:#e2e8f0;color:#1e293b;text-align:center;font-size:12px;font-weight:600;">维度信息</th>'
            for _ml, _mf, _fmt, _color in _p_metric_defs:
                _p_html += f'<th colspan="3" style="background:{_color};color:#1e293b;text-align:center;font-size:12px;font-weight:600;border-left:2px solid #fff;">{_ml}</th>'
            _p_html += '</tr>'
            _p_html += '<tr>'
            _p_html += f'<th style="min-width:110px;background:#e2e8f0;color:#1e293b;font-weight:600;">{_p_cmp_dim_label}</th>'
            _p_html += '<th style="min-width:72px;background:#e0f2fe;color:#1e293b;font-weight:600;">花费占比</th>'
            for _ml, _mf, _fmt, _color in _p_metric_defs:
                _p_html += f'<th style="min-width:80px;background:{_color};color:#1e293b;font-weight:600;">本期</th>'
                _p_html += f'<th style="min-width:80px;background:{_color};color:#1e293b;font-weight:600;">对比期</th>'
                _p_html += f'<th style="min-width:72px;background:{_color};color:#1e293b;font-weight:600;">变化率</th>'
            _p_html += '</tr></thead><tbody>'
            _p_total_row_idx = len(_p_cmp_tbl) - 1
            for i, r in enumerate(_p_cmp_tbl):
                is_total = (i == _p_total_row_idx)
                bg = '#fff7ed' if is_total else ('#fafafa' if i % 2 == 0 else 'white')
                fw = 'bold' if is_total else 'normal'
                _p_html += f'<tr style="background:{bg};font-weight:{fw};">'
                for j, c in enumerate(_p_cmp_cols):
                    val = r.get(c, '')
                    align = '' if j <= 1 else 'text-align:right;'
                    _p_html += f'<td style="{align}">{val}</td>'
                _p_html += '</tr>'
            _p_html += '</tbody></table></div>'
            st.markdown(_wrap_fullscreen(_p_html, title='📢 推广数据对比分析')[0], unsafe_allow_html=True)
            _render_download_panel(_p_cmp_tbl, _p_cmp_cols, f'promo_compare_{_p_cmp_dim_label}.csv', '📥 推广维度对比')

    # ─────────────────────────────────────────────────────────────
# TAB 3: 趋势分析
# ═══════════════════════════════════════════════════════════════
with tabs[3]:
    st.markdown('<div class="section-title">趋势分析</div>', unsafe_allow_html=True)
    st.markdown('### 📦 销售趋势')
    _s_c1, _s_c2, _s_c3 = st.columns([1, 2, 2])
    with _s_c1:
        _sales_dim = st.selectbox('分析维度', ['按日期', '按渠道', '按店铺', '按品类', '按型号'], key='tr_sales_dim')
    with _s_c2:
        gran = st.radio('粒度', ['月度', '周度', '日度'], horizontal=True, key='granularity')
    with _s_c3:
        _s_comp_mode = st.radio('对比模式', ['同比', '上月同期'], horizontal=True, key='s_comp_mode')

    # ── 销售维度字段映射 ──
    _s_dim_field = {'按渠道': '渠道', '按店铺': '店铺', '按品类': '品类', '按型号': '型号'}.get(_sales_dim, '')

    # ── 构建销售按日聚合 (dim_val, date) → {metrics} ──
    _sales_day = {}  # key: (dim_val, date_str)
    for r in daily:
        dt = r.get('日期', '')
        if not dt: continue
        if not (str(start) <= dt <= str(end)): continue
        dv = r.get(_s_dim_field, '') or '未标注' if _s_dim_field else ''
        key = (dv, dt)
        if key not in _sales_day:
            _sales_day[key] = {m: 0.0 for m in METRICS}
        for m in METRICS:
            _sales_day[key][m] += float(r.get(m, 0) or 0)

    # ── 全时段同维度日聚合（用于YoY） ──
    _sales_all_day = {}
    for r in daily_all_filtered:
        dt = r.get('日期', '')
        if not dt: continue
        dv = r.get(_s_dim_field, '') or '未标注' if _s_dim_field else ''
        key = (dv, dt)
        if key not in _sales_all_day:
            _sales_all_day[key] = {m: 0.0 for m in METRICS}
        for m in METRICS:
            _sales_all_day[key][m] += float(r.get(m, 0) or 0)

    # ── 上月同期日期映射（月份减1，天数不变）──
    def _date_mom(dt_str):
        """返回上月同日，如 2026-05-15 → 2026-04-15；跨年也正确处理"""
        try:
            dt_obj = datetime.date.fromisoformat(dt_str[:10])
            y, mo, d = dt_obj.year, dt_obj.month, dt_obj.day
            if mo == 1:
                nm, ny = 12, y - 1
            else:
                nm, ny = mo - 1, y
            import calendar
            max_day = calendar.monthrange(ny, nm)[1]
            return str(datetime.date(ny, nm, min(d, max_day)))
        except Exception:
            return None

    # ── 上月同期 day dict（从 _sales_all_day 取） ──
    _sales_mom_day = {}  # key: (dv, dt_str) → 上月同日对应数据
    for (dv, dt_str), v in _sales_day.items():
        mom_dt = _date_mom(dt_str)
        if mom_dt:
            _sales_mom_day[(dv, dt_str)] = _sales_all_day.get((dv, mom_dt))

    # ── 构建月度聚合 (dim_val, ym) → {metrics} ──
    _sales_month = {}
    for (dv, dt_str), v in _sales_day.items():
        ym = dt_str[:7]
        key = (dv, ym)
        if key not in _sales_month:
            _sales_month[key] = {m: 0.0 for m in METRICS}
        for m in METRICS:
            _sales_month[key][m] += v[m]

    # ── 构建 chart 用的 tr_data ──
    def _mk_label(dv, period):
        return f'{dv} | {period}' if _s_dim_field else period

    if gran == '月度':
        tr_data = [{
            '周期': _mk_label(dv, ym), '支付金额': v['支付金额'], '访客数': v['商品访客数'],
            '支付件数': v['支付件数'],
            '转化率': v['支付买家数'] / v['商品访客数'] if v['商品访客数'] else 0,
            '加购率': v['商品加购人数'] / v['商品访客数'] if v['商品访客数'] else 0,
        } for (dv, ym), v in sorted(_sales_month.items())]
    elif gran == '周度':
        week_dict = {}  # key: (dv, yr, wk)
        for (dv, dt_str), v in _sales_day.items():
            try:
                dt = datetime.date.fromisoformat(dt_str[:10])
                iso = dt.isocalendar()
                wkey = (dv, iso[0], iso[1])
            except Exception:
                continue
            if wkey not in week_dict:
                week_dict[wkey] = {m: 0.0 for m in METRICS}
                week_dict[wkey]['_dates'] = []
            for m in METRICS:
                week_dict[wkey][m] += v[m]
            week_dict[wkey]['_dates'].append(dt)
        tr_data = []
        for (dv, yr, wk), v in sorted(week_dict.items()):
            dates = v['_dates']
            w_start = min(dates).strftime('%m/%d')
            w_end = max(dates).strftime('%m/%d')
            tr_data.append({
                '周期': _mk_label(dv, f'{w_start}-{w_end}'),
                '支付金额': v['支付金额'], '访客数': v['商品访客数'],
                '支付件数': v['支付件数'],
                '转化率': v['支付买家数'] / v['商品访客数'] if v['商品访客数'] else 0,
                '加购率': v['商品加购人数'] / v['商品访客数'] if v['商品访客数'] else 0,
            })
    else:
        tr_data = [{
            '周期': _mk_label(dv, dt_str), '支付金额': v['支付金额'], '访客数': v['商品访客数'],
            '支付件数': v['支付件数'],
            '转化率': v['支付买家数'] / v['商品访客数'] if v['商品访客数'] else 0,
            '加购率': v['商品加购人数'] / v['商品访客数'] if v['商品访客数'] else 0,
        } for (dv, dt_str), v in sorted(_sales_day.items())]

    st.markdown('<div class="section-title">数据明细</div>', unsafe_allow_html=True)

    def _parse_num(v):
        if isinstance(v, (int, float)): return v
        try: return float(str(v).replace(',','').replace('%','').replace('¥',''))
        except: return 0

    def _date_yoy(dt_str):
        try:
            dt_obj = datetime.datetime.strptime(dt_str, '%Y-%m-%d').date()
            return str(dt_obj.replace(year=dt_obj.year - 1))
        except ValueError:
            return None

    # ── 根据对比模式决定对比日期查找函数和列名后缀 ──
    _s_cmp_suffix = '同比' if _s_comp_mode == '同比' else '上月同期'
    def _s_get_cmp_day(dv, dt_str):
        """返回对比期对应数据（同比：去年同日；上月同期：上月同日）"""
        if _s_comp_mode == '同比':
            ly_dt = _date_yoy(dt_str)
            return _sales_all_day.get((dv, ly_dt)) if ly_dt else None
        else:
            return _sales_mom_day.get((dv, dt_str))

    # ── 当选择了分析维度时，直接展示维度汇总表（不按日/月拆分）──
    if _s_dim_field:
        # 按维度值汇总：合并所有日期
        _dim_agg = {}
        for (dv, dt_str), v in _sales_day.items():
            if dv not in _dim_agg:
                _dim_agg[dv] = {m: 0.0 for m in METRICS}
            for m in METRICS:
                _dim_agg[dv][m] += v[m]
        # 对比期（同比/上月同期）：按维度汇总
        _dim_agg_ly = {}
        for (dv, dt_str), v in _sales_day.items():
            if _s_comp_mode == '同比':
                cmp_dt = _date_yoy(dt_str)
                cmp_v = _sales_all_day.get((dv, cmp_dt)) if cmp_dt else None
            else:
                cmp_v = _sales_mom_day.get((dv, dt_str))
            if not cmp_v: continue
            if dv not in _dim_agg_ly:
                _dim_agg_ly[dv] = {m: 0.0 for m in METRICS}
            for m in METRICS:
                _dim_agg_ly[dv][m] += cmp_v[m]
        # 推广花费按维度聚合（维度下所有日期求和，匹配渠道/店铺/品类/型号维度）
        _dim_promo = {}
        _dim_promo_amt = {}  # 推广总成交金额按维度聚合
        if promo_rows:
            for r in promo_filtered:
                d = r.get('_date', '')
                if not d or not (str(start) <= d <= str(end)): continue
                dv_p = r.get(_s_dim_field, '') or '未标注'
                _dim_promo[dv_p] = _dim_promo.get(dv_p, 0) + float(r.get('_花费', 0) or 0)
                _dim_promo_amt[dv_p] = _dim_promo_amt.get(dv_p, 0) + float(r.get('_总订单金额', 0) or 0)
        # 构建全时段推广日聚合字典（用于对比期推广数据查找）
        _p_fields_list = ['_花费','_展现数','_点击数','_直接订单金额','_总订单金额','_直接订单量','_总成交订单量']
        _promo_all_day = {}
        if promo_rows:
            for r in promo_rows:
                d = r.get('_date', '')
                if not d: continue
                if channel and r.get('_渠道', '') not in channel: continue
                if store and r.get('_店铺', '') not in store: continue
                if category and r.get('_品类', '') not in category: continue
                if model and r.get('_型号', '') not in model: continue
                dv_p = r.get(_s_dim_field, '') or '未标注'
                key = (dv_p, d)
                if key not in _promo_all_day:
                    _promo_all_day[key] = {fk: 0.0 for fk in _p_fields_list}
                for fk in _p_fields_list:
                    _promo_all_day[key][fk] += float(r.get(fk, 0) or 0)
        _dim_total_amt = sum(v['支付金额'] for v in _dim_agg.values()) or 1
        _dim_total_vis = sum(v['商品访客数'] for v in _dim_agg.values())
        # 对比期推广数据按维度聚合（花费 + 总成交金额）
        _dim_promo_ly = {}   # 对比期花费
        _dim_promo_amt_ly = {}  # 对比期推广总成交金额
        if promo_rows:
            for (dv, dt_str), v in _sales_day.items():
                if _s_comp_mode == '同比':
                    cmp_dt = _date_yoy(dt_str)
                else:
                    cmp_dt = _date_mom(dt_str)
                if not cmp_dt: continue
                dv_p = dv
                # 从 promo_all_day 找对比期推广数据
                for (pdv, pdt), pv in _promo_all_day.items():
                    if pdt == cmp_dt and pdv == dv_p:
                        _dim_promo_ly[dv] = _dim_promo_ly.get(dv, 0) + float(pv.get('_花费', 0) or 0)
                        _dim_promo_amt_ly[dv] = _dim_promo_amt_ly.get(dv, 0) + float(pv.get('_总订单金额', 0) or 0)
        _dim_tbl = []
        for dv, v in _dim_agg.items():
            vis = v['商品访客数']; amt = v['支付金额']
            ly_v = _dim_agg_ly.get(dv, {})
            ly_amt = ly_v.get('支付金额', 0); ly_vis = ly_v.get('商品访客数', 0)
            ly_buyers = ly_v.get('支付买家数', 0)
            cvr = v['支付买家数'] / vis if vis else 0
            ly_cvr = ly_buyers / ly_vis if ly_vis else 0
            yoy_amt = (amt - ly_amt) / ly_amt if ly_amt else None
            yoy_vis = (vis - ly_vis) / ly_vis if ly_vis else None
            yoy_cvr = (cvr - ly_cvr) / ly_cvr if ly_cvr else None
            atv_v = amt / v['支付买家数'] if v['支付买家数'] else 0
            ly_atv_v = ly_amt / ly_buyers if ly_buyers else None
            yoy_atv = (atv_v - ly_atv_v) / ly_atv_v if ly_atv_v else None
            spend = _dim_promo.get(dv, 0)
            promo_amt = _dim_promo_amt.get(dv, 0)  # 推广总成交金额
            prs = promo_amt / amt if amt else 0  # 推广成交占比
            # 对比期费率和推广成交占比
            ly_spend = _dim_promo_ly.get(dv, 0)
            ly_promo_amt = _dim_promo_amt_ly.get(dv, 0)
            ly_fee = ly_spend / ly_amt if ly_amt else None
            fee_yoy = (spend/amt - ly_fee) / ly_fee if amt and ly_fee and ly_fee > 0 else None
            ly_prs = ly_promo_amt / ly_amt if ly_amt else None
            prs_yoy = (prs - ly_prs) / ly_prs if ly_prs and ly_prs > 0 else None
            _dim_tbl.append({
                '维度': dv,
                '访客数': f"{int(vis):,}",
                '访客占比': f"{vis/_dim_total_vis*100:.2f}%" if _dim_total_vis else "0.00%",
                '买家数': f"{int(v['支付买家数']):,}",
                '支付件数': f"{int(v['支付件数']):,}",
                '成交金额(万)': round(amt/10000, 1),
                '成交占比': f"{amt/_dim_total_amt*100:.2f}%",
                '转化率': f"{cvr*100:.2f}%",
                '加购人数': f"{int(v['商品加购人数']):,}",
                '加购率': f"{v['商品加购人数']/vis*100:.2f}%" if vis else "0.00%",
                '客单价': round(amt/v['支付买家数'], 1) if v['支付买家数'] else 0,
                'UV价值': round(amt/vis, 1) if vis else 0,
                '费率': f"{spend/amt*100:.2f}%" if amt else '--',
                '推广成交占比': f"{prs*100:.2f}%",
                f'销额{_s_cmp_suffix}': f"{yoy_amt*100:+.2f}%" if yoy_amt is not None else '--',
                f'访客{_s_cmp_suffix}': f"{yoy_vis*100:+.2f}%" if yoy_vis is not None else '--',
                f'转化率{_s_cmp_suffix}': f"{yoy_cvr*100:+.2f}%" if yoy_cvr is not None else '--',
                f'客单价{_s_cmp_suffix}': f"{yoy_atv*100:+.2f}%" if yoy_atv is not None else '--',
                f'费率{_s_cmp_suffix}': f"{fee_yoy*100:+.2f}%" if fee_yoy is not None else '--',
                f'推广成交占比{_s_cmp_suffix}': f"{prs_yoy*100:+.2f}%" if prs_yoy is not None else '--',
            })
        # 合计行
        _dim_total_buyers = sum(v['支付买家数'] for v in _dim_agg.values())
        _dim_total_cart = sum(v['商品加购人数'] for v in _dim_agg.values())
        _dim_total_qty = sum(v['支付件数'] for v in _dim_agg.values())
        _ly_dim_vis = sum(v.get('商品访客数', 0) for v in _dim_agg_ly.values())
        _ly_dim_amt = sum(v.get('支付金额', 0) for v in _dim_agg_ly.values())
        _ly_dim_buyers = sum(v.get('支付买家数', 0) for v in _dim_agg_ly.values())
        _dim_total_spend = sum(_dim_promo.values())
        _dim_total_promo_amt = sum(_dim_promo_amt.values())
        _dim_prs = _dim_total_promo_amt / _dim_total_amt if _dim_total_amt else 0
        _dim_total_spend_ly = sum(_dim_promo_ly.values())
        _dim_total_promo_amt_ly = sum(_dim_promo_amt_ly.values())
        _dim_fee = _dim_total_spend / _dim_total_amt if _dim_total_amt else 0
        _dim_fee_ly = _dim_total_spend_ly / _ly_dim_amt if _ly_dim_amt else None
        _dim_fee_yoy = (_dim_fee - _dim_fee_ly) / _dim_fee_ly if _dim_fee_ly and _dim_fee_ly > 0 else None
        _dim_prs_ly = _dim_total_promo_amt_ly / _ly_dim_amt if _ly_dim_amt else None
        _dim_prs_yoy = (_dim_prs - _dim_prs_ly) / _dim_prs_ly if _dim_prs_ly and _dim_prs_ly > 0 else None
        _dim_cvr = _dim_total_buyers / _dim_total_vis if _dim_total_vis else 0
        _ly_dim_cvr = _ly_dim_buyers / _ly_dim_vis if _ly_dim_vis else 0
        _dim_atv = _dim_total_amt / _dim_total_buyers if _dim_total_buyers else 0
        _ly_dim_atv = _ly_dim_amt / _ly_dim_buyers if _ly_dim_buyers else None
        _yoy_dim_atv = (_dim_atv - _ly_dim_atv) / _ly_dim_atv if _ly_dim_atv else None
        _dim_tbl.append({
            '维度': '合计',
            '访客数': f"{int(_dim_total_vis):,}", '访客占比': '100.00%',
            '买家数': f"{int(_dim_total_buyers):,}", '支付件数': f"{int(_dim_total_qty):,}",
            '成交金额(万)': round(_dim_total_amt/10000, 1), '成交占比': '100.00%',
            '转化率': f"{_dim_cvr*100:.2f}%",
            '加购人数': f"{int(_dim_total_cart):,}",
            '加购率': f"{_dim_total_cart/_dim_total_vis*100:.2f}%" if _dim_total_vis else "0.00%",
            '客单价': round(_dim_total_amt/_dim_total_buyers, 1) if _dim_total_buyers else 0,
            'UV价值': round(_dim_total_amt/_dim_total_vis, 1) if _dim_total_vis else 0,
            '费率': f"{_dim_total_spend/_dim_total_amt*100:.2f}%" if _dim_total_amt else '--',
            '推广成交占比': f"{_dim_prs*100:.2f}%",
            f'销额{_s_cmp_suffix}': f"{(_dim_total_amt-_ly_dim_amt)/_ly_dim_amt*100:+.2f}%" if _ly_dim_amt else '--',
            f'访客{_s_cmp_suffix}': f"{(_dim_total_vis-_ly_dim_vis)/_ly_dim_vis*100:+.2f}%" if _ly_dim_vis else '--',
            f'转化率{_s_cmp_suffix}': f"{(_dim_cvr-_ly_dim_cvr)/_ly_dim_cvr*100:+.2f}%" if _ly_dim_cvr else '--',
            f'客单价{_s_cmp_suffix}': f"{_yoy_dim_atv*100:+.2f}%" if _yoy_dim_atv is not None else '--',
            f'费率{_s_cmp_suffix}': f"{_dim_fee_yoy*100:+.2f}%" if _dim_fee_yoy is not None else '--',
            f'推广成交占比{_s_cmp_suffix}': f"{_dim_prs_yoy*100:+.2f}%" if _dim_prs_yoy is not None else '--',
        })
        # 排序控件
        _ds_sort_cols = ['维度','访客数','买家数','支付件数','成交金额(万)','转化率','客单价','加购人数','加购率','UV价值','费率','推广成交占比']
        _dsc1, _dsc2 = st.columns([2, 1])
        with _dsc1:
            _ds_sort_by = st.selectbox('排序字段', _ds_sort_cols, index=4, key='dim_sort_col')
        with _dsc2:
            _ds_sort_desc = st.radio('', ['降序', '升序'], horizontal=True, key='dim_sort_dir', index=0)
        _dim_data_rows = [r for r in _dim_tbl if r.get('维度') != '合计']
        _dim_total_rows = [r for r in _dim_tbl if r.get('维度') == '合计']
        if _ds_sort_by == '维度':
            _dim_data_rows.sort(key=lambda r: r.get('维度', ''), reverse=(_ds_sort_desc == '降序'))
        else:
            _dim_data_rows.sort(key=lambda r: _parse_num(r.get(_ds_sort_by, 0)), reverse=(_ds_sort_desc == '降序'))
        _dim_headers = ['维度','访客数','访客占比','买家数','支付件数','成交金额(万)','成交占比','转化率','客单价','加购人数','加购率','UV价值','费率','推广成交占比',
                        f'销额{_s_cmp_suffix}',f'访客{_s_cmp_suffix}',f'转化率{_s_cmp_suffix}',f'客单价{_s_cmp_suffix}',f'费率{_s_cmp_suffix}',f'推广成交占比{_s_cmp_suffix}']
        _render_html_table(_dim_data_rows + _dim_total_rows, _dim_headers, _dim_headers, title=f'📦 销售{_sales_dim}汇总')
        _render_download_panel(_dim_data_rows + _dim_total_rows, _dim_headers, 'sales_dim_summary.csv', '📥 下载维度汇总')
    else:
        # ── 按日期模式：显示日度汇总/月度汇总 tab ──
        tab_daily, tab_monthly = st.tabs(['日度汇总', '月度汇总'])
        with tab_daily:
            # 构建每日推广花费字典（按日期累计，用于费率）
            _day_promo = {}
            _day_promo_amt = {}  # 每日推广总成交金额
            if promo_rows:
                for r in promo_filtered:
                    d = r.get('_date', '')
                    if not d or not (str(start) <= d <= str(end)): continue
                    _day_promo[d] = _day_promo.get(d, 0) + float(r.get('_花费', 0) or 0)
                    _day_promo_amt[d] = _day_promo_amt.get(d, 0) + float(r.get('_总订单金额', 0) or 0)
            # 构建全时段推广日聚合字典（用于对比期推广数据查找，按日期维度）
            _p_fields_list = ['_花费','_展现数','_点击数','_直接订单金额','_总订单金额','_直接订单量','_总成交订单量']
            _promo_all_day = {}
            if promo_rows:
                for r in promo_rows:
                    d = r.get('_date', '')
                    if not d: continue
                    if channel and r.get('_渠道', '') not in channel: continue
                    if store and r.get('_店铺', '') not in store: continue
                    if category and r.get('_品类', '') not in category: continue
                    if model and r.get('_型号', '') not in model: continue
                    dv_p = r.get(_s_dim_field, '') or '未标注'
                    key = (dv_p, d)
                    if key not in _promo_all_day:
                        _promo_all_day[key] = {fk: 0.0 for fk in _p_fields_list}
                    for fk in _p_fields_list:
                        _promo_all_day[key][fk] += float(r.get(fk, 0) or 0)
            # 使用 _sales_day 作为 day_dict (key: (dim_val, date))
            total_vis = sum(v['商品访客数'] for v in _sales_day.values())
            total_amt = sum(v['支付金额'] for v in _sales_day.values())
            total_buyers = sum(v['支付买家数'] for v in _sales_day.values())
            total_cart = sum(v['商品加购人数'] for v in _sales_day.values())
            # 预计算对比期：根据 _s_comp_mode 选同比或上月同期
            ly_day_dict = {}
            for (dv, dt_str), v in _sales_day.items():
                ly_day_dict[(dv, dt_str)] = _s_get_cmp_day(dv, dt_str)
            daily_tbl = []
            for (dv, dt_str), v in sorted(_sales_day.items()):
                vis = v['商品访客数']
                amt = v['支付金额']
                ly_v = ly_day_dict.get((dv, dt_str))
                yoy_amt = (amt - ly_v['支付金额']) / ly_v['支付金额'] if ly_v and ly_v['支付金额'] else None
                yoy_vis = (vis - ly_v['商品访客数']) / ly_v['商品访客数'] if ly_v and ly_v['商品访客数'] else None
                cvr = v['支付买家数'] / vis if vis else 0
                ly_cvr = ly_v['支付买家数'] / ly_v['商品访客数'] if ly_v and ly_v['商品访客数'] else None
                yoy_cvr = (cvr - ly_cvr) / ly_cvr if ly_cvr else None
                atv_d = amt / v['支付买家数'] if v['支付买家数'] else 0
                ly_atv_d = ly_v['支付金额'] / ly_v['支付买家数'] if ly_v and ly_v['支付买家数'] else None
                yoy_atv_d = (atv_d - ly_atv_d) / ly_atv_d if ly_atv_d else None
                day_spend = _day_promo.get(dt_str, 0)
                day_promo_amt = _day_promo_amt.get(dt_str, 0)
                day_prs = day_promo_amt / amt if amt else 0
                # 对比期推广数据（费率、推广成交占比的同比/上月同期）
                cmp_dt = _date_yoy(dt_str) if _s_comp_mode == '同比' else _date_mom(dt_str)
                cmp_spend = 0; cmp_promo_amt = 0
                if cmp_dt and promo_rows:
                    for (pdv, pdt), pv in _promo_all_day.items():
                        if pdt == cmp_dt and (not _s_dim_field or pdv == dv):
                            cmp_spend += float(pv.get('_花费', 0) or 0)
                            cmp_promo_amt += float(pv.get('_总订单金额', 0) or 0)
                cmp_fee_rate = cmp_spend / ly_v['支付金额'] if ly_v and ly_v['支付金额'] else None
                yoy_fee_rate = (day_spend/amt - cmp_fee_rate) / cmp_fee_rate if amt and cmp_fee_rate and cmp_fee_rate > 0 else None
                cmp_prs = cmp_promo_amt / ly_v['支付金额'] if ly_v and ly_v['支付金额'] else None
                yoy_prs = (day_prs - cmp_prs) / cmp_prs if cmp_prs and cmp_prs > 0 else None
                row = {
                    '日期': dt_str, '访客数': f"{int(vis):,}",
                    '访客占比': f"{vis/total_vis*100:.2f}%" if total_vis else "0.00%",
                    '买家数': f"{int(v['支付买家数']):,}", '支付件数': f"{int(v['支付件数']):,}",
                    '成交金额(万)': round(amt/10000, 1),
                    '成交占比': f"{amt/total_amt*100:.2f}%" if total_amt else "0.00%",
                    '转化率': f"{v['支付买家数']/vis*100:.2f}%" if vis else "0.00%",
                    '加购人数': f"{int(v['商品加购人数']):,}",
                    '加购率': f"{v['商品加购人数']/vis*100:.2f}%" if vis else "0.00%",
                    '客单价': round(amt/v['支付买家数'], 1) if v['支付买家数'] else 0,
                    'UV价值': round(amt/vis, 1) if vis else 0,
                    '费率': f"{day_spend/amt*100:.2f}%" if amt else '--',
                    '推广成交占比': f"{day_prs*100:.2f}%",
                    f'销额{_s_cmp_suffix}': f"{yoy_amt*100:+.2f}%" if yoy_amt is not None else '--',
                    f'访客{_s_cmp_suffix}': f"{yoy_vis*100:+.2f}%" if yoy_vis is not None else '--',
                    f'转化率{_s_cmp_suffix}': f"{yoy_cvr*100:+.2f}%" if yoy_cvr is not None else '--',
                    f'客单价{_s_cmp_suffix}': f"{yoy_atv_d*100:+.2f}%" if yoy_atv_d is not None else '--',
                    f'费率{_s_cmp_suffix}': f"{yoy_fee_rate*100:+.2f}%" if yoy_fee_rate is not None else '--',
                    f'推广成交占比{_s_cmp_suffix}': f"{yoy_prs*100:+.2f}%" if yoy_prs is not None else '--',
                }
                if _s_dim_field:
                    row['维度'] = dv
                daily_tbl.append(row)
            if daily_tbl:
                # 日度合计行：对比期与本期一一对应的日期求和
                _ly_total_vis_d = sum(v['商品访客数'] for v in ly_day_dict.values() if v)
                _ly_total_amt_d = sum(v['支付金额'] for v in ly_day_dict.values() if v)
                _ly_total_buyers_d = sum(v['支付买家数'] for v in ly_day_dict.values() if v)
                _total_yoy_amt_d = (total_amt - _ly_total_amt_d) / _ly_total_amt_d if _ly_total_amt_d else None
                _total_yoy_vis_d = (total_vis - _ly_total_vis_d) / _ly_total_vis_d if _ly_total_vis_d else None
                _total_cvr_d = total_buyers / total_vis if total_vis else 0
                _ly_total_cvr_d = _ly_total_buyers_d / _ly_total_vis_d if _ly_total_vis_d else 0
                _total_yoy_cvr_d = (_total_cvr_d - _ly_total_cvr_d) / _ly_total_cvr_d if _ly_total_cvr_d else None
                _total_atv_d = total_amt / total_buyers if total_buyers else 0
                _ly_total_atv_d = _ly_total_amt_d / _ly_total_buyers_d if _ly_total_buyers_d else None
                _total_yoy_atv_d = (_total_atv_d - _ly_total_atv_d) / _ly_total_atv_d if _ly_total_atv_d else None
                _total_promo_d = sum(_day_promo.values())
                _total_promo_amt_d = sum(_day_promo_amt.values())
                _total_rate_d = _total_promo_d / total_amt * 100 if total_amt else None
                _total_prs_d = _total_promo_amt_d / total_amt if total_amt else 0
                # 对比期推广数据聚合（按对比期日期）
                _ly_total_spend_d = 0
                _ly_total_promo_amt_d = 0
                if promo_rows:
                    for (dv, dt_str), v in _sales_day.items():
                        if _s_comp_mode == '同比':
                            cmp_dt = _date_yoy(dt_str)
                        else:
                            cmp_dt = _date_mom(dt_str)
                        if not cmp_dt: continue
                        for (pdv, pdt), pv in _promo_all_day.items():
                            if pdt == cmp_dt:
                                _ly_total_spend_d += float(pv.get('_花费', 0) or 0)
                                _ly_total_promo_amt_d += float(pv.get('_总订单金额', 0) or 0)
                _ly_total_rate_d = _ly_total_spend_d / _ly_total_amt_d * 100 if _ly_total_amt_d else None
                _total_yoy_rate_d = (_total_rate_d - _ly_total_rate_d) / _ly_total_rate_d if _ly_total_rate_d and _ly_total_rate_d > 0 else None
                _ly_total_prs_d = _ly_total_promo_amt_d / _ly_total_amt_d if _ly_total_amt_d else None
                _total_yoy_prs_d = (_total_prs_d - _ly_total_prs_d) / _ly_total_prs_d if _ly_total_prs_d and _ly_total_prs_d > 0 else None
                _total_qty = sum(v['支付件数'] for v in _sales_day.values())
                total_row = {
                    '日期': '总计', '访客数': f"{int(total_vis):,}", '访客占比': '100.00%',
                    '买家数': f"{int(total_buyers):,}", '支付件数': f"{int(_total_qty):,}",
                    '成交金额(万)': round(total_amt/10000, 1), '成交占比': '100.00%',
                    '转化率': f"{total_buyers/total_vis*100:.2f}%" if total_vis else "0.00%",
                    '加购人数': f"{int(total_cart):,}",
                    '加购率': f"{total_cart/total_vis*100:.2f}%" if total_vis else "0.00%",
                    '客单价': round(total_amt/total_buyers, 1) if total_buyers else 0,
                    'UV价值': round(total_amt/total_vis, 1) if total_vis else 0,
                    '费率': f"{_total_rate_d:.2f}%" if _total_rate_d is not None else '--',
                    '推广成交占比': f"{_total_prs_d*100:.2f}%",
                    f'销额{_s_cmp_suffix}': f"{_total_yoy_amt_d*100:+.2f}%" if _total_yoy_amt_d is not None else '--',
                    f'访客{_s_cmp_suffix}': f"{_total_yoy_vis_d*100:+.2f}%" if _total_yoy_vis_d is not None else '--',
                    f'转化率{_s_cmp_suffix}': f"{_total_yoy_cvr_d*100:+.2f}%" if _total_yoy_cvr_d is not None else '--',
                    f'客单价{_s_cmp_suffix}': f"{_total_yoy_atv_d*100:+.2f}%" if _total_yoy_atv_d is not None else '--',
                    f'费率{_s_cmp_suffix}': f"{_total_yoy_rate_d*100:+.2f}%" if _total_yoy_rate_d is not None else '--',
                    f'推广成交占比{_s_cmp_suffix}': f"{_total_yoy_prs_d*100:+.2f}%" if _total_yoy_prs_d is not None else '--',
                }
                if _s_dim_field:
                    total_row['维度'] = '合计'
                daily_tbl.append(total_row)
            # 排序控件
            _daily_sort_cols = ['日期','访客数','买家数','支付件数','成交金额(万)','转化率','客单价','加购人数','加购率','UV价值','费率','推广成交占比']
            if _s_dim_field:
                _daily_sort_cols = ['维度'] + _daily_sort_cols
            _dsc1, _dsc2 = st.columns([2, 1])
            with _dsc1:
                _daily_sort_by = st.selectbox('排序字段', _daily_sort_cols, index=0, key='daily_sort_col')
            with _dsc2:
                _daily_sort_desc = st.radio('', ['降序', '升序'], horizontal=True, key='daily_sort_dir', index=1)
            _daily_data_rows = [r for r in daily_tbl if r.get('日期') != '总计']
            _daily_total_row = [r for r in daily_tbl if r.get('日期') == '总计']
            if _daily_sort_by == '日期':
                _daily_data_rows.sort(key=lambda r: r.get('日期', ''), reverse=(_daily_sort_desc == '降序'))
            elif _daily_sort_by == '维度':
                _daily_data_rows.sort(key=lambda r: r.get('维度', ''), reverse=(_daily_sort_desc == '降序'))
            else:
                _daily_data_rows.sort(key=lambda r: _parse_num(r.get(_daily_sort_by, 0)), reverse=(_daily_sort_desc == '降序'))
            _daily_tbl_sorted = _daily_data_rows + _daily_total_row
            _daily_headers = ['日期','访客数','访客占比','买家数','支付件数','成交金额(万)','成交占比','转化率','客单价','加购人数','加购率','UV价值','费率','推广成交占比',
                              f'销额{_s_cmp_suffix}',f'访客{_s_cmp_suffix}',f'转化率{_s_cmp_suffix}',f'客单价{_s_cmp_suffix}',f'费率{_s_cmp_suffix}',f'推广成交占比{_s_cmp_suffix}']
            if _s_dim_field:
                _daily_headers = ['维度'] + _daily_headers
            _render_html_table(_daily_tbl_sorted, _daily_headers, _daily_headers, title='📦 销售日度趋势')
            # 下载原始数据
            _daily_dl = [{'维度': dv, '日期': dt, **{m: v[m] for m in METRICS}} for (dv, dt), v in _sales_day.items()]
            _render_download_panel(_daily_dl, (['维度'] if _s_dim_field else []) + ['日期'] + METRICS, 'daily_summary.csv')
        with tab_monthly:
            # 构建每月推广花费字典（按月份累计，用于费率）
            _month_promo = {}
            _month_promo_amt = {}  # 每月推广总成交金额
            if promo_rows:
                for r in promo_filtered:
                    d = r.get('_date', '')
                    if not d or len(d) < 7: continue
                    if not (str(start) <= d <= str(end)): continue
                    ym = d[:7]
                    _month_promo[ym] = _month_promo.get(ym, 0) + float(r.get('_花费', 0) or 0)
                    _month_promo_amt[ym] = _month_promo_amt.get(ym, 0) + float(r.get('_总订单金额', 0) or 0)
            # 使用 _sales_month (key: (dim_val, ym))
            total_vis_m = sum(v['商品访客数'] for v in _sales_month.values())
            total_amt_m = sum(v['支付金额'] for v in _sales_month.values())
            total_buyers_m = sum(v['支付买家数'] for v in _sales_month.values())
            total_cart_m = sum(v['商品加购人数'] for v in _sales_month.values())
            monthly_tbl = []
            for (dv, ym), v in sorted(_sales_month.items()):
                vis = v['商品访客数']
                amt = v['支付金额']
                # 对比期：从 ly_day_dict 聚合本月所有天对应的对比数据（支持同比/上月同期）
                ly_amt = 0; ly_vis = 0; ly_buyers = 0
                for (ly_dv, ly_dt_str), ly_v in ly_day_dict.items():
                    if ly_dt_str.startswith(ym) and ly_v:
                        if _s_dim_field and ly_dv != dv: continue
                        ly_amt += float(ly_v.get('支付金额', 0) or 0)
                        ly_vis += float(ly_v.get('商品访客数', 0) or 0)
                        ly_buyers += float(ly_v.get('支付买家数', 0) or 0)
                yoy_amt = (amt - ly_amt) / ly_amt if ly_amt else None
                yoy_vis = (vis - ly_vis) / ly_vis if ly_vis else None
                cvr = v['支付买家数'] / vis if vis else 0
                ly_cvr = ly_buyers / ly_vis if ly_vis else 0
                yoy_cvr = (cvr - ly_cvr) / ly_cvr if ly_cvr else None
                atv_m = amt / v['支付买家数'] if v['支付买家数'] else 0
                ly_atv_m = ly_amt / ly_buyers if ly_buyers else None
                yoy_atv_m = (atv_m - ly_atv_m) / ly_atv_m if ly_atv_m else None
                month_spend = _month_promo.get(ym, 0)
                month_promo_amt = _month_promo_amt.get(ym, 0)
                month_prs = month_promo_amt / amt if amt else 0
                # 对比期推广数据（按月份聚合对比期各天的推广花费和总成交金额）
                cmp_spend_m = 0; cmp_promo_amt_m = 0
                if promo_rows:
                    for (sdv, sdt), sv in _sales_day.items():
                        if not sdt.startswith(ym): continue
                        if _s_dim_field and sdv != dv: continue
                        cmp_dt_m = _date_yoy(sdt) if _s_comp_mode == '同比' else _date_mom(sdt)
                        if not cmp_dt_m: continue
                        for (pdv, pdt), pv in _promo_all_day.items():
                            if pdt == cmp_dt_m and (not _s_dim_field or pdv == sdv):
                                cmp_spend_m += float(pv.get('_花费', 0) or 0)
                                cmp_promo_amt_m += float(pv.get('_总订单金额', 0) or 0)
                cmp_fee_rate_m = cmp_spend_m / ly_amt if ly_amt else None
                yoy_fee_rate_m = (month_spend/amt - cmp_fee_rate_m) / cmp_fee_rate_m if amt and cmp_fee_rate_m and cmp_fee_rate_m > 0 else None
                cmp_prs_m = cmp_promo_amt_m / ly_amt if ly_amt else None
                yoy_prs_m = (month_prs - cmp_prs_m) / cmp_prs_m if cmp_prs_m and cmp_prs_m > 0 else None
                row = {
                    '月份': ym, '访客数': f"{int(vis):,}",
                    '访客占比': f"{vis/total_vis_m*100:.2f}%" if total_vis_m else "0.00%",
                    '买家数': f"{int(v['支付买家数']):,}", '支付件数': f"{int(v['支付件数']):,}",
                    '成交金额(万)': round(amt/10000, 1),
                    '成交占比': f"{amt/total_amt_m*100:.2f}%" if total_amt_m else "0.00%",
                    '转化率': f"{cvr*100:.2f}%" if vis else "0.00%",
                    '加购人数': f"{int(v['商品加购人数']):,}",
                    '加购率': f"{v['商品加购人数']/vis*100:.2f}%" if vis else "0.00%",
                    '客单价': round(amt/v['支付买家数'], 1) if v['支付买家数'] else 0,
                    'UV价值': round(amt/vis, 1) if vis else 0,
                    '费率': f"{month_spend/amt*100:.2f}%" if amt else '--',
                    '推广成交占比': f"{month_prs*100:.2f}%",
                    f'销额{_s_cmp_suffix}': f"{yoy_amt*100:+.2f}%" if yoy_amt is not None else '--',
                    f'访客{_s_cmp_suffix}': f"{yoy_vis*100:+.2f}%" if yoy_vis is not None else '--',
                    f'转化率{_s_cmp_suffix}': f"{yoy_cvr*100:+.2f}%" if yoy_cvr is not None else '--',
                    f'客单价{_s_cmp_suffix}': f"{yoy_atv_m*100:+.2f}%" if yoy_atv_m is not None else '--',
                    f'费率{_s_cmp_suffix}': f"{yoy_fee_rate_m*100:+.2f}%" if yoy_fee_rate_m is not None else '--',
                    f'推广成交占比{_s_cmp_suffix}': f"{yoy_prs_m*100:+.2f}%" if yoy_prs_m is not None else '--',
                }
                if _s_dim_field:
                    row['维度'] = dv
                monthly_tbl.append(row)
            if monthly_tbl:
                _mm_total_qty = sum(v['支付件数'] for v in _sales_month.values())
                _ly_daily_vis = sum(v['商品访客数'] for v in ly_day_dict.values() if v)
                _ly_daily_amt = sum(v['支付金额'] for v in ly_day_dict.values() if v)
                _ly_daily_buyers = sum(v['支付买家数'] for v in ly_day_dict.values() if v)
                _total_yoy_amt = (total_amt_m - _ly_daily_amt) / _ly_daily_amt if _ly_daily_amt else None
                _total_yoy_vis = (total_vis_m - _ly_daily_vis) / _ly_daily_vis if _ly_daily_vis else None
                _total_cvr = total_buyers_m / total_vis_m if total_vis_m else 0
                _ly_total_cvr = _ly_daily_buyers / _ly_daily_vis if _ly_daily_vis else 0
                _total_yoy_cvr = (_total_cvr - _ly_total_cvr) / _ly_total_cvr if _ly_total_cvr else None
                _total_atv_m = total_amt_m / total_buyers_m if total_buyers_m else 0
                _ly_total_atv_m = _ly_daily_amt / _ly_daily_buyers if _ly_daily_buyers else None
                _total_yoy_atv_m = (_total_atv_m - _ly_total_atv_m) / _ly_total_atv_m if _ly_total_atv_m else None
                _total_promo_m = sum(_month_promo.values())
                _total_promo_amt_m = sum(_month_promo_amt.values())
                _total_rate_m = _total_promo_m / total_amt_m * 100 if total_amt_m else None
                _total_prs_m = _total_promo_amt_m / total_amt_m if total_amt_m else 0
                # 对比期推广数据聚合（按对比期日期从 promo_all_day）
                _ly_total_spend_m = 0
                _ly_total_promo_amt_m = 0
                if promo_rows:
                    for (dv, dt_str), v in _sales_day.items():
                        if _s_comp_mode == '同比':
                            cmp_dt = _date_yoy(dt_str)
                        else:
                            cmp_dt = _date_mom(dt_str)
                        if not cmp_dt: continue
                        for (pdv, pdt), pv in _promo_all_day.items():
                            if pdt == cmp_dt:
                                _ly_total_spend_m += float(pv.get('_花费', 0) or 0)
                                _ly_total_promo_amt_m += float(pv.get('_总订单金额', 0) or 0)
                _ly_total_rate_m = _ly_total_spend_m / _ly_daily_amt * 100 if _ly_daily_amt else None
                _total_yoy_rate_m = (_total_rate_m - _ly_total_rate_m) / _ly_total_rate_m if _ly_total_rate_m and _ly_total_rate_m > 0 else None
                _ly_total_prs_m = _ly_total_promo_amt_m / _ly_daily_amt if _ly_daily_amt else None
                _total_yoy_prs_m = (_total_prs_m - _ly_total_prs_m) / _ly_total_prs_m if _ly_total_prs_m and _ly_total_prs_m > 0 else None
                total_row = {
                    '月份': '总计', '访客数': f"{int(total_vis_m):,}", '访客占比': '100.00%',
                    '买家数': f"{int(total_buyers_m):,}", '支付件数': f"{int(_mm_total_qty):,}",
                    '成交金额(万)': round(total_amt_m/10000, 1), '成交占比': '100.00%',
                    '转化率': f"{_total_cvr*100:.2f}%",
                    '加购人数': f"{int(total_cart_m):,}",
                    '加购率': f"{total_cart_m/total_vis_m*100:.2f}%" if total_vis_m else "0.00%",
                    '客单价': round(total_amt_m/total_buyers_m, 1) if total_buyers_m else 0,
                    'UV价值': round(total_amt_m/total_vis_m, 1) if total_vis_m else 0,
                    '费率': f"{_total_rate_m:.2f}%" if _total_rate_m is not None else '--',
                    '推广成交占比': f"{_total_prs_m*100:.2f}%",
                    f'销额{_s_cmp_suffix}': f"{_total_yoy_amt*100:+.2f}%" if _total_yoy_amt is not None else '--',
                    f'访客{_s_cmp_suffix}': f"{_total_yoy_vis*100:+.2f}%" if _total_yoy_vis is not None else '--',
                    f'转化率{_s_cmp_suffix}': f"{_total_yoy_cvr*100:+.2f}%" if _total_yoy_cvr is not None else '--',
                    f'客单价{_s_cmp_suffix}': f"{_total_yoy_atv_m*100:+.2f}%" if _total_yoy_atv_m is not None else '--',
                    f'费率{_s_cmp_suffix}': f"{_total_yoy_rate_m*100:+.2f}%" if _total_yoy_rate_m is not None else '--',
                    f'推广成交占比{_s_cmp_suffix}': f"{_total_yoy_prs_m*100:+.2f}%" if _total_yoy_prs_m is not None else '--',
                }
                if _s_dim_field:
                    total_row['维度'] = '合计'
                monthly_tbl.append(total_row)
                # 排序控件
                _mm_sort_cols = ['月份','访客数','买家数','支付件数','成交金额(万)','转化率','客单价','加购人数','加购率','UV价值','费率','推广成交占比']
                if _s_dim_field:
                    _mm_sort_cols = ['维度'] + _mm_sort_cols
                _mmc1, _mmc2 = st.columns([2, 1])
                with _mmc1:
                    _mm_sort_by = st.selectbox('排序字段', _mm_sort_cols, index=0, key='mm_sort_col')
                with _mmc2:
                    _mm_sort_desc = st.radio('', ['降序', '升序'], horizontal=True, key='mm_sort_dir', index=1)
                _mm_data_rows = [r for r in monthly_tbl if r.get('月份') != '总计']
                _mm_total_row = [r for r in monthly_tbl if r.get('月份') == '总计']
                def _parse_num_m(v):
                    if isinstance(v, (int, float)): return v
                    try: return float(str(v).replace(',','').replace('%','').replace('¥',''))
                    except: return 0
                if _mm_sort_by == '月份':
                    _mm_data_rows.sort(key=lambda r: r.get('月份', ''), reverse=(_mm_sort_desc == '降序'))
                elif _mm_sort_by == '维度':
                    _mm_data_rows.sort(key=lambda r: r.get('维度', ''), reverse=(_mm_sort_desc == '降序'))
                else:
                    _mm_data_rows.sort(key=lambda r: _parse_num_m(r.get(_mm_sort_by, 0)), reverse=(_mm_sort_desc == '降序'))
                _mm_sorted = _mm_data_rows + _mm_total_row
                _mm_headers = ['月份','访客数','访客占比','买家数','支付件数','成交金额(万)','成交占比','转化率','客单价','加购人数','加购率','UV价值','费率','推广成交占比',
                               f'销额{_s_cmp_suffix}',f'访客{_s_cmp_suffix}',f'转化率{_s_cmp_suffix}',f'客单价{_s_cmp_suffix}',f'费率{_s_cmp_suffix}',f'推广成交占比{_s_cmp_suffix}']
                if _s_dim_field:
                    _mm_headers = ['维度'] + _mm_headers
                _render_html_table(_mm_sorted, _mm_headers, _mm_headers, title='📦 销售月度趋势')
                # 下载原始月度数据
                _mm_dl = [{'维度': dv, '月份': ym, **{m: v[m] for m in METRICS}} for (dv, ym), v in _sales_month.items()]
                _render_download_panel(_mm_dl, (['维度'] if _s_dim_field else []) + ['月份'] + METRICS, 'monthly_summary.csv')
    
    st.markdown('---')
    t1, t2 = st.columns(2)
    with t1:
        fig = go.Figure()
        if tr_data:
            fig.add_trace(go.Bar(x=[r['周期'] for r in tr_data], y=[_wan(r['支付金额']) for r in tr_data],
                                  text=[f"{_wan(r['支付金额'])}万" for r in tr_data], textposition='outside',
                                  name='支付金额(万)', marker_color='#1d4ed8', opacity=0.85))
            fig.add_trace(go.Scatter(x=[r['周期'] for r in tr_data], y=[r['访客数'] for r in tr_data],
                                      name='访客数', yaxis='y2', line=dict(color='#06b6d4', width=2)))
        fig.update_layout(height=350, template='plotly_white', legend=dict(orientation='h'),
                        yaxis_title='支付金额(万)', yaxis2=dict(title='访客数', overlaying='y', side='right'))
        st.plotly_chart(fig, width="stretch")
        _render_download_panel(tr_data, ['周期','支付金额','访客数','转化率','加购率'], 'trend_amt_vs_vis.csv', '📥 趋势：金额/访客')
    with t2:
        fig = go.Figure()
        if tr_data:
            fig.add_trace(go.Scatter(x=[r['周期'] for r in tr_data], y=[r['转化率'] * 100 for r in tr_data],
                                      name='支付转化率(%)', line=dict(color='#22c55e', width=2)))
            fig.add_trace(go.Scatter(x=[r['周期'] for r in tr_data], y=[r['加购率'] * 100 for r in tr_data],
                                      name='加购率(%)', line=dict(color='#f59e0b', width=2)))
        fig.update_layout(height=350, template='plotly_white', legend=dict(orientation='h'), yaxis_title='比率(%)')
        st.plotly_chart(fig, width="stretch")
        _render_download_panel(tr_data, ['周期','支付金额','访客数','转化率','加购率'], 'trend_rate.csv', '📥 趋势：转化率/加购率')

    st.markdown('---')
    st.markdown('<div class="section-title">同比趋势叠加（月度）</div>', unsafe_allow_html=True)
    # 从 _sales_month 汇总为纯月度列表（跨维度聚合）
    _ym_agg = {}
    for (dv, ym), v in _sales_month.items():
        if ym not in _ym_agg:
            _ym_agg[ym] = {m: 0.0 for m in METRICS}
        for m in METRICS:
            _ym_agg[ym][m] += v[m]
    # 全时段月度聚合（用于同比查找）
    _all_ym_agg = {}
    for (dv, dt_str), v in _sales_all_day.items():
        ym = dt_str[:7]
        if ym not in _all_ym_agg:
            _all_ym_agg[ym] = {m: 0.0 for m in METRICS}
        for m in METRICS:
            _all_ym_agg[ym][m] += v[m]
    _yoy_monthly = [{'月份': ym, '支付金额': v['支付金额'], '商品访客数': v['商品访客数'],
                      '支付件数': v['支付件数'], '支付买家数': v['支付买家数'],
                      '商品加购人数': v['商品加购人数'],
                      '支付转化率': v['支付买家数'] / v['商品访客数'] if v['商品访客数'] else 0,
                      '加购率': v['商品加购人数'] / v['商品访客数'] if v['商品访客数'] else 0}
                     for ym, v in sorted(_ym_agg.items())]
    if _yoy_monthly:
        fig = go.Figure()
        fig.add_trace(go.Bar(x=[r['月份'] for r in _yoy_monthly], y=[_wan(r['支付金额']) for r in _yoy_monthly],
                              text=[f"{_wan(r['支付金额'])}万" for r in _yoy_monthly], textposition='outside',
                              name='本期月度金额', marker_color='#1d4ed8'))
        ly_data = [_all_ym_agg.get(month_shift(r['月份'], -12), {}).get('支付金额', 0) for r in _yoy_monthly]
        fig.add_trace(go.Scatter(x=[r['月份'] for r in _yoy_monthly], y=[_wan(v) for v in ly_data],
                                  name='去年同期金额', line=dict(color='#f59e0b', width=2, dash='dash')))
        fig.update_layout(height=380, template='plotly_white', legend=dict(orientation='h'), yaxis_title='支付金额(万)')
        st.plotly_chart(fig, width="stretch")

    st.markdown('---')
    st.markdown('<div class="section-title">周内趋势（每日均值）</div>', unsafe_allow_html=True)
    dow_map = {0: '周一', 1: '周二', 2: '周三', 3: '周四', 4: '周五', 5: '周六', 6: '周日'}
    dow_dict = {v: {m: [] for m in METRICS} for v in dow_map.values()}
    for (dv, dt_str), v in _sales_day.items():
        try:
            dt = datetime.date.fromisoformat(dt_str[:10])
            dow_name = dow_map[dt.weekday()]
            for m in METRICS:
                dow_dict[dow_name][m].append(v[m])
        except Exception:
            continue
    dow_avg = []
    for dow_name in ['周一', '周二', '周三', '周四', '周五', '周六', '周日']:
        d = dow_dict[dow_name]
        cnt = len(d['商品访客数'])
        if cnt:
            avg_amt = sum(d['支付金额']) / cnt
            avg_vis = sum(d['商品访客数']) / cnt
            avg_buyers = sum(d['支付买家数']) / cnt
            dow_avg.append({'星期': dow_name, '支付金额': _wan(avg_amt),
                             '访客数': int(avg_vis),
                             '转化率': avg_buyers / avg_vis if avg_vis else 0})
    if dow_avg:
        fig = go.Figure()
        # 柱状图 — 支付金额（左Y轴）
        fig.add_trace(go.Bar(
            x=[r['星期'] for r in dow_avg],
            y=[r['支付金额'] for r in dow_avg],
            name='支付金额',
            marker_color='#3B82F6',
            text=[f"{r['支付金额']:.1f}万" for r in dow_avg],
            textposition='outside',
            yaxis='y1'
        ))
        # 折线图 — 转化率（右Y轴，%）
        fig.add_trace(go.Scatter(
            x=[r['星期'] for r in dow_avg],
            y=[r['转化率'] * 100 for r in dow_avg],
            name='转化率',
            mode='lines+markers+text',
            text=[f"{r['转化率']*100:.2f}%" for r in dow_avg],
            textposition='top center',
            line=dict(color='#E6A817', width=3),
            marker=dict(size=10, color='#E6A817'),
            yaxis='y2'
        ))
        fig.update_layout(
            title='各星期日均支付金额与转化率',
            height=380,
            template='plotly_white',
            legend=dict(orientation='h', yanchor='bottom', y=1.02, xanchor='right', x=1),
            yaxis=dict(title='支付金额（万元）', side='left', showgrid=True),
            yaxis2=dict(title='转化率（%）', side='right', overlaying='y', showgrid=False),
            xaxis=dict(title='星期'),
            margin=dict(t=80)
        )
        st.plotly_chart(fig, width="stretch")

    # ═══════════════════════════════════════════════════
    # 📢 推广趋势（独立板块）
    # ═══════════════════════════════════════════════════
    if promo_rows:
        st.markdown('---')
        st.markdown('### 📢 推广趋势')
        _p_gran_c1, _p_gran_c2, _p_gran_c3 = st.columns([1, 2, 2])
        with _p_gran_c1:
            _p_gran = st.radio('粒度', ['月度', '周度', '日度'], horizontal=True, key='tr_promo_gran')
        with _p_gran_c2:
            _p_dim = st.selectbox('分析维度', ['按日期', '按渠道', '按店铺', '按品类', '按型号',
                                               '按产品线', '按营销场景', '按推广计划'], key='tr_promo_dim')
        with _p_gran_c3:
            _p_comp_mode = st.radio('对比模式', ['同比', '上月同期'], horizontal=True, key='p_comp_mode')
        _p_cmp_suffix = '同比' if _p_comp_mode == '同比' else '上月同期'
        # ── 推广维度字段映射 ──
        _p_dim_map = {'按渠道': '_渠道', '按店铺': '_店铺', '按品类': '_品类', '按型号': '_型号',
                       '按产品线': '产品线', '按营销场景': '_营销场景', '按推广计划': '_推广计划'}
        _p_dim_field = _p_dim_map.get(_p_dim, '')
        _p_use_dim = _p_dim_field and _p_dim != '按日期'

        # ── 构建推广按日聚合 (dim_val, date) → {promo_fields} ──
        _p_fields = ['_花费', '_展现数', '_点击数', '_直接订单金额', '_总订单金额', '_直接订单量', '_总成交订单量']
        _promo_day = {}
        for r in promo_filtered:
            d = r.get('_date', '')
            if not d or not (str(start) <= d <= str(end)): continue
            if _p_use_dim:
                if _p_dim_field == '产品线':
                    dv = r.get('产品线', '') or '未标注'
                elif _p_dim_field == '_推广计划':
                    dv = r.get('推广计划', '') or r.get('计划ID', '') or '未标注'
                else:
                    dv = r.get(_p_dim_field, '') or '未标注'
            else:
                dv = ''
            key = (dv, d)
            if key not in _promo_day:
                _promo_day[key] = {fk: 0.0 for fk in _p_fields}
            for fk in _p_fields:
                _promo_day[key][fk] += float(r.get(fk, 0) or 0)

        # ── 全时段推广日聚合（用于YoY） ──
        _promo_all_day = {}
        for r in promo_rows:
            d = r.get('_date', '')
            if not d: continue
            # 全局筛选
            if channel and r.get('_渠道', '') not in channel: continue
            if store and r.get('_店铺', '') not in store: continue
            if category and r.get('_品类', '') not in category: continue
            if model and r.get('_型号', '') not in model: continue
            if _p_use_dim:
                if _p_dim_field == '产品线':
                    dv = r.get('产品线', '') or '未标注'
                elif _p_dim_field == '_推广计划':
                    dv = r.get('推广计划', '') or r.get('计划ID', '') or '未标注'
                else:
                    dv = r.get(_p_dim_field, '') or '未标注'
            else:
                dv = ''
            key = (dv, d)
            if key not in _promo_all_day:
                _promo_all_day[key] = {fk: 0.0 for fk in _p_fields}
            for fk in _p_fields:
                _promo_all_day[key][fk] += float(r.get(fk, 0) or 0)

        # ── 对比期日查找（同比/上月同期均通过 _promo_all_day 查找）──
        def _p_date_yoy(dt_str):
            try:
                dt_obj = datetime.datetime.strptime(dt_str, '%Y-%m-%d').date()
                return str(dt_obj.replace(year=dt_obj.year - 1))
            except ValueError:
                return None
        def _p_date_mom(dt_str):
            """上月同日"""
            try:
                dt_obj = datetime.date.fromisoformat(dt_str[:10])
                y, mo, d = dt_obj.year, dt_obj.month, dt_obj.day
                if mo == 1:
                    nm, ny = 12, y - 1
                else:
                    nm, ny = mo - 1, y
                import calendar as _cal
                max_day = _cal.monthrange(ny, nm)[1]
                return str(datetime.date(ny, nm, min(d, max_day)))
            except Exception:
                return None
        _p_ly_day = {}
        for (dv, dt_str), v in _promo_day.items():
            if _p_comp_mode == '同比':
                cmp_dt = _p_date_yoy(dt_str)
            else:
                cmp_dt = _p_date_mom(dt_str)
            _p_ly_day[(dv, dt_str)] = _promo_all_day.get((dv, cmp_dt)) or {} if cmp_dt else {}

        # ── 销售成交金额预计算（用于推广费率：花费/销售支付金额）──
        # 按日期汇总销售成交金额（跨所有维度）
        _p_sales_day = {}
        for (dv_s, dt_str), v in _sales_day.items():
            _p_sales_day[dt_str] = _p_sales_day.get(dt_str, 0) + v['支付金额']
        # 按月汇总
        _p_sales_month = {}
        for dt_str, amt in _p_sales_day.items():
            ym = dt_str[:7]
            _p_sales_month[ym] = _p_sales_month.get(ym, 0) + amt
        # 按推广维度映射销售维度，汇总各维度值销售成交金额
        _sales_dim_map = {'_渠道': '渠道', '_店铺': '店铺', '_品类': '品类', '_型号': '型号'}
        _sales_dim_field = _sales_dim_map.get(_p_dim_field, None)
        _p_dim_sales = {}
        if _sales_dim_field:
            for row in daily:
                dt = row.get('日期', '')
                if not dt or not (str(start) <= dt <= str(end)): continue
                dv_s = row.get(_sales_dim_field, '') or '未标注'
                _p_dim_sales[dv_s] = _p_dim_sales.get(dv_s, 0) + float(row.get('支付金额', 0) or 0)
        # 全时段销售总额
        _p_total_sales = sum(v['支付金额'] for v in _sales_month.values()) if _sales_month else 0

        # ── 选择了分析维度时，显示维度汇总表（不按日/月拆分）──
        if _p_use_dim:
            # 按维度值汇总（合并所有日期）
            _p_dim_agg = {}
            for (dv, dt_str), v in _promo_day.items():
                if dv not in _p_dim_agg:
                    _p_dim_agg[dv] = {fk: 0.0 for fk in _p_fields}
                for fk in _p_fields:
                    _p_dim_agg[dv][fk] += v[fk]
            # 去年同期按维度汇总
            _p_dim_ly = {}
            for (dv, dt_str), v in _p_ly_day.items():
                if not v: continue
                if dv not in _p_dim_ly:
                    _p_dim_ly[dv] = {fk: 0.0 for fk in _p_fields}
                for fk in _p_fields:
                    _p_dim_ly[dv][fk] += v[fk]
            _pdim_total_amt = sum(v['_总订单金额'] for v in _p_dim_agg.values()) or 1
            _pdim_total_spend = sum(v['_花费'] for v in _p_dim_agg.values()) or 1
            _pdim_tbl = []
            for dv, v in sorted(_p_dim_agg.items()):
                spend = v['_花费']; clicks = v['_点击数']; impress = v['_展现数']
                direct_amt = v['_直接订单金额']; total_amt_v = v['_总订单金额']
                total_orders = v['_总成交订单量']
                cpc = spend / clicks if clicks else 0
                droi = direct_amt / spend if spend else 0
                troi = total_amt_v / spend if spend else 0
                tcvr = total_orders / clicks * 100 if clicks else 0
                ly = _p_dim_ly.get(dv, {})
                ly_spend = ly.get('_花费', 0); ly_clicks = ly.get('_点击数', 0)
                ly_direct = ly.get('_直接订单金额', 0); ly_total = ly.get('_总订单金额', 0)
                ly_torders = ly.get('_总成交订单量', 0)
                ly_cpc = ly_spend / ly_clicks if ly_clicks else 0
                ly_droi = ly_direct / ly_spend if ly_spend else 0
                ly_troi = ly_total / ly_spend if ly_spend else 0
                ly_tcvr = ly_torders / ly_clicks * 100 if ly_clicks else 0
                def _pdct(c, p):
                    if p and p != 0: return f'{(c-p)/p*100:+.1f}%'
                    return '--'
                _pdim_tbl.append({
                    '维度': dv,
                    '花费': f'¥{spend:,.0f}',
                    '花费占比': f'{spend/_pdim_total_spend*100:.2f}%' if _pdim_total_spend else '0.00%',
                    '费率': f'{spend/(_p_dim_sales.get(dv, _p_total_sales) if _sales_dim_field else _p_total_sales)*100:.2f}%' if (_p_dim_sales.get(dv, _p_total_sales) if _sales_dim_field else _p_total_sales) else '--',
                    'CPC': f'¥{cpc:.2f}',
                    '点击数': f'{int(clicks):,}',
                    '点击率': f'{clicks/impress*100:.2f}%' if impress else '0.00%',
                    '直接订单金额': f'¥{direct_amt:,.0f}',
                    '总订单金额': f'¥{total_amt_v:,.0f}',
                    '总金额占比': f'{total_amt_v/_pdim_total_amt*100:.2f}%',
                    '直接ROI': f'{droi:.2f}', '总ROI': f'{troi:.2f}',
                    '总转化率': f'{tcvr:.2f}%',
                    f'花费{_p_cmp_suffix}': _pdct(spend, ly_spend),
                    f'直接ROI{_p_cmp_suffix}': _pdct(droi, ly_droi),
                    f'总ROI{_p_cmp_suffix}': _pdct(troi, ly_troi),
                    f'CPC{_p_cmp_suffix}': _pdct(cpc, ly_cpc),
                    f'转化率{_p_cmp_suffix}': _pdct(tcvr, ly_tcvr),
                })
            # 合计行
            _pdt_spend = sum(v['_花费'] for v in _p_dim_agg.values())
            _pdt_clicks = sum(v['_点击数'] for v in _p_dim_agg.values())
            _pdt_impress = sum(v['_展现数'] for v in _p_dim_agg.values())
            _pdt_direct = sum(v['_直接订单金额'] for v in _p_dim_agg.values())
            _pdt_total = sum(v['_总订单金额'] for v in _p_dim_agg.values())
            _pdt_torders = sum(v['_总成交订单量'] for v in _p_dim_agg.values())
            _pdt_cpc = _pdt_spend / _pdt_clicks if _pdt_clicks else 0
            _pdt_droi = _pdt_direct / _pdt_spend if _pdt_spend else 0
            _pdt_troi = _pdt_total / _pdt_spend if _pdt_spend else 0
            _pdt_tcvr = _pdt_torders / _pdt_clicks * 100 if _pdt_clicks else 0
            _ly_pdt_spend = sum(v['_花费'] for v in _p_dim_ly.values())
            _ly_pdt_clicks = sum(v['_点击数'] for v in _p_dim_ly.values())
            _ly_pdt_direct = sum(v['_直接订单金额'] for v in _p_dim_ly.values())
            _ly_pdt_total = sum(v['_总订单金额'] for v in _p_dim_ly.values())
            _ly_pdt_torders = sum(v['_总成交订单量'] for v in _p_dim_ly.values())
            _ly_pdt_cpc = _ly_pdt_spend / _ly_pdt_clicks if _ly_pdt_clicks else 0
            _ly_pdt_droi = _ly_pdt_direct / _ly_pdt_spend if _ly_pdt_spend else 0
            _ly_pdt_troi = _ly_pdt_total / _ly_pdt_spend if _ly_pdt_spend else 0
            _ly_pdt_tcvr = _ly_pdt_torders / _ly_pdt_clicks * 100 if _ly_pdt_clicks else 0
            def _pdct2(c, p):
                if p and p != 0: return f'{(c-p)/p*100:+.1f}%'
                return '--'
            _pdim_tbl.append({
                '维度': '合计',
                '花费': f'¥{_pdt_spend:,.0f}',
                '花费占比': '100.00%',
                '费率': f'{_pdt_spend/_p_total_sales*100:.2f}%' if _p_total_sales else '--',
                'CPC': f'¥{_pdt_cpc:.2f}',
                '点击数': f'{int(_pdt_clicks):,}',
                '点击率': f'{_pdt_clicks/_pdt_impress*100:.2f}%' if _pdt_impress else '0.00%',
                '直接订单金额': f'¥{_pdt_direct:,.0f}',
                '总订单金额': f'¥{_pdt_total:,.0f}', '总金额占比': '100.00%',
                '直接ROI': f'{_pdt_droi:.2f}', '总ROI': f'{_pdt_troi:.2f}',
                '总转化率': f'{_pdt_tcvr:.2f}%',
                f'花费{_p_cmp_suffix}': _pdct2(_pdt_spend, _ly_pdt_spend),
                f'直接ROI{_p_cmp_suffix}': _pdct2(_pdt_droi, _ly_pdt_droi),
                f'总ROI{_p_cmp_suffix}': _pdct2(_pdt_troi, _ly_pdt_troi),
                f'CPC{_p_cmp_suffix}': _pdct2(_pdt_cpc, _ly_pdt_cpc),
                f'转化率{_p_cmp_suffix}': _pdct2(_pdt_tcvr, _ly_pdt_tcvr),
            })
            # 排序控件 — 默认按花费降序
            _ps_sort_cols = ['维度','花费','花费占比','费率','CPC','点击数','点击率','直接订单金额','总订单金额','直接ROI','总ROI','总转化率']
            _psc1, _psc2 = st.columns([2, 1])
            with _psc1:
                _ps_sort_by = st.selectbox('排序字段', _ps_sort_cols, index=1, key='promo_dim_sort_col')
            with _psc2:
                _ps_sort_desc = st.radio('', ['降序', '升序'], horizontal=True, key='promo_dim_sort_dir', index=0)
            _ps_data = [r for r in _pdim_tbl if r.get('维度') != '合计']
            _ps_total = [r for r in _pdim_tbl if r.get('维度') == '合计']
            if _ps_sort_by == '维度':
                _ps_data.sort(key=lambda r: r.get('维度', ''), reverse=(_ps_sort_desc == '降序'))
            else:
                _ps_data.sort(key=lambda r: _parse_num(r.get(_ps_sort_by, 0)), reverse=(_ps_sort_desc == '降序'))
            _pd_headers = ['维度','花费','花费占比','费率','CPC','点击数','点击率','直接订单金额','总订单金额','总金额占比','直接ROI','总ROI','总转化率',
                           f'花费{_p_cmp_suffix}',f'直接ROI{_p_cmp_suffix}',f'总ROI{_p_cmp_suffix}',f'CPC{_p_cmp_suffix}',f'转化率{_p_cmp_suffix}']
            _render_html_table(_ps_data + _ps_total, _pd_headers, _pd_headers, title=f'📢 推广{_p_dim}汇总')
            _render_download_panel(_ps_data + _ps_total, _pd_headers, 'promo_dim_summary.csv', '📥 下载推广维度汇总')
        else:
            # ── 推广日度表格 ──
            _pd_total_amt = sum(v['_总订单金额'] for v in _promo_day.values()) or 1
            _pd_total_spend = sum(v['_花费'] for v in _promo_day.values()) or 1
            _pd_tbl = []
            for (dv, dt_str), v in sorted(_promo_day.items()):
                ly = _p_ly_day.get((dv, dt_str), {})
                spend = v['_花费']; clicks = v['_点击数']; impress = v['_展现数']
                direct_amt = v['_直接订单金额']; total_amt_v = v['_总订单金额']
                total_orders = v['_总成交订单量']
                cpc = spend / clicks if clicks else 0
                droi = direct_amt / spend if spend else 0
                troi = total_amt_v / spend if spend else 0
                tcvr = total_orders / clicks * 100 if clicks else 0
                ly_spend = ly.get('_花费', 0); ly_clicks = ly.get('_点击数', 0)
                ly_direct_amt = ly.get('_直接订单金额', 0); ly_total_amt_v = ly.get('_总订单金额', 0)
                ly_total_orders = ly.get('_总成交订单量', 0)
                ly_cpc = ly_spend / ly_clicks if ly_clicks else 0
                ly_droi = ly_direct_amt / ly_spend if ly_spend else 0
                ly_troi = ly_total_amt_v / ly_spend if ly_spend else 0
                ly_tcvr = ly_total_orders / ly_clicks * 100 if ly_clicks else 0
                def _ppct(c, p):
                    if p and p != 0: return f"{(c-p)/p*100:+.1f}%"
                    return '--'
                row = {
                    '日期': dt_str,
                    '花费': f"¥{spend:,.0f}",
                    '花费占比': f"{spend/_pd_total_spend*100:.2f}%" if _pd_total_spend else '0.00%',
                    '费率': f"{spend/_p_sales_day.get(dt_str, 0)*100:.2f}%" if _p_sales_day.get(dt_str, 0) else '--',
                    'CPC': f"¥{cpc:.2f}",
                    '点击数': f"{int(clicks):,}",
                    '点击率': f"{clicks/impress*100:.2f}%" if impress else '0.00%',
                    '直接订单金额': f"¥{direct_amt:,.0f}",
                    '总订单金额': f"¥{total_amt_v:,.0f}",
                    '总金额占比': f"{total_amt_v/_pd_total_amt*100:.2f}%",
                    '直接ROI': f"{droi:.2f}",
                    '总ROI': f"{troi:.2f}",
                    '总转化率': f"{tcvr:.2f}%",
                    f'花费{_p_cmp_suffix}': _ppct(spend, ly_spend),
                    f'直接ROI{_p_cmp_suffix}': _ppct(droi, ly_droi),
                    f'总ROI{_p_cmp_suffix}': _ppct(troi, ly_troi),
                    f'CPC{_p_cmp_suffix}': _ppct(cpc, ly_cpc),
                    f'转化率{_p_cmp_suffix}': _ppct(tcvr, ly_tcvr),
                }
                if _p_use_dim:
                    row['维度'] = dv
                _pd_tbl.append(row)
            # 合计行
            if _pd_tbl:
                _pdt_spend = sum(v['_花费'] for v in _promo_day.values())
                _pdt_clicks = sum(v['_点击数'] for v in _promo_day.values())
                _pdt_direct = sum(v['_直接订单金额'] for v in _promo_day.values())
                _pdt_total = sum(v['_总订单金额'] for v in _promo_day.values())
                _pdt_torders = sum(v['_总成交订单量'] for v in _promo_day.values())
                _pdt_cpc = _pdt_spend / _pdt_clicks if _pdt_clicks else 0
                _pdt_droi = _pdt_direct / _pdt_spend if _pdt_spend else 0
                _pdt_troi = _pdt_total / _pdt_spend if _pdt_spend else 0
                _pdt_tcvr = _pdt_torders / _pdt_clicks * 100 if _pdt_clicks else 0
                _ly_pdt_spend = sum(v['_花费'] for v in _p_ly_day.values() if v)
                _ly_pdt_clicks = sum(v['_点击数'] for v in _p_ly_day.values() if v)
                _ly_pdt_direct = sum(v['_直接订单金额'] for v in _p_ly_day.values() if v)
                _ly_pdt_total = sum(v['_总订单金额'] for v in _p_ly_day.values() if v)
                _ly_pdt_torders = sum(v['_总成交订单量'] for v in _p_ly_day.values() if v)
                _ly_pdt_cpc = _ly_pdt_spend / _ly_pdt_clicks if _ly_pdt_clicks else 0
                _ly_pdt_droi = _ly_pdt_direct / _ly_pdt_spend if _ly_pdt_spend else 0
                _ly_pdt_troi = _ly_pdt_total / _ly_pdt_spend if _ly_pdt_spend else 0
                _ly_pdt_tcvr = _ly_pdt_torders / _ly_pdt_clicks * 100 if _ly_pdt_clicks else 0
                def _ppct2(c, p):
                    if p and p != 0: return f"{(c-p)/p*100:+.1f}%"
                    return '--'
                total_row = {
                    '日期': '总计',
                    '花费': f"¥{_pdt_spend:,.0f}",
                    '花费占比': '100.00%',
                    '费率': f"{_pdt_spend/_p_total_sales*100:.2f}%" if _p_total_sales else '--',
                    'CPC': f"¥{_pdt_cpc:.2f}", '点击数': f"{int(_pdt_clicks):,}", '点击率': '--',
                    '直接订单金额': f"¥{_pdt_direct:,.0f}",
                    '总订单金额': f"¥{_pdt_total:,.0f}", '总金额占比': '100.00%',
                    '直接ROI': f"{_pdt_droi:.2f}", '总ROI': f"{_pdt_troi:.2f}",
                    '总转化率': f"{_pdt_tcvr:.2f}%",
                    f'花费{_p_cmp_suffix}': _ppct2(_pdt_spend, _ly_pdt_spend),
                    f'直接ROI{_p_cmp_suffix}': _ppct2(_pdt_droi, _ly_pdt_droi),
                    f'总ROI{_p_cmp_suffix}': _ppct2(_pdt_troi, _ly_pdt_troi),
                    f'CPC{_p_cmp_suffix}': _ppct2(_pdt_cpc, _ly_pdt_cpc),
                    f'转化率{_p_cmp_suffix}': _ppct2(_pdt_tcvr, _ly_pdt_tcvr),
                }
                if _p_use_dim:
                    total_row['维度'] = '合计'
                _pd_tbl.append(total_row)
            _pd_headers = ['日期','花费','花费占比','费率','CPC','点击数','点击率','直接订单金额','总订单金额','总金额占比','直接ROI','总ROI','总转化率',
                           f'花费{_p_cmp_suffix}',f'直接ROI{_p_cmp_suffix}',f'总ROI{_p_cmp_suffix}',f'CPC{_p_cmp_suffix}',f'转化率{_p_cmp_suffix}']
            if _p_use_dim:
                _pd_headers = ['维度'] + _pd_headers
            _render_html_table(_pd_tbl, _pd_headers, _pd_headers, title='📢 推广日度趋势')
            _render_download_panel(_pd_tbl, _pd_headers, 'promo_daily_trend.csv', '📥 下载推广日度趋势')
    
            # ── 推广月度表格 ──
            _pm_dict = {}
            for (dv, dt_str), v in _promo_day.items():
                ym = dt_str[:7]
                key = (dv, ym)
                if key not in _pm_dict:
                    _pm_dict[key] = {fk: 0.0 for fk in _p_fields}
                for fk in _p_fields:
                    _pm_dict[key][fk] += v[fk]
            _pm_total_amt = sum(v['_总订单金额'] for v in _pm_dict.values()) or 1
            _pm_total_spend = sum(v['_花费'] for v in _pm_dict.values()) or 1
            _pm_tbl = []
            for (dv, ym), v in sorted(_pm_dict.items()):
                ly_vals = {}
                for (ly_dv, ly_dt), ly_v in _p_ly_day.items():
                    if ly_dt.startswith(ym) and ly_v:
                        if _p_use_dim and ly_dv != dv: continue
                        for fk in _p_fields:
                            ly_vals[fk] = ly_vals.get(fk, 0) + ly_v[fk]
                ly = ly_vals
                spend = v['_花费']; clicks = v['_点击数']; impress = v['_展现数']
                direct_amt = v['_直接订单金额']; total_amt_v = v['_总订单金额']
                total_orders = v['_总成交订单量']
                cpc = spend / clicks if clicks else 0
                droi = direct_amt / spend if spend else 0
                troi = total_amt_v / spend if spend else 0
                tcvr = total_orders / clicks * 100 if clicks else 0
                ly_spend = ly.get('_花费', 0); ly_clicks = ly.get('_点击数', 0)
                ly_direct_amt = ly.get('_直接订单金额', 0); ly_total_amt_v = ly.get('_总订单金额', 0)
                ly_total_orders = ly.get('_总成交订单量', 0)
                ly_cpc = ly_spend / ly_clicks if ly_clicks else 0
                ly_droi = ly_direct_amt / ly_spend if ly_spend else 0
                ly_troi = ly_total_amt_v / ly_spend if ly_spend else 0
                ly_tcvr = ly_total_orders / ly_clicks * 100 if ly_clicks else 0
                def _pmy(c, p):
                    if p and p != 0: return f"{(c-p)/p*100:+.1f}%"
                    return '--'
                row = {
                    '年月': ym,
                    '花费': f"¥{spend:,.0f}",
                    '花费占比': f"{spend/_pm_total_spend*100:.2f}%" if _pm_total_spend else '0.00%',
                    '费率': f"{spend/_p_sales_month.get(ym, 0)*100:.2f}%" if _p_sales_month.get(ym, 0) else '--',
                    'CPC': f"¥{cpc:.2f}",
                    '点击数': f"{int(clicks):,}",
                    '点击率': f"{clicks/impress*100:.2f}%" if impress else '0.00%',
                    '直接订单金额': f"¥{direct_amt:,.0f}",
                    '总订单金额': f"¥{total_amt_v:,.0f}",
                    '总金额占比': f"{total_amt_v/_pm_total_amt*100:.2f}%",
                    '直接ROI': f"{droi:.2f}", '总ROI': f"{troi:.2f}",
                    '总转化率': f"{tcvr:.2f}%",
                    f'花费{_p_cmp_suffix}': _pmy(spend, ly_spend),
                    f'直接ROI{_p_cmp_suffix}': _pmy(droi, ly_droi),
                    f'总ROI{_p_cmp_suffix}': _pmy(troi, ly_troi),
                    f'CPC{_p_cmp_suffix}': _pmy(cpc, ly_cpc),
                    f'转化率{_p_cmp_suffix}': _pmy(tcvr, ly_tcvr),
                }
                if _p_use_dim:
                    row['维度'] = dv
                _pm_tbl.append(row)
            # 推广月度合计行
            if _pm_tbl:
                _pmt_spend = sum(v['_花费'] for v in _pm_dict.values())
                _pmt_clicks = sum(v['_点击数'] for v in _pm_dict.values())
                _pmt_direct = sum(v['_直接订单金额'] for v in _pm_dict.values())
                _pmt_total = sum(v['_总订单金额'] for v in _pm_dict.values())
                _pmt_torders = sum(v['_总成交订单量'] for v in _pm_dict.values())
                _pmt_cpc = _pmt_spend / _pmt_clicks if _pmt_clicks else 0
                _pmt_droi = _pmt_direct / _pmt_spend if _pmt_spend else 0
                _pmt_troi = _pmt_total / _pmt_spend if _pmt_spend else 0
                _pmt_tcvr = _pmt_torders / _pmt_clicks * 100 if _pmt_clicks else 0
                _ly_pmt_spend = sum(v['_花费'] for v in _p_ly_day.values() if v)
                _ly_pmt_clicks = sum(v['_点击数'] for v in _p_ly_day.values() if v)
                _ly_pmt_direct = sum(v['_直接订单金额'] for v in _p_ly_day.values() if v)
                _ly_pmt_total = sum(v['_总订单金额'] for v in _p_ly_day.values() if v)
                _ly_pmt_torders = sum(v['_总成交订单量'] for v in _p_ly_day.values() if v)
                _ly_pmt_cpc = _ly_pmt_spend / _ly_pmt_clicks if _ly_pmt_clicks else 0
                _ly_pmt_droi = _ly_pmt_direct / _ly_pmt_spend if _ly_pmt_spend else 0
                _ly_pmt_troi = _ly_pmt_total / _ly_pmt_spend if _ly_pmt_spend else 0
                _ly_pmt_tcvr = _ly_pmt_torders / _ly_pmt_clicks * 100 if _ly_pmt_clicks else 0
                def _pmy2(c, p):
                    if p and p != 0: return f"{(c-p)/p*100:+.1f}%"
                    return '--'
                trow = {
                    '年月': '合计',
                    '花费': f"¥{_pmt_spend:,.0f}",
                    '花费占比': '100.00%',
                    '费率': f"{_pmt_spend/_p_total_sales*100:.2f}%" if _p_total_sales else '--',
                    'CPC': f"¥{_pmt_cpc:.2f}", '点击数': f"{int(_pmt_clicks):,}", '点击率': '--',
                    '直接订单金额': f"¥{_pmt_direct:,.0f}",
                    '总订单金额': f"¥{_pmt_total:,.0f}", '总金额占比': '100.00%',
                    '直接ROI': f"{_pmt_droi:.2f}", '总ROI': f"{_pmt_troi:.2f}",
                    '总转化率': f"{_pmt_tcvr:.2f}%",
                    f'花费{_p_cmp_suffix}': _pmy2(_pmt_spend, _ly_pmt_spend),
                    f'直接ROI{_p_cmp_suffix}': _pmy2(_pmt_droi, _ly_pmt_droi),
                    f'总ROI{_p_cmp_suffix}': _pmy2(_pmt_troi, _ly_pmt_troi),
                    f'CPC{_p_cmp_suffix}': _pmy2(_pmt_cpc, _ly_pmt_cpc),
                    f'转化率{_p_cmp_suffix}': _pmy2(_pmt_tcvr, _ly_pmt_tcvr),
                }
                if _p_use_dim:
                    trow['维度'] = '合计'
                _pm_tbl.append(trow)
            _pm_headers = ['年月','花费','花费占比','费率','CPC','点击数','点击率','直接订单金额','总订单金额','总金额占比','直接ROI','总ROI','总转化率',
                           f'花费{_p_cmp_suffix}',f'直接ROI{_p_cmp_suffix}',f'总ROI{_p_cmp_suffix}',f'CPC{_p_cmp_suffix}',f'转化率{_p_cmp_suffix}']
            if _p_use_dim:
                _pm_headers = ['维度'] + _pm_headers
            _render_html_table(_pm_tbl, _pm_headers, _pm_headers, title='📢 推广月度趋势')
            _render_download_panel(_pm_tbl, _pm_headers, 'promo_monthly_trend.csv', '📥 下载推广月度趋势')


# ═══════════════════════════════════════════════════════════════
# 智能诊断工具函数（华为方法论注入）
# ═══════════════════════════════════════════════════════════════

def _shapley_decompose_gmv(cur_sum, prev_sum):
    """
    Shapley值法分解GMV变化（华为式归因）
    GMV = 访客数 × 转化率 × 客单价
    返回各因子对ΔGMV的精确贡献（消除交互效应偏差）
    """
    V_cur  = cur_sum.get('商品访客数', 0)
    C_cur  = cur_sum.get('支付转化率', 0)
    A_cur  = cur_sum.get('客单价', 0)
    V_prev = prev_sum.get('商品访客数', 0)
    C_prev = prev_sum.get('支付转化率', 0)
    A_prev = prev_sum.get('客单价', 0)

    GMV_cur  = V_cur  * C_cur  * A_cur
    GMV_prev = V_prev * C_prev * A_prev
    delta = GMV_cur - GMV_prev

    if abs(delta) < 0.01 or V_prev == 0:
        return {'流量效应': 0, '转化效应': 0, '客单效应': 0, '交互效应': 0, 'delta': delta}

    # 6种排列求Shapley值（消除顺序偏差）
    # 排列1: V→C→A
    s1_v = (V_cur - V_prev) * C_prev * A_prev
    s1_c = V_cur * (C_cur - C_prev) * A_prev
    s1_a = V_cur * C_cur * (A_cur - A_prev)

    # 排列2: V→A→C
    s2_v = (V_cur - V_prev) * C_prev * A_prev
    s2_a = V_cur * (A_cur - A_prev) * C_prev
    s2_c = V_cur * A_cur * (C_cur - C_prev)

    # 排列3: C→V→A
    s3_c = (C_cur - C_prev) * V_prev * A_prev
    s3_v = C_cur * (V_cur - V_prev) * A_prev
    s3_a = C_cur * V_cur * (A_cur - A_prev)

    # 排列4: C→A→V
    s4_c = (C_cur - C_prev) * V_prev * A_prev
    s4_a = C_cur * (A_cur - A_prev) * V_prev
    s4_v = C_cur * A_cur * (V_cur - V_prev)

    # 排列5: A→V→C
    s5_a = (A_cur - A_prev) * V_prev * C_prev
    s5_v = A_cur * (V_cur - V_prev) * C_prev
    s5_c = A_cur * V_cur * (C_cur - C_prev)

    # 排列6: A→C→V
    s6_a = (A_cur - A_prev) * V_prev * C_prev
    s6_c = A_cur * (C_cur - C_prev) * V_prev
    s6_v = A_cur * C_cur * (V_cur - V_prev)

    shap_v = (s1_v + s2_v + s3_v + s4_v + s5_v + s6_v) / 6
    shap_c = (s1_c + s2_c + s3_c + s4_c + s5_c + s6_c) / 6
    shap_a = (s1_a + s2_a + s3_a + s4_a + s5_a + s6_a) / 6

    return {
        '流量效应': shap_v, '转化效应': shap_c, '客单效应': shap_a,
        '交互效应': delta - shap_v - shap_c - shap_a,
        'delta': delta
    }


def _gen_one_line_summary(gmv_g, shapley, ch_model_issues, vis_g, cvr_g, aov_g):
    """生成华为式一句话核心发现"""
    if gmv_g is None:
        return "⚠️ 数据不足，无法生成诊断总结。"
    direction = "增长" if gmv_g >= 0 else "下滑"
    pct_s = f"{gmv_g*100:+.1f}%"

    # 找最大Shapley贡献因子
    factors = [
        ('流量', shapley.get('流量效应', 0), vis_g),
        ('转化率', shapley.get('转化效应', 0), cvr_g),
        ('客单价', shapley.get('客单效应', 0), aov_g),
    ]
    main = max(factors, key=lambda x: abs(x[1]))
    main_pct = f"{main[2]*100:+.1f}%" if main[2] is not None else "--"

    # 找最大拖累型号
    worst_loss = 0; worst_name = ''
    for m in (ch_model_issues or []):
        loss = m.get('上期GMV', 0) - m.get('本期GMV', 0)
        if loss > worst_loss:
            worst_loss = loss
            worst_name = m.get('型号', '')

    summary = f"GMV环比{pct_s}，主因是<b>{main[0]}</b>变动{main_pct}"
    if worst_name and worst_loss > 0:
        summary += f"，其中<b>[{worst_name}]</b>单型号损失约¥{worst_loss:,.0f}"
    summary += "。"

    # 添加Shapley贡献度
    total_effect = abs(shapley.get('流量效应',0)) + abs(shapley.get('转化效应',0)) + abs(shapley.get('客单效应',0))
    if total_effect > 0:
        vp = abs(shapley['流量效应']) / total_effect * 100
        cp = abs(shapley['转化效应']) / total_effect * 100
        ap = abs(shapley['客单效应']) / total_effect * 100
        summary += f"<br><small style='color:#64748b;'>归因：流量 {vp:.0f}% | 转化 {cp:.0f}% | 客单价 {ap:.0f}%</small>"

    return summary


# ═══════════════════════════════════════════════════════════════
# TAB 4: 智能诊断 V4（华为复盘方法论注入）
# ═══════════════════════════════════════════════════════════════
# 构建筛选标签文字
_filter_parts = []
if channel: _filter_parts.append(f'渠道={"+".join(channel)}')
if store: _filter_parts.append(f'店铺={"+".join(store)}')
if category: _filter_parts.append(f'品类={"+".join(category)}')
if model: _filter_parts.append(f'型号={"+".join(model)}')
_filter_label = ' | '.join(_filter_parts) if _filter_parts else '全域'

with tabs[4]:
    st.markdown('<div class="section-title">🔍 智能诊断 V4 — 作战指挥室</div>', unsafe_allow_html=True)
    _cmp_label = f'上期 {_t2_prev_s} ~ {_t2_prev_e}'
    st.caption(f'诊断区间：{s} ~ {e} | 筛选范围：{_filter_label} | 对比区间：{_cmp_label}')

    # ══════════════════════════════════════
    # A. 核心数据准备（基于当前筛选条件）
    # ══════════════════════════════════════
    cur_rows_all = [r for r in daily if r.get('渠道', '') != '小豚天猫']
    cur_sum = summarize(cur_rows_all)  # 智能诊断模块排除小豚天猫后重新汇总

    def _agg_by_dims(rows, dims):
        out = {}
        for r in rows:
            key = tuple(r.get(d) or '未标注' for d in dims)
            out.setdefault(key, {m: 0.0 for m in METRICS})
            for m in METRICS:
                out[key][m] += float(r.get(m, 0) or 0)
        for k, v in out.items():
            v['支付转化率'] = v['支付买家数'] / v['商品访客数'] if v['商品访客数'] else 0
            v['客单价'] = v['支付金额'] / v['支付买家数'] if v['支付买家数'] else 0
            v['退款率'] = v['成功退款金额'] / v['支付金额'] if v['支付金额'] else 0
        return out

    prev_rows_all_raw = []
    for r in data['daily']:
        d = r.get('日期', '')
        if len(d) == 7: d = d + '-01'
        if _t2_prev_s <= d <= _t2_prev_e: prev_rows_all_raw.append(r)
    prev_rows_all = []
    for r in prev_rows_all_raw:
        if r.get('渠道', '') == '小豚天猫': continue
        if channel and r.get('渠道') not in channel: continue
        if store and r.get('店铺') not in store: continue
        if category and r.get('品类') not in category: continue
        if model and r.get('型号') not in model: continue
        prev_rows_all.append(r)

    cur_by_model   = _agg_by_dims(cur_rows_all, ['渠道','品类','型号'])
    cur_by_cat     = _agg_by_dims(cur_rows_all, ['渠道','品类'])
    cur_by_channel = _agg_by_dims(cur_rows_all, ['渠道'])
    prev_by_model  = _agg_by_dims(prev_rows_all, ['渠道','品类','型号'])
    prev_by_cat    = _agg_by_dims(prev_rows_all, ['渠道','品类'])
    prev_by_channel= _agg_by_dims(prev_rows_all, ['渠道'])

    prev_sum_all = summarize(prev_rows_all)
    def _global_chg(k):
        c = cur_sum.get(k, 0); p = prev_sum_all.get(k, 0)
        return (c - p) / p if p else None

    gmv_g = _global_chg('支付金额')
    vis_g = _global_chg('商品访客数')
    cvr_g = _global_chg('支付转化率')
    aov_g = _global_chg('客单价')
    ref_g = _global_chg('退款率')

    WARN_T = -0.05; DANGER_T = -0.15

    def _pct(v): return f'{v*100:+.1f}%' if v is not None else '--'

    # ── Shapley归因 ──
    shapley = _shapley_decompose_gmv(cur_sum, prev_sum_all)

    # ── 同比数据（GAP2）──
    yoy_sum = None
    gmv_yoy = vis_yoy = cvr_yoy = aov_yoy = ref_yoy = None
    try:
        yoy_rows_all = []
        for r in data['daily']:
            d = r.get('日期', '')
            if len(d) == 7: d = d + '-01'
            if yoy_s <= d <= yoy_e: yoy_rows_all.append(r)
        yoy_rows_filtered = []
        for r in yoy_rows_all:
            if r.get('渠道', '') == '小豚天猫': continue
            if channel and r.get('渠道') not in channel: continue
            if store and r.get('店铺') not in store: continue
            if category and r.get('品类') not in category: continue
            if model and r.get('型号') not in model: continue
            yoy_rows_filtered.append(r)
        if yoy_rows_filtered:
            yoy_sum = summarize(yoy_rows_filtered)
            def _yoy_chg(k):
                c = cur_sum.get(k, 0); p = yoy_sum.get(k, 0)
                return (c - p) / p if p else None
            gmv_yoy = _yoy_chg('支付金额')
            vis_yoy = _yoy_chg('商品访客数')
            cvr_yoy = _yoy_chg('支付转化率')
            aov_yoy = _yoy_chg('客单价')
            ref_yoy = _yoy_chg('退款率')
    except Exception:
        pass

    # ── 目标达成（GAP1）──
    gmv_target = None; gmv_target_gap = None
    try:
        if 'targets' in dir() and targets:
            from collections import Counter
            _ym_counter = Counter()
            for r in cur_rows_all:
                d = r.get('日期', '')[:7]
                if d: _ym_counter[d] += 1
            if _ym_counter:
                _main_ym = _ym_counter.most_common(1)[0][0]
                _tdata = targets.get(_main_ym, {})
                _shop_rows_t = _tdata.get('shop', [])
                if _shop_rows_t:
                    # 获取目标数据中属于当前筛选期间(s~e)的日期列
                    _tgt_dates = targets.get(_main_ym, {}).get('dates', [])
                    _period_dates = [d for d in _tgt_dates if s <= d <= e]
                    # 汇总所有店铺（排除天猫小豚）的成交金额目标
                    gmv_target = 0.0
                    for tr in _shop_rows_t:
                        _shop = str(tr.get('店铺', ''))
                        if _shop == '天猫小豚':
                            continue
                        if '成交金额' in str(tr.get('指标', '')) and '目标' in str(tr.get('指标', '')):
                            if _period_dates:
                                gmv_target += sum(tr.get(d, 0.0) for d in _period_dates)
                            else:
                                gmv_target += tr.get('合计', 0.0)
                    if gmv_target > 0:
                        gmv_target_gap = (cur_sum.get('支付金额', 0) - gmv_target) / gmv_target
    except Exception:
        pass

    # ── 健康评分 ──
    scores = {}
    for name, val in [('GMV',gmv_g),('流量',vis_g),('转化率',cvr_g),('客单价',aov_g),('退款率',ref_g)]:
        if val is None: scores[name] = 100
        elif val > WARN_T: scores[name] = 100
        elif val > DANGER_T: scores[name] = 60 + int((val - WARN_T) / (0 - WARN_T) * 40)
        else: scores[name] = max(0, int(val / DANGER_T * 60))
    if ref_g is not None and ref_g > 0:
        scores['退款率'] = max(0, 100 - abs(ref_g) * 300)
    health_score = sum(scores.values()) / len(scores)

    if health_score >= 90:   hv = ('🟢 整体健康', '#22c55e', '各项核心指标表现良好，继续保持现有经营策略。')
    elif health_score >= 70: hv = ('🟡 需要关注', '#f59e0b', f'部分指标出现波动（综合得分{health_score:.0f}/100），建议重点关注下方标红项。')
    elif health_score >= 50: hv = ('🔴 存在风险', '#ef4444', f'多项指标明显下滑（综合得分{health_score:.0f}/100），建议立即执行P0优先级行动。')
    else:                    hv = ('⚠️ 紧急告警', '#dc2626', f'整体经营状况堪忧（综合得分{health_score:.0f}/100），请优先处理所有P0任务！')

    # ══════════════════════════════════════════════════════════════
    # 🏛️ 30-40-30 首页结构：Layer1 一句话总结 → Layer2 三GAP速览 → Layer3 健康+KPI
    # ══════════════════════════════════════════════════════════════

    # ── Layer 1: 一句话核心发现（30% 定位差距）──
    _one_line = _gen_one_line_summary(gmv_g, shapley, [], vis_g, cvr_g, aov_g)
    _summary_bg = '#fef2f2' if (gmv_g is not None and gmv_g < DANGER_T) else ('#fff7ed' if (gmv_g is not None and gmv_g < WARN_T) else '#f0fdf4')
    _summary_border = '#fca5a5' if (gmv_g is not None and gmv_g < DANGER_T) else ('#fdba74' if (gmv_g is not None and gmv_g < WARN_T) else '#86efac')
    st.markdown(
        f"<div style='background:{_summary_bg};border:2px solid {_summary_border};border-radius:16px;"
        f"padding:18px 24px;margin:8px 0 16px 0;'>"
        f"<div style='font-size:12px;color:#64748b;margin-bottom:6px;'>📌 一句话核心发现</div>"
        f"<div style='font-size:17px;font-weight:800;color:#0f172a;line-height:1.6;'>{_one_line}</div>"
        f"</div>",
        unsafe_allow_html=True)

    # ── Layer 2: 三GAP速览（30% 差距定位）──
    st.markdown(
        "<div style='font-size:13px;color:#475569;font-weight:700;margin:4px 0 2px 0;'>"
        "📊 三GAP速览 — 目标达成 · 同比变化 · 环比变化</div>", unsafe_allow_html=True)
    st.markdown(
        "<div style='font-size:10px;color:#94a3b8;margin:0 0 8px 0;'>"
        "GAP=差距(Gap)：GAP1 目标vs实际 | GAP2 今年vs去年同期 | GAP3 本期vs上期</div>", unsafe_allow_html=True)

    gap1, gap2, gap3 = st.columns(3)

    # GAP 1: 目标 vs 实际
    with gap1:
        _gap1_color = '#22c55e'
        _gap1_icon = '✅'
        if gmv_target_gap is not None:
            if gmv_target_gap < -0.10: _gap1_color, _gap1_icon = '#dc2626', '🔴'
            elif gmv_target_gap < -0.05: _gap1_color, _gap1_icon = '#f59e0b', '🟡'
            elif gmv_target_gap < 0: _gap1_color, _gap1_icon = '#ea580c', '🟠'
        _gap1_val = f'{gmv_target_gap*100:+.1f}%' if gmv_target_gap is not None else '--'
        _gap1_target_str = f'¥{gmv_target:,.0f}' if gmv_target else '--'
        _gap1_actual_str = f'¥{cur_sum.get("支付金额",0):,.0f}'
        st.markdown(
            f"<div style='background:#f8fafc;border:1px solid #e2e8f0;border-radius:12px;padding:14px;text-align:center;'>"
            f"<div style='font-size:11px;color:#64748b;font-weight:700;margin-bottom:6px;'>🎯 GAP 1: 目标达成</div>"
            f"<div style='font-size:24px;font-weight:900;color:{_gap1_color};'>{_gap1_icon} {_gap1_val}</div>"
            f"<div style='font-size:10px;color:#94a3b8;margin-top:4px;'>实际 {_gap1_actual_str} / 目标 {_gap1_target_str}</div>"
            f"</div>", unsafe_allow_html=True)

    # GAP 2: 同比变化
    with gap2:
        _gap2_color = '#22c55e'; _gap2_icon = '📈'
        if gmv_yoy is not None:
            if gmv_yoy < -0.10: _gap2_color, _gap2_icon = '#dc2626', '📉'
            elif gmv_yoy < -0.05: _gap2_color, _gap2_icon = '#f59e0b', '📊'
        _gap2_val = f'{gmv_yoy*100:+.1f}%' if gmv_yoy is not None else '--'
        _yoy_label = f'去年同期 {yoy_s}~{yoy_e}' if yoy_sum else '去年同期'
        st.markdown(
            f"<div style='background:#f8fafc;border:1px solid #e2e8f0;border-radius:12px;padding:14px;text-align:center;'>"
            f"<div style='font-size:11px;color:#64748b;font-weight:700;margin-bottom:6px;'>📅 GAP 2: 同比变化</div>"
            f"<div style='font-size:24px;font-weight:900;color:{_gap2_color};'>{_gap2_icon} {_gap2_val}</div>"
            f"<div style='font-size:10px;color:#94a3b8;margin-top:4px;'>vs {_yoy_label}</div>"
            f"</div>", unsafe_allow_html=True)

    # GAP 3: 环比变化
    with gap3:
        _gap3_color = '#22c55e'; _gap3_icon = '📈'
        if gmv_g is not None:
            if gmv_g < DANGER_T: _gap3_color, _gap3_icon = '#dc2626', '📉'
            elif gmv_g < WARN_T: _gap3_color, _gap3_icon = '#f59e0b', '📊'
        _gap3_val = _pct(gmv_g)
        st.markdown(
            f"<div style='background:#f8fafc;border:1px solid #e2e8f0;border-radius:12px;padding:14px;text-align:center;'>"
            f"<div style='font-size:11px;color:#64748b;font-weight:700;margin-bottom:6px;'>🔄 GAP 3: 环比变化</div>"
            f"<div style='font-size:24px;font-weight:900;color:{_gap3_color};'>{_gap3_icon} {_gap3_val}</div>"
            f"<div style='font-size:10px;color:#94a3b8;margin-top:4px;'>vs 上期 {_t2_prev_s}~{_t2_prev_e}</div>"
            f"</div>", unsafe_allow_html=True)

    # ── Shapley归因条 ──
    _total_effect = abs(shapley['流量效应']) + abs(shapley['转化效应']) + abs(shapley['客单效应'])
    if _total_effect > 0:
        _vp = abs(shapley['流量效应']) / _total_effect * 100
        _cp = abs(shapley['转化效应']) / _total_effect * 100
        _ap = abs(shapley['客单效应']) / _total_effect * 100
        _delta_str = f'¥{shapley["delta"]:+,.0f}' if shapley['delta'] != 0 else '¥0'
        _v_color = '#dc2626' if shapley['流量效应'] < 0 else '#22c55e'
        _c_color = '#dc2626' if shapley['转化效应'] < 0 else '#22c55e'
        _a_color = '#dc2626' if shapley['客单效应'] < 0 else '#22c55e'
        st.markdown(
            f"<div style='background:#f1f5f9;border-radius:10px;padding:12px 18px;margin:12px 0 8px 0;'>"
            f"<div style='font-size:11px;color:#64748b;margin-bottom:6px;'>🔬 Shapley归因 — GMV变化分解（Δ={_delta_str}）</div>"
            f"<div style='display:flex;gap:12px;align-items:center;'>"
            f"<div style='flex:1;'><div style='font-size:10px;color:#475569;'>流量</div>"
            f"<div style='background:#e2e8f0;border-radius:6px;height:8px;margin:2px 0;'>"
            f"<div style='background:{_v_color};border-radius:6px;height:8px;width:{_vp:.0f}%;'></div></div>"
            f"<div style='font-size:10px;color:{_v_color};font-weight:700;'>{_vp:.0f}% (¥{shapley['流量效应']:+,.0f})</div></div>"
            f"<div style='flex:1;'><div style='font-size:10px;color:#475569;'>转化率</div>"
            f"<div style='background:#e2e8f0;border-radius:6px;height:8px;margin:2px 0;'>"
            f"<div style='background:{_c_color};border-radius:6px;height:8px;width:{_cp:.0f}%;'></div></div>"
            f"<div style='font-size:10px;color:{_c_color};font-weight:700;'>{_cp:.0f}% (¥{shapley['转化效应']:+,.0f})</div></div>"
            f"<div style='flex:1;'><div style='font-size:10px;color:#475569;'>客单价</div>"
            f"<div style='background:#e2e8f0;border-radius:6px;height:8px;margin:2px 0;'>"
            f"<div style='background:{_a_color};border-radius:6px;height:8px;width:{_ap:.0f}%;'></div></div>"
            f"<div style='font-size:10px;color:{_a_color};font-weight:700;'>{_ap:.0f}% (¥{shapley['客单效应']:+,.0f})</div></div>"
            f"</div></div>",
            unsafe_allow_html=True)

    # ── Layer 3: 健康评分 + 核心KPI（40% 深挖入口）──
    st.markdown('<hr style="margin:14px 0;border:none;border-top:1px dashed #cbd5e1;">', unsafe_allow_html=True)

    sc1, sc2, sc3 = st.columns([1, 4, 2])

    with sc1:
        st.metric('健康评分', f'{health_score:.0f}',
                  help=f"GMV:{scores['GMV']} | 流量:{scores['流量']} | 转化:{scores['转化率']} | 客单价:{scores['客单价']} | 退款:{scores['退款率']}")
    with sc2:
        st.markdown(
            f"<div style='background:{hv[1]}15;border-left:4px solid {hv[1]};border-radius:8px;"
            f"padding:12px 16px;margin-top:8px;'><strong>{hv[0]}</strong>&nbsp;&nbsp;{hv[2]}</div>",
            unsafe_allow_html=True)
    with sc3:
        _p0_cnt = sum(1 for k,v in [('gmv',gmv_g),('vis',vis_g),('cvr',cvr_g)] if v is not None and v < DANGER_T)
        st.metric('P0 级问题', _p0_cnt, help='GMV/流量/转化率中降幅>15%的指标数量')

    # 核心KPI卡片（精简版5卡）
    kpi5_cols = st.columns(5)
    _kpi5 = [
        ('💰 支付金额', gmv_g, cur_sum.get('支付金额',0), prev_sum_all.get('支付金额',0), '¥', False),
        ('👁 访客数',   vis_g, cur_sum.get('商品访客数',0), prev_sum_all.get('商品访客数',0), '', False),
        ('🔄 转化率',  cvr_g, cur_sum.get('支付转化率',0)*100, prev_sum_all.get('支付转化率',0)*100, '', True),
        ('🎫 客单价',  aov_g, cur_sum.get('客单价',0), prev_sum_all.get('客单价',0), '¥', False),
        ('↩️ 退款率',  ref_g, cur_sum.get('退款率',0)*100, prev_sum_all.get('退款率',0)*100, '', True),
    ]
    for col, (mname, mch, cv, pv, pre, ispct) in zip(kpi5_cols, _kpi5):
        with col:
            lvl = 'ok' if mch is None or mch > WARN_T else ('warn' if mch > DANGER_T else 'danger')
            bg = {'danger':'#fef2f2','warn':'#fff7ed','ok':'#f0fdf4'}[lvl]
            border = {'danger':'#fca5a5','warn':'#fdba74','ok':'#86efac'}[lvl]
            cv_s = f'{cv:.2f}%' if ispct else f'{cv:,.0f}'
            pv_s = f'{pv:.2f}%' if ispct else f'{pv:,.0f}'
            ch_s = _pct(mch)
            # Shapley归因提示（替代旧的简单归因）
            hint = ''
            if '支付金额' in mname and mch is not None and mch < 0 and _total_effect > 0:
                parts = []
                if shapley['流量效应'] < -1: parts.append(f"流量{_vp:.0f}%")
                if shapley['转化效应'] < -1: parts.append(f"转化{_cp:.0f}%")
                if shapley['客单效应'] < -1: parts.append(f"客单{_ap:.0f}%")
                hint = f'<div style="font-size:10px;color:#ea580c;margin-top:2px;">主因：{"+".join(parts) if parts else "多因子"}</div>' if parts else ''
            st.markdown(
                f'<div style="background:{bg};border:1px solid {border};border-radius:14px;'
                f'padding:10px;text-align:center;">'
                f'<div style="font-size:11px;color:#64748b;font-weight:700;">{mname}</div>'
                f'<div style="font-size:19px;font-weight:900;color:#0f172a;margin:3px 0;">{pre}{cv_s}</div>'
                f'<div style="font-size:10px;color:#94a3b8;">vs上期 {pre}{pv_s} ({ch_s})</div>{hint}</div>',
                unsafe_allow_html=True)

    # ══════════════════════════════════════════════════════════════
    # 人货场三大板块（子Tab）— Layer 4 深挖根因
    # ══════════════════════════════════════════════════════════════
    diag_tabs = st.tabs(['👥 人（流量&用户）', '📦 货（商品&转化）', '🏪 场（渠道&推广）', '🛠️ 执行清单'])

    # ────────────────────────────────────────────────────────────
    # 【人】：流量结构、渠道流量分布、加购漏斗的人侧、增长亮点
    # ────────────────────────────────────────────────────────────
    with diag_tabs[0]:
        st.markdown('<div class="section-title">👥 "人" — 流量来源 & 用户结构诊断</div>', unsafe_allow_html=True)
        st.caption('关注：访客规模变化 · 渠道流量分布 · 加购人群转化 · 增长亮点用户')

        # 1. 流量健康度概览
        h1c1, h1c2, h1c3, h1c4 = st.columns(4)
        _cur_vis  = cur_sum.get('商品访客数', 0)
        _prev_vis = prev_sum_all.get('商品访客数', 0)
        _cur_buyer  = cur_sum.get('支付买家数', 0)
        _prev_buyer = prev_sum_all.get('支付买家数', 0)
        _cur_cart   = cur_sum.get('商品加购人数', 0)
        _prev_cart  = prev_sum_all.get('商品加购人数', 0)
        _cur_cvr    = cur_sum.get('支付转化率', 0) * 100
        _prev_cvr   = prev_sum_all.get('支付转化率', 0) * 100

        for col, (lbl, cv, pv, is_pct) in zip(
            [h1c1, h1c2, h1c3, h1c4],
            [('访客数', _cur_vis, _prev_vis, False),
             ('成交买家数', _cur_buyer, _prev_buyer, False),
             ('加购人数', _cur_cart, _prev_cart, False),
             ('支付转化率', _cur_cvr, _prev_cvr, True)]):
            chg = (cv - pv) / pv if pv else None
            lvl = 'ok' if chg is None or chg > WARN_T else ('warn' if chg > DANGER_T else 'danger')
            bg = {'danger':'#fef2f2','warn':'#fff7ed','ok':'#f0fdf4'}[lvl]
            brd = {'danger':'#fca5a5','warn':'#fdba74','ok':'#86efac'}[lvl]
            cvs = f'{cv:.2f}%' if is_pct else f'{cv:,.0f}'
            pvs = f'{pv:.2f}%' if is_pct else f'{pv:,.0f}'
            col.markdown(
                f'<div style="background:{bg};border:1px solid {brd};border-radius:12px;padding:10px;text-align:center;">'
                f'<div style="font-size:11px;color:#64748b;font-weight:600;">{lbl}</div>'
                f'<div style="font-size:18px;font-weight:900;color:#0f172a;">{cvs}</div>'
                f'<div style="font-size:10px;color:#94a3b8;">vs上期 {pvs} ({_pct(chg)})</div></div>',
                unsafe_allow_html=True)

        st.markdown('#### 📊 渠道流量分布与变化')

        # 渠道流量分布表
        _ch_flow = []
        for ch_key, cv in cur_by_channel.items():
            pv = prev_by_channel.get(ch_key, {})
            _cv_vis  = cv.get('商品访客数', 0)
            _pv_vis  = pv.get('商品访客数', 0)
            _cv_buy  = cv.get('支付买家数', 0)
            _cv_cvr  = cv.get('支付转化率', 0) * 100
            _pv_cvr  = pv.get('支付转化率', 0) * 100
            _vis_chg = (_cv_vis - _pv_vis) / _pv_vis if _pv_vis else None
            _cvr_chg = (_cv_cvr - _pv_cvr)  # pp差
            _vis_share = _cv_vis / max(_cur_vis, 1) * 100
            _ch_flow.append({
                '渠道': ch_key[0],
                '访客数': f'{_cv_vis:,.0f}',
                '访客占比': f'{_vis_share:.1f}%',
                '访客变化': f"<span style='color:{'#22c55e' if (_vis_chg or 0) >= 0 else '#ef4444'};font-weight:700'>{_pct(_vis_chg)}</span>",
                '成交买家数': f'{_cv_buy:,.0f}',
                '转化率': f'{_cv_cvr:.2f}%',
                '转化率变化(pp)': f"<span style='color:{'#22c55e' if _cvr_chg >= 0 else '#ef4444'}'>{_cvr_chg:+.2f}pp</span>",
            })
        _ch_flow.sort(key=lambda x: float(x['访客占比'].rstrip('%')), reverse=True)
        if _ch_flow:
            _chf_html = _html_table(_ch_flow, height=min(340, len(_ch_flow)*36+50))
            st.markdown(_wrap_fullscreen(_chf_html, title='👥 渠道流量结构')[0], unsafe_allow_html=True)

        # 加购漏斗（人侧视角）
        st.markdown('#### 🔽 全域加购漏斗（人侧）')
        _funnel_data = [
            ('访客数',    _cur_vis,   _prev_vis),
            ('加购人数',  _cur_cart,  _prev_cart),
            ('成交买家数', _cur_buyer, _prev_buyer),
        ]
        _f_cols = st.columns(len(_funnel_data))
        for i, (name, cv, pv) in enumerate(_funnel_data):
            chg = (cv - pv) / pv if pv else None
            rate_to_vis = cv / max(_cur_vis, 1) * 100
            _f_cols[i].metric(
                label=f'{name}',
                value=f'{cv:,.0f}',
                delta=f'{_pct(chg)} (占访客{rate_to_vis:.1f}%)',
                delta_color='normal'
            )
        _add_rate_cur  = _cur_cart  / max(_cur_vis,  1) * 100
        _add_rate_prev = _prev_cart / max(_prev_vis, 1) * 100
        _pay_rate_cur  = _cur_buyer / max(_cur_cart,  1) * 100
        _pay_rate_prev = _prev_buyer / max(_prev_cart, 1) * 100
        st.caption(
            f'加购率：{_add_rate_cur:.2f}% (vs上期 {_add_rate_prev:.2f}%)｜'
            f'加购→成交率：{_pay_rate_cur:.1f}% (vs上期 {_pay_rate_prev:.1f}%)')

        # 增长亮点（新晋型号）
        st.markdown('#### 🌟 增长亮点 — 新晋爆发型号')
        cur_top20 = sorted(cur_by_model.items(), key=lambda x: x[1].get('支付金额',0), reverse=True)[:20]
        rising_stars = []
        for mk_key, mv in cur_top20:
            pv_m = prev_by_model.get(mk_key, {})
            mc = mv.get('支付金额', 0); mp = pv_m.get('支付金额', 0)
            if mp > 0:
                growth = (mc - mp) / mp
                if growth > 0.30 and mc > 500:
                    rising_stars.append({'渠道': mk_key[0], '品类': mk_key[1], '型号': mk_key[2],
                        '上期GMV': mp, '本期GMV': mc, '增速': growth,
                        '访客增幅': _pct((mv.get('商品访客数',0)-pv_m.get('商品访客数',0))/max(pv_m.get('商品访客数',1),1))})
            elif mc > 2000:
                rising_stars.append({'渠道': mk_key[0], '品类': mk_key[1], '型号': mk_key[2],
                    '上期GMV': 0, '本期GMV': mc, '增速': float('inf'), '访客增幅': '新上榜'})
        rising_stars.sort(key=lambda x: x['本期GMV'], reverse=True)
        if rising_stars:
            _rs_rows = []
            for r in rising_stars[:10]:
                sp = "<span style='color:#22c55e;font-weight:700'>🚀 新品</span>" if r['增速'] == float('inf') else f"<span style='color:#22c55e;font-weight:700'>+{r['增速']*100:.0f}%</span>"
                _rs_rows.append({'渠道': r['渠道'], '品类': r['品类'], '型号': r['型号'],
                    '上期GMV': f"¥{r['上期GMV']:,.0f}" if r['上期GMV'] > 0 else '新上榜',
                    '本期GMV': f"¥{r['本期GMV']:,.0f}", '增速': sp, '访客增幅': r['访客增幅']})
            _rs_html = _html_table(_rs_rows, height=min(320, len(_rs_rows)*34+40))
            st.markdown(_wrap_fullscreen(_rs_html, title='⭐ 黑马/明星单品')[0], unsafe_allow_html=True)
        else:
            st.info('ℹ️ 未发现增速>30%的型号（阈值：增速>30% 或 新上榜且GMV>¥500）。')

        _render_download_panel(
            [{'渠道': r['渠道'], '品类': r['品类'], '型号': r['型号'],
              '访客数': cur_by_model.get((r['渠道'],r['品类'],r['型号']),{}).get('商品访客数',0),
              '成交买家数': cur_by_model.get((r['渠道'],r['品类'],r['型号']),{}).get('支付买家数',0),
              '支付金额': cur_by_model.get((r['渠道'],r['品类'],r['型号']),{}).get('支付金额',0)}
             for r in rising_stars] if rising_stars else [],
            ['渠道','品类','型号','访客数','成交买家数','支付金额'],
            'diag_rising_stars.csv', '📥 增长亮点型号')

    # ────────────────────────────────────────────────────────────
    # 【货】：品类结构、SKU健康、转化漏斗、异常型号下钻
    # ────────────────────────────────────────────────────────────
    with diag_tabs[1]:
        st.markdown('<div class="section-title">📦 "货" — 商品结构 & 转化漏斗诊断</div>', unsafe_allow_html=True)
        st.caption('关注：品类GMV结构 · SKU转化漏斗 · 爆款健康度 · 客单价结构 · 退款率异常')

        # 品类贡献度分析
        st.markdown('#### 📊 品类销售结构（本期 vs 上期）')
        _cat_rows = []
        for cat_key, cv in cur_by_cat.items():
            pv = prev_by_cat.get(cat_key, {})
            _cv_gmv = cv.get('支付金额', 0); _pv_gmv = pv.get('支付金额', 0)
            _cv_cnt = cv.get('支付件数', 0)
            _cv_cvr = cv.get('支付转化率', 0) * 100
            _pv_cvr = pv.get('支付转化率', 0) * 100
            _cv_aov = cv.get('客单价', 0)
            _pv_aov = pv.get('客单价', 0)
            _gmv_chg = (_cv_gmv - _pv_gmv) / _pv_gmv if _pv_gmv else None
            _share = _cv_gmv / max(cur_sum.get('支付金额', 1), 1) * 100
            _cat_rows.append({
                '渠道': cat_key[0], '品类': cat_key[1],
                '本期GMV': f'¥{_cv_gmv:,.0f}',
                '成交占比': f'{_share:.1f}%',
                'GMV变化': f"<span style='color:{'#22c55e' if (_gmv_chg or 0) >= 0 else '#ef4444'};font-weight:700'>{_pct(_gmv_chg)}</span>",
                '销售件数': f'{_cv_cnt:,.0f}',
                '转化率': f'{_cv_cvr:.2f}%',
                '转化率变化': f"<span style='color:{'#22c55e' if _cv_cvr>=_pv_cvr else '#ef4444'}'>{_cv_cvr-_pv_cvr:+.2f}pp</span>",
                '客单价': f'¥{_cv_aov:,.0f}',
                '客单价变化': f"<span style='color:{'#22c55e' if _cv_aov>=_pv_aov else '#ef4444'}'>{_pct((_cv_aov-_pv_aov)/max(_pv_aov,1))}</span>",
            })
        _cat_rows.sort(key=lambda x: float(x['成交占比'].rstrip('%')), reverse=True)
        if _cat_rows:
            _cat_html = _html_table(_cat_rows, height=min(380, len(_cat_rows)*36+50))
            st.markdown(_wrap_fullscreen(_cat_html, title='📦 品类结构诊断')[0], unsafe_allow_html=True)
            _render_download_panel(_cat_rows, list(_cat_rows[0].keys()), 'diag_category.csv', '📥 品类结构')

        # 爆款健康度（上期TOP20 → 本期变化）
        st.markdown('#### ⚡ 爆款型号健康度（上期TOP20本期表现）')
        prev_top20 = sorted(prev_by_model.items(), key=lambda x: x[1].get('支付金额',0), reverse=True)[:20]
        drop_stars = []
        star_rows  = []
        for mk_key, pv in prev_top20:
            mv = cur_by_model.get(mk_key, {})
            pc = mv.get('支付金额', 0); pp = pv.get('支付金额', 0)
            if pp > 0:
                drop = (pc - pp) / pp
                drop_stars.append({'渠道': mk_key[0], '品类': mk_key[1], '型号': mk_key[2],
                    '上期GMV': pp, '本期GMV': pc, '缩水幅度': drop})
                _clr = '#22c55e' if drop >= 0 else ('#ef4444' if drop < -0.3 else '#f59e0b')
                star_rows.append({
                    '型号': mk_key[2], '品类': mk_key[1], '渠道': mk_key[0],
                    '上期GMV': f'¥{pp:,.0f}', '本期GMV': f'¥{pc:,.0f}',
                    '变化': f"<span style='color:{_clr};font-weight:700'>{_pct(drop)}</span>",
                    '状态': '🟢 稳健' if drop >= -0.05 else ('🟡 微降' if drop >= -0.3 else '🔴 断崖'),
                    '上期GMV份额': f"{pp/max(prev_sum_all.get('支付金额',1),1)*100:.1f}%"
                })
        drop_stars.sort(key=lambda x: x['缩水幅度'])
        star_rows.sort(key=lambda x: x['缩水幅度'] if isinstance(x.get('缩水幅度'), float) else 0)
        star_rows_disp = sorted(star_rows, key=lambda x: float(str(x.get('上期GMV份额','0%')).rstrip('%')), reverse=True)
        if star_rows_disp:
            st.markdown(_html_table(star_rows_disp[:20], height=min(520, len(star_rows_disp)*36+50)), unsafe_allow_html=True)

        # 转化率骤降型号
        st.markdown('#### 📉 转化率骤降型号（降幅>20%，访客>50）')
        cvr_drop_models = []
        for mk_key, mv in cur_by_model.items():
            pv_m = prev_by_model.get(mk_key, {})
            cvr_c = mv.get('支付转化率', 0); cvr_p = pv_m.get('支付转化率', 0)
            if cvr_p >= 0.005 and cvr_c < cvr_p:
                cvr_drop = (cvr_c - cvr_p) / cvr_p
                if cvr_drop < -0.20 and mv.get('商品访客数', 0) > 50:
                    cart_c = mv.get('商品加购人数', 0); cart_p = pv_m.get('商品加购人数', 0)
                    vis_c  = mv.get('商品访客数', 0)
                    if cart_c and cart_p and vis_c:
                        crc = cart_c / vis_c * 100; crp = cart_p / vis_c * 100
                        fn = f'加购率↓{crc-crp:.1f}pp→详情页吸引力下降' if crc < crp - 2 else '加购率正常→价格/评价/库存因素'
                    else:
                        fn = '数据不足'
                    cvr_drop_models.append({
                        '渠道': mk_key[0], '品类': mk_key[1], '型号': mk_key[2],
                        '上期转化': f'{cvr_p*100:.2f}%', '本期转化': f'{cvr_c*100:.2f}%',
                        '降幅': f"<span style='color:#dc2626;font-weight:700'>{_pct(cvr_drop)}</span>",
                        '本期访客': f"{mv.get('商品访客数',0):,.0f}",
                        '本期GMV': f"¥{mv.get('支付金额',0):,.0f}",
                        '漏斗判断': fn,
                    })
        cvr_drop_models.sort(key=lambda x: float(str(x.get('降幅','+0%')).replace('<span','').split('%')[0].split('>')[1] if '<span' in str(x.get('降幅','')) else '0') if False else 0)
        if cvr_drop_models:
            st.markdown(_html_table(cvr_drop_models[:15], height=min(450, len(cvr_drop_models)*36+40)), unsafe_allow_html=True)
            _render_download_panel(
                [{'渠道':r['渠道'],'品类':r['品类'],'型号':r['型号'],'本期访客':r['本期访客'],'本期GMV':r['本期GMV']} for r in cvr_drop_models],
                ['渠道','品类','型号','本期访客','本期GMV'], 'diag_cvr_drop.csv', '📥 转化骤降型号')
        else:
            st.info('✅ 未发现转化率骤降型号（阈值：降幅>20%，访客>50）。')

        # 客单价下跌型号
        st.markdown('#### 💰 客单价下跌型号（降幅>10%）')
        aov_drop_rows = []
        for mk_key, mv in cur_by_model.items():
            pv_m = prev_by_model.get(mk_key, {})
            ac = mv.get('客单价', 0); ap = pv_m.get('客单价', 0)
            if ap > 10 and ac < ap:
                ad = (ac - ap) / ap
                if ad < -0.10 and mv.get('支付件数', 0) > 10:
                    pc_c = mv.get('支付件数', 0); pc_p = pv_m.get('支付件数', 0)
                    if pc_c > pc_p * 1.3:      ar = '件数↑但均价↓→低价SKU占比提升/折扣加大'
                    elif pc_c < pc_p * 0.7:    ar = '件数↓且均价↓→高客单SKU销量萎缩'
                    else:                       ar = '件数持平→直接降价/促销力度加大'
                    aov_drop_rows.append({
                        '渠道': mk_key[0], '品类': mk_key[1], '型号': mk_key[2],
                        '上期客单价': f'¥{ap:,.0f}', '本期客单价': f'¥{ac:,.0f}',
                        '降幅': f"<span style='color:#ea580c;font-weight:700'>{_pct(ad)}</span>",
                        '本期件数': f'{mv.get("支付件数",0):,.0f}',
                        '上期件数': f'{pv_m.get("支付件数",0):,.0f}',
                        '初步判断': ar,
                    })
        if aov_drop_rows:
            st.markdown(_html_table(aov_drop_rows[:15], height=min(420, len(aov_drop_rows)*36+40)), unsafe_allow_html=True)
        else:
            st.info('✅ 未发现客单价明显下跌型号（阈值：降幅>10%，件数>10）。')

    # ────────────────────────────────────────────────────────────
    # 【场】：渠道 ROI 矩阵、推广诊断、流量效率
    # ────────────────────────────────────────────────────────────
    with diag_tabs[2]:
        st.markdown('<div class="section-title">🏪 "场" — 渠道效率 & 推广投放诊断</div>', unsafe_allow_html=True)
        st.caption('关注：各渠道GMV贡献与效率变化 · 推广ROI · 花费 vs 成交联动 · P0行动项')

        # 渠道效率矩阵
        st.markdown('#### 📊 渠道经营矩阵（本期 vs 上期）')
        _ch_matrix = []
        total_gmv_c = cur_sum.get('支付金额', 1)
        for ch_key, cv in cur_by_channel.items():
            pv = prev_by_channel.get(ch_key, {})
            _cv_gmv = cv.get('支付金额', 0); _pv_gmv = pv.get('支付金额', 0)
            _cv_vis = cv.get('商品访客数', 0)
            _cv_cvr = cv.get('支付转化率', 0) * 100
            _pv_cvr = pv.get('支付转化率', 0) * 100
            _cv_aov = cv.get('客单价', 0); _pv_aov = pv.get('客单价', 0)
            _gmv_chg = (_cv_gmv - _pv_gmv) / _pv_gmv if _pv_gmv else None
            _share = _cv_gmv / max(total_gmv_c, 1) * 100
            _status = '🟢' if (_gmv_chg or 0) > WARN_T else ('🟡' if (_gmv_chg or 0) > DANGER_T else '🔴')
            _ch_matrix.append({
                '状态': _status, '渠道': ch_key[0],
                '本期GMV': f'¥{_cv_gmv:,.0f}',
                'GMV占比': f'{_share:.1f}%',
                'GMV变化': f"<span style='color:{'#22c55e' if (_gmv_chg or 0)>=0 else '#ef4444'};font-weight:700'>{_pct(_gmv_chg)}</span>",
                '访客数': f'{_cv_vis:,.0f}',
                '转化率': f'{_cv_cvr:.2f}%',
                '转化率变化': f"<span style='color:{'#22c55e' if _cv_cvr>=_pv_cvr else '#ea580c'}'>{_cv_cvr-_pv_cvr:+.2f}pp</span>",
                '客单价': f'¥{_cv_aov:,.0f}',
                '客单价变化': f"<span style='color:{'#22c55e' if _cv_aov>=_pv_aov else '#ea580c'}'>{_pct((_cv_aov-_pv_aov)/max(_pv_aov,1))}</span>",
            })
        _ch_matrix.sort(key=lambda x: float(x['GMV占比'].rstrip('%')), reverse=True)
        if _ch_matrix:
            st.markdown(_html_table(_ch_matrix, height=min(340, len(_ch_matrix)*36+50)), unsafe_allow_html=True)
            _render_download_panel(_ch_matrix, list(_ch_matrix[0].keys()), 'diag_channel_matrix.csv', '📥 渠道矩阵')

        # 各渠道内GMV下滑最严重型号
        st.markdown('#### 🎯 渠道内 GMV 下滑型号 Top3（含多因子归因）')
        ch_model_issues = []
        total_cur_gmv = cur_sum.get('支付金额', 1)
        for ch_key, cur_v in cur_by_channel.items():
            ch_name = ch_key[0]
            prev_v = prev_by_channel.get(ch_key, {})
            ch_gmv_c = cur_v.get('支付金额', 0); ch_gmv_p = prev_v.get('支付金额', 0)
            ch_chg = (ch_gmv_c - ch_gmv_p) / ch_gmv_p if ch_gmv_p else None
            if ch_chg and ch_chg < 0:
                worst_models = []
                for mk_key, mv in cur_by_model.items():
                    if mk_key[0] != ch_name: continue
                    pv_m = prev_by_model.get(mk_key, {})
                    mc = mv.get('支付金额', 0); mp = pv_m.get('支付金额', 0)
                    m_chg = (mc - mp) / mp if mp else None
                    if m_chg is not None and m_chg < 0:
                        worst_models.append({
                            '渠道': mk_key[0], '品类': mk_key[1], '型号': mk_key[2],
                            '本期GMV': mc, '上期GMV': mp, '环比': m_chg,
                            '本期转化率': mv.get('支付转化率', 0), '上期转化率': pv_m.get('支付转化率', 0),
                            '本期访客': mv.get('商品访客数', 0), '上期访客': pv_m.get('商品访客数', 0),
                            '本期客单价': mv.get('客单价', 0), '上期客单价': pv_m.get('客单价', 0),
                        })
                worst_models.sort(key=lambda x: x['环比'])
                for wm in worst_models[:3]: ch_model_issues.append(wm)

        if ch_model_issues:
            _issue_rows = []
            for idx, w in enumerate(ch_model_issues[:25]):
                vis_chg_m = (w['本期访客']-w['上期访客'])/w['上期访客'] if w['上期访客'] else None
                cvr_diff  = (w['本期转化率']-w['上期转化率'])*100
                aov_chg_m = (w['本期客单价']-w['上期客单价'])/w['上期客单价'] if w['上期客单价'] else None
                factors = []
                if vis_chg_m is not None:
                    if vis_chg_m<=-0.20: factors.append(('流量崩塌', f'↓{_pct(vis_chg_m)}', 3))
                    elif vis_chg_m<=-0.10: factors.append(('流量大降', f'↓{_pct(vis_chg_m)}', 2))
                    elif vis_chg_m<=-0.05: factors.append(('流量微降', f'↓{_pct(vis_chg_m)}', 1))
                if cvr_diff<=-5:   factors.append(('转化崩溃', f'↓{cvr_diff:.1f}pp', 3))
                elif cvr_diff<=-2: factors.append(('转化下降', f'↓{cvr_diff:.1f}pp', 2))
                if aov_chg_m is not None:
                    if aov_chg_m<=-0.20: factors.append(('客单价暴跌', f'↓{_pct(aov_chg_m)}', 3))
                    elif aov_chg_m<=-0.10: factors.append(('客单价下跌', f'↓{_pct(aov_chg_m)}', 2))
                factors.sort(key=lambda x: x[2], reverse=True)
                if factors:
                    reason_parts = []
                    for fname, fdetail, fw in factors[:3]:
                        c = '#dc2626' if fw==3 else ('#ea580c' if fw==2 else '#f59e0b')
                        reason_parts.append(f"<span style='color:{c}'>● {fname}:{fdetail}</span>")
                    reason_str = ' '.join(reason_parts)
                else:
                    reason_str = "<span style='color:#64748b'>多因素平稳下滑</span>"
                impact_amt = max(0, w['上期GMV'] - w['本期GMV'])
                impact_pct = (impact_amt / total_cur_gmv * 100) if total_cur_gmv else 0
                _issue_rows.append({
                    '严重度': '🔴' if w['环比']<DANGER_T else ('🟠' if w['环比']<WARN_T else '🟡'),
                    '渠道': w['渠道'], '品类': w['品类'], '型号': w['型号'],
                    '本期GMV': f"¥{w['本期GMV']:,.0f}", '上期GMV': f"¥{w['上期GMV']:,.0f}",
                    '环比变化': f"<span style='color:#dc2626;font-weight:700'>{_pct(w['环比'])}</span>",
                    '拖累金额': f"¥{impact_amt:,.0f}({impact_pct:.1f}%)",
                    '主因归因': reason_str,
                })
            st.markdown(_html_table(_issue_rows, height=min(480, len(_issue_rows)*36+50)), unsafe_allow_html=True)
            top_drag = sorted(ch_model_issues, key=lambda x: x['环比'])[:5]
            st.caption(f"📌 最大拖累TOP5: {' | '.join([f'[{t['型号']}]{_pct(t['环比'])}' for t in top_drag])}")
        else:
            st.info('✅ 所有渠道各型号表现稳定，未发现显著异常下滑。')

        # 推广诊断（场：推广ROI分析）
        if promo_rows:
            st.markdown('<hr style="margin:16px 0;border:none;border-top:1px dashed #cbd5e1;">', unsafe_allow_html=True)
            st.markdown('#### 📢 推广效率诊断（场：推广投放）')

            promo_cur_diag  = [r for r in promo_rows if s            <= r.get('_date','') <= e           ]
            promo_prev_diag = [r for r in promo_rows if _t2_prev_s <= r.get('_date','') <= _t2_prev_e]
            def _promo_sum(rows):
                return {k: sum(r.get(f'_{k}',0) for r in rows) for k in ['花费','展现数','点击数','总订单金额','直接订单金额','总加购数']}
            p_cur  = _promo_sum(promo_cur_diag)
            p_prev = _promo_sum(promo_prev_diag)

            # 销售支付金额（本期和对比期）
            _diag_sales_cur = sum(float(r.get('支付金额', 0) or 0) for r in cur_rows_all if r.get('支付金额'))
            _diag_sales_prev = sum(float(r.get('支付金额', 0) or 0) for r in prev_rows_all if r.get('支付金额'))

            p_fc_g   = (p_cur['花费']       - p_prev['花费'])       / p_prev['花费']       if p_prev['花费'] else None
            p_roi_cur  = p_cur['总订单金额']  / p_cur['花费']         if p_cur['花费']  else 0
            p_roi_prev = p_prev['总订单金额'] / p_prev['花费']        if p_prev['花费'] else 0
            p_roi_g    = (p_roi_cur - p_roi_prev) / p_roi_prev       if p_roi_prev else None
            p_ctr_cur  = p_cur['点击数']  / p_cur['展现数']  * 100   if p_cur['展现数']  else 0
            p_ctr_prev = p_prev['点击数'] / p_prev['展现数'] * 100   if p_prev['展现数'] else 0
            p_ctr_g    = (p_ctr_cur - p_ctr_prev) / p_ctr_prev       if p_ctr_prev else None
            p_cpc_cur  = p_cur['花费']  / p_cur['点击数']            if p_cur['点击数']  else 0
            p_cpc_prev = p_prev['花费'] / p_prev['点击数']           if p_prev['点击数'] else 0
            p_cpc_g    = (p_cpc_cur - p_cpc_prev) / p_cpc_prev       if p_cpc_prev else None
            p_rate_cur  = p_cur['花费']  / _diag_sales_cur  * 100 if _diag_sales_cur  else 0
            p_rate_prev = p_prev['花费'] / _diag_sales_prev * 100 if _diag_sales_prev else 0
            p_rate_g    = (p_rate_cur - p_rate_prev) / p_rate_prev    if p_rate_prev else None

            # 推广KPI 5卡片
            p_kpi_cols = st.columns(5)
            _p_kpis = [
                ('💸 推广花费', p_fc_g,   f"¥{p_cur['花费']:,.0f}",     f"¥{p_prev['花费']:,.0f}"),
                ('📈 ROI',      p_roi_g,  f"{p_roi_cur:.2f}",            f"{p_roi_prev:.2f}"),
                ('🖱️ 点击率',  p_ctr_g,  f"{p_ctr_cur:.2f}%",           f"{p_ctr_prev:.2f}%"),
                ('💰 CPC',      p_cpc_g,  f"¥{p_cpc_cur:.2f}",           f"¥{p_cpc_prev:.2f}"),
                ('⚖️ 费率',    p_rate_g, f"{p_rate_cur:.1f}%",           f"{p_rate_prev:.1f}%"),
            ]
            for col, (name, chg, cur_s, prev_s_) in zip(p_kpi_cols, _p_kpis):
                with col:
                    _lvl = 'ok' if chg is None or chg > WARN_T else ('warn' if chg > DANGER_T else 'danger')
                    _bg  = {'danger':'#fef2f2','warn':'#fff7ed','ok':'#f0fdf4'}[_lvl]
                    _brd = {'danger':'#fca5a5','warn':'#fdba74','ok':'#86efac'}[_lvl]
                    st.markdown(
                        f'<div style="background:{_bg};border:1px solid {_brd};border-radius:12px;padding:10px;text-align:center;">'
                        f'<div style="font-size:11px;color:#64748b;font-weight:700;">{name}</div>'
                        f'<div style="font-size:18px;font-weight:900;color:#0f172a;">{cur_s}</div>'
                        f'<div style="font-size:10px;color:#94a3b8;">vs上期 {prev_s_} ({_pct(chg)})</div></div>',
                        unsafe_allow_html=True)

            # 推广诊断建议
            promo_suggestions = []
            if p_fc_g is not None and p_fc_g > 0.20:
                promo_suggestions.append(('P1','推广花费激增',
                    f'花费较上期增长{_pct(p_fc_g)}（当前¥{p_cur["花费"]:,.0f}）。'
                    f'① 检查各计划日预算上限 ② 核对ROI是否同步提升（当前{p_roi_cur:.2f} vs 上期{p_roi_prev:.2f}）'
                    f' ③ 若ROI下降，立即暂停低效计划'))
            elif p_fc_g is not None and p_fc_g < -0.20:
                promo_suggestions.append(('P0','推广花费大幅缩减',
                    f'花费较上期下降{_pct(p_fc_g)}（当前¥{p_cur["花费"]:,.0f}）。'
                    f'① 检查账户余额和计划状态 ② 若GMV同步下滑立即恢复核心计划预算 ③ 排查平台限流/违规'))
            if p_roi_g is not None and p_roi_g < -0.15:
                promo_suggestions.append(('P0','推广ROI显著恶化',
                    f'ROI从{p_roi_prev:.2f}→{p_roi_cur:.2f}（{_pct(p_roi_g)}）。'
                    f'① 筛选ROI<1的计划暂停或优化 ② 检查落地页转化路径 ③ 对比竞品价格确认是否涨价致转化下降'))
            elif p_roi_g is not None and p_roi_g > 0.15:
                promo_suggestions.append(('P2','推广ROI表现优异',
                    f'ROI从{p_roi_prev:.2f}→{p_roi_cur:.2f}（{_pct(p_roi_g)}）。'
                    f'① 加大高ROI计划预算（+30%测试）② 复制该计划定向和创意到其他SKU ③ 集中预算投放黄金时段'))
            if p_ctr_g is not None and p_ctr_g < -0.20:
                promo_suggestions.append(('P1','点击率大幅下滑',
                    f'CTR从{p_ctr_prev:.2f}%→{p_ctr_cur:.2f}%（{_pct(p_ctr_g)}）。'
                    f'① 更换主图/视频创意（A/B测试3组）② 检查定向人群是否过宽 ③ 若首屏占比下降需提高出价'))
            if p_cpc_g is not None and p_cpc_g > 0.30:
                promo_suggestions.append(('P1','点击成本快速上涨',
                    f'CPC从¥{p_cpc_prev:.2f}→¥{p_cpc_cur:.2f}（{_pct(p_cpc_g)}）。'
                    f'① 优化关键词/人群包质量分 ② 避开高峰时段竞价 ③ 测试长尾词降低竞争成本'))
            if p_rate_g is not None and p_rate_g > 0.20:
                promo_suggestions.append(('P1','推广费率过高',
                    f'费率从{p_rate_prev:.1f}%→{p_rate_cur:.1f}%（{_pct(p_rate_g)}）。'
                    f'① 设定费率红线（建议≤15%），超线计划立即优化 ② 提升客单价稀释费率 ③ 减少低转化时段投放'))
            if gmv_g is not None and gmv_g < 0 and p_fc_g is not None and p_fc_g < 0:
                promo_suggestions.append(('P0','销售&推广双降',
                    f'GMV{_pct(gmv_g)}且推广花费{_pct(p_fc_g)}，可能为系统性问题。'
                    f'① 对比行业大盘确认是否系统性下滑 ② 检查店铺DSR评分和违规记录 ③ 启动应急推广稳定基本盘'))
            elif gmv_g is not None and gmv_g < 0 and p_fc_g is not None and p_fc_g > 0:
                promo_suggestions.append(('P0','推广增但销售降（效率恶化）',
                    f'推广花费{_pct(p_fc_g)}但GMV{_pct(gmv_g)}，每多投1元推广反而亏损。'
                    f'① 立即暂停ROI<0.5的计划 ② 全面检查落地页和详情页 ③ 排查差评/缺货/涨价 ④ 修复转化再恢复投放'))

            if promo_suggestions:
                for pri, title, detail in sorted(promo_suggestions, key=lambda x: ['P0','P1','P2'].index(x[0])):
                    cls = {'P0':'tag-p0','P1':'tag-p1','P2':'tag-p2'}[pri]
                    with st.expander(f"<span class='action-tag {cls}'>{pri}</span> **{title}**", expanded=(pri=='P0')):
                        st.markdown(detail, unsafe_allow_html=True)
            else:
                st.info('✅ 推广数据表现平稳，未发现明显异常。')

    # ────────────────────────────────────────────────────────────
    # 【执行清单】：P0-P3 行动项（基于人货场诊断结果自动生成）
    # ────────────────────────────────────────────────────────────
    with diag_tabs[3]:
        st.markdown('<div class="section-title">🛠️ 作战任务令 — 三List体系（问题·风险·机会）</div>', unsafe_allow_html=True)
        st.caption(f'已识别 {len(ch_model_issues)} 个异常型号 | {len(cvr_drop_models)} 个转化骤降型号 | {len(drop_stars)} 个爆款掉量型号。每条措施绑定实际数据值。')

        actions = []
        def add_action(priority, title, detail, owner, timeline, metric_target, recover=''):
            actions.append({'p': priority, 't': title, 'd': detail, 'o': owner, 'tl': timeline, 'mt': metric_target, 'r': recover})

        # ── 人侧行动 ──
        if vis_g is not None and vis_g < -0.08:
            vis_loss_share = abs(vis_g) / (abs(vis_g) + abs(cvr_g or 0) + abs(aov_g or 0) + 0.001)
            gmv_loss = max(0, prev_sum_all.get('支付金额',0) - cur_sum.get('支付金额',0))
            recover_amt = gmv_loss * vis_loss_share
            add_action('P0', '【人·流量】紧急排查核心渠道流量断崖',
                f'<b>现状：</b>访客{_pct(vis_g)}（{prev_sum_all.get("商品访客数",0):,.0f}→{cur_sum.get("商品访客数",0):,.0f}），'
                f'拖累GMV约¥{recover_amt:,.0f}<br><br>'
                f'<b>🔍 根因排查（由粗到细）：</b><br>'
                f'<b>L1 流量断崖</b> → <b>L2 付费流量衰减？</b> 直通车后台→推广计划列表→展现量↓>30%→预算耗尽/质量分下降<br>'
                f'<b>L2 搜索排名下降？</b> 生意参谋→搜索分析→核心类目Top10词→对比搜索人气和CTR<br>'
                f'<b>L2 竞品活动冲击？</b> 平台搜索同品类→对比竞品是否有大促/新品上架<br>'
                f'<b>L3 直播间流量异常？</b> 查看流量来源→确认付费/免费比例变化<br><br>'
                f'<b>⚡ 应急措施 [快]：</b>表现最差的3个计划日预算+50%，观察3天<br>'
                f'<b>⚠️ 升级预案：</b>若3天无改善，升级为部门专项会议+申请额外推广预算',
                '运营负责人', '24小时内', f'访客≥{prev_sum_all.get("商品访客数",0)*0.95:,.0f}',
                recover=f'¥{recover_amt:,.0f}')

        # ── 货侧行动 ──
        if cvr_g is not None and cvr_g < -0.08:
            ccvr = cur_sum.get('支付转化率',0)*100; pcvr = prev_sum_all.get('支付转化率',0)*100
            lost_orders = cur_sum.get('商品访客数',0) * (prev_sum_all.get('支付转化率',0) - cur_sum.get('支付转化率',0))
            lost_gmv_val = lost_orders * cur_sum.get('客单价', 0)
            add_action('P0', '【货·转化】全店转化率紧急提升行动',
                f'<b>现状：</b>{pcvr:.2f}%→{ccvr:.2f}%（↓{pcvr-ccvr:.2f}pp），少成交约{lost_orders:,.0f}单，影响¥{lost_gmv_val:,.0f}<br><br>'
                f'<b>🔍 根因排查（由粗到细）：</b><br>'
                f'<b>L1 转化失效</b> → <b>L2 详情页吸引力下降？</b> 从「转化骤降型号」提取异常SKU→逐一做首屏3秒测试<br>'
                f'<b>L2 价格竞争力不足？</b> 平台搜索同款竞品前3名→对比售价，高于竞品8%则设限时9折<br>'
                f'<b>L2 差评/库存影响？</b> 导出近90天评价→词频统计→负面Top3→优化FAQ和卖点<br>'
                f'<b>L3 大促后遗症？</b> 检查大促是否刚结束→价格回调导致骤降→延长优惠3天<br><br>'
                f'<b>⚡ 应急措施 [快]：</b>转化骤降Top5型号临时设置限时折扣/赠品<br>'
                f'<b>⚠️ 升级预案：</b>若3天无改善，启动全店促销方案+美工详情页改版',
                '运营+美工', '3天内', f'转化率≥{pcvr*0.97:.2f}%',
                recover=f'¥{lost_gmv_val:,.0f}')

        if aov_g is not None and aov_g < -0.06:
            aov_loss = (prev_sum_all.get('客单价',0) - cur_sum.get('客单价',0)) * cur_sum.get('支付买家数', 1)
            add_action('P1', '【货·客单价】高客单价SKU曝光恢复',
                f'<b>现状：</b>¥{prev_sum_all.get("客单价",0):.0f}→¥{cur_sum.get("客单价",0):.0f}（{_pct(aov_g)}），损失约¥{aov_loss:,.0f}<br><br>'
                f'<b>🔍 根因排查：</b><br>'
                f'<b>L2 低价SKU占比提升？</b> 提取客单价前20 SKU→核对本周访客→圈出降幅最大的5个<br>'
                f'<b>L2 高客单SKU销量萎缩？</b> 检查高客单SKU是否有推广断流/竞品打压<br>'
                f'<b>L2 直接降价/促销力度加大？</b> 核对近期活动折扣力度变化<br><br>'
                f'<b>⚡ 优化措施 [中]：</b><br>'
                f'① 设置关联推荐：「搭配购买减X」「买二送一」，放在加购区下方<br>'
                f'② 满减门槛：均值¥{cur_sum.get("客单价",0):.0f}→满减线设¥{cur_sum.get("客单价",0)*1.3:,.0f}<br>'
                f'③ 直通车→「高消费力」人群溢价+20%',
                '运营', '1周内', f'客单价≥¥{prev_sum_all.get("客单价",0)*0.97:,.0f}',
                recover=f'¥{aov_loss:,.0f}')

        # ── 场侧行动 ──
        if ch_model_issues:
            ch_gmv_changes = {}
            for ck, cv_ch in cur_by_channel.items():
                pv_ch = prev_by_channel.get(ck, {})
                cc = cv_ch.get('支付金额',0); pp = pv_ch.get('支付金额',0)
                ch_gmv_changes[ck[0]] = (cc-pp)/pp if pp else None
            sorted_ch = sorted(ch_gmv_changes.items(), key=lambda x: x[1] if x[1] else 0)
            worst_ch  = sorted_ch[0] if sorted_ch else (None, None)
            if worst_ch[1] and worst_ch[1] < -0.05:
                ch_nm = worst_ch[0]; ch_pct = _pct(worst_ch[1])
                bad_in = [m for m in ch_model_issues if m['渠道'] == ch_nm][:3]
                ml = ', '.join([f"[{m['型号']}]({m['品类']})" for m in bad_in]) or '多个型号'
                add_action('P0', f'【场·渠道】{ch_nm}专项整改（GMV{ch_pct}）',
                    f'该渠道GMV{ch_pct}，集中在：{ml}<br><br>'
                    f'<b>🔍 根因排查：</b><br>'
                    f'<b>L2 推广计划异常？</b> 检查上述型号推广状态（停/降权/违规/预算耗尽）<br>'
                    f'<b>L2 DSR评分拖累？</b> DSR<4.7影响搜索权重，检查该渠道DSR近7天变化<br>'
                    f'<b>L2 活动报名遗漏？</b> 核对活动报名情况，重要会场补报<br>'
                    f'<b>L3 直播/内容影响？</b> 抖音渠道：查看7天直播时长和GMV/小时<br><br>'
                    f'<b>⚡ 应急措施 [快]：</b>上述型号临时增加推广预算30%+检查主图CTR<br>'
                    f'<b>⚠️ 升级预案：</b>若48小时无改善，启动渠道专项会议+调整渠道资源分配',
                    f'{ch_nm}渠道负责人', '48小时内', f'{ch_nm} GMV环比转正')

        # ── 型号级别具体措施（Top3）──
        if ch_model_issues:
            for bm in ch_model_issues[:3]:
                mod = bm['型号']; cat = bm['品类']; ch = bm['渠道']
                chg_pct = _pct(bm['环比'])
                vc = bm['本期访客']; vp = bm['上期访客']
                vm = (vc-vp)/vp if vp else None
                ccr = bm['本期转化率']*100; cpr = bm['上期转化率']*100
                cdr = ccr - cpr
                aov_cur = bm.get('本期客单价', 0); aov_prev = bm.get('上期客单价', 0)
                aov_chg = (aov_cur - aov_prev) / aov_prev if aov_prev else None
                # ── 3级根因判断 ──
                root_l1 = '流量断崖' if (vm and vm<-0.15) else ('转化失效' if cdr<-2 else '复合衰退')
                root_l2 = ''
                if root_l1 == '流量断崖':
                    if vm and vm < -0.30: root_l2 = '推广计划限流/预算耗尽'
                    elif vm and vm < -0.20: root_l2 = '搜索排名下降'
                    else: root_l2 = '竞品活动冲击'
                elif root_l1 == '转化失效':
                    if cdr < -5: root_l2 = '详情页吸引力下降'
                    elif aov_chg and aov_chg > 0.05: root_l2 = '价格竞争力不足'
                    else: root_l2 = '差评/库存影响'
                else:
                    root_l2 = '多因子叠加'
                pri = 'P0' if bm['环比'] < DANGER_T else 'P1'
                loss_amnt = bm['上期GMV'] - bm['本期GMV']
                detail = (
                    f'<b>数据：</b>¥{bm["本期GMV"]:,.0f} vs ¥{bm["上期GMV"]:,.0f}'
                    f'（损失¥{loss_amnt:,.0f}），转化{ccr:.2f}% vs {cpr:.2f}%<br>'
                    f'<b>根因：</b>{root_l1} → {root_l2}<br><br>'
                    f'<b>定向施策：</b><br>')
                if root_l1 == '流量断崖':
                    detail += (
                        f'① 检查该型号在{ch}的搜索排名和主图CTR'
                        f'（CTR={vc/(vp+0.001)*100:.1f}%，{"<3%需换主图" if vp and vc/vp<0.03 else "正常"}）<br>'
                        f'② 检查是否有推广计划被系统限流/预算耗尽<br>'
                        f'③ <b>⚡ 急救 [快]：</b>临时增加直通车日预算+50%，持续3天观察<br>'
                        f'④ <b>⚠️ 3天无改善：</b>启动竞品对标分析+主图A/B测试')
                else:
                    detail += (
                        f'① 打开该型号详情页模拟买家浏览——首屏3秒能否看清核心卖点？<br>'
                        f'② 评价审计：导出近60天评价→词频统计→负面Top3→优化话术<br>'
                        f'③ 价格对标：搜索同款竞品3家，'
                        f'{"高于竞品8%则设限时折 [快]" if bm["本期客单价"] >= bm["上期客单价"]*1.08 else "价格基本合理"}<br>'
                        f'④ 库存检查：确认该型号无缺货/预售状态<br>'
                        f'⑤ 差评处理：筛选出现≥2次的负面标签集中处理<br>'
                        f'⑥ <b>⚠️ 3天无改善：</b>启动全型号促销+详情页改版')
                add_action(pri, f'[{mod}]({cat}/{ch}) {root_l1}→{root_l2} — GMV{chg_pct}',
                    detail, f'运营-{cat}组', '3-5天见效', f'{mod} GMV环比>-5%',
                    recover=f'¥{loss_amnt:,.0f}')

        # ── 退款率措施 ──
        if ref_g is not None and ref_g > 0.05:
            crp = cur_sum.get('退款率', 0) * 100
            ref_loss = cur_sum.get('支付金额', 0) * ref_g
            if crp > 8 or ref_g > 0.10:
                add_action('P1', '【货·售后】退款率异常升高',
                    f'当前{crp:.1f}%，变化{_pct(ref_g)}，预估损失约¥{ref_loss:,.0f}<br>'
                    f'<b>🔍 根因排查：</b><br>'
                    f'① 导出近30天退款订单按「退款原因」归类，提取Top3原因<br>'
                    f'② 「质量问题/描述不符」>40%：质检团队抽检<br>'
                    f'③ 「物流慢/破损」>30%：改进包装+切换快递<br>'
                    f'④ 「不想要了」<说明详情页误导信息，需修正<br><br>'
                    f'<b>⚡ 应急措施 [快]：</b>退款率Top20%订单审核延长12小时+人工介入<br>'
                    f'<b>⚠️ 升级预案：</b>若7天无改善，启动供应链专项整改',
                    '客服+仓储+质检', '2周内', '退款率<5%',
                    recover=f'¥{ref_loss:,.0f}')

        # ── 常规措施 ──
        add_action('P3', '【常规】每周一上午健康检查',
            '每周一 10:00 完成：<br>'
            '① 打开本看板「智能诊断」Tab截图存档<br>'
            '② 对比上周同期标记变化±5%<br>'
            '③ 连续2周同一指标下滑→专项会议<br>'
            '④ 检查本周到期活动/优惠券续期',
            '运营负责人', '每周一固定', '周报存档', '')
        add_action('P3', '【常规】月度渠道ROI复盘',
            '每月5日前完成上月各渠道ROI：<br>'
            f'ROI=(渠道销售额-退货额)/渠道推广费用<br>'
            '① >5 加大投入 / 2-5 维持 / <2 缩减或优化<br>'
            '② ROI<2输出《XX渠道优化方案》',
            '运营+财务', '每月5号前', '全渠道均ROI>3', '')

        # ── 展示（三List体系）──
        actions_sorted = sorted(actions, key=lambda x: ['P0','P1','P2','P3'].index(x['p']))
        _p0_actions = [a for a in actions_sorted if a['p']=='P0']
        _p1_actions = [a for a in actions_sorted if a['p']=='P1']
        _p23_actions = [a for a in actions_sorted if a['p'] in ('P2','P3')]

        # ═══════════════════════════════════════
        # 🚨 List 1: 问题清单（P0/P1 立即行动）
        # ═══════════════════════════════════════
        st.markdown("<div style='background:#fef2f2;border:1px solid #fca5a5;border-radius:10px;padding:6px 14px;margin:12px 0 6px 0;'>"
                    "<b>🚨 问题清单</b> <small style='color:#94a3b8;'>— 已确认的异常，需立即处理</small></div>", unsafe_allow_html=True)

        if _p0_actions:
            st.markdown("<div style='font-size:12px;color:#dc2626;font-weight:700;margin:4px 0;'>P0 紧急行动</div>", unsafe_allow_html=True)
        for act in _p0_actions:
            cls = {'P0':'tag-p0'}[act['p']]
            _rec_str = f' | 预期挽回: {act["r"]}' if act.get('r') else ''
            with st.expander(f"<span class='action-tag {cls}'>{act['p']}</span> **{act['t']}** <small style='color:#94a3b8;'>| {act['o']} | 目标: {act['mt']} | 见效: {act['tl']}{_rec_str}</small>", expanded=True):
                st.markdown(act['d'], unsafe_allow_html=True)

        if _p1_actions:
            st.markdown("<div style='font-size:12px;color:#ea580c;font-weight:700;margin:8px 0 4px 0;'>P1 重点关注</div>", unsafe_allow_html=True)
        for act in _p1_actions:
            _rec_str = f' | 预期挽回: {act["r"]}' if act.get('r') else ''
            with st.expander(f"<span class='action-tag tag-p1'>{act['p']}</span> **{act['t']}** <small style='color:#94a3b8;'>| {act['o']} | 目标: {act['mt']}{_rec_str}</small>", expanded=False):
                st.markdown(act['d'], unsafe_allow_html=True)

        for act in _p23_actions:
            cls = {'P2':'tag-p2','P3':'tag-p3'}[act['p']]
            with st.expander(f"<span class='action-tag {cls}'>{act['p']}</span> **{act['t']}** <small style='color:#94a3b8;'>| {act['o']}</small>", expanded=False):
                st.markdown(act['d'], unsafe_allow_html=True)

        if not actions_sorted:
            st.success('✅ 当前所有核心指标健康，无额外干预。')

        # ═══════════════════════════════════════
        # ⚠️ List 2: 风险清单（预警监控）
        # ═══════════════════════════════════════
        risks = []
        # 1. 连续下滑但未达P0阈值的指标
        for name, val, thresh, suggestion in [
            ('GMV', gmv_g, -0.15, '检查核心渠道流量和爆款转化'),
            ('流量', vis_g, -0.15, '排查推广计划和搜索排名变化'),
            ('转化率', cvr_g, -0.15, '检查详情页优化和竞品价格'),
            ('客单价', aov_g, -0.10, '关注高客单SKU曝光量'),
        ]:
            if val is not None and val < WARN_T and val > thresh:
                level = '⚠️ 关注' if val < -0.10 else '👀 观察'
                icon = '🟡' if val < -0.10 else '🔵'
                risks.append({
                    '等级': f'{icon} {level}',
                    '指标': name,
                    '当前变化': _pct(val),
                    '阈值': f'跌破{_pct(thresh)}触发P0',
                    '建议': suggestion,
                })

        # 2. 渠道流量占比持续下降
        if len(cur_by_channel) >= 2:
            for ch_key, cv in cur_by_channel.items():
                ch_name = ch_key[0]
                cv_vis = cv.get('商品访客数', 0)
                pv = prev_by_channel.get(ch_key, {})
                pv_vis = pv.get('商品访客数', 0)
                if pv_vis > 100:
                    vis_chg = (cv_vis - pv_vis) / pv_vis
                    if vis_chg < -0.10:
                        cur_share = cv_vis / max(cur_sum.get('商品访客数', 1), 1) * 100
                        prev_share = pv_vis / max(prev_sum_all.get('商品访客数', 1), 1) * 100
                        risks.append({
                            '等级': '🟡 ⚠️ 关注',
                            '指标': f'{ch_name}渠道',
                            '当前变化': f'流量{_pct(vis_chg)}',
                            '阈值': f'占比 {prev_share:.1f}%→{cur_share:.1f}%',
                            '建议': '检查渠道推广预算和搜索权重',
                        })

        if risks:
            st.markdown("<div style='background:#fffbeb;border:1px solid #fcd34d;border-radius:10px;padding:6px 14px;margin:16px 0 6px 0;'>"
                        "<b>⚠️ 风险清单</b> <small style='color:#94a3b8;'>— 需持续监控，防止恶化</small></div>", unsafe_allow_html=True)
            _risk_rows = []
            for r in risks[:10]:
                _risk_rows.append({
                    '风险等级': r['等级'],
                    '监控指标': r['指标'],
                    '变化趋势': r['当前变化'],
                    '触发条件': r['阈值'],
                    '防控建议': r['建议'],
                })
            st.markdown(_html_table(_risk_rows, height=min(300, len(_risk_rows)*36+50)), unsafe_allow_html=True)
        else:
            st.markdown("<div style='background:#f0fdf4;border:1px solid #86efac;border-radius:10px;padding:6px 14px;margin:16px 0 6px 0;'>"
                        "<b>✅ 风险清单</b> <small style='color:#64748b;'>— 当前未发现明确风险信号</small></div>", unsafe_allow_html=True)

        # ═══════════════════════════════════════
        # 💡 List 3: 机会清单（增长杠杆）
        # ═══════════════════════════════════════
        opportunities = []

        # 1. 增长亮点型号（已有数据）
        if rising_stars:
            for r in rising_stars[:5]:
                opportunities.append({
                    '类型': '🚀 增长亮点',
                    '目标': f"{r['型号']}（{r['品类']}）",
                    '数据': f"本期GMV ¥{r['本期GMV']:,.0f}",
                    '机会点': '加大推广预算，复制成功模式到同类SKU',
                    '建议动作': '推广预算+30%，测试3天',
                })

        # 2. 高转化率但低流量的型号（隐藏的宝石）
        for mk_key, mv in cur_by_model.items():
            cvr_val = mv.get('支付转化率', 0) * 100
            vis_val = mv.get('商品访客数', 0)
            avg_cvr = cur_sum.get('支付转化率', 0) * 100 if cur_sum.get('支付转化率', 0) else 0
            if cvr_val > avg_cvr * 1.5 and vis_val > 30 and vis_val < 500:
                gmv_val = mv.get('支付金额', 0)
                if gmv_val > 500:
                    opportunities.append({
                        '类型': '💎 潜力型号',
                        '目标': f"{mk_key[2]}（{mk_key[1]}）",
                        '数据': f"转化{cvr_val:.1f}%(高于均值{avg_cvr:.1f}%) | 访客仅{vis_val:,.0f}",
                        '机会点': '转化率优异但流量不足，加大曝光可快速起量',
                        '建议动作': '加入推广计划，设置日预算¥200测试',
                    })

        if opportunities:
            opp_display = opportunities[:8]
            st.markdown("<div style='background:#eff6ff;border:1px solid #93c5fd;border-radius:10px;padding:6px 14px;margin:16px 0 6px 0;'>"
                        "<b>💡 机会清单</b> <small style='color:#94a3b8;'>— 可主动出击的增长机会</small></div>", unsafe_allow_html=True)
            _opp_rows = []
            for o in opp_display:
                _opp_rows.append({
                    '机会类型': o['类型'],
                    '目标': o['目标'],
                    '关键数据': o['数据'],
                    '机会分析': o['机会点'],
                    '建议动作': o['建议动作'],
                })
            st.markdown(_html_table(_opp_rows, height=min(300, len(_opp_rows)*36+50)), unsafe_allow_html=True)
        else:
            st.markdown("<div style='background:#f0fdf4;border:1px solid #86efac;border-radius:10px;padding:6px 14px;margin:16px 0 6px 0;'>"
                        "<b>💡 机会清单</b> <small style='color:#64748b;'>— 暂未发现显著增长机会</small></div>", unsafe_allow_html=True)

        # 下载诊断报告
        st.markdown("<hr style='margin:18px 0;border:none;border-top:1px dashed #cbd5e1;'>", unsafe_allow_html=True)
        dl_data = [{
            '诊断时间': datetime.datetime.now().strftime('%Y-%m-%d %H:%M'),
            '诊断区间': f'{s}~{e}', '对比区间': f'{prev_s}~{prev_e}',
            'GMV变化': _pct(gmv_g), '访客变化': _pct(vis_g),
            '转化率变化': f"{(cur_sum.get('支付转化率',0)-prev_sum_all.get('支付转化率',0))*100:+.2f}pp" if prev_sum_all.get('支付转化率',0) else '--',
            '客单价变化': _pct(aov_g), '退款率变化': _pct(ref_g),
            '健康评分': f'{health_score:.0f}/100', '健康结论': hv[0],
            '渠道异常型号数': len(ch_model_issues), '转化骤降数': len(cvr_drop_models),
            '爆款掉量数': len([d for d in (drop_stars if 'drop_stars' in dir() else []) if d.get('缩水幅度',0) < -0.30]),
            '增长亮点数': len(rising_stars) if 'rising_stars' in dir() else 0,
            'P0任务数': sum(1 for a in actions if a['p']=='P0'),
            'P1任务数': sum(1 for a in actions if a['p']=='P1'),
        }]
        for act in actions_sorted:
            dl_data.append({
                '诊断时间': '', '诊断区间': '', '对比区间': '',
                'GMV变化': '', '访客变化': '', '转化率变化': '', '客单价变化': '', '退款率变化': '',
                '健康评分': '', '健康结论': '', '渠道异常型号数': '', '转化骤降数': '',
                '爆款掉量数': '', '增长亮点数': '', 'P0任务数': '', 'P1任务数': '',
                '优先级': act['p'], '措施标题': act['t'], '负责人': act['o'],
                '见效周期': act['tl'], '量化目标': act['mt'],
            })
        _render_download_panel(
            dl_data if dl_data else [],
            list(dl_data[0].keys()) if dl_data else [],
            f'xiaotunbi_diagnosis_{s.replace("-","")}_{e.replace("-","")}.csv', '📥 完整诊断报告（含执行清单）')

    # ── 生成麦肯锡复盘PPT按钮（放在所有子Tab之后，确保变量已定义）──
    st.markdown('<hr style="margin:20px 0;border:none;border-top:1px dashed #cbd5e1;">', unsafe_allow_html=True)
    _ppt_col1, _ppt_col2 = st.columns([3, 1])
    with _ppt_col1:
        st.markdown(
            "<div style='font-size:13px;color:#64748b;'>📑 <b>一键生成麦肯锡风格复盘PPT</b> — "
            "包含封面、健康总览、人货场分析、执行清单共6页</div>",
            unsafe_allow_html=True)
    with _ppt_col2:
        _gen_ppt = st.button('🎯 生成复盘PPT', use_container_width=True, key='gen_mck_ppt')

    if _gen_ppt:
        with st.spinner('正在生成麦肯锡风格复盘PPT...'):
            _period_label_cur = f'{s} ~ {e}'
            _period_label_prev = f'{prev_s} ~ {prev_e}'
            _ppt_path = _generate_mckinsey_ppt(
                period_cur=_period_label_cur, period_prev=_period_label_prev,
                comp_mode=comp_mode, filter_label=_filter_label,
                health_score=health_score, health_status=hv[0], health_color=hv[1],
                gmv_g=gmv_g, vis_g=vis_g, cvr_g=cvr_g, aov_g=aov_g, ref_g=ref_g,
                cur_sum=cur_sum, prev_sum=prev_sum_all,
                cur_by_channel=cur_by_channel, prev_by_channel=prev_by_channel,
                cur_by_cat=cur_by_cat, prev_by_cat=prev_by_cat,
                cur_by_model=cur_by_model, prev_by_model=prev_by_model,
                rising_stars=rising_stars, drop_stars=drop_stars,
                cvr_drop_models=cvr_drop_models, aov_drop_rows=aov_drop_rows,
                ch_model_issues=ch_model_issues, promo_suggestions=promo_suggestions,
                actions=actions, WARN_T=WARN_T, DANGER_T=DANGER_T,
                s=s, e=e,
            )
            if _ppt_path:
                with open(_ppt_path, 'rb') as f:
                    st.download_button(
                        label=f'📥 下载复盘PPT ({_period_label_cur})',
                        data=f, file_name=f'xiaotunbi_复盘_{s.replace("-","")}_{e.replace("-","")}.pptx',
                        mime='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                        key='dl_mck_ppt')
                st.success(f'✅ PPT已生成（{_ppt_path}），点击上方按钮下载')

# ═══════════════════════════════════════════════════════════════
# TAB 5: 透视表分析（单期数据，无对比列）
# ═══════════════════════════════════════════════════════════════
with tabs[5]:
    st.markdown('<div class="section-title">透视表分析</div>', unsafe_allow_html=True)

    _pv_dim_opts = ['渠道', '店铺', '品类', '型号', '日期', '年月']
    _pv_promo_dim_opts = ['_渠道', '_店铺', '_品类', '_型号', '_产品线', '_营销场景', '_推广计划', '日期', '年月']

    # 销售可选指标（含占比/计算字段标记）
    _pv_sales_metrics_all = [
        '商品访客数', '访客占比', '支付买家数', '支付件数',
        '支付金额', '成交占比', '支付转化率', '商品加购人数',
        '加购率', 'UV价值', '客单价',
    ]
    # 推广可选指标（含占比/计算字段标记）
    _pv_promo_metrics_all = [
        '_花费', '_花费占比', '_展现数', '_CPC', '_点击数', '_点击率',
        '_直接订单金额', '_直接金额占比', '_总订单金额', '_总金额占比',
        '_直接ROI', '_总ROI', '_总转化率', '_直接转化率',
        '_总订单成本', '_直接订单成本',
    ]

    # ══════════════════════════════════════
    # 子专区1：销售数据透视表
    # ══════════════════════════════════════
    st.markdown('#### 销售数据透视表')

    # 字段配置行（置顶，水平排列）
    _p1_cfg_c1, _p1_cfg_c2, _p1_cfg_c3 = st.columns([1, 2, 1])
    with _p1_cfg_c1:
        _p1_row_dims = st.multiselect('行维度', _pv_dim_opts, default=['品类'], key='pv1_row')
    with _p1_cfg_c2:
        _p1_vals = st.multiselect('值指标', _pv_sales_metrics_all,
                                   default=_pv_sales_metrics_all, key='pv1_val')
    with _p1_cfg_c3:
        _p1_top_n = st.number_input('前N行（0=全部）', min_value=0, max_value=100,
                                     value=0, key='pv1_top')

    if _p1_vals and _p1_row_dims:
        _p1_raw = data['daily']
        if channel or store or category or model:
            _p1_filtered = []
            for r in _p1_raw:
                if channel and r.get('渠道') not in channel: continue
                if store and r.get('店铺') not in store: continue
                if category and r.get('品类') not in category: continue
                if model and r.get('型号') not in model: continue
                _p1_filtered.append(r)
            _p1_raw = _p1_filtered

        _p1_rows = get_period_rows(_p1_raw, today_s, today_e)

        # 注入时间维度字段（用于行维度筛选）
        for r in _p1_rows:
            dt = r.get('日期', '')
            r['日期'] = dt
            r['年月'] = dt[:7] if len(dt) >= 7 else dt

        # ── 去年同期数据 ──
        _yoy_start = start.replace(year=start.year - 1)
        _yoy_end = end.replace(year=end.year - 1)
        if start.month == 2 and start.day == 29:
            _yoy_start = _yoy_start.replace(day=28)
        if end.month == 2 and end.day == 29:
            _yoy_end = _yoy_end.replace(day=28)
        _p1_yoy_rows = get_period_rows(_p1_raw, str(_yoy_start), str(_yoy_end))

        # 同期数据也注入时间维度字段
        for r in _p1_yoy_rows:
            dt = r.get('日期', '')
            r['日期'] = dt
            r['年月'] = dt[:7] if len(dt) >= 7 else dt

        # 聚合：只汇总原始可加字段，计算字段后处理
        _P1_RAW_FIELDS = ['商品访客数', '支付买家数', '支付件数', '支付金额', '商品加购人数', '成功退款金额']

        def _pv1_group(rows, row_dims):
            agg = {}
            for r in rows:
                rk = tuple((r.get(d) or '未标注') for d in row_dims)
                if rk not in agg:
                    agg[rk] = {f: 0.0 for f in _P1_RAW_FIELDS}
                for f in _P1_RAW_FIELDS:
                    agg[rk][f] += float(r.get(f, 0) or 0)
            return agg

        _p1_agg = _pv1_group(_p1_rows, _p1_row_dims)
        _p1_yoy_agg = _pv1_group(_p1_yoy_rows, _p1_row_dims)
        _p1_total_visitors = sum(v['商品访客数'] for v in _p1_agg.values()) or 1
        _p1_total_amt = sum(v['支付金额'] for v in _p1_agg.values()) or 1

        # ── 按行维度聚合推广花费（用于费率计算）──
        _promo_dim_map = {'渠道': '_渠道', '店铺': '_店铺', '品类': '_品类', '型号': '_型号'}
        _p1_promo_spend = {}
        if promo_rows:
            for r in promo_filtered:
                rk = []
                for d in _p1_row_dims:
                    pd_key = _promo_dim_map.get(d, d)
                    val = r.get(pd_key)
                    if not val and d in ('日期', '年月'):
                        # 尝试从 _date 字段派生日期/年月
                        _dd = r.get('_date', '') or ''
                        if d == '年月' and len(_dd) >= 7:
                            val = _dd[:7]
                        elif d == '日期' and len(_dd) >= 10:
                            val = _dd[:10]
                    rk.append(val or '未标注')
                rk = tuple(rk)
                _p1_promo_spend[rk] = _p1_promo_spend.get(rk, 0) + float(r.get('_花费', 0) or 0)

        # 计算派生字段
        for rk, v in _p1_agg.items():
            v['支付转化率'] = v['支付买家数'] / v['商品访客数'] if v['商品访客数'] else 0
            v['客单价'] = v['支付金额'] / v['支付买家数'] if v['支付买家数'] else 0
            v['加购率'] = v['商品加购人数'] / v['商品访客数'] if v['商品访客数'] else 0
            v['退款率'] = v['成功退款金额'] / v['支付金额'] if v['支付金额'] else 0
            v['UV价值'] = v['支付金额'] / v['商品访客数'] if v['商品访客数'] else 0
            v['访客占比'] = v['商品访客数'] / _p1_total_visitors
            v['成交占比'] = v['支付金额'] / _p1_total_amt
            v['费率'] = _p1_promo_spend.get(rk, 0) / v['支付金额'] * 100 if v['支付金额'] else None

        # ── 计算同比 ──
        def _p1_yoy_pct(cur_v, ly_v):
            if ly_v and ly_v != 0:
                return (cur_v - ly_v) / ly_v * 100
            return None

        def _p1_yoy_key(rk):
            """将本期行键里的年月/日期字段往前推一年，用于在 yoy_agg 里查找对应 key"""
            new_rk = []
            for i, d in enumerate(_p1_row_dims):
                val = rk[i] if isinstance(rk, tuple) else rk
                if d == '年月' and isinstance(val, str) and len(val) >= 7:
                    try:
                        y, m = int(val[:4]), int(val[5:7])
                        new_rk.append(f'{y-1:04d}-{m:02d}')
                    except Exception:
                        new_rk.append(val)
                elif d == '日期' and isinstance(val, str) and len(val) >= 10:
                    try:
                        import datetime as _dt
                        _d = _dt.date.fromisoformat(val[:10])
                        new_rk.append(str(_d.replace(year=_d.year - 1)))
                    except Exception:
                        new_rk.append(val)
                else:
                    new_rk.append(val)
            return tuple(new_rk)

        for rk, v in _p1_agg.items():
            ly = _p1_yoy_agg.get(_p1_yoy_key(rk), _p1_yoy_agg.get(rk, {}))
            v['_YOY_访客数'] = _p1_yoy_pct(v.get('商品访客数', 0), ly.get('商品访客数', 0))
            v['_YOY_销售额'] = _p1_yoy_pct(v.get('支付金额', 0), ly.get('支付金额', 0))
            v['_YOY_销售量'] = _p1_yoy_pct(v.get('支付件数', 0), ly.get('支付件数', 0))
            # 转化率同比 = 本期转化率 vs 去年同期转化率
            _ly_buyers = ly.get('支付买家数', 0) or 0
            _ly_vis = ly.get('商品访客数', 0) or 0
            _ly_cvr = _ly_buyers / _ly_vis if _ly_vis else 0
            _cur_cvr = v.get('支付转化率', 0) or 0
            v['_YOY_转化率'] = _p1_yoy_pct(_cur_cvr, _ly_cvr)

        _p1_row_keys = sorted(_p1_agg.keys())
        if _p1_top_n > 0:
            _p1_row_keys = sorted(_p1_agg.keys(),
                                  key=lambda k: _p1_agg[k].get('支付金额', 0),
                                  reverse=True)[:_p1_top_n]

        # 合计行
        _p1_grand = {f: sum(v[f] for v in _p1_agg.values()) for f in _P1_RAW_FIELDS}
        _p1_grand['支付转化率'] = _p1_grand['支付买家数'] / _p1_grand['商品访客数'] if _p1_grand['商品访客数'] else 0
        _p1_grand['客单价'] = _p1_grand['支付金额'] / _p1_grand['支付买家数'] if _p1_grand['支付买家数'] else 0
        _p1_grand['加购率'] = _p1_grand['商品加购人数'] / _p1_grand['商品访客数'] if _p1_grand['商品访客数'] else 0
        _p1_grand['退款率'] = _p1_grand['成功退款金额'] / _p1_grand['支付金额'] if _p1_grand['支付金额'] else 0
        _p1_grand['UV价值'] = _p1_grand['支付金额'] / _p1_grand['商品访客数'] if _p1_grand['商品访客数'] else 0
        _p1_grand['访客占比'] = 1.0
        _p1_grand['成交占比'] = 1.0
        _p1_grand['费率'] = sum(_p1_promo_spend.values()) / _p1_grand['支付金额'] * 100 if _p1_grand['支付金额'] else None
        # Grand YoY
        _p1_grand_yoy = {f: sum(v[f] for v in _p1_yoy_agg.values()) for f in _P1_RAW_FIELDS}
        _p1_grand['_YOY_访客数'] = _p1_yoy_pct(_p1_grand.get('商品访客数', 0), _p1_grand_yoy.get('商品访客数', 0))
        _p1_grand['_YOY_销售额'] = _p1_yoy_pct(_p1_grand.get('支付金额', 0), _p1_grand_yoy.get('支付金额', 0))
        _p1_grand['_YOY_销售量'] = _p1_yoy_pct(_p1_grand.get('支付件数', 0), _p1_grand_yoy.get('支付件数', 0))
        _ly_g_buyers = _p1_grand_yoy.get('支付买家数', 0) or 0
        _ly_g_vis = _p1_grand_yoy.get('商品访客数', 0) or 0
        _ly_g_cvr = _ly_g_buyers / _ly_g_vis if _ly_g_vis else 0
        _p1_grand['_YOY_转化率'] = _p1_yoy_pct(_p1_grand.get('支付转化率', 0), _ly_g_cvr)

        def _fmt_s1(mc, val):
            if mc == '费率':
                return '{:.2f}%'.format(val) if val is not None else '--'
            if mc in ('支付金额', 'UV价值', '客单价'):
                return '¥{:,.0f}'.format(val)
            elif mc in ('支付转化率', '加购率', '退款率', '访客占比', '成交占比'):
                return '{:.2f}%'.format(val * 100)
            elif mc == '商品访客数':
                return '{:,.0f}'.format(int(val))
            else:
                return '{:,.0f}'.format(int(val))

        def _fmt_yoy(v):
            if v is None:
                return '--', '#94a3b8'
            sign = '+' if v >= 0 else ''
            color = '#22c55e' if v >= 0 else '#dc2626'
            return f'{sign}{v:.1f}%', color

        _p1_cols = _p1_vals
        # 排序控件：行维度字段（日期/年月等）+ 值指标
        _p1_dim_sort_opts = [d for d in _p1_row_dims if d in ('日期', '年月')]
        _p1_sort_opts = _p1_dim_sort_opts + list(_p1_vals)
        # 默认：若有日期/年月维度，默认选升序；否则默认降序
        _p1_sort_default_idx = 0
        _p1_sort_default_dir = '升序' if _p1_dim_sort_opts else '降序'
        _p1_sort_col, _p1_sort_dir = st.columns([2, 1])
        with _p1_sort_col:
            _p1_sort_key = st.selectbox('排序列', _p1_sort_opts, index=_p1_sort_default_idx, key='pv1_sort')
        with _p1_sort_dir:
            _p1_asc = st.radio('', ['降序', '升序'], horizontal=True, key='pv1_asc',
                               index=0 if _p1_sort_default_dir == '降序' else 1)
        # 按选定列排序：维度字段用字符串自然排序，值指标用数值排序
        def _get_p1_sort_val(rk, mc):
            if mc in _p1_row_dims:
                idx = _p1_row_dims.index(mc)
                return rk[idx] if isinstance(rk, tuple) else str(rk)
            return _p1_agg.get(rk, {}).get(mc, 0) or 0
        _p1_row_keys = sorted(_p1_row_keys,
                              key=lambda k: _get_p1_sort_val(k, _p1_sort_key),
                              reverse=(_p1_asc == '降序'))
        _yoy_cols = ['访客数同比', '转化率同比', '销售额同比', '销售量同比']
        _yoy_keys = ['_YOY_访客数', '_YOY_转化率', '_YOY_销售额', '_YOY_销售量']

        _th_html = '<thead><tr>'
        for d in _p1_row_dims:
            _th_html += '<th style="background:#e2e8f0;color:#1e293b;font-weight:600;padding:6px 10px;">' + d + '</th>'
        for mc in _p1_cols:
            _th_html += '<th style="background:#fef9c3;color:#1e293b;font-weight:600;text-align:right;padding:6px 10px;">' + mc + '</th>'
        _th_html += '<th style="background:#dbeafe;color:#1e293b;font-weight:600;text-align:right;padding:6px 10px;">费率</th>'
        for yc in _yoy_cols:
            _th_html += '<th style="background:#fce7f3;color:#1e293b;font-weight:600;text-align:right;padding:6px 10px;">' + yc + '</th>'
        _th_html += '</tr></thead>'

        _tb_html = '<tbody>'
        for rk in _p1_row_keys:
            _tb_html += '<tr>'
            for i, d in enumerate(_p1_row_dims):
                _tb_html += '<td style="padding:5px 10px;">' + str(rk[i]) + '</td>'
            for mc in _p1_cols:
                _v = _p1_agg.get(rk, {}).get(mc, 0) or 0
                _tb_html += '<td style="text-align:right;padding:5px 10px;">' + _fmt_s1(mc, _v) + '</td>'
            _rate_v = _p1_agg.get(rk, {}).get('费率')
            _tb_html += '<td style="text-align:right;padding:5px 10px;">' + _fmt_s1('费率', _rate_v) + '</td>'
            for yk in _yoy_keys:
                _yv = _p1_agg.get(rk, {}).get(yk)
                _ys, _yc = _fmt_yoy(_yv)
                _tb_html += '<td style="text-align:right;padding:5px 10px;color:' + _yc + ';">' + _ys + '</td>'
            _tb_html += '</tr>'
        _tb_html += '<tr style="background:#fff7ed;font-weight:bold;">'
        for i, d in enumerate(_p1_row_dims):
            _tb_html += '<td style="padding:5px 10px;">' + ('合计' if i == 0 else '') + '</td>'
        for mc in _p1_cols:
            _v = _p1_grand.get(mc, 0) or 0
            _tb_html += '<td style="text-align:right;padding:5px 10px;">' + _fmt_s1(mc, _v) + '</td>'
        _rate_g = _p1_grand.get('费率')
        _tb_html += '<td style="text-align:right;padding:5px 10px;">' + _fmt_s1('费率', _rate_g) + '</td>'
        for yk in _yoy_keys:
            _yv = _p1_grand.get(yk)
            _ys, _yc = _fmt_yoy(_yv)
            _tb_html += '<td style="text-align:right;padding:5px 10px;color:' + _yc + ';">' + _ys + '</td>'
        _tb_html += '</tr></tbody>'

        _html = '<div class="styled-table-wrap" style="max-height:600px;overflow-y:auto;overflow-x:auto;"><table class="styled-table">' + _th_html + _tb_html + '</table></div>'
        st.markdown(_html, unsafe_allow_html=True)
        # 下载原始数据
        _p1_dl = []
        for rk in _p1_row_keys:
            _dlr = {}
            for di, d in enumerate(_p1_row_dims):
                _dlr[d.lstrip('_')] = rk[di] if isinstance(rk, tuple) else str(rk)
            for mc in _p1_vals:
                _dlr[mc.lstrip('_')] = _p1_agg.get(rk, {}).get(mc, 0) or 0
            _dlr['费率'] = _p1_agg.get(rk, {}).get('费率')
            for yk in _yoy_keys:
                _dlr[yk.lstrip('_')] = _p1_agg.get(rk, {}).get(yk)
            _p1_dl.append(_dlr)
        if _p1_dl:
            _render_download_panel(_p1_dl, list(_p1_dl[0].keys()), 'pivot_sales.csv')
    else:
        st.info('请选择至少一个行维度和一个值指标')

    # ══════════════════════════════════════
    # 子专区2：推广数据透视表
    # ══════════════════════════════════════
    st.markdown('---')
    st.markdown('#### 推广数据透视表')

    # 字段配置行（置顶，水平排列）
    _p2_cfg_c1, _p2_cfg_c2, _p2_cfg_c3 = st.columns([1, 2, 1])
    with _p2_cfg_c1:
        _p2_row_dims = st.multiselect('行维度', _pv_promo_dim_opts, default=['_品类'], key='pv2_row')
    with _p2_cfg_c2:
        _p2_vals = st.multiselect('值指标', _pv_promo_metrics_all,
                                   default=_pv_promo_metrics_all, key='pv2_val')
    with _p2_cfg_c3:
        _p2_top_n = st.number_input('前N行（0=全部）', min_value=0, max_value=100,
                                     value=0, key='pv2_top')

    if _p2_vals and _p2_row_dims:
        _p2_raw = promo_rows if promo_rows else []
        # 按日期区间过滤（与销售透视表保持一致）
        _p2_raw = [r for r in _p2_raw if today_s <= (r.get('_date', '') or '') <= today_e]
        # 按全局筛选器过滤
        if channel or store or category or model:
            _p2_filtered = []
            for r in _p2_raw:
                if channel and r.get('_渠道') not in channel: continue
                if store and r.get('_店铺') not in store: continue
                if category and r.get('_品类') not in category: continue
                if model and r.get('_型号') not in model: continue
                _p2_filtered.append(r)
            _p2_raw = _p2_filtered

        # 注入时间维度字段（从 _date 派生）
        for r in _p2_raw:
            d = r.get('_date', '') or r.get('日期', '')
            r['日期'] = d
            r['年月'] = d[:7] if len(d) >= 7 else d
            # 注入产品线和推广计划字段（与趋势分析推广维度一致）
            r['_产品线'] = r.get('产品线', '') or '未标注'
            r['_推广计划'] = r.get('推广计划', '') or r.get('计划ID', '') or '未标注'

        _P2_RAW_FIELDS = ['_花费', '_展现数', '_点击数', '_总订单金额', '_直接订单金额', '_总成交订单量', '_直接订单量']

        def _pv2_group(rows, row_dims):
            agg = {}
            for r in rows:
                rk = tuple((r.get(d) or '未标注') for d in row_dims)
                if rk not in agg:
                    agg[rk] = {f: 0.0 for f in _P2_RAW_FIELDS}
                for f in _P2_RAW_FIELDS:
                    agg[rk][f] += float(r.get(f, 0) or 0)
            return agg

        _p2_agg = _pv2_group(_p2_raw, _p2_row_dims)
        # ── 去年同期：从 promo_rows 按去年日期区间重新过滤，而非用全局 promo_yoy ──
        try:
            import datetime as _pv2_dt
            _p2_yoy_s = str(start.replace(year=start.year - 1))
            _p2_yoy_e = str(end.replace(year=end.year - 1))
        except ValueError:
            _p2_yoy_s = str(start.replace(year=start.year - 1, day=28))
            _p2_yoy_e = str(end.replace(year=end.year - 1, day=28))
        _p2_yoy_raw = [r for r in (promo_rows if promo_rows else [])
                       if _p2_yoy_s <= (r.get('_date', '') or '') <= _p2_yoy_e]
        if channel or store or category or model:
            _p2_yoy_filtered = []
            for r in _p2_yoy_raw:
                if channel and r.get('_渠道') not in channel: continue
                if store and r.get('_店铺') not in store: continue
                if category and r.get('_品类') not in category: continue
                if model and r.get('_型号') not in model: continue
                _p2_yoy_filtered.append(r)
            _p2_yoy_raw = _p2_yoy_filtered
        # 同比数据注入时间维度字段（年月从 _date 的去年日期派生）
        for r in _p2_yoy_raw:
            d = r.get('_date', '') or r.get('日期', '')
            r['日期'] = d
            r['年月'] = d[:7] if len(d) >= 7 else d
            # 注入产品线和推广计划字段（与趋势分析推广维度一致）
            r['_产品线'] = r.get('产品线', '') or '未标注'
            r['_推广计划'] = r.get('推广计划', '') or r.get('计划ID', '') or '未标注'
        _p2_yoy_agg = _pv2_group(_p2_yoy_raw, _p2_row_dims)

        _p2_total_spend = sum(v['_花费'] for v in _p2_agg.values()) or 1
        _p2_total_total_amt = sum(v['_总订单金额'] for v in _p2_agg.values()) or 1
        _p2_total_direct_amt = sum(v['_直接订单金额'] for v in _p2_agg.values()) or 1

        # 计算派生字段
        for rk, v in _p2_agg.items():
            v['_点击率'] = v['_点击数'] / v['_展现数'] if v['_展现数'] else 0
            v['_CPC'] = v['_花费'] / v['_点击数'] if v['_点击数'] else 0
            v['_总ROI'] = v['_总订单金额'] / v['_花费'] if v['_花费'] else 0
            v['_直接ROI'] = v['_直接订单金额'] / v['_花费'] if v['_花费'] else 0
            v['_总转化率'] = v['_总成交订单量'] / v['_点击数'] if v['_点击数'] else 0
            v['_直接转化率'] = v['_直接订单量'] / v['_点击数'] if v['_点击数'] else 0
            v['_花费占比'] = v['_花费'] / _p2_total_spend
            v['_总金额占比'] = v['_总订单金额'] / _p2_total_total_amt
            v['_直接金额占比'] = v['_直接订单金额'] / _p2_total_direct_amt
            v['_总订单成本'] = v['_花费'] / v['_总订单金额'] * 100 if v['_总订单金额'] else 0
            v['_直接订单成本'] = v['_花费'] / v['_直接订单金额'] * 100 if v['_直接订单金额'] else 0

        # ── 计算同比 ──
        def _p2_yoy_pct(cur_v, ly_v):
            if ly_v and ly_v != 0:
                return (cur_v - ly_v) / ly_v * 100
            return None

        def _p2_yoy_key(rk):
            """将本期行键里的年月/日期字段往前推一年，用于在 yoy_agg 里查找对应 key"""
            new_rk = []
            for i, d in enumerate(_p2_row_dims):
                val = rk[i] if isinstance(rk, tuple) else rk
                d_strip = d.lstrip('_')
                if d_strip == '年月' and isinstance(val, str) and len(val) >= 7:
                    try:
                        y, m = int(val[:4]), int(val[5:7])
                        new_rk.append(f'{y-1:04d}-{m:02d}')
                    except Exception:
                        new_rk.append(val)
                elif d_strip == '日期' and isinstance(val, str) and len(val) >= 10:
                    try:
                        import datetime as _dt2
                        _d2 = _dt2.date.fromisoformat(val[:10])
                        new_rk.append(str(_d2.replace(year=_d2.year - 1)))
                    except Exception:
                        new_rk.append(val)
                else:
                    new_rk.append(val)
            return tuple(new_rk)

        for rk, v in _p2_agg.items():
            ly = _p2_yoy_agg.get(_p2_yoy_key(rk), _p2_yoy_agg.get(rk, {}))
            v['_YOY_花费'] = _p2_yoy_pct(v.get('_花费', 0), ly.get('_花费', 0))
            _ly_spend = ly.get('_花费', 0) or 0
            _ly_click = ly.get('_点击数', 0) or 0
            # 直接ROI同比
            _ly_droi = ly.get('_直接订单金额', 0) / _ly_spend if _ly_spend else 0
            _cur_droi = v.get('_直接ROI', 0) or 0
            v['_YOY_直接ROI'] = _p2_yoy_pct(_cur_droi, _ly_droi) if _ly_spend else None
            # 总ROI同比
            _ly_roi = ly.get('_总订单金额', 0) / _ly_spend if _ly_spend else 0
            _cur_roi = v.get('_总ROI', 0) or 0
            v['_YOY_总ROI'] = _p2_yoy_pct(_cur_roi, _ly_roi) if _ly_spend else None
            # CPC同比
            _ly_cpc = _ly_spend / _ly_click if _ly_click else 0
            _cur_cpc = v.get('_CPC', 0) or 0
            v['_YOY_CPC'] = _p2_yoy_pct(_cur_cpc, _ly_cpc) if _ly_click else None
            # 转化率同比（基于总成交订单量）
            _ly_orders = ly.get('_总成交订单量', 0) or 0
            _ly_cvr = _ly_orders / _ly_click if _ly_click else 0
            _cur_cvr = v.get('_总转化率', 0) or 0
            v['_YOY_转化率'] = _p2_yoy_pct(_cur_cvr, _ly_cvr) if _ly_click else None

        _p2_row_keys = sorted(_p2_agg.keys())
        # 排序控件（仅在前N行=0时启用自由排序）
        _p2_sort_enabled = (_p2_top_n == 0)
        if _p2_sort_enabled:
            # 行维度中的日期/年月字段（不带下划线前缀）
            _p2_dim_sort_opts = [d.lstrip('_') for d in _p2_row_dims if d.lstrip('_') in ('日期', '年月')]
            # 构建可排序的指标列表（日期维度在前）
            _p2_sort_metrics_base = ['花费', '展现数', '点击数', '点击率', 'CPC', '总ROI', '直接ROI',
                               '总转化率', '总订单金额', '直接订单金额', '花费占比']
            _p2_sort_metrics = _p2_dim_sort_opts + _p2_sort_metrics_base
            _p2_sort_default_dir = '升序' if _p2_dim_sort_opts else '降序'
            _p2_sort_col, _p2_sort_dir = st.columns([2, 1])
            with _p2_sort_col:
                _p2_sort_key = st.selectbox('排序列', _p2_sort_metrics, index=0, key='pv2_sort')
            with _p2_sort_dir:
                _p2_asc = st.radio('', ['降序', '升序'], horizontal=True, key='pv2_asc',
                                   index=0 if _p2_sort_default_dir == '降序' else 1)
            # 维度字段用字符串排序，值指标用数值排序
            def _get_p2_sort_val(rk, mc):
                if mc in ('日期', '年月'):
                    raw_dim = mc  # 推广透视表维度字段不带下划线
                    dim_list_stripped = [d.lstrip('_') for d in _p2_row_dims]
                    if mc in dim_list_stripped:
                        idx = dim_list_stripped.index(mc)
                        return rk[idx] if isinstance(rk, tuple) else str(rk)
                    return ''
                return _p2_agg.get(rk, {}).get('_' + mc, 0) or 0
            _p2_row_keys = sorted(_p2_row_keys,
                                   key=lambda k: _get_p2_sort_val(k, _p2_sort_key),
                                   reverse=(_p2_asc == '降序'))
        elif _p2_top_n > 0:
            _p2_row_keys = sorted(_p2_agg.keys(),
                                  key=lambda k: _p2_agg[k].get('_花费', 0),
                                  reverse=True)[:_p2_top_n]

        # 合计行
        _p2_grand = {f: sum(v[f] for v in _p2_agg.values()) for f in _P2_RAW_FIELDS}
        _p2_grand['_点击率'] = _p2_grand['_点击数'] / _p2_grand['_展现数'] if _p2_grand['_展现数'] else 0
        _p2_grand['_CPC'] = _p2_grand['_花费'] / _p2_grand['_点击数'] if _p2_grand['_点击数'] else 0
        _p2_grand['_总ROI'] = _p2_grand['_总订单金额'] / _p2_grand['_花费'] if _p2_grand['_花费'] else 0
        _p2_grand['_直接ROI'] = _p2_grand['_直接订单金额'] / _p2_grand['_花费'] if _p2_grand['_花费'] else 0
        _p2_grand['_总转化率'] = _p2_grand['_总成交订单量'] / _p2_grand['_点击数'] if _p2_grand['_点击数'] else 0
        _p2_grand['_直接转化率'] = _p2_grand['_直接订单量'] / _p2_grand['_点击数'] if _p2_grand['_点击数'] else 0
        _p2_grand['_花费占比'] = 1.0
        _p2_grand['_总金额占比'] = 1.0
        _p2_grand['_直接金额占比'] = 1.0
        _p2_grand['_总订单成本'] = _p2_grand['_花费'] / _p2_grand['_总订单金额'] * 100 if _p2_grand['_总订单金额'] else 0
        _p2_grand['_直接订单成本'] = _p2_grand['_花费'] / _p2_grand['_直接订单金额'] * 100 if _p2_grand['_直接订单金额'] else 0
        # Grand YoY
        _p2_grand_yoy = {f: sum(v[f] for v in _p2_yoy_agg.values()) for f in _P2_RAW_FIELDS}
        _p2_grand['_YOY_花费'] = _p2_yoy_pct(_p2_grand.get('_花费', 0), _p2_grand_yoy.get('_花费', 0))
        _g_ly_spend = _p2_grand_yoy.get('_花费', 0) or 1
        _g_ly_click = _p2_grand_yoy.get('_点击数', 0) or 1
        _p2_grand['_YOY_直接ROI'] = _p2_yoy_pct(_p2_grand['_直接ROI'], _p2_grand_yoy.get('_直接订单金额', 0) / _g_ly_spend if _g_ly_spend else 0)
        _p2_grand['_YOY_总ROI'] = _p2_yoy_pct(_p2_grand['_总ROI'], _p2_grand_yoy.get('_总订单金额', 0) / _g_ly_spend if _g_ly_spend else 0)
        _p2_grand['_YOY_CPC'] = _p2_yoy_pct(_p2_grand['_CPC'], _g_ly_spend / _g_ly_click if _g_ly_click else 0)
        _g_ly_orders = _p2_grand_yoy.get('_总成交订单量', 0) or 0
        _p2_grand['_YOY_转化率'] = _p2_yoy_pct(_p2_grand['_总转化率'], _g_ly_orders / _g_ly_click if _g_ly_click else 0)

        def _fmt_p2(mc, val):
            if mc in ('_花费', '_总订单金额', '_直接订单金额'):
                return '¥{:,.0f}'.format(val)
            elif mc == '_CPC':
                return '¥{:.2f}'.format(val)
            elif mc in ('_点击率', '_花费占比', '_总金额占比', '_直接金额占比', '_总转化率', '_直接转化率'):
                return '{:.2f}%'.format(val * 100)
            elif mc in ('_总ROI', '_直接ROI'):
                return '{:.2f}'.format(val)
            elif mc in ('_总订单成本', '_直接订单成本'):
                return '{:.2f}%'.format(val)
            elif mc in ('_展现数', '_点击数', '_总成交订单量', '_直接订单量'):
                return '{:,.0f}'.format(int(val))
            else:
                return '{:,.2f}'.format(val)

        def _fmt_p2_yoy(v):
            if v is None:
                return '--', '#94a3b8'
            sign = '+' if v >= 0 else ''
            color = '#22c55e' if v >= 0 else '#dc2626'
            return f'{sign}{v:.1f}%', color

        _p2_cols = _p2_vals
        _p2_yoy_cols = ['花费同比', '直接ROI同比', '总ROI同比', 'CPC同比', '转化率同比']
        _p2_yoy_keys = ['_YOY_花费', '_YOY_直接ROI', '_YOY_总ROI', '_YOY_CPC', '_YOY_转化率']

        _th_html = '<thead><tr>'
        for d in _p2_row_dims:
            _label = d.lstrip('_')
            _th_html += '<th style="background:#e2e8f0;color:#1e293b;font-weight:600;padding:6px 10px;">' + _label + '</th>'
        for mc in _p2_cols:
            _label = mc.lstrip('_')
            _th_html += '<th style="background:#dbeafe;color:#1e293b;font-weight:600;text-align:right;padding:6px 10px;">' + _label + '</th>'
        for yc in _p2_yoy_cols:
            _th_html += '<th style="background:#fce7f3;color:#1e293b;font-weight:600;text-align:right;padding:6px 10px;">' + yc + '</th>'
        _th_html += '</tr></thead>'

        _tb_html = '<tbody>'
        for rk in _p2_row_keys:
            _tb_html += '<tr>'
            for i, d in enumerate(_p2_row_dims):
                _tb_html += '<td style="padding:5px 10px;">' + str(rk[i]) + '</td>'
            for mc in _p2_cols:
                _v = _p2_agg.get(rk, {}).get(mc, 0) or 0
                _tb_html += '<td style="text-align:right;padding:5px 10px;">' + _fmt_p2(mc, _v) + '</td>'
            for yk in _p2_yoy_keys:
                _yv = _p2_agg.get(rk, {}).get(yk)
                _ys, _yc = _fmt_p2_yoy(_yv)
                _tb_html += '<td style="text-align:right;padding:5px 10px;color:' + _yc + ';">' + _ys + '</td>'
            _tb_html += '</tr>'
        _tb_html += '<tr style="background:#fff7ed;font-weight:bold;">'
        for i, d in enumerate(_p2_row_dims):
            _tb_html += '<td style="padding:5px 10px;">' + ('合计' if i == 0 else '') + '</td>'
        for mc in _p2_cols:
            _v = _p2_grand.get(mc, 0) or 0
            _tb_html += '<td style="text-align:right;padding:5px 10px;">' + _fmt_p2(mc, _v) + '</td>'
        for yk in _p2_yoy_keys:
            _yv = _p2_grand.get(yk)
            _ys, _yc = _fmt_p2_yoy(_yv)
            _tb_html += '<td style="text-align:right;padding:5px 10px;color:' + _yc + ';">' + _ys + '</td>'
        _tb_html += '</tr></tbody>'

        _html = '<div class="styled-table-wrap" style="max-height:600px;overflow-y:auto;overflow-x:auto;"><table class="styled-table">' + _th_html + _tb_html + '</table></div>'
        st.markdown(_html, unsafe_allow_html=True)
        # 下载原始数据
        _p2_dl = []
        for rk in _p2_row_keys:
            _dlr = {}
            for di, d in enumerate(_p2_row_dims):
                _dlr[d.lstrip('_')] = rk[di] if isinstance(rk, tuple) else str(rk)
            for mc in _p2_cols:
                _dlr[mc.lstrip('_')] = _p2_agg.get(rk, {}).get(mc, 0) or 0
            for yk in _p2_yoy_keys:
                _dlr[yk.lstrip('_')] = _p2_agg.get(rk, {}).get(yk)
            _p2_dl.append(_dlr)
        if _p2_dl:
            _render_download_panel(_p2_dl, list(_p2_dl[0].keys()), 'pivot_promo.csv')
    else:
        st.info('请选择至少一个行维度和一个值指标')

# ═══════════════════════════════════════════════════════════════
# TAB 6: 目标达成
# ═══════════════════════════════════════════════════════════════
with tabs[6]:
    st.markdown("""<div class="hero" style="padding:20px 28px 10px;"><div><span class="badge">🎯</span><span class="badge">目标达成追踪</span></div>
        <div class="hero-sub">店铺/单品日度目标 vs 实际达成，自动从销售数据填充</div></div>""", unsafe_allow_html=True)

    if not targets:
        st.info('📭 尚未上传目标数据。请在左侧「数据源更新」中选择「销售目标」上传 Excel，然后点击同步到云端。')
    else:
        _target_months = sorted(targets.keys(), reverse=True)
        _sel_ym = st.selectbox('选择目标月份', _target_months,
                               index=0, key='target_ym')

        if _sel_ym and _sel_ym in targets:
            tgt = targets[_sel_ym]
            shop_rows = tgt.get('shop', [])
            model_rows = tgt.get('model', [])
            date_list = tgt.get('dates', [])

            # 目标达成模块不受全局筛选影响，只受目标月份选择影响
            # 从原始销售数据中筛选目标月份的日期
            _cur_ym_prefix = _sel_ym + '-'
            _raw_daily_target = [r for r in data['daily'] if r.get('日期', '').startswith(_cur_ym_prefix)]
            # 推广数据同样不受全局筛选影响
            _all_promo_raw = load_promo_data(_promo_bytes) if _promo_bytes else []
            _raw_promo_target = [r for r in _all_promo_raw if r.get('_date', '').startswith(_cur_ym_prefix)]

            if not _raw_daily_target:
                st.warning('请先上传销售数据，才能自动计算达成率')
                st.stop()

            # ── 预计算销售日汇总（店铺+单品）──
            daily_by_shop_date = {}
            daily_by_model_date = {}
            for r in _raw_daily_target:
                shop = (r.get('店铺', '') or '').strip()
                model_name = (r.get('型号', '') or '').strip()
                d = r.get('日期', '')
                pay_amt = float(r.get('支付金额', 0) or 0)
                pay_qty = float(r.get('支付件数', 0) or 0)
                pay_buyers = float(r.get('支付买家数', 0) or 0)
                visitors = float(r.get('商品访客数', 0) or 0)

                key_sd = (shop, d)
                if key_sd not in daily_by_shop_date:
                    daily_by_shop_date[key_sd] = {'支付金额': 0.0, '支付件数': 0.0, '支付买家数': 0.0, '商品访客数': 0.0}
                daily_by_shop_date[key_sd]['支付金额'] += pay_amt
                daily_by_shop_date[key_sd]['支付件数'] += pay_qty
                daily_by_shop_date[key_sd]['支付买家数'] += pay_buyers
                daily_by_shop_date[key_sd]['商品访客数'] += visitors

                key_md = (shop, model_name, d)
                if key_md not in daily_by_model_date:
                    daily_by_model_date[key_md] = {'支付金额': 0.0, '支付件数': 0.0, '支付买家数': 0.0, '商品访客数': 0.0}
                daily_by_model_date[key_md]['支付金额'] += pay_amt
                daily_by_model_date[key_md]['支付件数'] += pay_qty
                daily_by_model_date[key_md]['支付买家数'] += pay_buyers
                daily_by_model_date[key_md]['商品访客数'] += visitors

            # ── 预计算推广花费日汇总（按店铺+日期）──
            promo_by_shop_date = {}
            if _raw_promo_target:
                for r in _raw_promo_target:
                    p_shop = (r.get('_店铺', '') or '').strip()
                    p_date = r.get('_date', '')
                    p_spend = r.get('_花费', 0.0)
                    key_ps = (p_shop, p_date)
                    if key_ps not in promo_by_shop_date:
                        promo_by_shop_date[key_ps] = 0.0
                    promo_by_shop_date[key_ps] += p_spend

            # ── 预计算推广花费日汇总（按店铺+型号+日期）──
            promo_by_model_date = {}
            if _raw_promo_target:
                for r in _raw_promo_target:
                    p_shop = (r.get('_店铺', '') or '').strip()
                    p_model = (r.get('_型号', '') or '').strip()
                    p_date = r.get('_date', '')
                    p_spend = r.get('_花费', 0.0)
                    if not p_model:
                        continue
                    key_pm = (p_shop, p_model, p_date)
                    if key_pm not in promo_by_model_date:
                        promo_by_model_date[key_pm] = 0.0
                    promo_by_model_date[key_pm] += p_spend

            # ── 预计算推广直接订单金额日汇总（按店铺+日期）──
            promo_direct_by_shop_date = {}
            if _raw_promo_target:
                for r in _raw_promo_target:
                    p_shop = (r.get('_店铺', '') or '').strip()
                    p_date = r.get('_date', '')
                    p_direct = r.get('_直接订单金额', 0.0)
                    key_ps = (p_shop, p_date)
                    if key_ps not in promo_direct_by_shop_date:
                        promo_direct_by_shop_date[key_ps] = 0.0
                    promo_direct_by_shop_date[key_ps] += p_direct

            # ── 预计算推广直接订单金额日汇总（按店铺+型号+日期）──
            promo_direct_by_model_date = {}
            if _raw_promo_target:
                for r in _raw_promo_target:
                    p_shop = (r.get('_店铺', '') or '').strip()
                    p_model = (r.get('_型号', '') or '').strip()
                    p_date = r.get('_date', '')
                    p_direct = r.get('_直接订单金额', 0.0)
                    if not p_model:
                        continue
                    key_pm = (p_shop, p_model, p_date)
                    if key_pm not in promo_direct_by_model_date:
                        promo_direct_by_model_date[key_pm] = 0.0
                    promo_direct_by_model_date[key_pm] += p_direct

            # ── 预计算去年同期数据（用于结构表同比）──
            _yoy_ym = None
            _yoy_date_list = []      # 同期天数（用于实际同比，基于今年实际有数据的日期）
            _yoy_date_list_full = [] # 全月（用于目标同比）
            _yoy_targets = None
            try:
                _cur_ym_parts = _sel_ym.split('-')
                _yoy_year = int(_cur_ym_parts[0]) - 1
                _yoy_ym = f'{_yoy_year}-{_cur_ym_parts[1]}'
                # 去年同期目标数据（可选，未上传时目标同比为空）
                if _yoy_ym in targets:
                    _yoy_targets = targets[_yoy_ym]
                # 实际同比日期：从今年销售数据中提取当月实际有数据的日期，映射到去年
                _cur_ym_prefix = _sel_ym  # e.g. '2026-06'
                _actual_dates_set = set()
                for r in _raw_daily_target:
                    d = r.get('日期', '')
                    if d and d.startswith(_cur_ym_prefix):
                        _actual_dates_set.add(d)
                # 按日期排序
                _actual_dates = sorted(_actual_dates_set)
                for d in _actual_dates:
                    try:
                        _dt = datetime.datetime.strptime(d, '%Y-%m-%d')
                        _ly_dt = _dt.replace(year=_dt.year - 1)
                        _yoy_date_list.append(_ly_dt.strftime('%Y-%m-%d'))
                    except ValueError:
                        pass
                # 全月：去年同期整月所有日期（用于目标同比）
                import calendar
                _last_day = calendar.monthrange(_yoy_year, int(_cur_ym_parts[1]))[1]
                _yoy_date_list_full = [f'{_yoy_year}-{int(_cur_ym_parts[1]):02d}-{day:02d}' for day in range(1, _last_day + 1)]
            except (ValueError, IndexError):
                pass

            # 去年同期销售日汇总（按店铺+日期）- 用于实际同比（同期天数）
            # 注意：去年同期数据必须从全量 data['daily'] 提取，不能从 _raw_daily_target（仅当前月份）
            _daily_shop_ly = {}
            _daily_model_ly = {}
            if _yoy_date_list:
                for r in data['daily']:
                    d = r.get('日期', '')
                    if d not in _yoy_date_list:
                        continue
                    shop = (r.get('店铺', '') or '').strip()
                    model_name = (r.get('型号', '') or '').strip()
                    pay_amt = float(r.get('支付金额', 0) or 0)
                    key_sd = (shop, d)
                    if key_sd not in _daily_shop_ly:
                        _daily_shop_ly[key_sd] = 0.0
                    _daily_shop_ly[key_sd] += pay_amt
                    key_md = (shop, model_name, d)
                    if key_md not in _daily_model_ly:
                        _daily_model_ly[key_md] = 0.0
                    _daily_model_ly[key_md] += pay_amt

            # 去年同期推广花费日汇总（按店铺+日期）- 用于实际同比（同期天数）
            # 注意：去年同期数据必须从全量 _all_promo_raw 提取，不能从 _raw_promo_target（仅当前月份）
            _promo_shop_ly = {}
            _promo_model_ly = {}
            if _yoy_date_list and _all_promo_raw:
                for r in _all_promo_raw:
                    d = r.get('_date', '')
                    if d not in _yoy_date_list:
                        continue
                    p_shop = (r.get('_店铺', '') or '').strip()
                    p_model = (r.get('_型号', '') or '').strip()
                    p_spend = r.get('_花费', 0.0)
                    key_ps = (p_shop, d)
                    if key_ps not in _promo_shop_ly:
                        _promo_shop_ly[key_ps] = 0.0
                    _promo_shop_ly[key_ps] += p_spend
                    if p_model:
                        key_pm = (p_shop, p_model, d)
                        if key_pm not in _promo_model_ly:
                            _promo_model_ly[key_pm] = 0.0
                        _promo_model_ly[key_pm] += p_spend

            # ── 去年全月数据（用于目标同比）──
            # 注意：去年同期数据必须从全量 data['daily'] 提取
            _daily_shop_ly_full = {}
            _daily_model_ly_full = {}
            if _yoy_date_list_full:
                for r in data['daily']:
                    d = r.get('日期', '')
                    if d not in _yoy_date_list_full:
                        continue
                    shop = (r.get('店铺', '') or '').strip()
                    model_name = (r.get('型号', '') or '').strip()
                    pay_amt = float(r.get('支付金额', 0) or 0)
                    key_sd = (shop, d)
                    if key_sd not in _daily_shop_ly_full:
                        _daily_shop_ly_full[key_sd] = 0.0
                    _daily_shop_ly_full[key_sd] += pay_amt
                    key_md = (shop, model_name, d)
                    if key_md not in _daily_model_ly_full:
                        _daily_model_ly_full[key_md] = 0.0
                    _daily_model_ly_full[key_md] += pay_amt

            # 注意：去年同期推广全月数据也必须从全量 _all_promo_raw 提取
            _promo_shop_ly_full = {}
            _promo_model_ly_full = {}
            if _yoy_date_list_full and _all_promo_raw:
                for r in _all_promo_raw:
                    d = r.get('_date', '')
                    if d not in _yoy_date_list_full:
                        continue
                    p_shop = (r.get('_店铺', '') or '').strip()
                    p_model = (r.get('_型号', '') or '').strip()
                    p_spend = r.get('_花费', 0.0)
                    key_ps = (p_shop, d)
                    if key_ps not in _promo_shop_ly_full:
                        _promo_shop_ly_full[key_ps] = 0.0
                    _promo_shop_ly_full[key_ps] += p_spend
                    if p_model:
                        key_pm = (p_shop, p_model, d)
                        if key_pm not in _promo_model_ly_full:
                            _promo_model_ly_full[key_pm] = 0.0
                        _promo_model_ly_full[key_pm] += p_spend

            def _yoy_pct_val(cur, ly):
                """同比变化率，返回数值（如 0.15 表示 +15%）"""
                if ly and ly != 0:
                    return (cur - ly) / ly
                return None

            def _fmt_yoy_val(v):
                """格式化同比值为字符串，如 +12.5% / -8.3% / --"""
                if v is None:
                    return '--'
                return f'{v*100:+.1f}%'

            def _yoy_color(v):
                """同比颜色"""
                if v is None:
                    return '#94a3b8'
                return '#22c55e' if v >= 0 else '#ef4444'

            import uuid as _uuid_mod

            def _indicator_type(indicator):
                """判断指标类型: 'target'=目标值, 'actual'=实际值(需填充), 'calc'=计算值(公式)"""
                # 计算行（需要公式计算）
                if '达成率' in indicator:
                    return 'calc'
                if indicator in ('实际费率', '实际费率(%)'):
                    return 'calc'
                # 实际值行（需从销售/推广数据填充）
                if '实际' in indicator or '达成' in indicator:
                    return 'actual'
                # 其余为目标行
                return 'target'

            def _get_actual_value(indicator, shop_name, d, model_name=''):
                """根据指标名称和店铺/单品，从销售或推广数据中获取实际值"""
                shop_name = (shop_name or '').strip()
                model_name = (model_name or '').strip()
                if model_name:
                    sd = daily_by_model_date.get((shop_name, model_name, d), {})
                else:
                    sd = daily_by_shop_date.get((shop_name, d), {})
                # 推广花费：单品模式用型号级，店铺模式用店铺级
                if model_name:
                    promo_val = promo_by_model_date.get((shop_name, model_name, d), 0.0)
                else:
                    promo_val = promo_by_shop_date.get((shop_name, d), 0.0)

                # 金额类实际值
                if '成交金额' in indicator or '销额' in indicator or '支付金额' in indicator:
                    return sd.get('支付金额', 0.0)
                # 件数类实际值
                if '件数' in indicator or '支付件数' in indicator or '销量' in indicator:
                    return sd.get('支付件数', 0.0)
                # 推广花费
                if '投入' in indicator or '花费' in indicator:
                    return promo_val
                # 默认
                return 0.0

            def _get_target_row_by_indicator(rows_data, target_indicator):
                """从目标行列表中查找对应目标值行（按指标名匹配）"""
                for tr in rows_data:
                    if tr['指标'] == target_indicator:
                        return tr
                return None

            def _fmt_val(v, is_pct=False):
                """格式化数值：金额用千分位，百分比保留1位"""
                if is_pct:
                    return f'{v:.1f}%'
                if v == 0:
                    return '--'
                return f'{v:,.0f}'

            def _fmt_rate_cell(rate_val):
                """格式化达成率单元格，带颜色"""
                if rate_val is None:
                    return '--'
                color = '#22c55e' if rate_val >= 100 else '#ef4444' if rate_val < 80 else '#f59e0b'
                return f'<span style="color:{color};font-weight:bold;">{rate_val:.1f}%</span>'

            def _build_fullscreen_js(tbl_id, title):
                """为目标达成表格生成全屏 JS"""
                return f"""
<script>
(function() {{
    var overlay = null;
    window['_fsOpen_{tbl_id}'] = function() {{
        if (overlay) return;
        var tblWrap = document.getElementById('{tbl_id}');
        var content = tblWrap.innerHTML;
        overlay = document.createElement('div');
        overlay.id = '{tbl_id}_fs';
        overlay.innerHTML = '<style>.styled-table{{width:100%;border-collapse:collapse;font-size:13px;}}.styled-table th{{background:#1e293b;color:#e2e8f0;border-bottom:2px solid #334155;padding:10px 8px;position:sticky;top:0;z-index:1;}}.styled-table td{{padding:7px 10px;border-bottom:1px solid #e5e7eb;white-space:nowrap;color:#1e293b;}}.styled-table tbody tr:nth-child(even){{background:#f8fafc;}}.styled-table tbody tr:hover{{background:#e2e8f0;}}</style>' +
            '<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:10px;flex-shrink:0;padding:0 8px;"><span style="color:#fff;font-size:18px;font-weight:700;">{title}</span><button onclick="window._fsClose_{tbl_id}()" style="background:#ef4444;color:#fff;border:none;border-radius:6px;padding:6px 18px;cursor:pointer;font-size:14px;font-weight:600;">✕ 关闭</button></div><div style="flex:1;overflow:auto;background:#fff;border-radius:8px;min-height:0;">' + content + '</div>';
        overlay.style.cssText = 'display:flex;position:fixed;top:0;left:0;width:100vw;height:100vh;background:rgba(0,0,0,0.82);z-index:2147483647;flex-direction:column;padding:20px;box-sizing:border-box;';
        document.body.appendChild(overlay);
    }};
    window['_fsClose_{tbl_id}'] = function() {{
        if (overlay) {{ overlay.remove(); overlay = null; }}
    }};
}})();
</script>
"""

            def _render_target_table(header_cols, title, table_data):
                """渲染带全屏按钮的目标达成表格"""
                tbl_id = 'tgt_' + _uuid_mod.uuid4().hex[:8]
                fs_js = _build_fullscreen_js(tbl_id, title)
                fs_btn = (
                    f'<button onclick="window._fsOpen_{tbl_id}()" '
                    f'style="float:right;margin-bottom:4px;padding:3px 10px;font-size:12px;'
                    f'background:#1d4ed8;color:#fff;border:none;border-radius:4px;cursor:pointer;">⛶ 全屏</button>'
                )
                html = fs_js
                html += f'<div style="display:flex;justify-content:space-between;align-items:center;"><span style="font-weight:700;font-size:14px;">{title}</span>{fs_btn}</div>'
                html += f'<div id="{tbl_id}" class="styled-table-wrap"><table class="styled-table">'
                html += '<thead><tr>'
                for h in header_cols:
                    html += f'<th style="text-align:center;white-space:nowrap;">{h}</th>'
                html += '</tr></thead><tbody>'
                for row in table_data:
                    html += '<tr>'
                    for j, cell in enumerate(row):
                        align = 'left' if j == 0 else 'right'
                        html += f'<td style="text-align:{align};white-space:nowrap;">{cell}</td>'
                    html += '</tr>'
                html += '</tbody></table></div>'
                st.markdown(html, unsafe_allow_html=True)

            def _is_rate_indicator(indicator):
                """判断是否为费率类指标（需用%格式化）"""
                return '费率' in indicator

            def _build_rows_for_entity(rows_data, shop_name, model_name=''):
                """
                为一个店铺或单品构建表格数据行。
                两遍遍历：第一遍收集 target/actual 值，第二遍生成 calc 行。
                返回 (table_data, actual_summary) 其中 actual_summary 用于全店铺合计。
                """
                shop_name = (shop_name or '').strip()
                model_name = (model_name or '').strip()
                table_data = []
                actual_summary = {}  # {indicator: {'合计': val, date: val}}

                # ── 第一遍：处理 target 和 actual 行，构建 actual_summary ──
                calc_rows = []  # 暂存 calc 行，第二遍处理
                for sr in rows_data:
                    indicator = sr['指标']
                    # 跳过纯数字的垃圾指标名（如 "0.4"）
                    try:
                        float(indicator)
                        continue
                    except ValueError:
                        pass
                    itype = _indicator_type(indicator)

                    if itype == 'target':
                        is_rate = _is_rate_indicator(indicator)
                        row_vals = [indicator]
                        total = sr.get('合计', 0) or 0
                        if is_rate:
                            row_vals.append(_fmt_val(total * 100, is_pct=True))
                        else:
                            row_vals.append(f'{total:,.0f}')
                        for d in date_list:
                            v = sr.get(d, 0) or 0
                            if is_rate:
                                row_vals.append(_fmt_val(v * 100, is_pct=True))
                            else:
                                row_vals.append(f'{v:,.0f}' if v else '--')
                        table_data.append(row_vals)

                    elif itype == 'actual':
                        # 单品模式：跳过Excel中的actual行（实际值由代码自动计算），避免与统一指标行重复
                        if model_name:
                            continue
                        actual_total = 0.0
                        for d in date_list:
                            av = _get_actual_value(indicator, shop_name, d, model_name)
                            actual_total += av
                        row_vals = [indicator]
                        row_vals.append(f'{actual_total:,.0f}' if actual_total else '--')
                        for d in date_list:
                            av = _get_actual_value(indicator, shop_name, d, model_name)
                            row_vals.append(f'{av:,.0f}' if av else '--')
                        table_data.append(row_vals)
                        # 记录实际值供计算行和全店铺合计使用
                        actual_summary[indicator] = {'合计': actual_total}
                        for d in date_list:
                            actual_summary[indicator][d] = _get_actual_value(indicator, shop_name, d, model_name)

                    elif itype == 'calc':
                        calc_rows.append(sr)

                # ── 第一遍后：自动追加缺失的 actual 行 ──
                # 1. 店铺模式：追加"实际支付件数"（如果Excel中没有）
                if not model_name and '实际支付件数' not in actual_summary:
                    qty_total = 0.0
                    qty_row = ['实际支付件数']
                    for d in date_list:
                        qv = _get_actual_value('实际支付件数', shop_name, d)
                        qty_total += qv
                    qty_row.append(f'{qty_total:,.0f}' if qty_total else '--')
                    for d in date_list:
                        qv = _get_actual_value('实际支付件数', shop_name, d)
                        qty_row.append(f'{qv:,.0f}' if qv else '--')
                    table_data.append(qty_row)
                    actual_summary['实际支付件数'] = {'合计': qty_total}
                    for d in date_list:
                        actual_summary['实际支付件数'][d] = _get_actual_value('实际支付件数', shop_name, d)

                # 2. 单品模式：填充 actual_summary 供 calc 行计算使用，并追加"实际投入金额"行到表格
                if model_name:
                    # 实际成交金额：用于达成率计算
                    amt_total = 0.0
                    for d in date_list:
                        av = _get_actual_value('成交金额', shop_name, d, model_name)
                        amt_total += av
                    actual_summary['实际成交金额'] = {'合计': amt_total}
                    for d in date_list:
                        actual_summary['实际成交金额'][d] = _get_actual_value('成交金额', shop_name, d, model_name)
                    # 实际投入金额：用于费率计算 + 追加到表格展示
                    spend_total = 0.0
                    for d in date_list:
                        sv = promo_by_model_date.get((shop_name, model_name, d), 0.0)
                        spend_total += sv
                    actual_summary['实际投入金额'] = {'合计': spend_total}
                    for d in date_list:
                        sv = promo_by_model_date.get((shop_name, model_name, d), 0.0)
                        actual_summary['实际投入金额'][d] = sv
                    # 追加"实际投入金额"行到表格（显示每日推广花费）
                    spend_row = ['实际推广花费']
                    spend_row.append(f'{spend_total:,.0f}' if spend_total else '--')
                    for d in date_list:
                        sv = promo_by_model_date.get((shop_name, model_name, d), 0.0)
                        spend_row.append(f'{sv:,.0f}' if sv else '--')
                    table_data.append(spend_row)
                    # 实际支付件数
                    qty_total = 0.0
                    for d in date_list:
                        qv = _get_actual_value('实际支付件数', shop_name, d, model_name)
                        qty_total += qv
                    actual_summary['实际支付件数'] = {'合计': qty_total}
                    for d in date_list:
                        actual_summary['实际支付件数'][d] = _get_actual_value('实际支付件数', shop_name, d, model_name)
                    # 追加"实际成交金额"行到表格
                    amt_row = ['实际成交金额']
                    amt_row.append(f'{amt_total:,.0f}' if amt_total else '--')
                    for d in date_list:
                        av = _get_actual_value('成交金额', shop_name, d, model_name)
                        amt_row.append(f'{av:,.0f}' if av else '--')
                    table_data.append(amt_row)
                    # 追加"实际支付件数"行到表格
                    qty_row = ['实际支付件数']
                    qty_row.append(f'{qty_total:,.0f}' if qty_total else '--')
                    for d in date_list:
                        qv = _get_actual_value('实际支付件数', shop_name, d, model_name)
                        qty_row.append(f'{qv:,.0f}' if qv else '--')
                    table_data.append(qty_row)

                # ── 第二遍：处理 calc 行（此时 actual_summary 已完整）──
                for sr in calc_rows:
                    indicator = sr['指标']

                    if '达成率' in indicator:
                        # 成交金额达成率 = 实际成交金额 / 成交金额目标 × 100
                        actual_key = '实际成交金额' if model_name else '成交金额达成'
                        target_key = '成交金额目标'
                        actual_row = actual_summary.get(actual_key, {})
                        target_row = _get_target_row_by_indicator(rows_data, target_key)
                        if target_row is None:
                            continue

                        row_vals = [indicator]
                        calc_total_target = 0.0
                        calc_total_actual = 0.0
                        for d in date_list:
                            a = actual_row.get(d, 0)
                            if a > 0:
                                t = target_row.get(d, 0) or 0
                                calc_total_target += t
                                calc_total_actual += a
                        rate = calc_total_actual / calc_total_target * 100 if calc_total_target > 0 else 0
                        row_vals.append(_fmt_rate_cell(rate))
                        for d in date_list:
                            t = target_row.get(d, 0) or 0
                            a = actual_row.get(d, 0)
                            if t > 0:
                                r = a / t * 100
                                row_vals.append(_fmt_rate_cell(r))
                            else:
                                row_vals.append('--')
                        table_data.append(row_vals)

                    elif '费率' in indicator and '目标' not in indicator:
                        # 实际费率 = 推广花费 / 实际成交金额 × 100
                        actual_key = '实际成交金额' if model_name else '成交金额达成'
                        actual_row = actual_summary.get(actual_key, {})
                        spend_row_data = actual_summary.get('实际投入金额', {})
                        row_vals = [indicator]
                        # 直接汇总后相除，与结构表保持一致（不逐日过滤无销额的天）
                        calc_total_spend = spend_row_data.get('合计', 0.0)
                        calc_total_actual_amt = actual_row.get('合计', 0.0)
                        rate = calc_total_spend / calc_total_actual_amt * 100 if calc_total_actual_amt > 0 else 0
                        row_vals.append(_fmt_val(rate, is_pct=True))
                        for d in date_list:
                            spend = spend_row_data.get(d, 0)
                            actual_amt = actual_row.get(d, 0)
                            if actual_amt > 0:
                                r = spend / actual_amt * 100
                                row_vals.append(_fmt_val(r, is_pct=True))
                            else:
                                row_vals.append('--')
                        table_data.append(row_vals)

                # ── 预先收集买家数/访客数（后续两行都需要）──
                total_buyers = 0.0
                total_visitors = 0.0
                daily_visitors = {}
                daily_buyers = {}
                for d in date_list:
                    if model_name:
                        sd = daily_by_model_date.get((shop_name, model_name, d), {})
                    else:
                        sd = daily_by_shop_date.get((shop_name, d), {})
                    b = sd.get('支付买家数', 0.0)
                    v = sd.get('商品访客数', 0.0)
                    total_buyers += b
                    total_visitors += v
                    daily_buyers[d] = b
                    daily_visitors[d] = v

                # ── 追加「实际访客数」行 ──
                vis_row = ['实际访客数']
                vis_row.append(f'{total_visitors:,.0f}' if total_visitors else '--')
                for d in date_list:
                    v = daily_visitors[d]
                    vis_row.append(f'{v:,.0f}' if v else '--')
                table_data.append(vis_row)

                # ── 追加「实际转化率」行 = 买家数 / 访客数 ──
                cr_row = ['实际转化率']
                cr_total = total_buyers / total_visitors * 100 if total_visitors > 0 else 0
                cr_row.append(_fmt_val(cr_total, is_pct=True) if total_visitors > 0 else '--')
                for d in date_list:
                    buyers = daily_buyers[d]
                    visitors = daily_visitors[d]
                    if visitors > 0:
                        cr_row.append(_fmt_val(buyers / visitors * 100, is_pct=True))
                    else:
                        cr_row.append('--')
                table_data.append(cr_row)
                # 存入 actual_summary 供全店铺合计汇总使用
                actual_summary['_buyers'] = {'合计': total_buyers}
                actual_summary['_visitors'] = {'合计': total_visitors}
                for d in date_list:
                    actual_summary['_buyers'][d] = daily_buyers[d]
                    actual_summary['_visitors'][d] = daily_visitors[d]

                # ── 追加「直接ROI」行 = 推广直接订单金额 / 推广花费 ──
                total_direct = 0.0
                total_spend = 0.0
                daily_direct = {}
                daily_spend_roi = {}
                for d in date_list:
                    if model_name:
                        dv = promo_direct_by_model_date.get((shop_name, model_name, d), 0.0)
                        sv = promo_by_model_date.get((shop_name, model_name, d), 0.0)
                    else:
                        dv = promo_direct_by_shop_date.get((shop_name, d), 0.0)
                        sv = promo_by_shop_date.get((shop_name, d), 0.0)
                    total_direct += dv
                    total_spend += sv
                    daily_direct[d] = dv
                    daily_spend_roi[d] = sv
                roi_row = ['直接ROI']
                roi_total = total_direct / total_spend if total_spend > 0 else 0
                roi_row.append(f'{roi_total:.2f}' if total_spend > 0 else '--')
                for d in date_list:
                    dv = daily_direct[d]
                    sv = daily_spend_roi[d]
                    if sv > 0:
                        roi_row.append(f'{dv / sv:.2f}')
                    else:
                        roi_row.append('--')
                table_data.append(roi_row)
                # 存入 actual_summary 供全店铺合计汇总使用
                actual_summary['_direct_amt'] = {'合计': total_direct}
                actual_summary['_spend_for_roi'] = {'合计': total_spend}
                for d in date_list:
                    actual_summary['_direct_amt'][d] = daily_direct[d]
                    actual_summary['_spend_for_roi'][d] = daily_spend_roi[d]

                return table_data, actual_summary

            # ═══════════════════════════════════════════════
            # 店铺目标达成表
            # ═══════════════════════════════════════════════

            # ── 店铺销售结构 & 费用结构分析（前置）──
            # 构建 shops_order_pre
            shops_order_pre = []
            seen_shops_pre = set()
            for sr in shop_rows:
                s = sr['店铺']
                if s not in seen_shops_pre:
                    shops_order_pre.append(s)
                    seen_shops_pre.add(s)
            if shops_order_pre and any(s != '天猫小豚' for s in shops_order_pre):
                st.subheader('📊 店铺销售结构 & 费用结构')
                struct_rows = []
                total_amt_target = 0.0
                total_amt_actual = 0.0
                total_spend_budget = 0.0
                total_spend_actual = 0.0
                # 合计同比
                total_amt_target_ly = 0.0
                total_amt_actual_ly = 0.0
                total_spend_budget_ly = 0.0
                total_spend_actual_ly = 0.0

                shop_struct = {}
                for shop_name in shops_order_pre:
                    if shop_name == '天猫小豚':
                        continue
                    shop_data = [sr for sr in shop_rows if sr['店铺'] == shop_name]
                    amt_target_row = _get_target_row_by_indicator(shop_data, '成交金额目标')
                    rate_target_row = _get_target_row_by_indicator(shop_data, '目标费率')
                    amt_target = sum(amt_target_row.get(d, 0) or 0 for d in date_list) if amt_target_row else 0
                    spend_budget = 0.0
                    if amt_target_row and rate_target_row:
                        for d in date_list:
                            t = amt_target_row.get(d, 0) or 0
                            r = rate_target_row.get(d, 0) or 0
                            spend_budget += t * r
                    amt_actual = 0.0
                    for d in date_list:
                        sd = daily_by_shop_date.get((shop_name, d), {})
                        amt_actual += sd.get('支付金额', 0.0)
                    spend_actual = 0.0
                    for d in date_list:
                        spend_actual += promo_by_shop_date.get((shop_name, d), 0.0)

                    # ── 去年同期数据 ──
                    # 目标同比基准 = 去年同期全月实际值
                    amt_target_ly = 0.0
                    spend_budget_ly = 0.0
                    if _yoy_date_list_full:
                        for d in _yoy_date_list_full:
                            amt_target_ly += _daily_shop_ly_full.get((shop_name, d), 0.0)
                        for d in _yoy_date_list_full:
                            spend_budget_ly += _promo_shop_ly_full.get((shop_name, d), 0.0)
                    # 实际同比：用去年同期同期天数
                    amt_actual_ly = 0.0
                    spend_actual_ly = 0.0
                    if _yoy_date_list:
                        for d in _yoy_date_list:
                            amt_actual_ly += _daily_shop_ly.get((shop_name, d), 0.0)
                        for d in _yoy_date_list:
                            spend_actual_ly += _promo_shop_ly.get((shop_name, d), 0.0)

                    shop_struct[shop_name] = {
                        'amt_target': amt_target,
                        'spend_budget': spend_budget,
                        'amt_actual': amt_actual,
                        'spend_actual': spend_actual,
                        'amt_target_ly': amt_target_ly,
                        'spend_budget_ly': spend_budget_ly,
                        'amt_actual_ly': amt_actual_ly,
                        'spend_actual_ly': spend_actual_ly,
                    }
                    total_amt_target += amt_target
                    total_amt_actual += amt_actual
                    total_spend_budget += spend_budget
                    total_spend_actual += spend_actual
                    total_amt_target_ly += amt_target_ly
                    total_amt_actual_ly += amt_actual_ly
                    total_spend_budget_ly += spend_budget_ly
                    total_spend_actual_ly += spend_actual_ly

                for shop_name in shops_order_pre:
                    if shop_name == '天猫小豚':
                        continue
                    s = shop_struct[shop_name]
                    sales_pct_target = s['amt_target'] / total_amt_target * 100 if total_amt_target > 0 else 0
                    sales_pct_actual = s['amt_actual'] / total_amt_actual * 100 if total_amt_actual > 0 else 0
                    spend_pct_budget = s['spend_budget'] / total_spend_budget * 100 if total_spend_budget > 0 else 0
                    spend_pct_actual = s['spend_actual'] / total_spend_actual * 100 if total_spend_actual > 0 else 0
                    budget_rate = s['spend_budget'] / s['amt_target'] * 100 if s['amt_target'] > 0 else 0
                    actual_rate = s['spend_actual'] / s['amt_actual'] * 100 if s['amt_actual'] > 0 else 0
                    # 去年同期费率（目标用全月，实际用同期天数）
                    budget_rate_ly = s['spend_budget_ly'] / s['amt_target_ly'] * 100 if s['amt_target_ly'] > 0 else 0
                    actual_rate_ly = s['spend_actual_ly'] / s['amt_actual_ly'] * 100 if s['amt_actual_ly'] > 0 else 0
                    # 销额进度条: 实际/目标
                    sales_prog = min(s['amt_actual'] / s['amt_target'] * 100, 100) if s['amt_target'] > 0 else 0
                    sales_prog_color = '#22c55e' if sales_prog >= 100 else '#f59e0b' if sales_prog >= 70 else '#ef4444'
                    sales_prog_html = f'<div style="background:#e5e7eb;border-radius:4px;height:10px;width:80px;display:inline-block;vertical-align:middle;"><div style="width:{sales_prog:.0f}%;background:{sales_prog_color};height:10px;border-radius:4px;"></div></div> <span style="font-size:11px;color:{sales_prog_color};">{sales_prog:.1f}%</span>'
                    # 花费进度条: 实际/预算
                    spend_prog = min(s['spend_actual'] / s['spend_budget'] * 100, 100) if s['spend_budget'] > 0 else 0
                    spend_prog_color = '#22c55e' if spend_prog <= 105 else '#f59e0b' if spend_prog <= 120 else '#ef4444'
                    spend_prog_html = f'<div style="background:#e5e7eb;border-radius:4px;height:10px;width:80px;display:inline-block;vertical-align:middle;"><div style="width:{spend_prog:.0f}%;background:{spend_prog_color};height:10px;border-radius:4px;"></div></div> <span style="font-size:11px;color:{spend_prog_color};">{spend_prog:.1f}%</span>'
                    # ── 同比计算 ──
                    yoy_amt_target = _yoy_pct_val(s['amt_target'], s['amt_target_ly'])
                    yoy_amt_actual = _yoy_pct_val(s['amt_actual'], s['amt_actual_ly'])
                    yoy_spend_budget = _yoy_pct_val(s['spend_budget'], s['spend_budget_ly'])
                    yoy_spend_actual = _yoy_pct_val(s['spend_actual'], s['spend_actual_ly'])
                    yoy_budget_rate = _yoy_pct_val(budget_rate, budget_rate_ly)
                    yoy_actual_rate = _yoy_pct_val(actual_rate, actual_rate_ly)
                    # 新列序：同比放在对应数据旁边
                    struct_rows.append({
                        '_sort_key': s['amt_target'],
                        'is_total': False,
                        'data': [
                            shop_name,
                            f'{s["amt_target"]:,.0f}', _fmt_yoy_val(yoy_amt_target), f'{sales_pct_target:.1f}%',
                            f'{s["amt_actual"]:,.0f}', _fmt_yoy_val(yoy_amt_actual), f'{sales_pct_actual:.1f}%',
                            sales_prog_html,
                            f'{s["spend_budget"]:,.0f}', _fmt_yoy_val(yoy_spend_budget), f'{spend_pct_budget:.1f}%',
                            f'{s["spend_actual"]:,.0f}', _fmt_yoy_val(yoy_spend_actual), f'{spend_pct_actual:.1f}%',
                            spend_prog_html,
                            f'{budget_rate:.1f}%', _fmt_yoy_val(yoy_budget_rate),
                            f'{actual_rate:.1f}%', _fmt_yoy_val(yoy_actual_rate),
                            f'{actual_rate - budget_rate:+.1f}%',
                        ]
                    })

                # 按目标销额从高到低排序
                struct_rows.sort(key=lambda x: x['_sort_key'], reverse=True)

                # ── 合计行 ──
                if total_amt_target > 0:
                    _tt_sales_pct_target = 100.0
                    _tt_sales_pct_actual = 100.0
                    _tt_spend_pct_budget = 100.0
                    _tt_spend_pct_actual = 100.0
                    _tt_budget_rate = total_spend_budget / total_amt_target * 100 if total_amt_target > 0 else 0
                    _tt_actual_rate = total_spend_actual / total_amt_actual * 100 if total_amt_actual > 0 else 0
                    _tt_budget_rate_ly = total_spend_budget_ly / total_amt_target_ly * 100 if total_amt_target_ly > 0 else 0
                    _tt_actual_rate_ly = total_spend_actual_ly / total_amt_actual_ly * 100 if total_amt_actual_ly > 0 else 0
                    _tt_sales_prog = min(total_amt_actual / total_amt_target * 100, 100) if total_amt_target > 0 else 0
                    _tt_sales_prog_color = '#22c55e' if _tt_sales_prog >= 100 else '#f59e0b' if _tt_sales_prog >= 70 else '#ef4444'
                    _tt_sales_prog_html = f'<div style="background:#e5e7eb;border-radius:4px;height:10px;width:80px;display:inline-block;vertical-align:middle;"><div style="width:{_tt_sales_prog:.0f}%;background:{_tt_sales_prog_color};height:10px;border-radius:4px;"></div></div> <span style="font-size:11px;color:{_tt_sales_prog_color};">{_tt_sales_prog:.1f}%</span>'
                    _tt_spend_prog = min(total_spend_actual / total_spend_budget * 100, 100) if total_spend_budget > 0 else 0
                    _tt_spend_prog_color = '#22c55e' if _tt_spend_prog <= 105 else '#f59e0b' if _tt_spend_prog <= 120 else '#ef4444'
                    _tt_spend_prog_html = f'<div style="background:#e5e7eb;border-radius:4px;height:10px;width:80px;display:inline-block;vertical-align:middle;"><div style="width:{_tt_spend_prog:.0f}%;background:{_tt_spend_prog_color};height:10px;border-radius:4px;"></div></div> <span style="font-size:11px;color:{_tt_spend_prog_color};">{_tt_spend_prog:.1f}%</span>'
                    _tt_yoy_amt_target = _yoy_pct_val(total_amt_target, total_amt_target_ly)
                    _tt_yoy_amt_actual = _yoy_pct_val(total_amt_actual, total_amt_actual_ly)
                    _tt_yoy_spend_budget = _yoy_pct_val(total_spend_budget, total_spend_budget_ly)
                    _tt_yoy_spend_actual = _yoy_pct_val(total_spend_actual, total_spend_actual_ly)
                    _tt_yoy_budget_rate = _yoy_pct_val(_tt_budget_rate, _tt_budget_rate_ly)
                    _tt_yoy_actual_rate = _yoy_pct_val(_tt_actual_rate, _tt_actual_rate_ly)
                    total_row = [
                        '<b>合计</b>',
                        f'<b>{total_amt_target:,.0f}</b>', _fmt_yoy_val(_tt_yoy_amt_target), '<b>100.0%</b>',
                        f'<b>{total_amt_actual:,.0f}</b>', _fmt_yoy_val(_tt_yoy_amt_actual), '<b>100.0%</b>',
                        _tt_sales_prog_html,
                        f'<b>{total_spend_budget:,.0f}</b>', _fmt_yoy_val(_tt_yoy_spend_budget), '<b>100.0%</b>',
                        f'<b>{total_spend_actual:,.0f}</b>', _fmt_yoy_val(_tt_yoy_spend_actual), '<b>100.0%</b>',
                        _tt_spend_prog_html,
                        f'<b>{_tt_budget_rate:.1f}%</b>', _fmt_yoy_val(_tt_yoy_budget_rate),
                        f'<b>{_tt_actual_rate:.1f}%</b>', _fmt_yoy_val(_tt_yoy_actual_rate),
                        f'<b>{_tt_actual_rate - _tt_budget_rate:+.1f}%</b>',
                    ]
                    struct_rows.append({'_sort_key': -1, 'is_total': True, 'data': total_row})

                if struct_rows:
                    struct_header = ['店铺',
                                     '目标销额', '同比', '销额占比',
                                     '实际销额', '同比', '实际占比', '销额进度',
                                     '预算花费', '同比', '预算占比',
                                     '实际花费', '同比', '实际占比', '花费进度',
                                     '预算费率', '同比', '实际费率', '同比', '费率差异']
                    # 差异列索引: 费率差异=18
                    diff_cols = {18}
                    # 同比列索引: 2,5,9,12,16,17
                    yoy_cols = {2, 5, 9, 12, 16, 17}
                    html = '<div class="styled-table-wrap"><table class="styled-table"><thead><tr>'
                    for h in struct_header:
                        html += f'<th style="text-align:center;white-space:nowrap;">{h}</th>'
                    html += '</tr></thead><tbody>'
                    for item in struct_rows:
                        row = item['data']
                        is_total = item.get('is_total', False)
                        html += '<tr style="font-weight:bold;background:#f1f5f9;">' if is_total else '<tr>'
                        for j, cell in enumerate(row):
                            style = 'text-align:center;white-space:nowrap;'
                            if j in diff_cols:
                                try:
                                    v = float(str(cell).rstrip('%').replace('<b>','').replace('</b>',''))
                                    color = '#22c55e' if v > 0 else '#ef4444' if v < 0 else '#888'
                                    style += f'color:{color};font-weight:bold;'
                                except:
                                    pass
                            elif j in yoy_cols:
                                try:
                                    raw = str(cell).replace('<b>','').replace('</b>','').rstrip('%')
                                    v = float(raw)
                                    color = '#22c55e' if v >= 0 else '#ef4444'
                                    style += f'color:{color};font-weight:bold;'
                                except:
                                    style += 'color:#94a3b8;'
                            html += f'<td style="{style}">{cell}</td>'
                        html += '</tr>'
                    html += '</tbody></table></div>'
                    st.markdown(html, unsafe_allow_html=True)
                    st.caption('💡 费率差异 = 实际费率 − 预算费率；目标同比=去年同期全月，实际同比=去年同期同期天数')

            # ── 单品销售结构 & 费用结构分析（前置）──
            # 先收集所有 (shop, model) 的唯一组合，并为每个组合找到成交金额目标行
            # 元组: (model, amt_target, spend_budget, amt_actual, spend_actual, amt_target_ly, spend_budget_ly, amt_actual_ly, spend_actual_ly)
            model_struct_pre = {}  # {shop: [(model, ...)]}
            seen_models = set()  # (shop, model)
            for mr in model_rows:
                shop_name = mr['店铺']
                model_name = mr['型号']
                key = (shop_name, model_name)
                if key in seen_models:
                    continue
                seen_models.add(key)
                # 从 model_rows 中查找该 (shop, model) 的成交金额目标行
                _mr_filtered = [r for r in model_rows if r['店铺'] == shop_name and r['型号'] == model_name]
                amt_target_row = _get_target_row_by_indicator(_mr_filtered, '成交金额目标')
                amt_target = sum(amt_target_row.get(d, 0) or 0 for d in date_list) if amt_target_row else 0
                # 预算花费 = 成交金额目标 × 目标费率（逐日加权）
                rate_target_row = _get_target_row_by_indicator(_mr_filtered, '目标费率')
                spend_budget = 0.0
                if amt_target_row and rate_target_row:
                    for d in date_list:
                        t = amt_target_row.get(d, 0) or 0
                        r = rate_target_row.get(d, 0) or 0
                        spend_budget += t * r
                amt_actual = 0.0
                for d in date_list:
                    sd = daily_by_model_date.get((shop_name, model_name, d), {})
                    amt_actual += sd.get('支付金额', 0.0)
                spend_actual = 0.0
                for d in date_list:
                    sv = promo_by_model_date.get((shop_name, model_name, d), 0.0)
                    spend_actual += sv
                # ── 去年同期数据 ──
                # 目标同比基准 = 去年同期全月实际值
                amt_target_ly = 0.0
                spend_budget_ly = 0.0
                if _yoy_date_list_full:
                    for d in _yoy_date_list_full:
                        amt_target_ly += _daily_model_ly_full.get((shop_name, model_name, d), 0.0)
                    for d in _yoy_date_list_full:
                        sv = _promo_model_ly_full.get((shop_name, model_name, d), 0.0)
                        spend_budget_ly += sv
                # 实际同比：用去年同期同期天数
                amt_actual_ly = 0.0
                spend_actual_ly = 0.0
                if _yoy_date_list:
                    for d in _yoy_date_list:
                        amt_actual_ly += _daily_model_ly.get((shop_name, model_name, d), 0.0)
                    for d in _yoy_date_list:
                        sv = _promo_model_ly.get((shop_name, model_name, d), 0.0)
                        spend_actual_ly += sv
                if shop_name not in model_struct_pre:
                    model_struct_pre[shop_name] = []
                model_struct_pre[shop_name].append((model_name, amt_target, spend_budget, amt_actual, spend_actual,
                                                    amt_target_ly, spend_budget_ly, amt_actual_ly, spend_actual_ly))

            if model_struct_pre:
                allowed_shops = ['华为京东自营', '天猫华为官旗', '天猫智选']
                st.subheader('📊 单品销售结构 & 费用结构（按店铺分组）')
                for shop_name in allowed_shops:
                    if shop_name not in model_struct_pre:
                        continue
                    models = model_struct_pre[shop_name]
                    shop_total_target = sum(m[1] for m in models)
                    shop_total_budget = sum(m[2] for m in models)
                    shop_total_actual = sum(m[3] for m in models)
                    shop_total_spend = sum(m[4] for m in models)

                    mrows = []
                    for (model_name, amt_target, spend_budget, amt_actual, spend_actual,
                         amt_target_ly, spend_budget_ly, amt_actual_ly, spend_actual_ly) in models:
                        sales_pct_target = amt_target / shop_total_target * 100 if shop_total_target > 0 else 0
                        sales_pct_actual = amt_actual / shop_total_actual * 100 if shop_total_actual > 0 else 0
                        budget_pct = spend_budget / shop_total_budget * 100 if shop_total_budget > 0 else 0
                        spend_pct = spend_actual / shop_total_spend * 100 if shop_total_spend > 0 else 0
                        budget_rate = spend_budget / amt_target * 100 if amt_target > 0 else 0
                        actual_rate = spend_actual / amt_actual * 100 if amt_actual > 0 else 0
                        # 销额进度条
                        sales_prog = min(amt_actual / amt_target * 100, 100) if amt_target > 0 else 0
                        sp_color = '#22c55e' if sales_prog >= 100 else '#f59e0b' if sales_prog >= 70 else '#ef4444'
                        sales_prog_html = f'<div style="background:#e5e7eb;border-radius:4px;height:10px;width:70px;display:inline-block;vertical-align:middle;"><div style="width:{sales_prog:.0f}%;background:{sp_color};height:10px;border-radius:4px;"></div></div> <span style="font-size:11px;color:{sp_color};">{sales_prog:.1f}%</span>'
                        # 花费进度条
                        spend_prog = min(spend_actual / spend_budget * 100, 100) if spend_budget > 0 else 0
                        ep_color = '#22c55e' if spend_prog <= 105 else '#f59e0b' if spend_prog <= 120 else '#ef4444'
                        spend_prog_html = f'<div style="background:#e5e7eb;border-radius:4px;height:10px;width:70px;display:inline-block;vertical-align:middle;"><div style="width:{spend_prog:.0f}%;background:{ep_color};height:10px;border-radius:4px;"></div></div> <span style="font-size:11px;color:{ep_color};">{spend_prog:.1f}%</span>'
                        # 去年同期费率
                        budget_rate_ly = spend_budget_ly / amt_target_ly * 100 if amt_target_ly > 0 else 0
                        actual_rate_ly = spend_actual_ly / amt_actual_ly * 100 if amt_actual_ly > 0 else 0
                        # ── 同比计算 ──
                        yoy_amt_target = _yoy_pct_val(amt_target, amt_target_ly)
                        yoy_amt_actual = _yoy_pct_val(amt_actual, amt_actual_ly)
                        yoy_spend_budget = _yoy_pct_val(spend_budget, spend_budget_ly)
                        yoy_spend_actual = _yoy_pct_val(spend_actual, spend_actual_ly)
                        yoy_budget_rate = _yoy_pct_val(budget_rate, budget_rate_ly)
                        yoy_actual_rate = _yoy_pct_val(actual_rate, actual_rate_ly)
                        # 新列序：同比放在对应数据旁边
                        mrows.append({
                            '_sort_key': amt_target,
                            'data': [
                                model_name,
                                f'{amt_target:,.0f}', _fmt_yoy_val(yoy_amt_target), f'{sales_pct_target:.1f}%',
                                f'{amt_actual:,.0f}', _fmt_yoy_val(yoy_amt_actual), f'{sales_pct_actual:.1f}%',
                                sales_prog_html,
                                f'{spend_budget:,.0f}', _fmt_yoy_val(yoy_spend_budget), f'{budget_pct:.1f}%',
                                f'{spend_actual:,.0f}', _fmt_yoy_val(yoy_spend_actual), f'{spend_pct:.1f}%',
                                spend_prog_html,
                                f'{budget_rate:.1f}%', _fmt_yoy_val(yoy_budget_rate),
                                f'{actual_rate:.1f}%', _fmt_yoy_val(yoy_actual_rate),
                                f'{actual_rate - budget_rate:+.1f}%',
                            ]
                        })

                    # 按目标销额从高到低排序
                    mrows.sort(key=lambda x: x['_sort_key'], reverse=True)

                    if mrows:
                        st.markdown(f'**{shop_name}**')
                        mheader = ['型号', '目标销额', '同比', '销额占比',
                                   '实际销额', '同比', '实际占比', '销额进度',
                                   '预算花费', '同比', '预算占比',
                                   '实际花费', '同比', '实际占比', '花费进度',
                                   '预算费率', '同比', '实际费率', '同比', '费率差异']
                        diff_cols_m = {18}
                        yoy_cols_m = {2, 5, 9, 12, 16, 17}
                        html = '<div class="styled-table-wrap"><table class="styled-table"><thead><tr>'
                        for h in mheader:
                            html += f'<th style="text-align:center;white-space:nowrap;">{h}</th>'
                        html += '</tr></thead><tbody>'
                        for item in mrows:
                            row = item['data']
                            html += '<tr>'
                            for j, cell in enumerate(row):
                                style = 'text-align:center;white-space:nowrap;'
                                if j in diff_cols_m:
                                    try:
                                        v = float(str(cell).rstrip('%'))
                                        color = '#22c55e' if v > 0 else '#ef4444' if v < 0 else '#888'
                                        style += f'color:{color};font-weight:bold;'
                                    except:
                                        pass
                                elif j in yoy_cols_m:
                                    try:
                                        raw = str(cell).rstrip('%')
                                        v = float(raw)
                                        color = '#22c55e' if v >= 0 else '#ef4444'
                                        style += f'color:{color};font-weight:bold;'
                                    except:
                                        style += 'color:#94a3b8;'
                                html += f'<td style="{style}">{cell}</td>'
                            html += '</tr>'
                        html += '</tbody></table></div>'
                        st.markdown(html, unsafe_allow_html=True)
                    st.caption('💡 销额/花费进度 = 实际/目标（预算）；费率差异 = 实际费率 − 预算费率；目标同比=去年同期全月，实际同比=去年同期同期天数')

            st.subheader('🏪 店铺目标达成')
            if shop_rows:
                shops_order = []
                seen_shops = set()
                for sr in shop_rows:
                    s = sr['店铺']
                    if s not in seen_shops:
                        shops_order.append(s)
                        seen_shops.add(s)

                # 全店铺合计（排除天猫小豚）
                all_shop_target = {}    # {指标: {date: val}}
                all_shop_actual = {}    # {指标: {date: val}}
                # 加权费率分子: {date: Σ(成交金额目标 × 目标费率)}
                all_shop_weighted_rate = {}
                # 保存每个店铺自己的 actual_summary，供结构分析用
                shop_actual_summaries = {}  # {shop_name: actual_summary}

                for shop_name in shops_order:
                    shop_data = [sr for sr in shop_rows if sr['店铺'] == shop_name]
                    if not shop_data:
                        continue

                    header_cols = ['指标', '合计'] + date_list
                    table_data, actual_summary = _build_rows_for_entity(shop_data, shop_name)
                    shop_actual_summaries[shop_name] = actual_summary

                    if table_data:
                        _render_target_table(header_cols, shop_name, table_data)

                    # 累计全店铺数据（排除天猫小豚）
                    if shop_name != '天猫小豚':
                        for sr in shop_data:
                            indicator = sr['指标']
                            itype = _indicator_type(indicator)
                            if itype == 'target':
                                if indicator not in all_shop_target:
                                    all_shop_target[indicator] = {}
                                for d in date_list:
                                    v = sr.get(d, 0) or 0
                                    all_shop_target[indicator][d] = all_shop_target[indicator].get(d, 0.0) + v
                                # 如果是目标费率，累积加权分子: Σ(成交金额目标 × 目标费率)
                                if '费率' in indicator:
                                    amt_target_row = _get_target_row_by_indicator(shop_data, '成交金额目标')
                                    for d in date_list:
                                        amt_t = (amt_target_row.get(d, 0) or 0) if amt_target_row else 0
                                        rate_v = sr.get(d, 0) or 0
                                        all_shop_weighted_rate[d] = all_shop_weighted_rate.get(d, 0.0) + amt_t * rate_v
                            elif itype == 'actual':
                                if indicator not in all_shop_actual:
                                    all_shop_actual[indicator] = {}
                                for d in date_list:
                                    all_shop_actual[indicator][d] = all_shop_actual[indicator].get(d, 0.0) + actual_summary.get(indicator, {}).get(d, 0)
                        # 追加代码自动生成的 actual 指标（Excel中没有的才补，避免重复累加）
                        excel_actual_indicators = set()
                        for sr in shop_data:
                            if _indicator_type(sr['指标']) == 'actual':
                                excel_actual_indicators.add(sr['指标'])
                        for indicator, vals in actual_summary.items():
                            if indicator not in excel_actual_indicators:
                                if indicator not in all_shop_actual:
                                    all_shop_actual[indicator] = {}
                                for d in date_list:
                                    all_shop_actual[indicator][d] = all_shop_actual[indicator].get(d, 0.0) + vals.get(d, 0)

                # ── 全店铺合计（排除天猫小豚）──
                if all_shop_target:
                    st.subheader('🏢 全店铺合计达成（不含天猫小豚）')
                    header_cols = ['指标', '合计'] + date_list
                    table_data = []

                    # 固定指标顺序（与单店铺表格一致）
                    # 收集所有存在的指标类型
                    existing_indicators = set(list(all_shop_target.keys()) + list(all_shop_actual.keys()))
                    fixed_order = ['成交金额目标', '成交金额达成', '实际投入金额', '目标费率', '实际支付件数', '成交金额达成率', '实际费率']
                    always_show = {'成交金额达成率', '实际费率', '实际支付件数'}
                    all_indicators = [ind for ind in fixed_order if ind in existing_indicators or ind in always_show]

                    for indicator in all_indicators:
                        itype = _indicator_type(indicator)

                        if itype == 'target':
                            is_rate = _is_rate_indicator(indicator)
                            row_vals = [indicator]
                            if is_rate and '目标费率' in indicator:
                                # 全店铺目标费率 = Σ(成交金额目标 × 各店铺目标费率) / Σ(成交金额目标)
                                # 按每日计算加权平均
                                grand_sum_weighted = sum(all_shop_weighted_rate.values())
                                grand_sum_amt = sum(all_shop_target.get('成交金额目标', {}).values())
                                if grand_sum_amt > 0:
                                    row_vals.append(_fmt_val(grand_sum_weighted / grand_sum_amt * 100, is_pct=True))
                                else:
                                    row_vals.append('--')
                                for d in date_list:
                                    amt_target = all_shop_target.get('成交金额目标', {}).get(d, 0)
                                    weighted_v = all_shop_weighted_rate.get(d, 0)
                                    if amt_target > 0:
                                        r = weighted_v / amt_target * 100
                                        row_vals.append(_fmt_val(r, is_pct=True))
                                    else:
                                        row_vals.append('--')
                            else:
                                grand_sum = sum(all_shop_target.get(indicator, {}).values())
                                row_vals.append(f'{grand_sum:,.0f}')
                                for d in date_list:
                                    v = all_shop_target.get(indicator, {}).get(d, 0)
                                    row_vals.append(f'{v:,.0f}' if v else '--')
                            table_data.append(row_vals)

                        elif itype == 'actual':
                            row_vals = [indicator]
                            grand_sum = sum(all_shop_actual.get(indicator, {}).values())
                            row_vals.append(f'{grand_sum:,.0f}' if grand_sum else '--')
                            for d in date_list:
                                v = all_shop_actual.get(indicator, {}).get(d, 0)
                                row_vals.append(f'{v:,.0f}' if v else '--')
                            table_data.append(row_vals)

                        elif itype == 'calc':
                            if '达成率' in indicator:
                                target_key = '成交金额目标'
                                actual_key = '成交金额达成'
                                target_data = all_shop_target.get(target_key, {})
                                actual_data = all_shop_actual.get(actual_key, {})
                                row_vals = [indicator]
                                total_t = 0.0
                                total_a = 0.0
                                for d in date_list:
                                    a = actual_data.get(d, 0)
                                    if a > 0:
                                        total_t += target_data.get(d, 0)
                                        total_a += a
                                rate = total_a / total_t * 100 if total_t > 0 else 0
                                row_vals.append(_fmt_rate_cell(rate))
                                for d in date_list:
                                    t = target_data.get(d, 0)
                                    a = actual_data.get(d, 0)
                                    if t > 0:
                                        r = a / t * 100
                                        row_vals.append(_fmt_rate_cell(r))
                                    else:
                                        row_vals.append('--')
                                table_data.append(row_vals)

                            elif '费率' in indicator and '目标' not in indicator:
                                actual_amt_key = '成交金额达成'
                                spend_key = '实际投入金额'
                                actual_amt_data = all_shop_actual.get(actual_amt_key, {})
                                spend_data = all_shop_actual.get(spend_key, {})
                                row_vals = [indicator]
                                total_amt = 0.0
                                total_spend = 0.0
                                for d in date_list:
                                    amt = actual_amt_data.get(d, 0)
                                    if amt > 0:
                                        total_amt += amt
                                        total_spend += spend_data.get(d, 0)
                                rate = total_spend / total_amt * 100 if total_amt > 0 else 0
                                row_vals.append(_fmt_val(rate, is_pct=True))
                                for d in date_list:
                                    amt = actual_amt_data.get(d, 0)
                                    spend = spend_data.get(d, 0)
                                    if amt > 0:
                                        r = spend / amt * 100
                                        row_vals.append(_fmt_val(r, is_pct=True))
                                    else:
                                        row_vals.append('--')
                                table_data.append(row_vals)

                    # ── 全店铺合计「实际访客数」+「实际转化率」行 ──
                    all_buyers_data = all_shop_actual.get('_buyers', {})
                    all_visitors_data = all_shop_actual.get('_visitors', {})
                    total_buyers_all = sum(v for k, v in all_buyers_data.items() if k != '合计')
                    total_visitors_all = sum(v for k, v in all_visitors_data.items() if k != '合计')
                    # 实际访客数行
                    vis_row_all = ['实际访客数']
                    vis_row_all.append(f'{total_visitors_all:,.0f}' if total_visitors_all else '--')
                    for d in date_list:
                        v = all_visitors_data.get(d, 0)
                        vis_row_all.append(f'{v:,.0f}' if v else '--')
                    table_data.append(vis_row_all)
                    # 实际转化率行
                    cr_row = ['实际转化率']
                    cr_row.append(_fmt_val(total_buyers_all / total_visitors_all * 100, is_pct=True) if total_visitors_all > 0 else '--')
                    for d in date_list:
                        buyers = all_buyers_data.get(d, 0)
                        visitors = all_visitors_data.get(d, 0)
                        if visitors > 0:
                            cr_row.append(_fmt_val(buyers / visitors * 100, is_pct=True))
                        else:
                            cr_row.append('--')
                    table_data.append(cr_row)

                    # ── 全店铺合计「直接ROI」行 ──
                    all_direct_data = all_shop_actual.get('_direct_amt', {})
                    all_spend_roi_data = all_shop_actual.get('_spend_for_roi', {})
                    total_direct_all = sum(v for k, v in all_direct_data.items() if k != '合计')
                    total_spend_roi_all = sum(v for k, v in all_spend_roi_data.items() if k != '合计')
                    roi_row_all = ['直接ROI']
                    roi_row_all.append(f'{total_direct_all / total_spend_roi_all:.2f}' if total_spend_roi_all > 0 else '--')
                    for d in date_list:
                        dv = all_direct_data.get(d, 0)
                        sv = all_spend_roi_data.get(d, 0)
                        if sv > 0:
                            roi_row_all.append(f'{dv / sv:.2f}')
                        else:
                            roi_row_all.append('--')
                    table_data.append(roi_row_all)

                    if table_data:
                        _render_target_table(header_cols, '全店铺合计（不含天猫小豚）', table_data)

            else:
                st.info('该月份无店铺目标数据')

            # ═══════════════════════════════════════════════
            # 单品目标达成表
            # ═══════════════════════════════════════════════
            st.subheader('📦 单品目标达成')
            if model_rows:
                model_groups = {}
                for mr in model_rows:
                    key = (mr['店铺'], mr['型号'])
                    if key not in model_groups:
                        model_groups[key] = []
                    model_groups[key].append(mr)

                for (shop_name, model_name), mdata in model_groups.items():
                    header_cols = ['指标', '合计'] + date_list
                    table_data, actual_summary = _build_rows_for_entity(mdata, shop_name, model_name)

                    # 单品表格去重：按指标名只保留首次出现的行（防止Excel中重复录入）
                    if model_name and table_data:
                        seen_indicators = set()
                        deduped = []
                        for row in table_data:
                            indicator = row[0]
                            if indicator not in seen_indicators:
                                seen_indicators.add(indicator)
                                deduped.append(row)
                        table_data = deduped

                    if table_data:
                        # 跳过全空表格：检查是否所有数据行（除第一列指标名外）都是 '--'
                        all_empty = True
                        for row in table_data:
                            for cell in row[1:]:  # 跳过指标名列
                                if cell != '--':
                                    all_empty = False
                                    break
                            if not all_empty:
                                break
                        if not all_empty:
                            _render_target_table(header_cols, f'{shop_name} · {model_name}', table_data)
            else:
                st.info('该月份无单品目标数据')